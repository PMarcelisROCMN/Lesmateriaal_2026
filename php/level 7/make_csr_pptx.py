from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Colors ────────────────────────────────────────────────────────────────
C_CTRL   = RGBColor(0x10, 0xB9, 0x81)  # emerald green  – Controller
C_SVC    = RGBColor(0xF5, 0x9E, 0x0B)  # amber          – Service
C_REPO   = RGBColor(0x3B, 0x82, 0xF6)  # blue           – Repository
C_GRAY   = RGBColor(0x64, 0x74, 0x8B)  # slate          – Client / DB
C_BG     = RGBColor(0x0F, 0x17, 0x2A)  # dark navy
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT   = RGBColor(0x1E, 0x29, 0x3B)
C_MUTED  = RGBColor(0x6B, 0x72, 0x80)
C_SLIDE  = RGBColor(0xF8, 0xFA, 0xFC)
C_CODE   = RGBColor(0x1A, 0x23, 0x37)
C_CODE_T = RGBColor(0xE2, 0xE8, 0xF0)
C_GREEN  = RGBColor(0x05, 0x7A, 0x55)
C_RED    = RGBColor(0xB9, 0x1C, 0x1C)
C_CTRL_L = RGBColor(0xD1, 0xFA, 0xE5)
C_SVC_L  = RGBColor(0xFE, 0xF3, 0xC7)
C_REPO_L = RGBColor(0xDB, 0xEA, 0xFE)
C_CTRL_D = RGBColor(0x06, 0x57, 0x3D)
C_SVC_D  = RGBColor(0x78, 0x35, 0x00)
C_REPO_D = RGBColor(0x1D, 0x4E, 0xD8)

W, H = 13.33, 7.5

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)


# ── Low-level helpers ─────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def bg(slide, c):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = c


def box(slide, l, t, w, h, fill, bdr=None, bw=1.5):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if bdr:
        s.line.color.rgb = bdr
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s


def _set_run(run, text, sz, bold, italic, color, font):
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color


def tb(slide, text, l, t, w, h, sz=16, bold=False, c=None,
       al=PP_ALIGN.LEFT, italic=False, font="Calibri", wrap=True):
    if c is None:
        c = C_TEXT
    shape = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = al
    _set_run(p.add_run(), text, sz, bold, italic, c, font)
    return shape


def boxtb(slide, text, l, t, w, h, fill, tc=None, sz=16, bold=False,
          al=PP_ALIGN.CENTER, bdr=None, bw=1.5, font="Calibri", wrap=True,
          vcenter=True):
    if tc is None:
        tc = C_WHITE
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if bdr:
        s.line.color.rgb = bdr
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    tf = s.text_frame
    tf.word_wrap = wrap
    if vcenter:
        body_pr = tf._txBody.find(qn('a:bodyPr'))
        if body_pr is not None:
            body_pr.set('anchor', 'ctr')
    p = tf.paragraphs[0]
    p.alignment = al
    _set_run(p.add_run(), text, sz, bold, False, tc, font)
    return s


def blist(slide, items, l, t, w, h, sz=15, c=None, sp=6, font="Calibri"):
    if c is None:
        c = C_TEXT
    shape = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(sp if i > 0 else 0)
        _set_run(p.add_run(), item, sz, False, False, c, font)
    return shape


def code_box(slide, lines, l, t, w, h, sz=13):
    box(slide, l, t, w, h, C_CODE)
    shape = slide.shapes.add_textbox(
        Inches(l + 0.18), Inches(t + 0.2), Inches(w - 0.35), Inches(h - 0.35))
    shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0)
        _set_run(p.add_run(), line, sz, False, False, C_CODE_T, "Consolas")
    return shape


def divider(slide, y=0.9, c=None):
    if c is None:
        c = C_MUTED
    box(slide, 0.4, y, W - 0.8, 0.015, c)


# ── Mini stack indicator (shown on layer slides) ──────────────────────────

def mini_stack(slide, active):
    """Draw small C→S→R stack at top-right, highlighting 'active' layer."""
    colors = {"Controller": C_CTRL, "Service": C_SVC, "Repository": C_REPO}
    labels = ["Controller", "Service", "Repository"]
    x, y, bw2, bh2, gap = 10.55, 0.15, 0.88, 0.42, 0.06
    for i, label in enumerate(labels):
        fill = colors[label] if label == active else C_GRAY
        alpha = 1.0 if label == active else 0.5
        boxtb(slide, label, x + i * (bw2 + gap), y, bw2, bh2, fill,
              sz=10, bold=(label == active))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ═══════════════════════════════════════════════════════════════════════════
s1 = new_slide()
bg(s1, C_BG)

box(s1, 0, 0, W, 0.14, C_CTRL)          # top accent
box(s1, 0, H - 0.14, W, 0.14, C_REPO)  # bottom accent

tb(s1, "Controller  –  Service  –  Repository",
   0.6, 1.5, 12.1, 1.5, sz=42, bold=True, c=C_WHITE, al=PP_ALIGN.CENTER)

tb(s1, "Een gelaagde architectuur voor PHP webapplicaties",
   0.6, 3.0, 12.1, 0.75, sz=22, c=RGBColor(0x94, 0xA3, 0xB8),
   al=PP_ALIGN.CENTER)

# Three preview boxes
bw3, bh3 = 3.6, 0.85
gap3 = 0.3
sx3 = (W - 3 * bw3 - 2 * gap3) / 2
by3 = 4.6
boxtb(s1, "Controller", sx3,                  by3, bw3, bh3, C_CTRL, sz=22, bold=True)
boxtb(s1, "Service",    sx3 + bw3 + gap3,     by3, bw3, bh3, C_SVC,  tc=C_TEXT, sz=22, bold=True)
boxtb(s1, "Repository", sx3 + 2*(bw3+gap3),   by3, bw3, bh3, C_REPO, sz=22, bold=True)

tb(s1, "→", sx3 + bw3 - 0.1,           by3 + 0.2, 0.5, 0.45, sz=22, c=C_WHITE, al=PP_ALIGN.CENTER)
tb(s1, "→", sx3 + 2*bw3 + gap3 - 0.1,  by3 + 0.2, 0.5, 0.45, sz=22, c=C_WHITE, al=PP_ALIGN.CENTER)

tb(s1, "CSR Pattern", 0.6, 6.2, 12.1, 0.5, sz=13,
   c=RGBColor(0x47, 0x55, 0x69), al=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Overview: De drie lagen
# ═══════════════════════════════════════════════════════════════════════════
s2 = new_slide()
bg(s2, C_SLIDE)

box(s2, 0, 0, W, 0.9, C_BG)
tb(s2, "De drie lagen", 0.5, 0.15, 10, 0.65, sz=30, bold=True, c=C_WHITE)

# ── Left: vertical stack diagram ──────────────────────────────────────────
cx, cw, bh_d = 0.55, 4.5, 0.68
gap_d = 0.38

layers_d = [
    (C_GRAY,  "Client  (browser)",  1.15),
    (C_CTRL,  "Controller",         1.15 + bh_d + gap_d),
    (C_SVC,   "Service",            1.15 + 2*(bh_d + gap_d)),
    (C_REPO,  "Repository",         1.15 + 3*(bh_d + gap_d)),
    (C_GRAY,  "Database",           1.15 + 4*(bh_d + gap_d)),
]

for i, (col, lbl, ly) in enumerate(layers_d):
    tc = C_TEXT if col == C_SVC else C_WHITE
    boxtb(s2, lbl, cx, ly, cw, bh_d, col, tc=tc, sz=17, bold=(col not in (C_GRAY,)))

    if i < len(layers_d) - 1:
        labels_between = [
            "HTTP request / response",
            "methodeaanroep",
            "methodeaanroep",
            "SQL query",
        ]
        tb(s2, "▼", cx + cw/2 - 0.12, ly + bh_d + 0.03, 0.3, 0.3,
           sz=14, c=C_MUTED, al=PP_ALIGN.CENTER)
        tb(s2, labels_between[i], cx + cw/2 + 0.22, ly + bh_d + 0.05, 2.4, 0.28,
           sz=11, c=C_MUTED)

# ── Right: description cards ───────────────────────────────────────────────
dx = 5.7
desc_data = [
    (C_GRAY,  "Client",       "Stuurt HTTP-verzoeken naar de server (browser, app, Postman...)"),
    (C_CTRL,  "Controller",   "Ontvangt het verzoek, valideert de input, stuurt een HTTP-response terug"),
    (C_SVC,   "Service",      "Bevat de businesslogica — de regels die gelden voor de applicatie"),
    (C_REPO,  "Repository",   "Voert SQL-queries uit en zet databaserijen om naar objecten"),
    (C_GRAY,  "Database",     "Bewaart alle data persistent (MySQL, SQLite, ...)"),
]

for i, (col, name, desc) in enumerate(desc_data):
    ly = 1.15 + i * (bh_d + gap_d)
    box(s2, dx,       ly + 0.1, 0.22, 0.48, col)
    tb(s2, name, dx + 0.35, ly,       W - dx - 0.4, 0.42, sz=15, bold=True, c=C_TEXT)
    tb(s2, desc, dx + 0.35, ly + 0.38, W - dx - 0.4, 0.48, sz=13, c=C_MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2b – We beginnen met de Repository (transition)
# ═══════════════════════════════════════════════════════════════════════════
s2b = new_slide()
bg(s2b, C_BG)

box(s2b, 0, 0, W, 0.14, C_REPO)
box(s2b, 0, H - 0.14, W, 0.14, C_REPO)

tb(s2b, "We beginnen met...", 0.7, 0.6, 8.0, 0.65, sz=22,
   c=RGBColor(0x94, 0xA3, 0xB8))
tb(s2b, "de Repository", 0.7, 1.2, 8.0, 1.15, sz=54, bold=True, c=C_REPO)

tb(s2b, "Waarom hier starten?", 0.7, 2.52, 8.0, 0.44, sz=16, bold=True,
   c=RGBColor(0x7D, 0xD3, 0xFC))
blist(s2b, [
    "→  SQL kennen jullie al  —  de Repository is het meest herkenbaar",
    "→  Meest concreet: een databaserij wordt omgezet naar een PHP-object",
    "→  We bouwen van beneden naar boven: data eerst, logica daarna",
    "→  Service en Controller worden pas duidelijk als Repository bekend is",
], 0.82, 3.0, 7.8, 2.5, sz=16, c=RGBColor(0xCB, 0xD5, 0xE0), sp=10)

# Stack on right — Repository highlighted, others dimmed
C_DIM = RGBColor(0x1C, 0x28, 0x3A)
rx_t, rw_t, rbh_t, rgap_t = 9.3, 3.7, 0.62, 0.3

stack_t = [
    (C_DIM,  C_MUTED, "Client"),
    (C_DIM,  C_MUTED, "Controller"),
    (C_DIM,  C_MUTED, "Service"),
    (C_REPO, C_WHITE, "Repository"),
    (C_DIM,  C_MUTED, "Database"),
]
sy_t = 1.2
for i, (fill_t, tc_t, lbl_t) in enumerate(stack_t):
    y_t = sy_t + i * (rbh_t + rgap_t)
    boxtb(s2b, lbl_t, rx_t, y_t, rw_t, rbh_t, fill_t,
          tc=tc_t, sz=17, bold=(lbl_t == "Repository"))
    if i < len(stack_t) - 1:
        arr_c = C_REPO if lbl_t == "Service" else RGBColor(0x2D, 0x3C, 0x4E)
        tb(s2b, "▼", rx_t + rw_t / 2 - 0.1, y_t + rbh_t + 0.04, 0.28, 0.24,
           sz=12, c=arr_c, al=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Repository
# ═══════════════════════════════════════════════════════════════════════════
s3 = new_slide()
bg(s3, C_SLIDE)

box(s3, 0, 0, W, 1.0, C_REPO)
tb(s3, "Repository", 0.5, 0.08, 9, 0.55, sz=34, bold=True, c=C_WHITE)
tb(s3, "Datalaag  —  de enige plek waar SQL staat",
   0.5, 0.6, 9, 0.38, sz=15, c=RGBColor(0xBF, 0xDB, 0xFE))
mini_stack(s3, "Repository")

lx, lw = 0.45, 5.95

tb(s3, "Verantwoordelijkheid", lx, 1.18, lw, 0.4, sz=15, bold=True, c=C_REPO_D)
blist(s3, [
    "• Alle SQL uitvoeren (SELECT, INSERT, UPDATE, DELETE)",
    "• Rijen ophalen en omzetten naar PHP-objecten (hydration)",
    "• findAll()  ·  findById()  ·  save()  ·  delete()",
    "• De PDO-verbinding gebruiken",
], lx + 0.1, 1.58, lw - 0.15, 1.5, sz=15, c=C_TEXT)

tb(s3, "Wat hoort WEL hier?", lx, 3.18, lw, 0.38, sz=15, bold=True, c=C_GREEN)
blist(s3, [
    "✓  hydrate(array $row): Message  →  row omzetten naar object",
    "✓  PDO::prepare() en execute()",
    "✓  lastInsertId() bij INSERT",
], lx + 0.1, 3.56, lw - 0.15, 1.1, sz=15, c=C_GREEN)

tb(s3, "Wat hoort NIET hier?", lx, 4.76, lw, 0.38, sz=15, bold=True, c=C_RED)
blist(s3, [
    "✗  Business rules  (\"mag de user dit?\")",
    "✗  Inputvalidatie  (\"is dit veld ingevuld?\")",
    "✗  HTTP-statuscodes of JSON-responses",
], lx + 0.1, 5.14, lw - 0.15, 1.1, sz=15, c=C_RED)

# Code example
code_box(s3, [
    "// MessageRepository.php",
    "",
    "public function findAll(): array",
    "{",
    "    $stmt = $this->pdo",
    '        ->prepare("SELECT * FROM messages");',
    "    $stmt->execute();",
    "    $rows = $stmt->fetchAll();",
    "",
    "    // Zet elke row om naar een Message-object",
    "    return array_map($this->hydrate(...), $rows);",
    "}",
    "",
    "private function hydrate(array $row): Message",
    "{",
    "    return new Message(",
    "        id:        (int) $row['id'],",
    "        message:   $row['message'],",
    "        timestamp: new DateTimeImmutable(),",
    "        userId:    (int) $row['user_id']",
    "    );",
    "}",
], 6.6, 1.15, 6.5, 6.15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3b – Repository: Hydration & Domain Objects (in-depth)
# ═══════════════════════════════════════════════════════════════════════════
s3b = new_slide()
bg(s3b, C_SLIDE)

box(s3b, 0, 0, W, 1.0, C_REPO)
tb(s3b, "Repository — Hydration & Domain Objects",
   0.5, 0.1, 9.5, 0.55, sz=28, bold=True, c=C_WHITE)
tb(s3b, "Waarom geeft de Repository objecten terug en geen arrays?",
   0.5, 0.63, 9.5, 0.35, sz=14, c=RGBColor(0xBF, 0xDB, 0xFE))
mini_stack(s3b, "Repository")

lx_h, lw_h = 0.45, 5.7

# What is hydration
tb(s3b, "Wat is hydration?", lx_h, 1.18, lw_h, 0.38, sz=15, bold=True, c=C_REPO_D)
box(s3b, lx_h, 1.56, lw_h, 1.0, C_REPO_L, bdr=C_REPO)
tb(s3b, "Een databaserij is een associatieve array:",
   lx_h + 0.15, 1.62, lw_h - 0.25, 0.28, sz=13, c=C_TEXT)
tb(s3b, "['id'=>1,  'message'=>'Hallo!',  'user_id'=>3]",
   lx_h + 0.15, 1.9, lw_h - 0.25, 0.28, sz=12, c=C_REPO_D, font="Consolas")
tb(s3b, "Hydration = die array omzetten naar een getypeerd PHP-object.",
   lx_h + 0.15, 2.22, lw_h - 0.25, 0.3, sz=13, c=C_TEXT)

# Why domain objects
tb(s3b, "Waarom Domain Objects en geen arrays?", lx_h, 2.72, lw_h, 0.38,
   sz=15, bold=True, c=C_REPO_D)
blist(s3b, [
    "✓  Type safety  —  $msg->message  vs  $row['mesage']  (typo!)",
    "✓  IDE-autocomplete en type hints per property",
    "✓  Encapsulatie: service / controller weet niet hoe de DB eruitziet",
    "✓  Kolomnaam wijzigt? Alleen de Repository en het model aanpassen",
    "✓  Methoden op object mogelijk: $msg->isOwnedBy($user)",
], lx_h + 0.1, 3.12, lw_h - 0.15, 2.12, sz=14, c=C_TEXT, sp=5)

box(s3b, lx_h, 5.35, lw_h, 0.45, C_REPO_L, bdr=C_REPO)
tb(s3b, "Gouden regel:  alleen de Repository kent de databasestructuur.",
   lx_h + 0.15, 5.41, lw_h - 0.25, 0.34, sz=13, bold=True, c=C_REPO_D)

# Right side — code comparison
rx_h, rw_h = 6.3, 6.75

box(s3b, rx_h, 1.18, rw_h, 0.35, RGBColor(0xFF, 0xEC, 0xEC), bdr=C_RED)
tb(s3b, "Zonder hydration — kwetsbaar", rx_h + 0.14, 1.2, rw_h - 0.22, 0.3,
   sz=13, bold=True, c=C_RED)
code_box(s3b, [
    "$rows = $stmt->fetchAll();",
    "",
    "// Typo  →  PHP geeft géén foutmelding",
    "echo $rows[0]['mesage'];",
    "",
    "// Wat voor type is dit?  int?  string?",
    "echo $rows[0]['timestamp'];",
], rx_h, 1.53, rw_h, 1.6, sz=12)

box(s3b, rx_h, 3.23, rw_h, 0.35, RGBColor(0xEC, 0xFF, 0xEC), bdr=C_GREEN)
tb(s3b, "Met hydration — veilig en duidelijk", rx_h + 0.14, 3.25, rw_h - 0.22, 0.3,
   sz=13, bold=True, c=C_GREEN)
code_box(s3b, [
    "$messages = $this->repository->findAll();",
    "",
    "// IDE-autocomplete  +  volledig type-safe",
    "echo $messages[0]->message;",
    "",
    "// Altijd een DateTimeImmutable  —  nooit raden",
    "echo $messages[0]->timestamp->format('d-m-Y H:i');",
    "",
    "// Methoden op het object zijn mogelijk",
    "if ($message->isOwnedBy($currentUser)) { ... }",
], rx_h, 3.58, rw_h, 2.72, sz=12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Service
# ═══════════════════════════════════════════════════════════════════════════
s4 = new_slide()
bg(s4, C_SLIDE)

box(s4, 0, 0, W, 1.0, C_SVC)
tb(s4, "Service", 0.5, 0.08, 9, 0.55, sz=34, bold=True, c=C_TEXT)
tb(s4, "Businesslogica-laag  —  de regels van de applicatie",
   0.5, 0.6, 9, 0.38, sz=15, c=RGBColor(0x78, 0x35, 0x00))
mini_stack(s4, "Service")

lx, lw = 0.45, 5.95

tb(s4, "Verantwoordelijkheid", lx, 1.18, lw, 0.4, sz=15, bold=True, c=C_SVC_D)
blist(s4, [
    "• Bevat de businesslogica van de applicatie",
    "• Stelt regels vast: \"mag dit? klopt dit logisch?\"",
    "• Coördineert aanroepen naar de repository",
    "• Geeft het resultaat terug aan de controller",
], lx + 0.1, 1.58, lw - 0.15, 1.5, sz=15, c=C_TEXT)

tb(s4, "Wat hoort WEL hier?", lx, 3.18, lw, 0.38, sz=15, bold=True, c=C_GREEN)
blist(s4, [
    "✓  Autorisatiecheck  (\"is dit jouw bericht?\")",
    "✓  password_hash()  —  business rule, niet opslag",
    "✓  Aanroepen van repository-methoden",
], lx + 0.1, 3.56, lw - 0.15, 1.1, sz=15, c=C_GREEN)

tb(s4, "Wat hoort NIET hier?", lx, 4.76, lw, 0.38, sz=15, bold=True, c=C_RED)
blist(s4, [
    "✗  SQL-queries  (dat doet de repository)",
    "✗  HTTP-statuscodes of json_encode()",
    "✗  Directe inputvalidatie van raw user-input",
], lx + 0.1, 5.14, lw - 0.15, 1.1, sz=15, c=C_RED)

code_box(s4, [
    "// MessageService.php",
    "",
    "public function delete(",
    "    int $id,",
    "    User $currentUser",
    "): void {",
    "",
    "    // Repository haalt het bericht op",
    "    $msg = $this->repository->findById($id)",
    "        ?? throw new RuntimeException(",
    "               'Bericht niet gevonden', 404",
    "           );",
    "",
    "    // Business rule: alleen eigenaar mag verwijderen",
    "    if ($msg->userId !== $currentUser->id)",
    "        throw new RuntimeException(",
    "            'Niet jouw bericht!', 403",
    "        );",
    "",
    "    $this->repository->delete($id);",
    "}",
], 6.6, 1.15, 6.5, 6.15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Controller
# ═══════════════════════════════════════════════════════════════════════════
s5 = new_slide()
bg(s5, C_SLIDE)

box(s5, 0, 0, W, 1.0, C_CTRL)
tb(s5, "Controller", 0.5, 0.08, 9, 0.55, sz=34, bold=True, c=C_WHITE)
tb(s5, "HTTP-laag  —  het contactpunt met de buitenwereld",
   0.5, 0.6, 9, 0.38, sz=15, c=RGBColor(0xA7, 0xF3, 0xD0))
mini_stack(s5, "Controller")

lx, lw = 0.45, 5.95

tb(s5, "Verantwoordelijkheid", lx, 1.18, lw, 0.4, sz=15, bold=True, c=C_CTRL_D)
blist(s5, [
    "• Ontvangt het HTTP-verzoek (GET, POST, PUT, DELETE)",
    "• Leest de input (body, URL-parameters, querystring)",
    "• Valideert de invoer (lengte, verplicht veld, format)",
    "• Roept de juiste service-methode aan",
    "• Stuurt een HTTP-response terug (JSON + statuscode)",
], lx + 0.1, 1.58, lw - 0.15, 1.8, sz=15, c=C_TEXT)

tb(s5, "Wat hoort WEL hier?", lx, 3.48, lw, 0.38, sz=15, bold=True, c=C_GREEN)
blist(s5, [
    "✓  Inputvalidatie  (\"is het veld leeg? juiste lengte?\")",
    "✓  HTTP-statuscodes  (200, 201, 404, 422 ...)",
    "✓  json_encode() en headers",
], lx + 0.1, 3.86, lw - 0.15, 1.1, sz=15, c=C_GREEN)

tb(s5, "Wat hoort NIET hier?", lx, 5.06, lw, 0.38, sz=15, bold=True, c=C_RED)
blist(s5, [
    "✗  SQL-queries of database-aanroepen",
    "✗  Business rules  (\"mag dit?\", \"klopt dit logisch?\")",
    "✗  Password hashing of andere domain logic",
], lx + 0.1, 5.44, lw - 0.15, 1.1, sz=15, c=C_RED)

code_box(s5, [
    "// MessageController.php",
    "",
    "public function store(array $body): void",
    "{",
    "    // 1. Valideer de invoer",
    "    $error = $this->validateMessage($body);",
    "    if ($error) {",
    "        $this->json(['error' => $error], 422);",
    "        return;",
    "    }",
    "",
    "    // 2. Service aanroepen",
    "    $now = new DateTimeImmutable();",
    "    $msg = $this->messageService->create(",
    "        $body['message'], $now, $currentUser",
    "    );",
    "",
    "    // 3. HTTP-response sturen",
    "    $this->json($msg, 201);",
    "}",
], 6.6, 1.15, 6.5, 6.15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Data flow
# ═══════════════════════════════════════════════════════════════════════════
s6 = new_slide()
bg(s6, C_SLIDE)

box(s6, 0, 0, W, 0.88, C_BG)
tb(s6, "Samenwerking: POST /messages  (bericht plaatsen)",
   0.5, 0.14, 12, 0.62, sz=26, bold=True, c=C_WHITE)

steps_fwd = [
    (C_GRAY, "1. Client",      "Stuurt POST\n/messages met\nberichttekst als JSON"),
    (C_CTRL, "2. Controller",  "Valideert invoer\nRoept messageService\n.create() aan"),
    (C_SVC,  "3. Service",     "Controleert:\nis user ingelogd?\nMag dit?"),
    (C_REPO, "4. Repository",  "Voert INSERT\nINTO messages\nuit via PDO"),
    (C_GRAY, "5. Database",    "Bewaart de\nrij en geeft\nlastInsertId terug"),
]

bw_f, bh_head, bh_body = 2.1, 0.5, 1.0
gap_f = 0.27
sx_f = (W - len(steps_fwd) * bw_f - (len(steps_fwd)-1) * gap_f) / 2
sy_f = 1.08

for i, (col, lbl, desc) in enumerate(steps_fwd):
    x = sx_f + i * (bw_f + gap_f)
    tc = C_TEXT if col == C_SVC else C_WHITE
    boxtb(s6, lbl, x, sy_f, bw_f, bh_head, col, tc=tc, sz=14, bold=True)
    box(s6, x, sy_f + bh_head, bw_f, bh_body,
        RGBColor(0xF1, 0xF5, 0xF9), bdr=col)
    tb(s6, desc, x + 0.1, sy_f + bh_head + 0.1,
       bw_f - 0.2, bh_body - 0.15, sz=12, c=C_MUTED)
    if i < len(steps_fwd) - 1:
        ax = x + bw_f + gap_f/2 - 0.12
        tb(s6, "→", ax, sy_f + bh_head*0.2, 0.28, 0.38,
           sz=18, c=C_MUTED, al=PP_ALIGN.CENTER)

# Return path
tb(s6, "Response terugsturen (omgekeerd):",
   0.4, 2.78, 8, 0.38, sz=13, bold=True, c=C_MUTED)

steps_back = [
    (C_GRAY, "Database",    "Geeft rij terug"),
    (C_REPO, "Repository",  "hydrate() →\nMessage object"),
    (C_SVC,  "Service",     "Geeft Message\nobject terug"),
    (C_CTRL, "Controller",  "json_encode() +\nstatus 201"),
    (C_GRAY, "Client",      "Ontvangt JSON\nmet nieuw bericht"),
]

bh_r = 0.88
sy_r = 3.18

for i, (col, lbl, desc) in enumerate(steps_back):
    x = sx_f + i * (bw_f + gap_f)
    tc = C_TEXT if col == C_SVC else C_WHITE
    boxtb(s6, lbl, x, sy_r, bw_f, 0.46, col, tc=tc, sz=13, bold=True)
    box(s6, x, sy_r + 0.46, bw_f, 0.68,
        RGBColor(0xF1, 0xF5, 0xF9), bdr=col)
    tb(s6, desc, x + 0.1, sy_r + 0.52, bw_f - 0.2, 0.56, sz=11, c=C_MUTED)
    if i < len(steps_back) - 1:
        ax = x + bw_f + gap_f/2 - 0.12
        tb(s6, "→", ax, sy_r + 0.08, 0.28, 0.32,
           sz=18, c=C_MUTED, al=PP_ALIGN.CENTER)

# Bottom callout
box(s6, 0.4, 4.62, W - 0.8, 0.7, RGBColor(0xF0, 0xF9, 0xFF),
    bdr=RGBColor(0xBA, 0xE6, 0xFD))
tb(s6, "Elke laag kent alleen zijn buren.  "
       "De Controller kent geen SQL.  "
       "De Repository weet niets van HTTP.",
   0.65, 4.72, W - 1.3, 0.5, sz=14, bold=True,
   c=RGBColor(0x07, 0x53, 0x85), al=PP_ALIGN.CENTER)

# Dependency arrow
tb(s6, "Controller  →  roept aan  →  Service  →  roept aan  →  Repository",
   0.4, 5.52, W - 0.8, 0.4, sz=13, c=C_MUTED, al=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Samenvatting (table)
# ═══════════════════════════════════════════════════════════════════════════
s7 = new_slide()
bg(s7, C_SLIDE)

box(s7, 0, 0, W, 0.88, C_BG)
tb(s7, "Samenvatting", 0.5, 0.15, 10, 0.62, sz=30, bold=True, c=C_WHITE)

# Column positions
c1x, c1w = 0.35, 2.0
c2x, c2w = 2.45, 3.9
c3x, c3w = 6.45, 3.3
c4x, c4w = 9.85, 3.2

# Header row
hdr_y = 1.06
for cx2, cw2, lbl in [
    (c1x, c1w, "Laag"),
    (c2x, c2w, "Verantwoordelijkheid"),
    (c3x, c3w, "Weet van..."),
    (c4x, c4w, "Weet NIET van..."),
]:
    boxtb(s7, lbl, cx2, hdr_y, cw2, 0.48, C_BG, sz=13, bold=True)

table_rows = [
    (C_CTRL, C_CTRL_L, "Controller",
     "HTTP-verzoek ontvangen, invoer valideren, JSON-response sturen",
     "GET / POST / PUT / DELETE, status codes, json_encode()",
     "SQL, business rules, password hashing"),
    (C_SVC, C_SVC_L, "Service",
     "Businesslogica en domain rules uitvoeren, aanroepen van repository",
     "PHP-objecten, autorisatieregels, domain logic",
     "SQL, HTTP, status codes, json_encode()"),
    (C_REPO, C_REPO_L, "Repository",
     "SQL-queries uitvoeren, rijen ophalen en omzetten naar objecten",
     "Database, PDO, SQL, hydration",
     "HTTP, business rules, inputvalidatie"),
]

rh = 1.55
for i, (col, light, name, resp, knows, not_knows) in enumerate(table_rows):
    ry = hdr_y + 0.52 + i * (rh + 0.07)

    # Layer name cell
    box(s7, c1x, ry, c1w, rh, light)
    box(s7, c1x, ry, 0.18, rh, col)
    tb(s7, name, c1x + 0.28, ry + 0.55, c1w - 0.35, 0.5,
       sz=16, bold=True, c=col, al=PP_ALIGN.CENTER)

    # Responsibility
    box(s7, c2x, ry, c2w, rh, RGBColor(0xF8, 0xFA, 0xFC),
        bdr=RGBColor(0xE2, 0xE8, 0xF0))
    tb(s7, resp, c2x + 0.12, ry + 0.18, c2w - 0.22, rh - 0.28,
       sz=13, c=C_TEXT, wrap=True)

    # Knows
    box(s7, c3x, ry, c3w, rh,
        RGBColor(0xF0, 0xFD, 0xF4), bdr=RGBColor(0xBB, 0xF7, 0xD0))
    tb(s7, "✓  " + knows, c3x + 0.12, ry + 0.18, c3w - 0.22, rh - 0.28,
       sz=13, c=C_GREEN, wrap=True)

    # Does NOT know
    box(s7, c4x, ry, c4w, rh,
        RGBColor(0xFF, 0xF1, 0xF2), bdr=RGBColor(0xFE, 0xCA, 0xCA))
    tb(s7, "✗  " + not_knows, c4x + 0.12, ry + 0.18, c4w - 0.22, rh - 0.28,
       sz=13, c=C_RED, wrap=True)

# Footer
box(s7, 0.35, 6.82, W - 0.7, 0.52, RGBColor(0xF0, 0xF9, 0xFF),
    bdr=RGBColor(0xBA, 0xE6, 0xFD))
tb(s7, "Separation of Concerns  —  één verantwoordelijkheid per laag",
   0.5, 6.88, W - 1.0, 0.4, sz=15, bold=True,
   c=RGBColor(0x07, 0x53, 0x85), al=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Cheatsheet (printable)
# ═══════════════════════════════════════════════════════════════════════════
s8 = new_slide()
bg(s8, C_WHITE)

# Title bar
box(s8, 0, 0, W, 0.62, C_BG)
tb(s8, "Cheatsheet  —  Controller · Service · Repository",
   0.35, 0.08, 12.5, 0.38, sz=20, bold=True, c=C_WHITE)
tb(s8, "Eén verantwoordelijkheid per laag  ·  Separation of Concerns",
   0.35, 0.44, 12.5, 0.18, sz=10, c=RGBColor(0x94, 0xA3, 0xB8))

col_w_cs = 4.1
col_gap_cs = 0.22
col_xs_cs = [
    0.3,
    0.3 + col_w_cs + col_gap_cs,
    0.3 + 2 * (col_w_cs + col_gap_cs),
]

cheat_data = [
    (C_CTRL, C_CTRL_L, C_CTRL_D, "Controller",
     "HTTP-verzoeken ontvangen en responses sturen",
     [
         "→  Invoer valideren (lengte, verplicht veld...)",
         "→  HTTP status codes (200, 201, 404, 422...)",
         "→  json_encode() + Content-Type header",
         "→  Service-methode aanroepen",
         "→  Request body en URL-params inlezen",
     ],
     [
         "→  SQL of database-aanroepen",
         "→  Business rules  (\"mag dit?\")",
         "→  Password hashing of domain logic",
         "→  Aanroepen van de repository direct",
     ]),
    (C_SVC, C_SVC_L, C_SVC_D, "Service",
     "Businesslogica en domain rules uitvoeren",
     [
         "→  Autorisatiecheck  (\"is dit jouw item?\")",
         "→  password_hash() / password_verify()",
         "→  Domain rules en business logic",
         "→  Repository-methoden aanroepen",
     ],
     [
         "→  SQL of directe database-toegang",
         "→  HTTP status codes of headers",
         "→  json_encode()",
         "→  Raw userinput valideren",
     ]),
    (C_REPO, C_REPO_L, C_REPO_D, "Repository",
     "SQL uitvoeren en rijen omzetten naar objecten",
     [
         "→  SELECT / INSERT / UPDATE / DELETE",
         "→  hydrate(array $row): DomainObject",
         "→  findAll() · findById() · save() · delete()",
         "→  PDO, prepare(), execute(), lastInsertId()",
     ],
     [
         "→  Business rules of autorisatie",
         "→  HTTP status codes of JSON",
         "→  Inputvalidatie van userinput",
         "→  Aanroepen van services of controllers",
     ]),
]

for i, (col, light, dark, name, taak, goed, fout) in enumerate(cheat_data):
    cx = col_xs_cs[i]
    tc_hdr = C_TEXT if col == C_SVC else C_WHITE

    # Column header
    boxtb(s8, name, cx, 0.68, col_w_cs, 0.6, col, tc=tc_hdr, sz=18, bold=True)

    # Taak section
    box(s8, cx, 1.3, col_w_cs, 0.78, light)
    tb(s8, "Taak:", cx + 0.12, 1.32, col_w_cs - 0.18, 0.26, sz=11, bold=True, c=dark)
    tb(s8, taak,   cx + 0.12, 1.55, col_w_cs - 0.18, 0.5,  sz=12, c=C_TEXT, wrap=True)

    # Hoort WEL hier
    box(s8, cx, 2.1, col_w_cs, 0.3, RGBColor(0xF0, 0xFD, 0xF4), bdr=C_GREEN)
    tb(s8, "✓  Hoort WEL hier:", cx + 0.1, 2.12, col_w_cs - 0.15, 0.26,
       sz=11, bold=True, c=C_GREEN)
    blist(s8, goed, cx + 0.14, 2.42, col_w_cs - 0.22, 1.62, sz=12, c=C_TEXT, sp=3)

    # Hoort NIET hier
    box(s8, cx, 4.07, col_w_cs, 0.3, RGBColor(0xFF, 0xF1, 0xF2), bdr=C_RED)
    tb(s8, "✗  Hoort NIET hier:", cx + 0.1, 4.09, col_w_cs - 0.15, 0.26,
       sz=11, bold=True, c=C_RED)
    blist(s8, fout, cx + 0.14, 4.39, col_w_cs - 0.22, 1.28, sz=12, c=C_RED, sp=3)

# Separator
box(s8, 0.3, 5.73, W - 0.6, 0.02, RGBColor(0xCB, 0xD5, 0xE0))

# Data flow
tb(s8, "Data flow:", 0.35, 5.81, 1.5, 0.36, sz=12, bold=True, c=C_TEXT)

flow_items = [
    ("Client",      C_GRAY, 1.1),
    ("→",           None,   0.28),
    ("Controller",  C_CTRL, 1.45),
    ("→",           None,   0.28),
    ("Service",     C_SVC,  1.1),
    ("→",           None,   0.28),
    ("Repository",  C_REPO, 1.45),
    ("→",           None,   0.28),
    ("Database",    C_GRAY, 1.1),
]
# center the flow row
flow_total = sum(fw for _, _, fw in flow_items) + (len(flow_items) - 1) * 0.05
fx_cs = (W - flow_total) / 2
for lbl_f, col_f, fw_f in flow_items:
    if col_f is None:
        tb(s8, lbl_f, fx_cs, 5.84, fw_f, 0.3, sz=14, c=C_MUTED, al=PP_ALIGN.CENTER)
    else:
        tc_f = C_TEXT if col_f == C_SVC else C_WHITE
        boxtb(s8, lbl_f, fx_cs, 5.81, fw_f, 0.36, col_f, tc=tc_f, sz=12, bold=True)
    fx_cs += fw_f + 0.05

# Footer callout
box(s8, 0.3, 6.28, W - 0.6, 0.85, RGBColor(0xF8, 0xFA, 0xFC),
    bdr=RGBColor(0xCB, 0xD5, 0xE0))
tb(s8, "Hydration:", 0.48, 6.36, 1.3, 0.28, sz=12, bold=True, c=C_REPO_D)
tb(s8, "de Repository zet elke databaserij om naar een Domain Object (Message, User, ...).",
   1.75, 6.36, W - 2.1, 0.3, sz=12, c=C_TEXT, wrap=True)
tb(s8, "Gouden regel:", 0.48, 6.66, 1.6, 0.28, sz=12, bold=True, c=C_REPO_D)
tb(s8, "alleen de Repository weet hoe de databasetabel eruitziet — "
       "service en controller werken uitsluitend met objecten.",
   1.75, 6.66, W - 2.1, 0.3, sz=12, c=C_TEXT, wrap=True)


# ── Save ──────────────────────────────────────────────────────────────────
output = (r"C:\Users\Peter Marcelis\Documents\Lesmateriaal"
          r"\tweedejaars_periode_twee\CSR_Pattern_Presentatie_v2.pptx")
prs.save(output)
print("Saved:", output)
