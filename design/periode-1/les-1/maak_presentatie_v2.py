"""
Les 1: Wat is Design?
Design regels:
- Geen rounded cards met flat headers erop (hoekconflict)
- Rounded alleen voor kleine elementen: badges, bullets
- Alles wat in een container staat, past er ook in
- Geen em-dashes in zichtbare tekst
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image, ImageDraw
import numpy as np, io

# Kleuren
NAVY   = (6,  14,  44)
NAVY2  = (30, 58, 138)
BLUE   = (37, 99, 235)
BLUE2  = (59,130, 246)
CYAN   = (6, 182, 212)
WHITE  = (255,255,255)
OFFWH  = (248,250,252)
DARK   = (15, 23,  42)
BODY   = (51, 65,  85)
MUTED  = (100,116,135)
GREEN  = (22,163, 74)
RED    = (220, 38, 38)

def rgb(c): return RGBColor(c[0],c[1],c[2])

SW, SH = Inches(13.33), Inches(7.5)
BW, BH = 1920, 1080

# Achtergronden
def grad_arr(w,h,c1,c2,d='diag'):
    x,y = np.linspace(0,1,w,dtype=np.float32), np.linspace(0,1,h,dtype=np.float32)
    if d=='h': t=np.tile(x,(h,1))
    elif d=='v': t=np.tile(y.reshape(-1,1),(1,w))
    else:
        xx,yy=np.meshgrid(x,y); t=xx*.55+yy*.45
    a=np.zeros((h,w,3),dtype=np.float32)
    for i in range(3): a[:,:,i]=c1[i]+(c2[i]-c1[i])*t
    return np.clip(a,0,255).astype(np.uint8)

def dark_bg(circles=True):
    img=Image.fromarray(grad_arr(BW,BH,NAVY,NAVY2)).convert('RGBA')
    if circles:
        ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
        d.ellipse([BW-500,-180,BW+200,520],fill=(59,130,246,20))
        d.ellipse([-60,BH-360,380,BH+100],fill=(6,182,212,25))
        img=Image.alpha_composite(img,ov)
    return img.convert('RGB')

def blue_bg():
    img=Image.fromarray(grad_arr(BW,BH,(29,78,216),(37,99,235))).convert('RGBA')
    ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
    d.ellipse([BW-380,-120,BW+120,380],fill=(255,255,255,10))
    img=Image.alpha_composite(img,ov)
    return img.convert('RGB')

def light_bg():
    img=Image.fromarray(grad_arr(BW,BH,OFFWH,WHITE)).convert('RGBA')
    ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
    d.ellipse([BW-580,-220,BW+180,540],fill=(59,130,246,8))
    img=Image.alpha_composite(img,ov)
    return img.convert('RGB')

def to_buf(img):
    b=io.BytesIO(); img.save(b,format='PNG'); b.seek(0); return b

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
BLANK=prs.slide_layouts[6]

def new_slide(bg=None):
    s=prs.slides.add_slide(BLANK)
    if bg:
        p=s.shapes.add_picture(to_buf(bg),0,0,SW,SH)
        t=s.shapes._spTree; t.remove(p._element); t.insert(2,p._element)
    else:
        s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(WHITE)
    return s

# Basis shapes
def box(s,l,t,w,h,c,border=None):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c)
    if border: sh.line.color.rgb=rgb(border); sh.line.width=Pt(0.75)
    else: sh.line.fill.background()
    return sh

def gbox(s,l,t,w,h,c1,c2,angle=90):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.gradient(); stops=sh.fill.gradient_stops
    stops[0].position=0.0; stops[0].color.rgb=rgb(c1)
    stops[1].position=1.0; stops[1].color.rgb=rgb(c2)
    sh.fill.gradient_angle=angle; sh.line.fill.background(); return sh

def dot(s,cx,cy,r,c):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,cx-r,cy-r,r*2,r*2)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c); sh.line.fill.background(); return sh

def label(s,text,l,t,w,h,size=18,c=BODY,bold=False,italic=False,align=PP_ALIGN.LEFT):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=text; r.font.name="Calibri"; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=rgb(c)
    return bx

def blist(s,items,l,t,w,h,size=18,c=BODY,gap=7):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(gap); r=p.add_run()
        r.text=f"  {chr(8226)}  {item}"; r.font.name="Calibri"; r.font.size=Pt(size); r.font.color.rgb=rgb(c)
    return bx

def notitie(s,t): s.notes_slide.notes_text_frame.text=t

def paginanr(s,n,tot=18,light=False):
    label(s,f"{n} / {tot}",Inches(12.5),Inches(7.1),Inches(0.8),Inches(0.35),
          size=11,c=(180,180,200) if light else MUTED,align=PP_ALIGN.RIGHT)

# Pagina-opmaak hulpfuncties
def hero_header(s,title,sub=None):
    label(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(1.0),size=36,c=WHITE,bold=True)
    gbox(s,Inches(1.0),Inches(1.42),Inches(2.0),Inches(0.08),BLUE,CYAN,angle=0)
    if sub: label(s,sub,Inches(1.0),Inches(1.55),Inches(11.3),Inches(0.42),size=15,c=(160,185,230),italic=True)

def page_header(s,title,sub=None):
    label(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(0.95),size=34,c=DARK,bold=True)
    gbox(s,Inches(1.0),Inches(1.38),Inches(1.8),Inches(0.08),BLUE,CYAN,angle=0)
    if sub: label(s,sub,Inches(1.0),Inches(1.5),Inches(11.3),Inches(0.38),size=14,c=MUTED,italic=True)

# Kaart: plat vlak + smalle gekleurde linkerstrip. Geen rounded corners.
def card(s,l,t,w,h,accent=BLUE,bg=OFFWH,border=None):
    STRIP=Inches(0.07)
    box(s,l,t,w,h,bg,border=border if border else (215,225,240))
    gbox(s,l,t,STRIP,h,accent,CYAN,angle=270)
    # Geeft terug: startpunt voor content
    return l+STRIP+Inches(0.18), t+Inches(0.14), w-STRIP-Inches(0.35)

# Kolomkop (flat, geen rounded probleem)
def col_kop(s,l,t,w,h,text,c):
    box(s,l,t,w,h,c)
    label(s,text,l,t+Inches(0.07),w,h-Inches(0.07),size=20,c=WHITE,bold=True,align=PP_ALIGN.CENTER)


# ── SLIDE 1: Titel ────────────────────────────────────────
s=new_slide(dark_bg(True)); paginanr(s,1,light=True)
gbox(s,0,0,Inches(0.5),SH,BLUE,CYAN,angle=270)
label(s,"Wat is Design?",Inches(1.0),Inches(1.8),Inches(9.5),Inches(1.7),size=60,c=WHITE,bold=True)
gbox(s,Inches(1.0),Inches(3.5),Inches(2.8),Inches(0.1),BLUE,CYAN,angle=0)
label(s,"Periode 1  ·  Les 1",Inches(1.0),Inches(3.7),Inches(6.0),Inches(0.5),size=18,c=(160,185,230))
label(s,"Design voor het web",Inches(1.0),Inches(4.3),Inches(6.0),Inches(0.4),size=14,c=(100,130,180),italic=True)
notitie(s,"Welkom bij les 1. Vandaag starten we breed: wat is design? Daarna maken we kennis met Figma.")


# ── SLIDE 2: Agenda ───────────────────────────────────────
s=new_slide(light_bg()); paginanr(s,2)
page_header(s,"Vandaag")
items=["Wat is design en wat niet","Hoe denkt een gebruiker?","Het kanaal bepaalt de boodschap","Kennismaking met Figma","Opdracht van vandaag"]
bkl=[BLUE,(79,130,246),(99,150,246),(6,182,212),(16,163,110)]
for i,(item,bc) in enumerate(zip(items,bkl)):
    t=Inches(1.9)+i*Inches(0.97)
    dot(s,Inches(1.35),t+Inches(0.22),Inches(0.22),bc)
    label(s,str(i+1),Inches(1.13),t+Inches(0.06),Inches(0.44),Inches(0.38),size=14,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    label(s,item,Inches(1.85),t+Inches(0.03),Inches(9.0),Inches(0.5),size=21,c=DARK)
notitie(s,"Loop de agenda door. We beginnen met theorie en eindigen concreet in Figma.")


# ── SLIDE 3: Design != Mooi ───────────────────────────────
s=new_slide(); paginanr(s,3)
# Linker donker paneel
gbox(s,0,0,Inches(5.0),SH,NAVY,(20,40,100),angle=270)
gbox(s,0,0,Inches(0.5),SH,BLUE,CYAN,angle=270)
label(s,"Design",Inches(0.75),Inches(1.5),Inches(4.0),Inches(0.8),size=42,c=WHITE,bold=True)
label(s,"≠",Inches(0.75),Inches(2.28),Inches(4.0),Inches(0.75),size=54,c=CYAN,bold=True)
label(s,"Iets mooi\nmaken",Inches(0.75),Inches(2.95),Inches(4.0),Inches(1.4),size=32,c=WHITE,bold=True)
# Rechter inhoud
blist(s,["Design lost een probleem op voor een gebruiker","Mooi zijn is een bijproduct, niet het doel","Slecht design kan er prachtig uitzien","Goed design kan sober zijn en perfect werken"],
      Inches(5.4),Inches(1.6),Inches(7.4),Inches(2.8),size=18,c=DARK)
# Quote: platte kaart met linkerstrip
cl,ct,cw=card(s,Inches(5.4),Inches(4.85),Inches(7.5),Inches(2.1),accent=NAVY)
label(s,'"Design is not just what it looks like. Design is how it works."',
      cl,ct,cw,Inches(1.3),size=17,c=DARK,italic=True)
label(s,"Steve Jobs",cl,ct+Inches(1.35),cw,Inches(0.4),size=13,c=MUTED)
notitie(s,"Gebruik twee concrete voorbeelden: een mooie maar onbruikbare site, en Google homepage 1998.")


# ── SLIDE 4: Design vs Kunst ──────────────────────────────
s=new_slide(light_bg()); paginanr(s,4)
page_header(s,"Design vs. Kunst")
KW=Inches(5.5); GAP=Inches(0.45); KH=Inches(5.05)
KTOP=Inches(1.75); HDR=Inches(0.55)
data=[
    ("Design",NAVY,NAVY2,["In dienst van een doel","Altijd een gebruiker centraal","Meetbaar: werkt het?","Beperkingen zijn het spel"]),
    ("Kunst",(80,30,120),(140,50,180),["Vrije expressie","Publiek is optioneel","Subjectief","Beperkingen zijn een hindernis"]),
]
for i,(title_,c1,c2,items) in enumerate(data):
    L=Inches(1.0)+i*(KW+GAP)
    col_kop(s,L,KTOP,KW,HDR,title_,c1)
    box(s,L,KTOP+HDR,KW,KH-HDR,OFFWH,border=(215,225,240))
    for j,item in enumerate(items):
        bg=OFFWH if j%2==0 else WHITE
        box(s,L,KTOP+HDR+j*Inches(1.1),KW,Inches(1.05),bg,border=(220,228,245))
        label(s,item,L+Inches(0.2),KTOP+HDR+j*Inches(1.1)+Inches(0.12),KW-Inches(0.4),Inches(0.82),size=17,c=BODY)
label(s,"Zelfde vaardigheden, andere intentie",Inches(1.0),Inches(6.88),Inches(11.3),Inches(0.38),size=14,c=MUTED,italic=True,align=PP_ALIGN.CENTER)
notitie(s,"De vaardigheden overlappen: kleur, compositie, typografie. Maar design werkt altijd in dienst van iemand anders.")


# ── SLIDE 5: Centrale vraag ───────────────────────────────
s=new_slide(blue_bg()); paginanr(s,5,light=True)
gbox(s,0,0,Inches(0.5),SH,BLUE2,CYAN,angle=270)
label(s,"De vraag die je altijd moet stellen:",Inches(1.0),Inches(1.1),Inches(11.3),Inches(0.5),size=17,c=(180,210,255),italic=True)
label(s,'"Kan mijn gebruiker zonder na te denken\ndoen wat ze wil doen?"',
      Inches(1.2),Inches(1.9),Inches(10.9),Inches(2.4),size=38,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
gbox(s,Inches(4.5),Inches(4.4),Inches(4.3),Inches(0.09),BLUE2,CYAN,angle=0)
label(s,"Steve Krug,  Don't Make Me Think  (2000)",Inches(1.0),Inches(4.6),Inches(11.3),Inches(0.5),size=16,c=(180,210,255),italic=True,align=PP_ALIGN.CENTER)
notitie(s,"Laat deze quote even staan. Vraag studenten: wanneer heb jij voor het laatst gefrustreerd een website verlaten?")


# ── SLIDE 6: DMMT ────────────────────────────────────────
s=new_slide(light_bg()); paginanr(s,6)
page_header(s,"Don't Make Me Think")
box(s,Inches(1.0),Inches(1.72),Inches(4.8),Inches(0.5),NAVY)
label(s,"Steve Krug  ·  2000, nog steeds actueel",Inches(1.18),Inches(1.78),Inches(4.42),Inches(0.42),size=14,c=WHITE,bold=True)
blist(s,["Gebruikers scannen pagina's, ze lezen niet","Ze klikken op het eerste dat redelijk lijkt","Elk vraagteken kost mentale energie","Jouw design moet zo vanzelfsprekend zijn dat niemand nadenkt","De gebruiker is druk, gehaast en minder zorgvuldig dan jij hoopt"],
      Inches(1.0),Inches(2.35),Inches(7.5),Inches(4.0),size=19,c=DARK)
# Rechter donkere paneel: plat vlak, content ruim inset
box(s,Inches(9.2),Inches(1.65),Inches(3.7),Inches(5.4),NAVY)
gbox(s,Inches(9.2),Inches(1.65),Inches(3.7),Inches(0.08),BLUE,CYAN,angle=0)
label(s,"Vragen van je gebruiker",Inches(9.4),Inches(1.85),Inches(3.3),Inches(0.4),size=13,c=(160,185,230),italic=True)
vragen=["Waar moet ik klikken?","Wat doet dit?","Ben ik op de goede pagina?","Hoe ga ik terug?"]
for i,v in enumerate(vragen):
    box(s,Inches(9.4),Inches(2.45)+i*Inches(1.0),Inches(3.3),Inches(0.82),(30,50,100))
    label(s,f"?  {v}",Inches(9.58),Inches(2.55)+i*Inches(1.0),Inches(3.05),Inches(0.6),size=13,c=(200,220,255))
label(s,"Elk vraagteken is een gefrustreerde gebruiker",Inches(9.4),Inches(6.6),Inches(3.3),Inches(0.38),size=11,c=(100,130,180),italic=True,align=PP_ALIGN.CENTER)
notitie(s,"Dit boek is van 2000 maar de principes zijn tijdloos. Niets aan ons scangedrag is veranderd, eerder erger geworden door social media.")


# ── SLIDE 7: F-patroon ───────────────────────────────────
s=new_slide(light_bg()); paginanr(s,7)
page_header(s,"Zo lezen gebruikers een pagina")
blist(s,["Gebruikers volgen een F-patroon (Nielsen Norman Group, 2006)","Bovenaan links krijgt de meeste aandacht","Eerste woorden van elke zin zijn het belangrijkst","Koppen worden gelezen, lopende tekst nauwelijks","Zet kritieke informatie bovenaan en links"],
      Inches(1.0),Inches(1.75),Inches(7.3),Inches(4.5),size=18,c=DARK)
label(s,"nngroup.com  ·  f-shaped-pattern-reading-web-content",Inches(1.0),Inches(6.5),Inches(7.5),Inches(0.42),size=11,c=MUTED)
# Heatmap: donker kader
CARD_L=Inches(8.7); CARD_T=Inches(1.55); CARD_W=Inches(4.25); CARD_H=Inches(5.3)
box(s,CARD_L,CARD_T,CARD_W,CARD_H,DARK)
gbox(s,CARD_L,CARD_T,CARD_W,Inches(0.08),BLUE,CYAN,angle=0)
# Max breedte voor balken = CARD_W - 2*padding
MAX_BAR=CARD_W-Inches(0.5)  # = 3.75"
BAR_L=CARD_L+Inches(0.25)
bars=[(MAX_BAR,(220,38,38)),(Inches(2.5),(234,88,12)),(Inches(1.3),(217,119,6)),(Inches(1.3),(180,160,0)),(Inches(1.3),(180,160,0))]
for i,(bw,bc) in enumerate(bars):
    BT=CARD_T+Inches(0.35)+i*Inches(0.92)
    box(s,BAR_L,BT,bw,Inches(0.6),bc)
label(s,"Simulatie heatmap  ·  rood = meest bekeken",CARD_L,CARD_T+CARD_H+Inches(0.08),CARD_W,Inches(0.35),size=10,c=MUTED,align=PP_ALIGN.CENTER)
notitie(s,"Open nngroup.com in de browser en laat de echte heatmaps zien. Vraag: hoe schrijf je koppen zodat ze ook los van de rest betekenis hebben?")


# ── SLIDE 8: McLuhan ──────────────────────────────────────
s=new_slide(dark_bg(True)); paginanr(s,8,light=True)
gbox(s,0,0,Inches(0.5),SH,BLUE,CYAN,angle=270)
hero_header(s,"The Medium is the Message")
blist(s,["Marshall McLuhan, filosoof en mediatheoreticus (1964)","Het kanaal dat je kiest communiceert al iets, voor je eerste woord","Zelfde bericht via SMS, brief of telefoongesprek: compleet anders","Desktop en mobiel zijn twee verschillende media"],
      Inches(1.0),Inches(2.1),Inches(11.3),Inches(2.7),size=18,c=(200,220,255))
# Quote: platte kaart op donkere achtergrond
box(s,Inches(1.0),Inches(5.0),Inches(11.3),Inches(1.95),(20,35,80))
gbox(s,Inches(1.0),Inches(5.0),Inches(0.08),Inches(1.95),BLUE,CYAN,angle=270)
label(s,'"The medium is the message."',Inches(1.25),Inches(5.2),Inches(10.7),Inches(0.75),size=26,c=WHITE,bold=True,italic=True,align=PP_ALIGN.CENTER)
label(s,"Marshall McLuhan,  Understanding Media  (1964)",Inches(1.25),Inches(5.95),Inches(10.7),Inches(0.5),size=14,c=(160,185,230),italic=True,align=PP_ALIGN.CENTER)
notitie(s,"Geef het voorbeeld: een sollicitatiebrief per WhatsApp vs. per e-mail vs. op papier. Identieke inhoud, compleet andere boodschap.")


# ── SLIDE 9: Responsive ───────────────────────────────────
s=new_slide(light_bg()); paginanr(s,9)
page_header(s,"Hetzelfde design, andere ervaring")
KW=Inches(3.75); GAP=Inches(0.24); KH=Inches(5.3); KTOP=Inches(1.75); HDR=Inches(0.55)
data=[
    ("Desktop",NAVY,["Muis en toetsenbord","Grote viewport","Hover states werken","Gefocuste sessie","Meer ruimte voor inhoud"]),
    ("Tablet",(30,100,180),["Aanraking","Middelgrote viewport","Beperkte hover","Half zittend","Flexibele sessie"]),
    ("Mobiel",BLUE,["Vinger, thumb zone","Kleine viewport","Geen hover states","Onderweg, gehaast","Versnipperde aandacht"]),
]
for i,(title_,c1,items) in enumerate(data):
    L=Inches(1.0)+i*(KW+GAP)
    col_kop(s,L,KTOP,KW,HDR,title_,c1)
    box(s,L,KTOP+HDR,KW,KH-HDR,OFFWH,border=(215,225,240))
    for j,item in enumerate(items):
        bg=OFFWH if j%2==0 else WHITE
        box(s,L,KTOP+HDR+j*Inches(0.92),KW,Inches(0.88),bg,border=(220,228,245))
        label(s,item,L+Inches(0.15),KTOP+HDR+j*Inches(0.92)+Inches(0.12),KW-Inches(0.3),Inches(0.65),size=15,c=BODY)
label(s,"Responsive design is niet alleen technisch, het is een designvraag",Inches(1.0),Inches(7.12),Inches(11.3),Inches(0.35),size=14,c=MUTED,italic=True,align=PP_ALIGN.CENTER)
notitie(s,"Vraag studenten: gebruik jij dezelfde apps anders op je telefoon dan op een laptop? Wat doe je anders?")


# ── SLIDE 10: Onzichtbaar ─────────────────────────────────
s=new_slide(light_bg()); paginanr(s,10)
page_header(s,"Goed design valt niet op")
blist(s,["Als iets goed werkt, merk je het niet","Als iets slecht werkt, merk je het meteen","Don Norman (1988): een deur met trekhandvat die je moet duwen","Het design geeft de verkeerde aanwijzing en dat leidt tot frustratie"],
      Inches(1.0),Inches(1.75),Inches(11.3),Inches(2.8),size=19,c=DARK)
# Uitleg kaart (plat, linkerstrip)
cl,ct,cw=card(s,Inches(1.0),Inches(4.7),Inches(11.3),Inches(2.42),accent=BLUE,bg=OFFWH)
label(s,"Norman Door",cl,ct,cw,Inches(0.45),size=16,c=DARK,bold=True)
label(s,"Een ontwerp dat de verkeerde aanwijzing geeft",cl,ct+Inches(0.48),cw,Inches(0.38),size=14,c=MUTED)
voorbeelden=[("Knop die er niet uitziet als knop",RED),("Link zonder kleur of onderstreping",(180,100,0)),("Formulierveld zonder label",BLUE),("Navigatie die niemand begrijpt",(120,40,160))]
for j,(ex,ec) in enumerate(voorbeelden):
    EL=cl+j*Inches(2.77); ET=ct+Inches(1.0)
    box(s,EL,ET,Inches(2.6),Inches(0.88),ec+(12,) if len(ec)==3 else ec,border=ec)
    # Simpele platte box met gekleurde rand
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,EL,ET,Inches(2.6),Inches(0.88))
    sh.fill.solid(); fc=sh.fill.fore_color
    fc.rgb=rgb(ec); sh.line.fill.background()
    # Pas transparantie toe via XML
    spPr=sh._element.find(qn('p:spPr'))
    sf=spPr.find('.//' + qn('a:solidFill'))
    sc=sf.find(qn('a:srgbClr'))
    al=etree.SubElement(sc,qn('a:alpha')); al.set('val','15000')
    label(s,ex,EL+Inches(0.1),ET+Inches(0.08),Inches(2.4),Inches(0.72),size=12,c=DARK)
notitie(s,"Zoek 'Norman door Vox' op YouTube, 4 minuten. Vraag daarna: welk webdesign-equivalent heb jij deze week meegemaakt?")


# ── SLIDE 11: Live analyse ────────────────────────────────
s=new_slide(); paginanr(s,11)
gbox(s,0,0,Inches(6.55),SH,(15,45,15),(20,65,20),angle=270)
gbox(s,Inches(6.78),0,Inches(6.55),SH,(60,15,15),(85,20,20),angle=270)
box(s,Inches(6.55),0,Inches(0.23),SH,WHITE)
label(s,"Goed design",Inches(0.6),Inches(0.5),Inches(5.5),Inches(0.65),size=26,c=WHITE,bold=True)
box(s,Inches(0.6),Inches(1.15),Inches(2.2),Inches(0.08),(90,210,90))
label(s,"Open live: stripe.com",Inches(0.6),Inches(1.4),Inches(5.5),Inches(0.55),size=20,c=(180,240,180),bold=True)
for i,q in enumerate(["Wat zie je als eerste?","Weet je direct wat je hier kunt doen?","Hoe snel vind je wat je zoekt?","Hoeveel moeite kost navigeren?"]):
    box(s,Inches(0.6),Inches(2.15)+i*Inches(1.0),Inches(5.5),Inches(0.82),(25,70,25))
    label(s,q,Inches(0.78),Inches(2.25)+i*Inches(1.0),Inches(5.15),Inches(0.62),size=16,c=WHITE)
label(s,"Slecht design",Inches(7.15),Inches(0.5),Inches(5.5),Inches(0.65),size=26,c=WHITE,bold=True)
box(s,Inches(7.15),Inches(1.15),Inches(2.2),Inches(0.08),(235,90,90))
label(s,"Open live: een verouderde bedrijfssite",Inches(7.15),Inches(1.4),Inches(5.5),Inches(0.55),size=20,c=(255,180,180),bold=True)
for i,q in enumerate(["Waar is de navigatie?","Wat wil deze site dat je doet?","Hoe snel raak je gefrustreerd?","Vind je wat je zoekt?"]):
    box(s,Inches(7.15),Inches(2.15)+i*Inches(1.0),Inches(5.5),Inches(0.82),(75,20,20))
    label(s,q,Inches(7.33),Inches(2.25)+i*Inches(1.0),Inches(5.15),Inches(0.62),size=16,c=WHITE)
label(s,"Kijk niet of het mooi is, kijk of het werkt",0,Inches(6.82),SW,Inches(0.42),size=16,c=(200,200,200),italic=True,align=PP_ALIGN.CENTER)
notitie(s,"Interactief moment, geen slides nodig. Open de browser. Stel vragen aan de klas. Houd het bij 5 minuten per site.")


# ── SLIDE 12: Figma intro ─────────────────────────────────
s=new_slide(dark_bg(True)); paginanr(s,12,light=True)
gbox(s,0,0,Inches(0.5),SH,BLUE,CYAN,angle=270)
hero_header(s,"Kennismaking met Figma")
feats=[("Browser-based","Geen installatie vereist, ook desktop app beschikbaar"),("Real-time samenwerken","Meerdere mensen tegelijk in hetzelfde bestand"),("Gratis voor studenten","Education plan via figma.com/education"),("Industriestandaard","Gebruikt door Spotify, Airbnb, Uber en Google")]
for i,(title_,desc) in enumerate(feats):
    L=Inches(1.0)+(i%2)*Inches(5.95); T=Inches(2.1)+(i//2)*Inches(2.15)
    box(s,L,T,Inches(5.7),Inches(1.85),(20,35,80))
    gbox(s,L,T,Inches(5.7),Inches(0.08),BLUE,CYAN,angle=0)
    label(s,title_,L+Inches(0.2),T+Inches(0.2),Inches(5.3),Inches(0.5),size=18,c=WHITE,bold=True)
    label(s,desc,L+Inches(0.2),T+Inches(0.78),Inches(5.3),Inches(0.75),size=14,c=(160,185,230))
notitie(s,"Geef studenten 5 minuten om een account aan te maken. Gebruik schoolmailadres voor het gratis Education plan.")


# ── SLIDE 13: Interface ───────────────────────────────────
s=new_slide(light_bg()); paginanr(s,13)
page_header(s,"De interface van Figma")
zones=[("1","Toolbar","Bovenaan","Gereedschappen: frame, shapes, tekst, pen",BLUE),("2","Layers panel","Links","Alles wat je maakt staat hier",( 124,58,237)),("3","Canvas","Midden","Jouw werkruimte, oneindig groot, begin altijd met een frame",(5,150,105)),("4","Properties panel","Rechts","Afmetingen, kleuren, tekst en effecten van het geselecteerde element",RED)]
for i,(num,name,loc,desc,c) in enumerate(zones):
    T=Inches(1.82)+i*Inches(1.28)
    dot(s,Inches(1.2),T+Inches(0.28),Inches(0.25),c)
    label(s,num,Inches(0.97),T+Inches(0.1),Inches(0.46),Inches(0.45),size=16,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    label(s,name,Inches(1.62),T+Inches(0.07),Inches(2.2),Inches(0.48),size=20,c=DARK,bold=True)
    box(s,Inches(3.9),T+Inches(0.1),Inches(1.15),Inches(0.4),c)
    label(s,loc,Inches(3.92),T+Inches(0.13),Inches(1.11),Inches(0.35),size=12,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    label(s,desc,Inches(5.22),T+Inches(0.09),Inches(7.7),Inches(0.48),size=16,c=BODY)
    if i<3: box(s,Inches(1.0),T+Inches(0.64),Inches(11.9),Inches(0.02),(225,230,242))
notitie(s,"Laat dit nu zien in Figma. Wijs de 4 zones aan, klik erop. Studenten openen tegelijk hun eigen Figma.")


# ── SLIDE 14: Frames ──────────────────────────────────────
s=new_slide(light_bg()); paginanr(s,14)
page_header(s,"Frame = Jouw scherm")
blist(s,["Een frame stelt een schermformaat voor, vergelijkbaar met een artboard","Sneltoets: F, klik en sleep of kies een preset rechts in het panel","Presets: iPhone 14 (390x844)  en  Desktop (1440x1024)","Frames kunnen in frames, dit wordt later belangrijk bij components","Geef frames altijd een naam: 'home-desktop', 'hero-mobiel'"],
      Inches(1.0),Inches(1.75),Inches(11.3),Inches(3.85),size=19,c=DARK)
cl,ct,cw=card(s,Inches(1.0),Inches(5.75),Inches(11.3),Inches(1.38),accent=BLUE,bg=OFFWH)
label(s,"Gouden regel:",cl,ct,cw,Inches(0.42),size=14,c=DARK,bold=True)
label(s,"Begin altijd met een frame. Nooit direct op de canvas ontwerpen.",cl,ct+Inches(0.45),cw,Inches(0.55),size=19,c=DARK,bold=True)
notitie(s,"Demo: druk F, kies Desktop, geef de naam 'home'. Druk daarna F opnieuw en kies iPhone 14. Laat zien hoe je beide naast elkaar kunt hebben.")


# ── SLIDE 15: Layers ──────────────────────────────────────
s=new_slide(light_bg()); paginanr(s,15)
page_header(s,"Layers panel")
blist(s,["Elke shape, tekst of frame verschijnt in de layers","Volgorde is diepte: bovenste laag in de lijst staat vooraan","Ctrl+G of Cmd+G om elementen te groeperen","Naamgeving is cruciaal, een developer of collega werkt met jouw bestand"],
      Inches(1.0),Inches(1.75),Inches(6.8),Inches(3.2),size=18,c=DARK)
label(s,"Gebruik: sectie / element als naamconventie",Inches(1.0),Inches(5.05),Inches(6.8),Inches(0.45),size=15,c=BLUE,bold=True)
# Vergelijking: twee platte panelen
PW=Inches(2.8); PH=Inches(4.85); PT=Inches(1.65)
for ci,(label_,items,bg,tc) in enumerate([("Slecht",["Rectangle 47","Group 3","Text 12","Ellipse 2","Rectangle 48"],(80,15,15),(255,160,160)),("Goed",["hero / achtergrond","nav / logo","nav / links","hero / kop","hero / knop"],(15,50,15),(140,230,140))]):
    L=Inches(8.4)+ci*Inches(3.0)
    box(s,L,PT,PW,PH,bg)
    gbox(s,L,PT,PW,Inches(0.08),RED if ci==0 else GREEN,RED if ci==0 else GREEN,angle=0)
    label(s,label_,L+Inches(0.15),PT+Inches(0.15),PW-Inches(0.3),Inches(0.38),size=14,c=tc,bold=True)
    for j,item in enumerate(items):
        box(s,L+Inches(0.15),PT+Inches(0.65)+j*Inches(0.78),PW-Inches(0.3),Inches(0.65),(bg[0]+12,bg[1]+12,bg[2]+12))
        label(s,item,L+Inches(0.25),PT+Inches(0.72)+j*Inches(0.78),PW-Inches(0.5),Inches(0.5),size=12,c=tc)
notitie(s,"Laat studenten 3 rechthoeken plaatsen en benoemen. Laat dan zien wat er gebeurt als je een bestand deelt met verkeerde namen.")


# ── SLIDE 16: Shapes ──────────────────────────────────────
s=new_slide(light_bg()); paginanr(s,16)
page_header(s,"Je eerste gereedschappen")
tools=[("R","Rechthoek","Klik en sleep  ·  Shift ingedrukt voor vierkant",NAVY),("O","Ellipse","Klik en sleep  ·  Shift ingedrukt voor cirkel",BLUE),("T","Tekst","Klik op de canvas en begin te typen",(5,150,105)),("V","Select","Altijd terug naar V na gebruik van een ander gereedschap",(124,58,237)),("F","Frame","Klik en sleep of kies preset rechts in het panel",RED)]
for i,(key,name,desc,c) in enumerate(tools):
    T=Inches(1.75)+i*Inches(1.02)
    box(s,Inches(1.0),T,Inches(0.72),Inches(0.72),c)
    label(s,key,Inches(1.0),T+Inches(0.05),Inches(0.72),Inches(0.6),size=24,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    label(s,name,Inches(1.9),T+Inches(0.1),Inches(2.5),Inches(0.52),size=20,c=DARK,bold=True)
    label(s,desc,Inches(4.5),T+Inches(0.12),Inches(7.9),Inches(0.48),size=15,c=BODY)
notitie(s,"Demo: maak een button. Rechthoek (R), tekst (T), centreer de tekst. Verander kleur via het properties panel rechts.")


# ── SLIDE 17: Opdracht ────────────────────────────────────
s=new_slide(dark_bg(False)); paginanr(s,17,light=True)
gbox(s,0,0,Inches(0.5),SH,BLUE,CYAN,angle=270)
label(s,"Opdracht  ·  Les 1",Inches(1.0),Inches(0.45),Inches(11.3),Inches(0.7),size=32,c=WHITE,bold=True)
label(s,"Goed Doel Website  ·  Analyse",Inches(1.0),Inches(1.12),Inches(8.0),Inches(0.45),size=17,c=(160,185,230),italic=True)
gbox(s,Inches(1.0),Inches(1.55),Inches(2.2),Inches(0.08),BLUE,CYAN,angle=0)
stappen=[("1","Maak een nieuw Figma bestand: [Naam] Goed Doel Project"),("2","Zoek 3 websites van goede doelen, bijvoorbeeld Rode Kruis, Unicef of Amnesty"),("3","Maak screenshots en plak ze op aparte frames in Figma"),("4","Annoteer per website: hero, navigatie en call-to-action"),("5","Schrijf 2 tot 3 pluspunten en 2 tot 3 verbeterpunten per site"),("6","Deel via Share, Anyone with link, Can view")]
bkl2=[BLUE,(59,130,246),(99,150,246),(6,182,212),(5,150,105),(16,163,110)]
for i,((num,step),bc) in enumerate(zip(stappen,bkl2)):
    T=Inches(1.82)+i*Inches(0.84)
    dot(s,Inches(1.2),T+Inches(0.2),Inches(0.2),bc)
    label(s,num,Inches(0.98),T+Inches(0.04),Inches(0.4),Inches(0.36),size=13,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    label(s,step,Inches(1.6),T+Inches(0.02),Inches(11.1),Inches(0.48),size=16,c=(210,225,255))
box(s,Inches(1.0),Inches(6.88),Inches(11.3),Inches(0.52),(20,35,80))
label(s,"Inleveren: Figma link voor het begin van les 2, twee ontwerpen bespreken we klassikaal",Inches(1.25),Inches(6.96),Inches(10.8),Inches(0.38),size=13,c=(180,210,255),align=PP_ALIGN.CENTER)
notitie(s,"Geef de laatste 15 minuten om te starten. Loop rond. Veel studenten hebben moeite met tekst in Figma, leg dat even extra uit.")


# ── SLIDE 18: Volgende les ────────────────────────────────
s=new_slide(dark_bg(True)); paginanr(s,18,light=True)
gbox(s,0,0,Inches(0.5),SH,BLUE,CYAN,angle=270)
label(s,"Volgende les",Inches(1.0),Inches(1.5),Inches(11.3),Inches(0.7),size=42,c=WHITE,bold=True)
label(s,"Whitespace en Alignment",Inches(1.0),Inches(2.2),Inches(8.0),Inches(0.6),size=26,c=CYAN)
gbox(s,Inches(1.0),Inches(2.82),Inches(2.5),Inches(0.09),BLUE,CYAN,angle=0)
for i,item in enumerate(["Het 8pt grid systeem als basis voor consistent spacing","Micro en macro whitespace: hoe ademruimte werkt","Je eerste wireframe in Figma voor het goed doel project"]):
    box(s,Inches(1.0),Inches(3.1)+i*Inches(0.92),(SW-Inches(1.5)),Inches(0.75),(20,35,80))
    gbox(s,Inches(1.0),Inches(3.1)+i*Inches(0.92),Inches(0.07),Inches(0.75),BLUE,CYAN,angle=270)
    label(s,item,Inches(1.25),Inches(3.18)+i*Inches(0.92),(SW-Inches(2.0)),Inches(0.55),size=17,c=(210,225,255))
box(s,Inches(1.0),Inches(6.1),(SW-Inches(1.5)),Inches(1.05),(15,28,70))
label(s,"Tip: Installeer de Figma desktop app via figma.com/downloads voor betere performance",Inches(1.25),Inches(6.42),(SW-Inches(2.0)),Inches(0.42),size=15,c=(160,185,230),align=PP_ALIGN.CENTER)
notitie(s,"Sluit de les af. Vraag nog of iedereen een Figma account heeft en de opdracht begrijpt.")

prs.save(r"C:\Users\Peter Marcelis\Documents\Lesmateriaal\design\periode-1\les-1\01_intro-wat-is-design.pptx")
print(f"Klaar: {len(prs.slides)} slides")
