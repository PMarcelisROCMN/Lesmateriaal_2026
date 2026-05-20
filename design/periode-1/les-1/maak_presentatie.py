from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Kleurenpalet ──────────────────────────────────────────
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x11, 0x18, 0x27)
BLUE   = RGBColor(0x23, 0x63, 0xEB)
BODY   = RGBColor(0x37, 0x41, 0x51)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
LBLUE  = RGBColor(0xEF, 0xF6, 0xFF)
QUOTE_BG = RGBColor(0xF0, 0xF4, 0xFF)

# ── Afmetingen ────────────────────────────────────────────
SW   = Inches(13.33)
SH   = Inches(7.5)
ML   = Inches(1.0)      # margin left (na blauwe balk)
MT   = Inches(0.55)     # margin top
CW   = Inches(11.5)     # content width
BAR  = Inches(0.45)     # breedte blauwe balk

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ── Hulpfuncties ──────────────────────────────────────────

def new_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return slide


def rect(slide, l, t, w, h, color, border=False):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if not border:
        s.line.fill.background()
    else:
        s.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        s.line.width = Pt(0.75)
    return s


def txt(slide, text, l, t, w, h, size=20, color=BODY,
        bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.name      = "Calibri"
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    return box, tf


def bullets(slide, items, l, t, w, h, size=19, color=BODY, spacing=7):
    """items = list of str  OR  list of (str, size, bold)"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        run = p.add_run()
        if isinstance(item, tuple):
            text, sz = item[0], item[1]
            bd = item[2] if len(item) > 2 else False
        else:
            text, sz, bd = item, size, False
        run.text           = f"  •  {text}"
        run.font.name      = "Calibri"
        run.font.size      = Pt(sz)
        run.font.bold      = bd
        run.font.color.rgb = color
    return box, tf


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def bar(slide):
    rect(slide, 0, 0, BAR, SH, BLUE)


def title_block(slide, title_text, sub=None):
    txt(slide, title_text, ML, MT, CW, Inches(0.9),
        size=34, color=DARK, bold=True)
    rect(slide, ML, Inches(1.38), Inches(1.6), Inches(0.07), BLUE)
    if sub:
        txt(slide, sub, ML, Inches(1.5), CW, Inches(0.45),
            size=15, color=GRAY, italic=True)


def slide_number(slide, n, total=18):
    txt(slide, f"{n} / {total}",
        Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.35),
        size=12, color=GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════
#  SLIDE 1 — Titelslide
# ══════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, 0, BAR * 3, SH, BLUE)                    # brede blauwe balk links
txt(s, "Wat is Design?",
    Inches(1.8), Inches(2.2), Inches(9), Inches(1.4),
    size=52, color=DARK, bold=True)
txt(s, "Periode 1  ·  Les 1",
    Inches(1.8), Inches(3.55), Inches(6), Inches(0.5),
    size=18, color=GRAY)
rect(s, Inches(1.8), Inches(3.45), Inches(1.8), Inches(0.06), BLUE)
txt(s, "Design voor het web",
    Inches(1.8), Inches(4.05), Inches(6), Inches(0.4),
    size=14, color=GRAY, italic=True)
notes(s, "Welkom bij les 1. Vandaag starten we breed: wat is design eigenlijk? Daarna maken we kennis met Figma.")

# ══════════════════════════════════════════════════════════
#  SLIDE 2 — Agenda
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Vandaag"); slide_number(s, 2)
bullets(s, [
    "Wat is design — en wat niet",
    "Hoe denkt een gebruiker?",
    "Het kanaal bepaalt de boodschap",
    "Kennismaking met Figma",
    "Opdracht van vandaag",
], ML, Inches(1.8), CW, Inches(4.5), size=22)
notes(s, "Loop de agenda door. Benadruk dat we beginnen met theorie en eindigen met een concrete opdracht in Figma.")

# ══════════════════════════════════════════════════════════
#  SLIDE 3 — Design ≠ Mooi maken
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Design ≠ Iets mooi maken"); slide_number(s, 3)
bullets(s, [
    "Design lost een probleem op voor een specifieke gebruiker",
    "Mooi is een bijproduct — niet het doel",
    "Slecht design kan er prachtig uitzien en toch niet werken",
    "Goed design kan sober zijn en perfect werken",
], ML, Inches(1.8), CW, Inches(2.8), size=20)
# Quote box
rect(s, ML, Inches(4.4), CW, Inches(1.8), QUOTE_BG, border=True)
txt(s, '"Design is not just what it looks like and feels like.\nDesign is how it works."',
    Inches(1.3), Inches(4.55), Inches(9.5), Inches(1.0),
    size=19, color=DARK, italic=True, align=PP_ALIGN.CENTER)
txt(s, "— Steve Jobs",
    ML, Inches(5.6), CW, Inches(0.4),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)
notes(s, "Gebruik twee concrete voorbeelden: een esthetisch mooie website die onbruikbaar is (laat live zien), en de eerste Google homepage uit 1998 — lelijk maar briljant werkend.")

# ══════════════════════════════════════════════════════════
#  SLIDE 4 — Design vs. Kunst
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Design vs. Kunst"); slide_number(s, 4)

col_w = Inches(5.0)
gap   = Inches(0.4)
col1_l = ML
col2_l = ML + col_w + gap

# Kolom 1
rect(s, col1_l, Inches(1.8), col_w, Inches(0.55), BLUE)
txt(s, "Design", col1_l, Inches(1.83), col_w, Inches(0.5),
    size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
for i, item in enumerate([
    "In dienst van een doel",
    "Altijd een gebruiker centraal",
    "Meetbaar: werkt het?",
    "Beperkingen zijn het spel",
]):
    rect(s, col1_l, Inches(2.45 + i * 0.82), col_w, Inches(0.72),
         RGBColor(0xF8, 0xFA, 0xFF) if i % 2 == 0 else WHITE, border=True)
    txt(s, item, Inches(1.15), Inches(2.55 + i * 0.82), Inches(4.7), Inches(0.6),
        size=17, color=BODY)

# Kolom 2
rect(s, col2_l, Inches(1.8), col_w, Inches(0.55), DARK)
txt(s, "Kunst", col2_l, Inches(1.83), col_w, Inches(0.5),
    size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
for i, item in enumerate([
    "Vrije expressie",
    "Publiek is optioneel",
    "Subjectief",
    "Beperkingen zijn een hindernis",
]):
    rect(s, col2_l, Inches(2.45 + i * 0.82), col_w, Inches(0.72),
         RGBColor(0xF8, 0xF8, 0xF8) if i % 2 == 0 else WHITE, border=True)
    txt(s, item, Inches(6.6), Inches(2.55 + i * 0.82), Inches(4.7), Inches(0.6),
        size=17, color=BODY)

txt(s, "Zelfde vaardigheden — andere intentie",
    ML, Inches(6.7), CW, Inches(0.4),
    size=15, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
notes(s, "Benadruk dat de vaardigheden overlappen: kleur, compositie, typografie. Maar de intentie is anders. Jullie werk is altijd in dienst van iemand anders — dat maakt het vak interessant én uitdagend.")

# ══════════════════════════════════════════════════════════
#  SLIDE 5 — De centrale vraag
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); slide_number(s, 5)
# Grote quote, gecentreerd
rect(s, BAR, 0, SW - BAR, SH, LBLUE)
txt(s, "De vraag die je altijd moet stellen:",
    ML, Inches(1.5), CW, Inches(0.5),
    size=18, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, '"Kan mijn gebruiker zonder na te denken\ndoen wat ze wil doen?"',
    ML, Inches(2.3), CW, Inches(2.2),
    size=34, color=DARK, bold=True, align=PP_ALIGN.CENTER)
rect(s, Inches(4.5), Inches(4.6), Inches(4.3), Inches(0.07), BLUE)
txt(s, "— Steve Krug,  Don't Make Me Think  (2000)",
    ML, Inches(4.8), CW, Inches(0.5),
    size=16, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
notes(s, "Laat deze quote even staan. Vraag studenten: wanneer heb jij voor het laatst gefrustreerd een website verlaten? Wat was er mis?")

# ══════════════════════════════════════════════════════════
#  SLIDE 6 — Don't Make Me Think
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Don't Make Me Think"); slide_number(s, 6)
rect(s, ML, Inches(1.8), Inches(3.8), Inches(0.38), BLUE)
txt(s, "Steve Krug — 2000, nog steeds actueel",
    ML, Inches(1.82), Inches(6), Inches(0.35),
    size=14, color=WHITE, bold=True)
bullets(s, [
    "Gebruikers scannen pagina's — ze lezen niet",
    "Ze klikken op het eerste dat redelijk lijkt",
    "Elk vraagteken kost mentale energie",
    "Jouw design moet zo vanzelfsprekend zijn dat niemand nadenkt",
    "De gebruiker is druk, gehaast, en niet zo zorgvuldig als jij hoopt",
], ML, Inches(2.3), CW, Inches(4.0), size=20)
notes(s, "Dit boek is van 2000 maar de principes zijn tijdloos. Niets aan ons scangedrag is veranderd — het is eerder erger geworden door social media. Tip: laat een eye-tracking heatmap zien (te vinden op nngroup.com).")

# ══════════════════════════════════════════════════════════
#  SLIDE 7 — Gebruikers scannen (F-patroon)
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Zo lezen gebruikers een pagina"); slide_number(s, 7)
bullets(s, [
    "Gebruikers volgen een F-patroon (Nielsen Norman Group, 2006)",
    "Bovenaan links krijgt de meeste aandacht",
    "Eerste woorden van elke zin zijn het belangrijkst",
    "Koppen worden gelezen — alinea's nauwelijks",
    "Zet kritieke informatie bovenaan en links",
], ML, Inches(1.8), Inches(7.5), Inches(4.0), size=20)
# Rechts: gesimuleerde F-patroon visualisatie
for i, (w, color_) in enumerate([
    (Inches(4.5), RGBColor(0xFF, 0x4D, 0x4D)),
    (Inches(3.2), RGBColor(0xFF, 0x7A, 0x00)),
    (Inches(1.5), RGBColor(0xFF, 0xAA, 0x00)),
    (Inches(1.5), RGBColor(0xFF, 0xAA, 0x00)),
    (Inches(1.5), RGBColor(0xFF, 0xAA, 0x00)),
]):
    rect(s, Inches(8.8), Inches(1.85 + i * 0.92), w, Inches(0.6), color_)
txt(s, "Heatmap simulatie — rood = meest bekeken",
    Inches(8.8), Inches(6.55), Inches(4.0), Inches(0.4),
    size=12, color=GRAY, italic=True)
txt(s, "Bron: nngroup.com",
    ML, Inches(6.7), CW, Inches(0.4), size=12, color=GRAY)
notes(s, "Open nngroup.com/articles/f-shaped-pattern in de browser en laat de echte heatmaps zien. Vraag: hoe schrijf je koppen zodat ze ook los van de rest betekenis hebben?")

# ══════════════════════════════════════════════════════════
#  SLIDE 8 — The Medium is the Message
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "The Medium is the Message"); slide_number(s, 8)
bullets(s, [
    "Marshall McLuhan — filosoof en mediatheoreticus (1964)",
    "Het kanaal dat je kiest communiceert al iets, vóór je eerste woord",
    "Hetzelfde bericht via SMS, brief of telefoongesprek — compleet andere ervaring",
    "Voor design: desktop en mobiel zijn twee verschillende media",
], ML, Inches(1.8), CW, Inches(3.2), size=20)
rect(s, ML, Inches(5.0), CW, Inches(1.7), QUOTE_BG, border=True)
txt(s, '"The medium is the message."',
    Inches(1.3), Inches(5.15), Inches(9.5), Inches(0.6),
    size=22, color=DARK, bold=True, italic=True, align=PP_ALIGN.CENTER)
txt(s, "— Marshall McLuhan,  Understanding Media  (1964)",
    ML, Inches(5.85), CW, Inches(0.4),
    size=14, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
notes(s, "Geef het voorbeeld: een sollicitatiebrief per WhatsApp vs. per e-mail vs. op papier. De inhoud is identiek, maar de boodschap is compleet anders. Vertaal dit daarna naar responsive design.")

# ══════════════════════════════════════════════════════════
#  SLIDE 9 — Responsive design als mediumvraag
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Hetzelfde design, andere ervaring"); slide_number(s, 9)

col_w = Inches(3.4)
headers = ["Desktop", "Tablet", "Mobiel"]
colors  = [BLUE, RGBColor(0x60, 0xA5, 0xFA), RGBColor(0x93, 0xC5, 0xFD)]
items_  = [
    ["Muis, ruimte, zittend", "Grotere viewport", "Hover states werken", "Gefocuste sessie"],
    ["Aanraking, half zittend", "Middelgrote viewport", "Beperkte hover", "Flexibele sessie"],
    ["Vinger, onderweg, gehaast", "Kleine viewport", "Geen hover states", "Versnipperde aandacht"],
]
for col, (header, color_, its) in enumerate(zip(headers, colors, items_)):
    lft = ML + col * (col_w + Inches(0.25))
    rect(s, lft, Inches(1.8), col_w, Inches(0.5), color_)
    txt(s, header, lft, Inches(1.83), col_w, Inches(0.45),
        size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for r, item in enumerate(its):
        bg = RGBColor(0xF0, 0xF4, 0xFF) if r % 2 == 0 else WHITE
        rect(s, lft, Inches(2.38 + r * 0.88), col_w, Inches(0.82), bg, border=True)
        txt(s, item, Inches(lft / 914400 + 0.1), Inches(2.48 + r * 0.88),
            Inches(3.2), Inches(0.7), size=15, color=BODY)

txt(s, "Responsive design is niet alleen technisch — het is een designvraag",
    ML, Inches(6.7), CW, Inches(0.4),
    size=15, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
notes(s, "Vraag studenten: gebruik jij dezelfde apps anders op je telefoon dan op een laptop? Wat doe je anders en waarom? Dit is precies het mediumprobleem dat McLuhan beschreef.")

# ══════════════════════════════════════════════════════════
#  SLIDE 10 — Goed design is onzichtbaar
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Goed design valt niet op"); slide_number(s, 10)
bullets(s, [
    "Als iets goed werkt, merk je het niet",
    "Als iets slecht werkt, merk je het meteen",
    "Don Norman (1988): een deur met een trekhandvat die je toch moet duwen",
    "Het design communiceert verkeerd → frustratie",
], ML, Inches(1.8), CW, Inches(3.0), size=20)
rect(s, ML, Inches(4.9), CW, Inches(1.8), QUOTE_BG, border=True)
txt(s, "Norman Door  =  een ontwerp dat de verkeerde aanwijzing geeft",
    Inches(1.3), Inches(5.0), Inches(9.5), Inches(0.5),
    size=18, color=DARK, bold=True, align=PP_ALIGN.CENTER)
txt(s, "Webequivalenten: een knop die er niet uitziet als knop · een link zonder kleur · een formulier zonder labels",
    Inches(1.3), Inches(5.55), Inches(9.5), Inches(0.8),
    size=16, color=BODY, align=PP_ALIGN.CENTER)
txt(s, "Bron: Don Norman — The Design of Everyday Things (1988)",
    ML, Inches(6.75), CW, Inches(0.35), size=12, color=GRAY)
notes(s, "Zoek 'Norman door' op YouTube — de Vox video (4 minuten) is perfect. Laat hem zien als er tijd is. Vraag daarna: welk webdesign-equivalent heb jij deze week meegemaakt?")

# ══════════════════════════════════════════════════════════
#  SLIDE 11 — Live analyse
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Goed vs. Slecht — Live analyse"); slide_number(s, 11)
rect(s, ML, Inches(1.8), Inches(5.0), Inches(4.4), RGBColor(0xF0, 0xFD, 0xF4), border=True)
txt(s, "Goed", ML + Inches(0.15), Inches(1.9), Inches(4.7), Inches(0.4),
    size=16, color=RGBColor(0x16, 0xA3, 0x4A), bold=True)
bullets(s, [
    "Open: stripe.com",
    "Wat zie je als eerste?",
    "Weet je meteen wat je hier kunt doen?",
    "Hoeveel moeite kost het om iets te vinden?",
], ML + Inches(0.15), Inches(2.35), Inches(4.7), Inches(3.5), size=17,
   color=BODY, spacing=10)

rect(s, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.4), RGBColor(0xFF, 0xF1, 0xF2), border=True)
txt(s, "Slecht", Inches(6.95), Inches(1.9), Inches(5.3), Inches(0.4),
    size=16, color=RGBColor(0xDC, 0x26, 0x26), bold=True)
bullets(s, [
    "Zoek een verouderde bedrijfssite",
    "Waar is de navigatie?",
    "Wat wil deze site dat je doet?",
    "Hoe snel raak je gefrustreerd?",
], Inches(6.95), Inches(2.35), Inches(5.3), Inches(3.5), size=17,
   color=BODY, spacing=10)

txt(s, "Kijk niet of het mooi is — kijk of het werkt",
    ML, Inches(6.6), CW, Inches(0.45),
    size=16, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
notes(s, "Dit is een interactief moment — geen slides nodig, open de browser. Stel vragen aan de klas. Houd het bij 5 minuten per site. Bespreek na: wat maakt het verschil?")

# ══════════════════════════════════════════════════════════
#  SLIDE 12 — Kennismaking Figma
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Kennismaking met Figma"); slide_number(s, 12)
bullets(s, [
    "Browser-gebaseerd — geen installatie verplicht (desktop app beschikbaar)",
    "Real-time samenwerking — meerdere mensen tegelijk in één bestand",
    "Gratis voor studenten via het Education plan",
    "Industriestandaard — gebruikt door Spotify, Airbnb, Uber, Google",
    "Vervangt Sketch (alleen Mac) en Adobe XD (gestopt in 2023)",
], ML, Inches(1.8), CW, Inches(4.0), size=20)
rect(s, ML, Inches(5.95), CW, Inches(0.95), LBLUE, border=True)
txt(s, "Account aanmaken:  figma.com  →  'Get started for free'  →  gebruik je schoolmailadres  →  vraag Education plan aan",
    Inches(1.3), Inches(6.05), Inches(10.8), Inches(0.75),
    size=16, color=DARK, align=PP_ALIGN.CENTER)
notes(s, "Geef studenten nu 5 minuten om een account aan te maken als ze dat nog niet hebben gedaan. Laat daarna zelf Figma openen en leid ze door de interface.")

# ══════════════════════════════════════════════════════════
#  SLIDE 13 — Figma interface: de vier zones
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "De interface van Figma"); slide_number(s, 13)

zones = [
    ("1", "Toolbar",          "Bovenaan",  "Gereedschappen: frame, shapes, tekst, pen, hand",    BLUE),
    ("2", "Layers panel",     "Links",     "Alles wat je maakt staat hier — jouw structuur",       RGBColor(0x7C, 0x3A, 0xED)),
    ("3", "Canvas",           "Midden",    "Jouw werkruimte — oneindig groot, begin altijd met een frame", RGBColor(0x05, 0x96, 0x69)),
    ("4", "Properties panel", "Rechts",    "Afmetingen, kleuren, tekst, effecten van het geselecteerde element", RGBColor(0xDC, 0x26, 0x26)),
]
for i, (num, name, loc, desc, color_) in enumerate(zones):
    t = Inches(1.85 + i * 1.25)
    rect(s, ML, t, Inches(0.5), Inches(0.5), color_)
    txt(s, num, ML, t, Inches(0.5), Inches(0.5),
        size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, name, Inches(1.7), t, Inches(2.2), Inches(0.5),
        size=18, color=DARK, bold=True)
    txt(s, f"({loc})", Inches(3.9), t, Inches(1.5), Inches(0.5),
        size=14, color=GRAY, italic=True)
    txt(s, desc, Inches(5.4), t, Inches(7.2), Inches(0.5),
        size=16, color=BODY)

notes(s, "Laat dit nu zien in Figma zelf. Wijs de 4 zones aan, klik erop, laat zien wat er verandert. Studenten openen tegelijk hun eigen Figma.")

# ══════════════════════════════════════════════════════════
#  SLIDE 14 — Frames
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Frame = Jouw scherm"); slide_number(s, 14)
bullets(s, [
    "Een frame stelt een schermformaat voor (zoals een artboard)",
    "Sneltoets:  F  → klik en sleep, of kies een preset rechts",
    "Presets:  iPhone 14  (390×844)  ·  Desktop  (1440×1024)",
    "Frames kunnen in frames — dit wordt later belangrijk bij components",
    "Naam je frames altijd:  'home-desktop', 'home-mobiel'",
], ML, Inches(1.8), CW, Inches(4.0), size=20)
rect(s, ML, Inches(5.95), CW, Inches(1.0), QUOTE_BG, border=True)
txt(s, "Gouden regel:  begin ALTIJD met een frame — nooit direct op de canvas ontwerpen",
    Inches(1.3), Inches(6.1), Inches(10.8), Inches(0.7),
    size=17, color=DARK, bold=True, align=PP_ALIGN.CENTER)
notes(s, "Demo in Figma: druk F, kies Desktop, geef het de naam 'home'. Druk daarna F opnieuw en kies iPhone 14. Laat zien hoe je beide frames naast elkaar kunt hebben.")

# ══════════════════════════════════════════════════════════
#  SLIDE 15 — Layers panel
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Layers panel — Jouw structuur"); slide_number(s, 15)
bullets(s, [
    "Elke shape, tekst of frame verschijnt in de layers",
    "Volgorde = diepte:  bovenste laag in de lijst = meest vooraan",
    "Ctrl+G  (of Cmd+G)  om elementen te groeperen",
    "Naamgeving is cruciaal — een developer of collega werkt met jouw bestand",
], ML, Inches(1.8), Inches(6.5), Inches(3.8), size=20)

# Visueel: goed vs. slecht layers voorbeeld
rect(s, Inches(7.8), Inches(1.7), Inches(4.8), Inches(4.8), RGBColor(0xF9, 0xFA, 0xFB), border=True)
txt(s, "Slecht", Inches(7.95), Inches(1.75), Inches(2.0), Inches(0.4),
    size=14, color=RGBColor(0xDC, 0x26, 0x26), bold=True)
txt(s, "Goed", Inches(10.2), Inches(1.75), Inches(2.0), Inches(0.4),
    size=14, color=RGBColor(0x16, 0xA3, 0x4A), bold=True)
bad  = ["Rectangle 47", "Group 3", "Text 12", "Ellipse 2", "Rectangle 48"]
good = ["hero / achtergrond", "nav / logo", "nav / links", "hero / kop", "hero / knop"]
for i, (b, g) in enumerate(zip(bad, good)):
    txt(s, b, Inches(7.95), Inches(2.25 + i * 0.68), Inches(2.1), Inches(0.6),
        size=13, color=RGBColor(0xDC, 0x26, 0x26))
    txt(s, g, Inches(10.2), Inches(2.25 + i * 0.68), Inches(2.3), Inches(0.6),
        size=13, color=RGBColor(0x16, 0xA3, 0x4A))

txt(s, "Gebruik:  sectie / element  als naamconventie",
    ML, Inches(5.7), Inches(6.5), Inches(0.5),
    size=16, color=BLUE, bold=True)
notes(s, "Laat studenten hun eigen Figma bestand openen. Vraag ze om 3 rechthoeken te plaatsen en te noemen. Laat daarna zien wat er gebeurt als je een bestand deelt met verkeerde namen — niemand snapt er iets van.")

# ══════════════════════════════════════════════════════════
#  SLIDE 16 — Shapes & tekst
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Je eerste elementen"); slide_number(s, 16)

tools = [
    ("R", "Rechthoek",  "Klik en sleep — + Shift voor perfect vierkant"),
    ("O", "Ellipse",    "Klik en sleep — + Shift voor perfecte cirkel"),
    ("T", "Tekst",      "Klik op de canvas en begin te typen"),
    ("V", "Select",     "Altijd terug naar V na gebruik van een ander gereedschap"),
    ("F", "Frame",      "Klik en sleep, of kies preset rechts in het panel"),
]
for i, (key, name, desc) in enumerate(tools):
    t = Inches(1.75 + i * 1.02)
    rect(s, ML, t, Inches(0.55), Inches(0.55), BLUE)
    txt(s, key, ML, t, Inches(0.55), Inches(0.55),
        size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, name, Inches(1.75), t, Inches(2.0), Inches(0.55),
        size=18, color=DARK, bold=True)
    txt(s, desc, Inches(3.8), t, Inches(8.8), Inches(0.55),
        size=16, color=BODY)

notes(s, "Demo: maak live een eenvoudige button — een rechthoek (R), dan een tekstlabel (T), centreer de tekst over de rechthoek. Verander de kleur via het properties panel rechts. Studenten doen dit tegelijk mee.")

# ══════════════════════════════════════════════════════════
#  SLIDE 17 — Opdracht
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Opdracht — Les 1"); slide_number(s, 17)
rect(s, ML, Inches(1.75), Inches(0.45), Inches(4.6), BLUE)
txt(s, "Goed Doel Website — Analyse",
    Inches(1.65), Inches(1.75), Inches(9.5), Inches(0.45),
    size=16, color=WHITE, bold=True)

steps = [
    "Maak een nieuw Figma bestand:  [Naam] — Goed Doel Project",
    "Zoek 3 websites van goede doelen (bijv. Rode Kruis, Unicef, Amnesty)",
    "Maak screenshots en plak ze op aparte frames in Figma",
    "Annoteer met de teksttool per website:",
    "     →  Hero section  |  Navigatie  |  Call-to-action",
    "     →  2–3 dingen die goed werken",
    "     →  2–3 dingen die beter kunnen",
    "Deel het bestand:  Share → Anyone with link → Can view",
]
for i, step in enumerate(steps):
    t = Inches(2.3 + i * 0.54)
    if not step.startswith("  "):
        rect(s, ML, t + Inches(0.08), Inches(0.3), Inches(0.3),
             BLUE if i < 4 else RGBColor(0x93, 0xC5, 0xFD))
        txt(s, str(i + 1) if i < 4 else "", ML, t + Inches(0.06),
            Inches(0.3), Inches(0.32), size=12, color=WHITE,
            bold=True, align=PP_ALIGN.CENTER)
    txt(s, step, Inches(1.65), t, Inches(10.5), Inches(0.52),
        size=16, color=DARK if not step.startswith("  ") else BODY)

rect(s, ML, Inches(6.65), CW, Inches(0.5), LBLUE, border=True)
txt(s, "Inleveren:  begin volgende les — 2 ontwerpen bespreken we klassikaal",
    Inches(1.3), Inches(6.7), Inches(10.8), Inches(0.38),
    size=14, color=DARK, align=PP_ALIGN.CENTER)
notes(s, "Geef studenten de laatste 15 minuten van de les om hiermee te starten. Loop rond. Veel studenten hebben moeite met het plaatsen van tekst in Figma — leg dat even extra uit.")

# ══════════════════════════════════════════════════════════
#  SLIDE 18 — Volgende les
# ══════════════════════════════════════════════════════════
s = new_slide(); bar(s); title_block(s, "Volgende les: Whitespace & Alignment"); slide_number(s, 18)
bullets(s, [
    "Het 8pt grid systeem — de basis van consistent spacing",
    "Micro vs. macro whitespace — hoe ademruimte werkt",
    "Je eerste wireframe in Figma (voor het goed doel project)",
], ML, Inches(1.8), CW, Inches(3.0), size=21)
rect(s, ML, Inches(4.9), CW, Inches(1.15), LBLUE, border=True)
txt(s, "Tip voor thuis:",
    Inches(1.3), Inches(5.0), Inches(2.0), Inches(0.4),
    size=15, color=DARK, bold=True)
txt(s, "Installeer de Figma desktop app voor betere performance →  figma.com/downloads",
    Inches(1.3), Inches(5.42), Inches(10.5), Inches(0.5),
    size=16, color=BODY)
rect(s, ML, Inches(6.2), CW, Inches(1.0), QUOTE_BG, border=True)
txt(s, "Zorg dat je opdracht klaar is: 3 geanalyseerde websites in Figma",
    Inches(1.3), Inches(6.3), Inches(10.5), Inches(0.65),
    size=17, color=DARK, bold=True, align=PP_ALIGN.CENTER)
notes(s, "Sluit de les af. Vraag nog één keer of iedereen een Figma account heeft en de opdracht begrijpt. Zorg dat de gedeelde link klaarstaat in Teams of de leeromgeving.")


# ── Opslaan ───────────────────────────────────────────────
output = r"C:\Users\Peter Marcelis\Documents\Lesmateriaal\design\periode-1\les-1\01_intro-wat-is-design.pptx"
prs.save(output)
print(f"Klaar: {output}")
