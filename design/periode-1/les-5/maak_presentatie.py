"""
Les 5: Componenten en Design Systems
Periode 1 - Design Curriculum Jaar 1
Goed doel website project
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image, ImageDraw
import numpy as np, io

NAVY=(6,14,44); NAVY2=(30,58,138); BLUE=(37,99,235); BLUE2=(59,130,246)
CYAN=(6,182,212); TEAL=(13,148,136); WHITE=(255,255,255); OFFWH=(248,250,252)
DARK=(15,23,42); BODY=(51,65,85); MUTED=(100,116,135)
GREEN=(22,163,74); RED=(220,38,38)
ACCENT1=(22,163,74); ACCENT2=(74,222,128)

def rgb(c): return RGBColor(c[0],c[1],c[2])
SW,SH=Inches(13.33),Inches(7.5); BW,BH=1920,1080

def garr(w,h,c1,c2,d='diag'):
    x,y=np.linspace(0,1,w,dtype=np.float32),np.linspace(0,1,h,dtype=np.float32)
    if d=='h': t=np.tile(x,(h,1))
    elif d=='v': t=np.tile(y.reshape(-1,1),(1,w))
    else: xx,yy=np.meshgrid(x,y); t=xx*.55+yy*.45
    a=np.zeros((h,w,3),dtype=np.float32)
    for i in range(3): a[:,:,i]=c1[i]+(c2[i]-c1[i])*t
    return np.clip(a,0,255).astype(np.uint8)

def dark_bg(circ=True):
    img=Image.fromarray(garr(BW,BH,NAVY,NAVY2)).convert('RGBA')
    if circ:
        ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
        d.ellipse([BW-500,-180,BW+200,520],fill=(ACCENT1[0],ACCENT1[1],ACCENT1[2],20))
        d.ellipse([-60,BH-360,380,BH+100],fill=(ACCENT2[0],ACCENT2[1],ACCENT2[2],22))
        img=Image.alpha_composite(img,ov)
    return img.convert('RGB')

def light_bg():
    img=Image.fromarray(garr(BW,BH,OFFWH,WHITE)).convert('RGBA')
    ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
    d.ellipse([BW-580,-220,BW+180,540],fill=(ACCENT1[0],ACCENT1[1],ACCENT1[2],8))
    img=Image.alpha_composite(img,ov); return img.convert('RGB')

def buf(img):
    b=io.BytesIO(); img.save(b,format='PNG'); b.seek(0); return b

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
BLANK=prs.slide_layouts[6]

def ns(bg=None):
    s=prs.slides.add_slide(BLANK)
    if bg:
        p=s.shapes.add_picture(buf(bg),0,0,SW,SH)
        t=s.shapes._spTree; t.remove(p._element); t.insert(2,p._element)
    else:
        s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(WHITE)
    return s

def box(s,l,t,w,h,c,border=None):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c)
    if border: sh.line.color.rgb=rgb(border); sh.line.width=Pt(0.75)
    else: sh.line.fill.background()
    return sh

def gbox(s,l,t,w,h,c1,c2,angle=90):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.gradient(); st=sh.fill.gradient_stops
    st[0].position=0.0; st[0].color.rgb=rgb(c1)
    st[1].position=1.0; st[1].color.rgb=rgb(c2)
    sh.fill.gradient_angle=angle; sh.line.fill.background(); return sh

def dot(s,cx,cy,r,c):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,cx-r,cy-r,r*2,r*2)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c); sh.line.fill.background()

def lbl(s,text,l,t,w,h,size=18,c=BODY,bold=False,italic=False,align=PP_ALIGN.LEFT):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=text; r.font.name="Calibri"; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=rgb(c)

def blist(s,items,l,t,w,h,size=18,c=BODY,gap=7):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(gap); r=p.add_run()
        r.text=f"  {chr(8226)}  {item}"; r.font.name="Calibri"; r.font.size=Pt(size); r.font.color.rgb=rgb(c)

def notitie(s,t): s.notes_slide.notes_text_frame.text=t

def pnr(s,n,tot=15,light=False):
    lbl(s,f"{n} / {tot}",Inches(12.5),Inches(7.1),Inches(0.8),Inches(0.35),size=11,c=(180,180,200) if light else MUTED,align=PP_ALIGN.RIGHT)

def hero_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(1.0),size=36,c=WHITE,bold=True)
    gbox(s,Inches(1.0),Inches(1.42),Inches(2.0),Inches(0.08),ACCENT1,ACCENT2,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.55),Inches(11.3),Inches(0.42),size=15,c=(180,240,200),italic=True)

def page_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(0.95),size=34,c=DARK,bold=True)
    gbox(s,Inches(1.0),Inches(1.38),Inches(1.8),Inches(0.08),ACCENT1,ACCENT2,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.5),Inches(11.3),Inches(0.38),size=14,c=MUTED,italic=True)

def card(s,l,t,w,h,accent=None,bg=OFFWH):
    a=accent if accent else ACCENT1
    STRIP=Inches(0.07)
    box(s,l,t,w,h,bg,border=(215,225,240))
    gbox(s,l,t,STRIP,h,a,ACCENT2,angle=270)
    return l+STRIP+Inches(0.18), t+Inches(0.14), w-STRIP-Inches(0.35)

def col_kop(s,l,t,w,h,text,c):
    box(s,l,t,w,h,c)
    lbl(s,text,l,t+Inches(0.07),w,h-Inches(0.07),size=20,c=WHITE,bold=True,align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1: TITEL
# ─────────────────────────────────────────────────────────────────────────────
s1=ns(dark_bg(True))
# Green accent strip left
gbox(s1,Inches(0),Inches(1.8),Inches(0.12),Inches(3.8),ACCENT1,ACCENT2,angle=270)
lbl(s1,"Les 5",Inches(0.4),Inches(1.3),Inches(8),Inches(0.55),size=16,c=(180,240,200),bold=False,italic=True)
lbl(s1,"Componenten\nen Design Systems",Inches(0.4),Inches(1.8),Inches(10),Inches(2.2),size=52,c=WHITE,bold=True)
lbl(s1,"Periode 1  |  Les 5",Inches(0.4),Inches(3.95),Inches(8),Inches(0.5),size=16,c=(180,240,200))
gbox(s1,Inches(0.4),Inches(4.55),Inches(3.0),Inches(0.07),ACCENT1,ACCENT2,angle=0)
lbl(s1,"Ontwerpen als systeem, niet als losse onderdelen",Inches(0.4),Inches(4.7),Inches(10),Inches(0.5),size=17,c=(200,230,255),italic=True)
dot(s1,Inches(11.5),Inches(1.5),Inches(0.55),(22,163,74,))
dot(s1,Inches(12.2),Inches(2.1),Inches(0.3),(74,222,128))
dot(s1,Inches(10.8),Inches(5.8),Inches(0.2),(22,163,74))
pnr(s1,1,light=True)
notitie(s1,"Welkom bij les 5. Vandaag leer je hoe je ontwerpen opbouwt als een systeem van herbruikbare componenten - net zoals developers dat doen met code. Dit is een cruciale stap in professioneel ontwerpen.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2: AGENDA
# ─────────────────────────────────────────────────────────────────────────────
s2=ns(light_bg())
page_hdr(s2,"Vandaag","Les 5 overzicht")
items=[
    "Wat is een design system en waarom gebruik je het",
    "Atomic design: atoms, molecules, organisms",
    "Figma: je eerste component aanmaken",
    "Varianten en properties instellen",
    "Auto Layout voor flexibele componenten",
    "Goed doel: hero en navigatie als component",
]
CL=Inches(1.0); CT=Inches(1.85); CW=Inches(11.0); CH=Inches(0.68); GAP=Inches(0.08)
colors=[ACCENT1,BLUE,(124,58,237),TEAL,(217,119,6),RED]
for i,item in enumerate(items):
    t=CT+i*(CH+GAP)
    cl,_,_ = card(s2,CL,t,CW,CH,accent=colors[i])
    lbl(s2,f"{i+1}.",CL+Inches(0.1),t+Inches(0.14),Inches(0.35),CH-Inches(0.14),size=17,c=colors[i],bold=True)
    lbl(s2,item,CL+Inches(0.55),t+Inches(0.14),CW-Inches(0.9),CH-Inches(0.14),size=17,c=DARK)
pnr(s2,2)
notitie(s2,"Loop de agenda door. Leg uit dat alles wat we in de vorige lessen hebben gedaan (wit ruimte, typografie, kleur) nu samenkomt in een systeem. Vandaag leren studenten denken als een design system engineer.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3: WAT IS EEN DESIGN SYSTEM
# ─────────────────────────────────────────────────────────────────────────────
s3=ns()
# Left dark panel
box(s3,Inches(0),Inches(0),Inches(5.2),SH,NAVY)
gbox(s3,Inches(0),Inches(0),Inches(0.12),SH,ACCENT1,ACCENT2,angle=270)
lbl(s3,"Wat is een design system?",Inches(0.25),Inches(0.4),Inches(4.7),Inches(0.7),size=22,c=WHITE,bold=True)
gbox(s3,Inches(0.25),Inches(1.12),Inches(1.5),Inches(0.07),ACCENT1,ACCENT2,angle=0)
lbl(s3,"Een design system is een collectie van herbruikbare componenten, richtlijnen en beslissingen die samenwerken als een taal.",
    Inches(0.25),Inches(1.28),Inches(4.65),Inches(2.2),size=19,c=(210,230,255))
lbl(s3,"Zonder design system:",Inches(0.25),Inches(3.55),Inches(4.7),Inches(0.4),size=14,c=(160,180,220),italic=True)
blist(s3,["Elke pagina ziet er anders uit","Aanpassingen kosten veel tijd","Inconsistente stijlen en maten"],
      Inches(0.25),Inches(3.95),Inches(4.7),Inches(1.3),size=14,c=(180,200,240))
# Right light panel
lbl(s3,"Voordelen van een design system",Inches(5.5),Inches(0.45),Inches(7.5),Inches(0.55),size=22,c=DARK,bold=True)
gbox(s3,Inches(5.5),Inches(1.0),Inches(1.5),Inches(0.07),ACCENT1,ACCENT2,angle=0)
benefits=[
    "Consistentie over alle paginas van je website",
    "Sneller bouwen door hergebruik van componenten",
    "Makkelijker aanpassen: wijzig eenmalig in de master",
    "Betere samenwerking tussen ontwerpers en developers",
]
blist(s3,benefits,Inches(5.5),Inches(1.2),Inches(7.5),Inches(2.4),size=17,c=BODY)
# Card at bottom with examples
cl,ct,cw=card(s3,Inches(5.5),Inches(3.8),Inches(7.5),Inches(2.9),accent=ACCENT1)
lbl(s3,"Bekende design systems:",cl,ct,cw,Inches(0.38),size=15,c=DARK,bold=True)
examples=[
    "Material Design (Google) - m3.material.io",
    "Human Interface Guidelines (Apple) - developer.apple.com",
    "Tailwind UI (web) - tailwindui.com",
]
blist(s3,examples,cl,ct+Inches(0.4),cw,Inches(2.1),size=14,c=BODY)
pnr(s3,3)
notitie(s3,"Laat een van de bekende design systems zien op het scherm. Material Design is goed toegankelijk en gratis. Vraag: 'Welke websites gebruiken denk je een design system?' - bijna alle grote platforms.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4: ATOMIC DESIGN
# ─────────────────────────────────────────────────────────────────────────────
s4=ns(light_bg())
page_hdr(s4,"Atomic Design: van klein naar groot","Brad Frost (2013) - ontwerpen als chemische bouwstenen")
PW=Inches(3.7); PH=Inches(4.8); PT=Inches(1.72); GAP=Inches(0.22)
starts=[Inches(1.0), Inches(1.0)+PW+GAP, Inches(1.0)+2*(PW+GAP)]
headers=["Atoms","Molecules","Organisms"]
hcols=[ACCENT1,BLUE,NAVY2]
bodies=[
    ["de kleinste bouwstenen","Voorbeelden:","knop","invoerveld","label","icoon","kleurvlak"],
    ["combinaties van atoms","Voorbeelden:","zoekveld (invoer + knop)","nav-item (icoon + label)","formulierveld (label + input + fout)"],
    ["complete secties","Voorbeelden:","navigatiebalk (logo + links + knop)","hero sectie (heading + tekst + CTA)","testimonial kaart"],
]
KH=Inches(0.5)
for i,(xl,hc,hdr,bdy) in enumerate(zip(starts,hcols,headers,bodies)):
    col_kop(s4,xl,PT,PW,KH,hdr,hc)
    box(s4,xl,PT+KH,PW,PH-KH,WHITE,border=(215,225,240))
    lbl(s4,bdy[0],xl+Inches(0.15),PT+KH+Inches(0.12),PW-Inches(0.3),Inches(0.35),size=13,c=MUTED,italic=True)
    lbl(s4,bdy[1],xl+Inches(0.15),PT+KH+Inches(0.5),PW-Inches(0.3),Inches(0.35),size=14,c=hc,bold=True)
    for j,item in enumerate(bdy[2:]):
        lbl(s4,f"  {chr(8226)}  {item}",xl+Inches(0.15),PT+KH+Inches(0.88)+j*Inches(0.38),PW-Inches(0.3),Inches(0.38),size=14,c=BODY)
    # Visual cues at bottom of each panel
    vis_t=PT+PH-Inches(0.85)
    if i==0:
        box(s4,xl+Inches(0.3),vis_t,Inches(0.9),Inches(0.32),hc)
        box(s4,xl+Inches(1.35),vis_t,Inches(0.6),Inches(0.32),OFFWH,border=(200,210,220))
        dot(s4,xl+Inches(2.5),vis_t+Inches(0.16),Inches(0.16),hc)
    elif i==1:
        box(s4,xl+Inches(0.2),vis_t,Inches(1.5),Inches(0.32),OFFWH,border=(200,210,220))
        box(s4,xl+Inches(1.8),vis_t,Inches(0.72),Inches(0.32),hc)
    else:
        box(s4,xl+Inches(0.2),vis_t-Inches(0.1),PW-Inches(0.4),Inches(0.55),OFFWH,border=(200,210,220))
        box(s4,xl+Inches(0.2),vis_t-Inches(0.1),Inches(0.6),Inches(0.55),hc)
pnr(s4,4)
notitie(s4,"Atomic Design is een methode van Brad Frost uit 2013. Het helpt ontwerpers denken in herbruikbare eenheden. Vergelijk met LEGO: losse blokjes (atoms) vormen samen grotere structuren. Vraag studenten: welke atoms kennen ze al van hun goed doel ontwerp?")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5: COMPONENT IN FIGMA
# ─────────────────────────────────────────────────────────────────────────────
s5=ns(dark_bg(False))
hero_hdr(s5,"Figma: je eerste component","Van gewone shapes naar herbruikbare bouwstenen")
steps=[
    ("1","Maak een frame aan met de gewenste afmetingen, bijvoorbeeld 180x48 voor een knop"),
    ("2","Stijl het element met kleur, typografie en afstanden uit je stijlen"),
    ("3","Selecteer het frame en druk op Ctrl+Alt+K (Mac: Cmd+Alt+K) om een component te maken"),
    ("4","Het icoon wordt een ruitvorm: dit is je component (de master)"),
    ("5","Sleep het component vanuit het Assets panel om instances te maken"),
    ("6","Wijzig de master: alle instances worden automatisch bijgewerkt"),
]
ST=Inches(1.95); SW2=Inches(10.5); SH2=Inches(0.54); SGAP=Inches(0.08)
for i,(nr,txt) in enumerate(steps):
    t=ST+i*(SH2+SGAP)
    box(s5,Inches(1.0),t,Inches(0.45),SH2,ACCENT1)
    lbl(s5,nr,Inches(1.0),t+Inches(0.06),Inches(0.45),SH2-Inches(0.06),size=17,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    box(s5,Inches(1.46),t,SW2-Inches(0.46),SH2,(25,35,60),border=(40,55,90))
    lbl(s5,txt,Inches(1.65),t+Inches(0.08),SW2-Inches(0.7),SH2-Inches(0.08),size=15,c=(210,230,255))
# Tip card at bottom
cl,ct,cw=card(s5,Inches(1.0),Inches(6.2),Inches(11.0),Inches(0.88),accent=ACCENT1,bg=(20,30,55))
lbl(s5,"Tip:",cl,ct,Inches(0.45),Inches(0.5),size=14,c=ACCENT2,bold=True)
lbl(s5,"Geef componenten een beschrijvende naam: 'button/primary/large' is beter dan 'rectangle 1'",
    cl+Inches(0.52),ct,cw-Inches(0.52),Inches(0.5),size=14,c=(200,230,255),italic=True)
pnr(s5,5,light=True)
notitie(s5,"Doe dit live voor in Figma. Maak een eenvoudige knop-shape, stijl hem, en gebruik dan Ctrl+Alt+K. Laat zien hoe de ruitvorm verschijnt. Maak daarna een instance door Alt+slepen en verander de master. Studenten zien het effect direct.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6: VARIANTEN
# ─────────────────────────────────────────────────────────────────────────────
s6=ns(light_bg())
page_hdr(s6,"Varianten en Properties","Meerdere versies van hetzelfde component beheren")
# Left explanation
lbl(s6,"Wat zijn varianten?",Inches(1.0),Inches(1.72),Inches(5.5),Inches(0.45),size=18,c=DARK,bold=True)
blist(s6,[
    "Varianten zijn meerdere versies van hetzelfde component",
    "Gebruik ze voor states: default, hover, active, disabled",
    "Gebruik ze voor sizes: small, medium, large",
    "Gebruik ze voor types: primary, secondary, ghost",
],Inches(1.0),Inches(2.18),Inches(5.5),Inches(2.2),size=16,c=BODY)
# How to add variants
lbl(s6,"Varianten toevoegen in Figma:",Inches(1.0),Inches(4.45),Inches(5.5),Inches(0.4),size=16,c=DARK,bold=True)
blist(s6,[
    "Selecteer het component in Figma",
    "Klik op '+ Add variant' in het rechter paneel",
    "Geef elke variant een Property en Value naam",
],Inches(1.0),Inches(4.85),Inches(5.5),Inches(1.6),size=15,c=BODY)
# Right visual demo: button variant table
lbl(s6,"Voorbeeld: knop varianten",Inches(7.0),Inches(1.72),Inches(5.8),Inches(0.45),size=18,c=DARK,bold=True)
# Table header
cols=["primary","secondary","ghost"]
hcols=[ACCENT1,BLUE,(124,58,237)]
TL=Inches(7.0); TT=Inches(2.22); TW=Inches(1.6); TH=Inches(0.45); TGAP=Inches(0.08)
for j,(col,hc) in enumerate(zip(cols,hcols)):
    xl=TL+j*(TW+TGAP)
    box(s6,xl,TT,TW,TH,hc)
    lbl(s6,col,xl,TT+Inches(0.06),TW,TH-Inches(0.06),size=13,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
# Row 2: large
lbl(s6,"large",TL-Inches(0.82),TT+TH+TGAP+Inches(0.1),Inches(0.8),Inches(0.5),size=13,c=MUTED,italic=True)
row2_colors=[ACCENT1,OFFWH,WHITE]
row2_borders=[None,(200,210,220),(200,210,220)]
row2_texts=["Doneer","Doneer","Doneer"]
row2_tc=[WHITE,BLUE,BLUE]
for j,(rc,rb,rt,rtc) in enumerate(zip(row2_colors,row2_borders,row2_texts,row2_tc)):
    xl=TL+j*(TW+TGAP)
    box(s6,xl,TT+TH+TGAP,TW,Inches(0.52),rc,border=rb)
    lbl(s6,rt,xl,TT+TH+TGAP+Inches(0.1),TW,Inches(0.38),size=14,c=rtc,bold=True,align=PP_ALIGN.CENTER)
# Row 3: small
lbl(s6,"small",TL-Inches(0.82),TT+2*(TH+TGAP)+Inches(0.52)+Inches(0.2),Inches(0.8),Inches(0.35),size=13,c=MUTED,italic=True)
for j,(rc,rb,rt,rtc) in enumerate(zip(row2_colors,row2_borders,row2_texts,row2_tc)):
    xl=TL+j*(TW+TGAP)
    box(s6,xl,TT+2*(TH+TGAP)+Inches(0.52),TW,Inches(0.38),rc,border=rb)
    lbl(s6,rt,xl,TT+2*(TH+TGAP)+Inches(0.52)+Inches(0.04),TW,Inches(0.32),size=12,c=rtc,bold=True,align=PP_ALIGN.CENTER)
# Disabled row
lbl(s6,"disabled",TL-Inches(0.82),TT+3*(TH+TGAP)+Inches(0.9)+Inches(0.2),Inches(0.82),Inches(0.35),size=13,c=MUTED,italic=True)
for j in range(3):
    xl=TL+j*(TW+TGAP)
    box(s6,xl,TT+3*(TH+TGAP)+Inches(0.9),(TW),Inches(0.38),(200,210,220),border=(215,225,235))
    lbl(s6,"Doneer",xl,TT+3*(TH+TGAP)+Inches(0.9)+Inches(0.04),TW,Inches(0.32),size=12,c=(160,170,180),bold=True,align=PP_ALIGN.CENTER)
pnr(s6,6)
notitie(s6,"Varianten zijn een van de krachtigste functies van Figma. Laat zien hoe je in Figma een component set maakt en varianten toevoegt. Leg uit dat dit in code overeenkomt met CSS klassen of props in React.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7: AUTO LAYOUT
# ─────────────────────────────────────────────────────────────────────────────
s7=ns(dark_bg(False))
hero_hdr(s7,"Auto Layout: flexibele componenten","Shift+A in Figma - zoals flexbox in CSS")
# Bullets
blist(s7,[
    "Auto Layout laat elementen automatisch omgaan met veranderende inhoud",
    "Zoals flexbox in CSS: richting, gap, padding en uitlijning instellen",
    "Knoppen groeien mee met de tekst erin",
    "Lijsten herschikken zich als items worden toegevoegd",
    "Essentieel voor responsief ontwerpen in Figma",
],Inches(1.0),Inches(2.05),Inches(7.2),Inches(2.5),size=16,c=(210,230,255))
# Steps to enable
lbl(s7,"Auto Layout inschakelen:",Inches(1.0),Inches(4.65),Inches(7.2),Inches(0.42),size=16,c=ACCENT2,bold=True)
blist(s7,[
    "Selecteer een frame of component in Figma",
    "Druk op Shift+A om Auto Layout in te schakelen",
    "Stel in: richting (H/V), gap en padding (gebruik 8pt waarden: 8, 16, 24, 32)",
],Inches(1.0),Inches(5.1),Inches(7.2),Inches(1.5),size=15,c=(200,220,255))
# Visual demo: two boxes side by side
DT=Inches(2.2); DL=Inches(8.5); DW=Inches(2.2); DH=Inches(2.6)
# Without Auto Layout
box(s7,DL,DT,DW,DH,(20,30,60),border=(50,70,100))
lbl(s7,"Zonder Auto Layout",DL,DT+Inches(0.1),DW,Inches(0.38),size=12,c=(160,180,220),bold=True,align=PP_ALIGN.CENTER)
box(s7,DL+Inches(0.2),DT+Inches(0.6),DW-Inches(0.4),Inches(0.42),ACCENT1)
lbl(s7,"Knop met lange tekst...",DL+Inches(0.2),DT+Inches(0.6)+Inches(0.06),DW-Inches(0.4),Inches(0.32),size=11,c=WHITE,align=PP_ALIGN.CENTER)
# Clip indicator
box(s7,DL+DW-Inches(0.3),DT+Inches(0.6),Inches(0.3),Inches(0.42),RED)
lbl(s7,"Clip",DL+DW-Inches(0.32),DT+Inches(0.66),Inches(0.32),Inches(0.3),size=9,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
lbl(s7,"Tekst wordt afgeknipt",DL,DT+Inches(1.18),DW,Inches(0.38),size=11,c=(180,100,100),align=PP_ALIGN.CENTER)
# With Auto Layout
DL2=Inches(11.0)
box(s7,DL2,DT,DW,DH,(20,30,60),border=(50,70,100))
lbl(s7,"Met Auto Layout",DL2,DT+Inches(0.1),DW,Inches(0.38),size=12,c=ACCENT2,bold=True,align=PP_ALIGN.CENTER)
box(s7,DL2+Inches(0.15),DT+Inches(0.6),DW-Inches(0.3),Inches(0.42),ACCENT1)
lbl(s7,"Knop past zich aan",DL2+Inches(0.15),DT+Inches(0.6)+Inches(0.06),DW-Inches(0.3),Inches(0.32),size=11,c=WHITE,align=PP_ALIGN.CENTER)
lbl(s7,"Inhoud past altijd",DL2,DT+Inches(1.18),DW,Inches(0.38),size=11,c=ACCENT2,align=PP_ALIGN.CENTER)
pnr(s7,7,light=True)
notitie(s7,"Demo in Figma: maak een knop zonder Auto Layout en typ een langere tekst. Laat zien dat het overloopt. Zet dan Auto Layout aan en doe hetzelfde. Studenten zien direct het verschil. Vergelijk met CSS flexbox die ze later zullen gebruiken.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8: COMPONENT STATES
# ─────────────────────────────────────────────────────────────────────────────
s8=ns(light_bg())
page_hdr(s8,"Component states: interactie vastleggen","Hoe een component er uitziet in verschillende situaties")
states=[
    ("Default"  , MUTED              , "De basisstaat"        , "Hoe het component er altijd uitziet als er niets mee gebeurt"),
    ("Hover"    , ACCENT1            , "Muiscursor eroverheen", "Visuele feedback als de gebruiker naar het element wijst"),
    ("Active"   , BLUE               , "Klikken of drukken"   , "Het moment van interactie, korte visuele bevestiging"),
    ("Disabled" , (150,160,170)      , "Niet beschikbaar"     , "Verlaagd contrast, laat zien dat het element inactief is"),
]
CW=Inches(2.95); CH=Inches(3.6); CT=Inches(1.78); CL=Inches(1.0); CGAP=Inches(0.22)
KH=Inches(0.5)
for i,(name,col,short,desc) in enumerate(states):
    xl=CL+i*(CW+CGAP)
    col_kop(s8,xl,CT,CW,KH,name,col)
    box(s8,xl,CT+KH,CW,CH-KH,WHITE,border=(215,225,240))
    lbl(s8,short,xl+Inches(0.15),CT+KH+Inches(0.15),CW-Inches(0.3),Inches(0.38),size=14,c=col,bold=True)
    lbl(s8,desc,xl+Inches(0.15),CT+KH+Inches(0.55),CW-Inches(0.3),Inches(0.9),size=13,c=BODY)
    # Visual cue: a simulated button in each state
    VT=CT+KH+Inches(1.6)
    VW=CW-Inches(0.6); VH=Inches(0.5)
    VL=xl+Inches(0.3)
    if name=="Default":
        box(s8,VL,VT,VW,VH,col)
        lbl(s8,"Doneer",VL,VT+Inches(0.08),VW,Inches(0.38),size=14,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    elif name=="Hover":
        box(s8,VL,VT,VW,VH,(10,130,55))
        lbl(s8,"Doneer",VL,VT+Inches(0.08),VW,Inches(0.38),size=14,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
        gbox(s8,VL,VT,VW,Inches(0.08),(255,255,255),(255,255,255),angle=0)
    elif name=="Active":
        box(s8,VL,VT+Inches(0.05),VW,VH-Inches(0.05),col)
        lbl(s8,"Doneer",VL,VT+Inches(0.12),VW,Inches(0.38),size=14,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    else:
        box(s8,VL,VT,VW,VH,(210,215,220),border=(200,205,210))
        lbl(s8,"Doneer",VL,VT+Inches(0.08),VW,Inches(0.38),size=14,c=(150,160,170),bold=True,align=PP_ALIGN.CENTER)
    lbl(s8,f"Variant: {name.lower()}",xl+Inches(0.15),CT+KH+Inches(2.35),CW-Inches(0.3),Inches(0.35),size=12,c=MUTED,italic=True)
pnr(s8,8)
notitie(s8,"States zijn essentieel voor een goede gebruikerservaring. Gebruikers verwachten visuele feedback bij interactie. Leg uit dat in Figma states als varianten worden gedefinieerd en dat developers later de CSS :hover, :active en :disabled pseudo-classes gebruiken.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9: GOED DOEL NAV COMPONENT
# ─────────────────────────────────────────────────────────────────────────────
s9=ns(dark_bg(False))
hero_hdr(s9,"Goed doel: navigatie als component","Stap voor stap je nav bouwen in Figma")
steps=[
    ("1","Maak een frame aan van 1440x80px, geef het de naam 'nav'"),
    ("2","Voeg toe: logo-placeholder (rechthoek), 3 menu-links (tekst), CTA knop"),
    ("3","Selecteer alles en maak Auto Layout aan: horizontaal, space between"),
    ("4","Maak het frame een component (Ctrl+Alt+K)"),
    ("5","Maak een variant 'scrolled' met een lichtere achtergrondkleur"),
]
ST=Inches(1.95); SH2=Inches(0.54); SGAP=Inches(0.08)
for i,(nr,txt) in enumerate(steps):
    t=ST+i*(SH2+SGAP)
    box(s9,Inches(1.0),t,Inches(0.45),SH2,ACCENT1)
    lbl(s9,nr,Inches(1.0),t+Inches(0.06),Inches(0.45),SH2-Inches(0.06),size=17,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    box(s9,Inches(1.46),t,Inches(10.0),SH2,(25,35,60),border=(40,55,90))
    lbl(s9,txt,Inches(1.65),t+Inches(0.08),Inches(9.8),SH2-Inches(0.08),size=15,c=(210,230,255))
# Nav visual mockup
NL=Inches(1.0); NT=Inches(5.1); NW=Inches(11.5); NH=Inches(0.7)
box(s9,NL,NT,NW,NH,NAVY2,border=(50,80,160))
box(s9,NL+Inches(0.2),NT+Inches(0.15),Inches(1.0),Inches(0.4),ACCENT1)
lbl(s9,"LOGO",NL+Inches(0.2),NT+Inches(0.18),Inches(1.0),Inches(0.36),size=13,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
for j,link in enumerate(["Over ons","Projecten","Contact"]):
    lx=NL+Inches(2.0)+j*Inches(1.8)
    lbl(s9,link,lx,NT+Inches(0.18),Inches(1.6),Inches(0.36),size=13,c=(200,220,255),align=PP_ALIGN.CENTER)
box(s9,NL+NW-Inches(1.8),NT+Inches(0.12),Inches(1.5),Inches(0.46),ACCENT1)
lbl(s9,"Doneer nu",NL+NW-Inches(1.8),NT+Inches(0.16),Inches(1.5),Inches(0.38),size=13,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
# Note below mockup
lbl(s9,"Hergebruik de nav op alle paginas door instances te plaatsen",
    Inches(1.0),Inches(5.88),Inches(11.5),Inches(0.4),size=14,c=(180,200,240),italic=True)
pnr(s9,9,light=True)
notitie(s9,"Doe dit live in Figma samen met studenten. De navigatie is een perfect voorbeeld van een organism in atomic design: het bestaat uit logo (atom), links (atoms) en een knop-component (molecule). Benadruk dat de 'scrolled' variant nuttig is voor sticky navigatie.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10: GOED DOEL HERO COMPONENT
# ─────────────────────────────────────────────────────────────────────────────
s10=ns(dark_bg(False))
hero_hdr(s10,"Goed doel: hero als component","De bovenkant van je pagina als herbruikbaar onderdeel")
steps=[
    ("1","Maak een frame aan van 1440x600px (of full viewport), naam 'hero'"),
    ("2","Voeg toe: hero-afbeelding placeholder, H1 titel, subtitel, CTA knop"),
    ("3","Maak de CTA knop een instance van je knop-component"),
    ("4","Gebruik je typografie- en kleurstijlen uit les 3 en 4"),
    ("5","Maak het hele frame een component voor hergebruik over paginas"),
]
ST=Inches(1.95); SH2=Inches(0.54); SGAP=Inches(0.08)
for i,(nr,txt) in enumerate(steps):
    t=ST+i*(SH2+SGAP)
    box(s10,Inches(1.0),t,Inches(0.45),SH2,ACCENT1)
    lbl(s10,nr,Inches(1.0),t+Inches(0.06),Inches(0.45),SH2-Inches(0.06),size=17,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    box(s10,Inches(1.46),t,Inches(10.0),SH2,(25,35,60),border=(40,55,90))
    lbl(s10,txt,Inches(1.65),t+Inches(0.08),Inches(9.8),SH2-Inches(0.08),size=15,c=(210,230,255))
# Hero visual mockup
HL=Inches(1.0); HT=Inches(5.12); HW=Inches(11.5); HH=Inches(1.7)
box(s10,HL,HT,HW,HH,NAVY2,border=(50,80,160))
box(s10,HL,HT,HW,Inches(0.04),ACCENT1)
box(s10,HL,HT,Inches(3.5),HH,(30,45,80))
lbl(s10,"[ Afbeelding ]",HL+Inches(0.8),HT+Inches(0.62),Inches(2.0),Inches(0.45),size=13,c=(140,160,200),align=PP_ALIGN.CENTER)
lbl(s10,"Help ons het verschil maken",HL+Inches(3.7),HT+Inches(0.15),Inches(5.5),Inches(0.55),size=20,c=WHITE,bold=True)
lbl(s10,"Samen bouwen we aan een betere toekomst voor iedereen",HL+Inches(3.7),HT+Inches(0.72),Inches(5.5),Inches(0.45),size=13,c=(180,210,255))
box(s10,HL+Inches(3.7),HT+Inches(1.22),Inches(1.6),Inches(0.38),ACCENT1)
lbl(s10,"Doneer nu",HL+Inches(3.7),HT+Inches(1.24),Inches(1.6),Inches(0.34),size=13,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
pnr(s10,10,light=True)
notitie(s10,"De hero is het gezicht van de goed doel website. Benadruk dat studenten nu hun eerder gemaakte knop-component hergebruiken als instance. Dit is het moment waarop het design system echt gaat werken: een stijlwijziging in de master knop past automatisch de hero bij.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11: DESIGN SYSTEM OVERZICHT
# ─────────────────────────────────────────────────────────────────────────────
s11=ns(light_bg())
page_hdr(s11,"Design system: wat heb je nu","Een overzicht van je mini design system na 5 lessen")
checklist=[
    (ACCENT1, "Kleurstijlen (les 4)", "Primair, neutrals, semantisch (succes, fout, waarschuwing)"),
    (BLUE,    "Tekststijlen (les 3)", "H1 t/m caption, met juiste font, grootte en regelafstand"),
    (TEAL,    "Spacing systeem (les 2)", "8pt grid consequent toegepast op alle componenten"),
    (GREEN,   "Basiscomponenten (les 5)", "Knop, invoerveld, navigatiebalk, hero sectie"),
]
CT=Inches(1.95); CW=Inches(11.0); CH=Inches(0.88); CGAP=Inches(0.14)
for i,(ac,title,desc) in enumerate(checklist):
    t=CT+i*(CH+CGAP)
    cl,ct2,cw2=card(s11,Inches(1.0),t,CW,CH,accent=ac)
    dot(s11,Inches(1.42),t+CH/2,Inches(0.16),ac)
    lbl(s11,title,cl+Inches(0.15),ct2,cw2*0.38,CH-Inches(0.14),size=17,c=DARK,bold=True)
    lbl(s11,desc,cl+Inches(0.15)+cw2*0.38,ct2,cw2*0.58,CH-Inches(0.14),size=15,c=BODY)
# Conclusion card
cl2,ct3,cw3=card(s11,Inches(1.0),Inches(5.55),Inches(11.0),Inches(1.28),accent=ACCENT1)
lbl(s11,"Dit is je mini design system.",cl2,ct3,cw3,Inches(0.45),size=18,c=DARK,bold=True)
lbl(s11,"Professionele teams bouwen hierop door: ze voegen documentatie, do's and don'ts, en developer handoffs toe. Jouw systeem is de basis.",
    cl2,ct3+Inches(0.47),cw3,Inches(0.65),size=15,c=BODY)
pnr(s11,11)
notitie(s11,"Moment van terugblik. Studenten hebben in 5 lessen de fundamenten gelegd van een design system. Vergelijk dit met professionele systemen als Material Design - die zijn ook ooit klein begonnen. Moedig studenten aan trots te zijn op wat ze hebben opgebouwd.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12: BRONNEN
# ─────────────────────────────────────────────────────────────────────────────
s12=ns(light_bg())
page_hdr(s12,"Bronnen","Verder lezen en ontdekken over design systems en componenten")
bronnen=[
    ("Atomic Design boek",           "atomicdesign.bradfrost.com",                                        ACCENT1),
    ("Figma componenten docs",        "figma.com/best-practices/components-styles-and-shared-libraries",   BLUE),
    ("Storybook component library",   "storybook.js.org",                                                  (124,58,237)),
    ("Material Design 3",             "m3.material.io",                                                    RED),
    ("Supernova design system tool",  "supernova.io",                                                      (217,119,6)),
]
CT=Inches(1.72); CW=Inches(11.0); CH=Inches(0.72); CGAP=Inches(0.12)
for i,(title,url,ac) in enumerate(bronnen):
    t=CT+i*(CH+CGAP)
    cl,ct2,cw2=card(s12,Inches(1.0),t,CW,CH,accent=ac)
    lbl(s12,title,cl,ct2,cw2*0.42,CH-Inches(0.14),size=17,c=DARK,bold=True)
    lbl(s12,url,cl+cw2*0.42,ct2+Inches(0.04),cw2*0.56,CH-Inches(0.18),size=14,c=ac,italic=True)
pnr(s12,12)
notitie(s12,"Stuur de links na de les via de cursuspagina. Atomic Design van Brad Frost is gratis online beschikbaar. Storybook is handig om later te laten zien als studenten verder gaan met development. Material Design geeft goed inzicht in hoe grote bedrijven design systems documenteren.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13: LIVE OEFENING
# ─────────────────────────────────────────────────────────────────────────────
s13=ns(light_bg())
page_hdr(s13,"Oefening: je eerste component","Doe het zelf in Figma - 20 minuten")
stappen=[
    ("Stap 1", "Open je goed doel wireframe in Figma", ACCENT1),
    ("Stap 2", "Maak een knopcomponent: 'button/primary/large' met Auto Layout, padding 16x48px", BLUE),
    ("Stap 3", "Voeg drie varianten toe: primary, secondary, disabled", (124,58,237)),
    ("Stap 4", "Gebruik de knop als instance in je hero en navigatie", TEAL),
]
CT=Inches(1.85); CW=Inches(11.0); CH=Inches(1.0); CGAP=Inches(0.18)
for i,(step,desc,ac) in enumerate(stappen):
    t=CT+i*(CH+CGAP)
    cl,ct2,cw2=card(s13,Inches(1.0),t,CW,CH,accent=ac)
    lbl(s13,step,cl,ct2,Inches(1.2),CH-Inches(0.14),size=16,c=ac,bold=True)
    lbl(s13,desc,cl+Inches(1.25),ct2,cw2-Inches(1.25),CH-Inches(0.14),size=16,c=DARK)
# Goal card
cl2,ct3,cw3=card(s13,Inches(1.0),Inches(6.08),Inches(11.0),Inches(0.9),accent=ACCENT1)
lbl(s13,"Doel:",cl2,ct3+Inches(0.05),Inches(0.6),Inches(0.55),size=15,c=ACCENT1,bold=True)
lbl(s13,"Na deze les heb je een werkende knop-component met varianten in je goed doel project",
    cl2+Inches(0.65),ct3+Inches(0.05),cw3-Inches(0.65),Inches(0.55),size=15,c=DARK)
pnr(s13,13)
notitie(s13,"Geef studenten 20 minuten voor deze oefening. Loop rond en help individueel. Veel studenten zullen moeite hebben met de variant property namen - wijs ze op de Figma documentatie. Stap 3 (instances gebruiken) is het meest waardevol: dat is het aha-moment van design systems.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14: OPDRACHT
# ─────────────────────────────────────────────────────────────────────────────
s14=ns(dark_bg(False))
# Accent strip left
gbox(s14,Inches(0),Inches(0),Inches(0.12),SH,ACCENT1,ACCENT2,angle=270)
lbl(s14,"Opdracht  |  Les 5",Inches(0.25),Inches(0.35),Inches(11),Inches(0.55),size=16,c=(160,200,180),italic=True)
lbl(s14,"Goed Doel Website  |  Componenten",Inches(0.25),Inches(0.9),Inches(12.5),Inches(0.75),size=28,c=WHITE,bold=True)
gbox(s14,Inches(0.25),Inches(1.65),Inches(2.0),Inches(0.08),ACCENT1,ACCENT2,angle=0)
stappen=[
    ("1","Maak een knopcomponent aan met minimaal 3 varianten (primary, secondary, disabled)"),
    ("2","Maak een navigatiecomponent aan met Auto Layout (horizontaal, space between)"),
    ("3","Maak een testimonialkaart als component en plaats 3 instances naast elkaar"),
    ("4","Zorg dat alle componenten je kleur- en typografiestijlen gebruiken uit les 3 en 4"),
    ("5","Documenteer je componenten met een beschrijving in het Figma component panel"),
]
ST=Inches(1.85); SH2=Inches(0.62); SGAP=Inches(0.1)
for i,(nr,txt) in enumerate(stappen):
    t=ST+i*(SH2+SGAP)
    box(s14,Inches(0.25),t,Inches(0.45),SH2,ACCENT1)
    lbl(s14,nr,Inches(0.25),t+Inches(0.08),Inches(0.45),SH2-Inches(0.08),size=17,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    box(s14,Inches(0.71),t,Inches(11.8),SH2,(22,35,65),border=(40,60,100))
    lbl(s14,txt,Inches(0.9),t+Inches(0.1),Inches(11.6),SH2-Inches(0.1),size=15,c=(210,230,255))
pnr(s14,14,light=True)
notitie(s14,"Inleverdeadline: voor de volgende les. Controleer of studenten: (1) echte Figma components gebruiken, niet gewoon gekopieerde shapes, (2) Auto Layout correct hebben ingesteld, (3) instances gebruiken in plaats van kopieeen. Dit zijn de drie grootste valkuilen.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15: VOLGENDE LES
# ─────────────────────────────────────────────────────────────────────────────
s15=ns(dark_bg(True))
gbox(s15,Inches(0),Inches(0),Inches(0.12),SH,ACCENT1,ACCENT2,angle=270)
lbl(s15,"Volgende les",Inches(0.25),Inches(1.2),Inches(12),Inches(0.8),size=44,c=WHITE,bold=True)
gbox(s15,Inches(0.25),Inches(2.0),Inches(2.5),Inches(0.08),ACCENT1,ACCENT2,angle=0)
lbl(s15,"Project Afronding Periode 1",Inches(0.25),Inches(2.15),Inches(12),Inches(0.55),size=22,c=(180,240,200))
preview=[
    "Terugblik: wat hebben we geleerd in periode 1",
    "Goed doel website: volledige review checklist",
    "Peer review: feedback geven en ontvangen",
    "Presentatie en oplevering van het project",
]
blist(s15,preview,Inches(0.5),Inches(2.88),Inches(11.5),Inches(2.2),size=18,c=(200,230,255))
dot(s15,Inches(11.5),Inches(1.8),Inches(0.5),ACCENT1)
dot(s15,Inches(12.2),Inches(2.6),Inches(0.28),ACCENT2)
dot(s15,Inches(10.8),Inches(5.5),Inches(0.18),ACCENT1)
cl,ct2,cw2=card(s15,Inches(0.5),Inches(5.58),Inches(12.4),Inches(1.2),accent=ACCENT1,bg=(20,30,55))
lbl(s15,"Zorg dat je Figma project klaar is voor review:",cl,ct2,cw2,Inches(0.42),size=16,c=ACCENT2,bold=True)
lbl(s15,"Alle paginas, componenten en stijlen moeten aanwezig en netjes georganiseerd zijn.",
    cl,ct2+Inches(0.44),cw2,Inches(0.55),size=16,c=(200,230,255))
pnr(s15,15,light=True)
notitie(s15,"Sluit de les af met een korte samenvatting: vandaag hebben studenten geleerd wat een design system is, hoe atomic design werkt, hoe je Figma components maakt met varianten en Auto Layout, en hoe je deze toepast op het goed doel project. Volgende les is de finale van periode 1.")

# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
OUT=r"C:\Users\Peter Marcelis\Documents\Lesmateriaal\design\periode-1\les-5\05_componenten-design-systems.pptx"
prs.save(OUT)
print(f"Opgeslagen: {OUT}")
print(f"Totaal slides: {len(prs.slides)}")
