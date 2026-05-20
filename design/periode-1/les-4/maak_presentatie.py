"""
Les 4: Kleur en Contrast
Periode 1 - Ontwerp voor het Web
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image, ImageDraw
import numpy as np, io

NAVY=(6,14,44); NAVY2=(30,58,138); BLUE=(37,99,235); BLUE2=(59,130,246)
CYAN=(6,182,212); TEAL=(13,148,136); WHITE=(255,255,255); OFFWH=(248,250,252)
DARK=(15,23,42); BODY=(51,65,85); MUTED=(100,116,135)
GREEN=(22,163,74); RED=(220,38,38)
ACCENT1=(234,88,12); ACCENT2=(251,191,36)  # Orange/Amber for les 4

def rgb(c): return RGBColor(c[0],c[1],c[2])
SW,SH=Inches(13.33),Inches(7.5); BW,BH=1920,1080

def garr(w,h,c1,c2,d='diag'):
    x,y=np.linspace(0,1,w,dtype=np.float32),np.linspace(0,1,h,dtype=np.float32)
    if d=='h': t=np.tile(x,(h,1))
    elif d=='v': t=np.tile(y.reshape(-1,1),(1,w))
    else: xx,yy=np.meshgrid(x,y); t=xx*.55+yy*.45
    a=np.zeros((h,w,3),dtype=np.float32)
    for i in range(3): a[:,:,i]=c1[i]+(c2[i]-c1[i])*t
    return np.clip(a,0,255).astype(np.uint8)

def dark_bg(circ=True):
    img=Image.fromarray(garr(BW,BH,NAVY,NAVY2)).convert('RGBA')
    if circ:
        ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
        d.ellipse([BW-500,-180,BW+200,520],fill=(ACCENT1[0],ACCENT1[1],ACCENT1[2],20))
        d.ellipse([-60,BH-360,380,BH+100],fill=(ACCENT2[0],ACCENT2[1],ACCENT2[2],22))
        img=Image.alpha_composite(img,ov)
    return img.convert('RGB')

def light_bg():
    img=Image.fromarray(garr(BW,BH,OFFWH,WHITE)).convert('RGBA')
    ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
    d.ellipse([BW-580,-220,BW+180,540],fill=(ACCENT1[0],ACCENT1[1],ACCENT1[2],8))
    img=Image.alpha_composite(img,ov); return img.convert('RGB')

def buf(img):
    b=io.BytesIO(); img.save(b,format='PNG'); b.seek(0); return b

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
BLANK=prs.slide_layouts[6]

def ns(bg=None):
    s=prs.slides.add_slide(BLANK)
    if bg:
        p=s.shapes.add_picture(buf(bg),0,0,SW,SH)
        t=s.shapes._spTree; t.remove(p._element); t.insert(2,p._element)
    else:
        s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(WHITE)
    return s

def box(s,l,t,w,h,c,border=None):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c)
    if border: sh.line.color.rgb=rgb(border); sh.line.width=Pt(0.75)
    else: sh.line.fill.background()
    return sh

def gbox(s,l,t,w,h,c1,c2,angle=90):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.gradient(); st=sh.fill.gradient_stops
    st[0].position=0.0; st[0].color.rgb=rgb(c1)
    st[1].position=1.0; st[1].color.rgb=rgb(c2)
    sh.fill.gradient_angle=angle; sh.line.fill.background(); return sh

def dot(s,cx,cy,r,c):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,cx-r,cy-r,r*2,r*2)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c); sh.line.fill.background()

def lbl(s,text,l,t,w,h,size=18,c=BODY,bold=False,italic=False,align=PP_ALIGN.LEFT):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=text; r.font.name="Calibri"; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=rgb(c)

def blist(s,items,l,t,w,h,size=18,c=BODY,gap=7):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(gap); r=p.add_run()
        r.text=f"  {chr(8226)}  {item}"; r.font.name="Calibri"; r.font.size=Pt(size); r.font.color.rgb=rgb(c)

def notitie(s,t): s.notes_slide.notes_text_frame.text=t

def pnr(s,n,tot=15,light=False):
    lbl(s,f"{n} / {tot}",Inches(12.5),Inches(7.1),Inches(0.8),Inches(0.35),size=11,c=(180,180,200) if light else MUTED,align=PP_ALIGN.RIGHT)

def hero_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(1.0),size=36,c=WHITE,bold=True)
    gbox(s,Inches(1.0),Inches(1.42),Inches(2.0),Inches(0.08),ACCENT1,ACCENT2,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.55),Inches(11.3),Inches(0.42),size=15,c=(240,210,180),italic=True)

def page_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(0.95),size=34,c=DARK,bold=True)
    gbox(s,Inches(1.0),Inches(1.38),Inches(1.8),Inches(0.08),ACCENT1,ACCENT2,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.5),Inches(11.3),Inches(0.38),size=14,c=MUTED,italic=True)

def card(s,l,t,w,h,accent=None,bg=OFFWH):
    a=accent if accent else ACCENT1
    STRIP=Inches(0.07)
    box(s,l,t,w,h,bg,border=(215,225,240))
    gbox(s,l,t,STRIP,h,a,ACCENT2,angle=270)
    return l+STRIP+Inches(0.18), t+Inches(0.14), w-STRIP-Inches(0.35)

def col_kop(s,l,t,w,h,text,c):
    box(s,l,t,w,h,c)
    lbl(s,text,l,t+Inches(0.07),w,h-Inches(0.07),size=20,c=WHITE,bold=True,align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1: TITEL
# ─────────────────────────────────────────────────────────────────────────────
s1=ns(dark_bg(True))
# Left accent strip
gbox(s1,Inches(0.0),Inches(0.0),Inches(0.18),Inches(7.5),ACCENT1,ACCENT2,angle=270)
# Title
lbl(s1,"Kleur\nen Contrast",Inches(1.0),Inches(1.2),Inches(8.5),Inches(2.5),size=56,c=WHITE,bold=True)
# Gradient underline
gbox(s1,Inches(1.0),Inches(3.65),Inches(3.0),Inches(0.1),ACCENT1,ACCENT2,angle=0)
# Periode label
lbl(s1,"Periode 1 . Les 4",Inches(1.0),Inches(3.85),Inches(8.0),Inches(0.5),size=18,c=(200,200,230),bold=False)
# Tagline
lbl(s1,"Hoe kleur emotie en hierarchie stuurt",Inches(1.0),Inches(4.45),Inches(9.0),Inches(0.6),size=20,c=(240,210,180),italic=True)
# Decorative dots (bottom right corner, subtle)
dot(s1,Inches(11.8),Inches(6.0),Inches(0.18),(ACCENT1[0],ACCENT1[1],ACCENT1[2]))
dot(s1,Inches(12.4),Inches(5.6),Inches(0.12),(ACCENT2[0],ACCENT2[1],ACCENT2[2]))
pnr(s1,1,light=True)
notitie(s1,"Welkom bij les 4. Kleur is een van de krachtigste tools in design. Vandaag leren studenten hoe kleurtheorie werkt en hoe ze dat toepassen in hun goed doel project.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2: AGENDA
# ─────────────────────────────────────────────────────────────────────────────
s2=ns(light_bg())
page_hdr(s2,"Vandaag","Les 4: Kleur en Contrast")

agenda_items=[
    ("Kleurmodel: RGB, HSL en het kleurwiel", BLUE),
    ("Kleurharmonie: welke combinaties werken", (124,58,237)),
    ("WCAG contrast ratios voor toegankelijkheid", GREEN),
    ("Semantische kleuren: success, error, warning", RED),
    ("Kleurpaletten opzetten met shades", ACCENT1),
    ("Figma: kleurstijlen instellen voor het project", TEAL),
]

col1_items = agenda_items[:3]
col2_items = agenda_items[3:]

for idx,(item,col) in enumerate(col1_items):
    ty=Inches(2.0)+idx*Inches(1.45)
    dot(s2,Inches(1.35),ty+Inches(0.22),Inches(0.16),col)
    lbl(s2,f"{idx+1}. {item}",Inches(1.65),ty,Inches(4.8),Inches(1.2),size=17,c=DARK)

for idx,(item,col) in enumerate(col2_items):
    ty=Inches(2.0)+idx*Inches(1.45)
    dot(s2,Inches(7.15),ty+Inches(0.22),Inches(0.16),col)
    lbl(s2,f"{idx+4}. {item}",Inches(7.45),ty,Inches(4.8),Inches(1.2),size=17,c=DARK)

pnr(s2,2)
notitie(s2,"Loop de agenda kort door. Kleur is een breed onderwerp - vandaag richten we ons op praktische toepassing voor het goed doel project.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3: HET KLEURWIEL EN HSL
# ─────────────────────────────────────────────────────────────────────────────
s3=ns(dark_bg(False))
# Left dark panel
box(s3,Inches(0.0),Inches(0.0),Inches(6.0),Inches(7.5),(12,20,56))

# Title spanning full slide
lbl(s3,"Het kleurwiel en HSL",Inches(0.4),Inches(0.4),Inches(12.0),Inches(0.8),size=32,c=WHITE,bold=True)
gbox(s3,Inches(0.4),Inches(1.18),Inches(2.0),Inches(0.07),ACCENT1,ACCENT2,angle=0)

# Colorful dots arranged in a circular pattern on the left panel (color wheel simulation)
cx=Inches(2.9); cy=Inches(4.2); base_r=Inches(1.55)
hue_colors=[
    (220,38,38),    # Red
    (234,88,12),    # Orange
    (234,179,8),    # Yellow
    (22,163,74),    # Green
    (37,99,235),    # Blue
    (124,58,237),   # Purple
]
import math
for i,hc in enumerate(hue_colors):
    angle=math.radians(i*60 - 90)
    px=cx+base_r*math.cos(angle)
    py=cy+base_r*math.sin(angle)
    dot(s3,px,py,Inches(0.38),hc)

# Inner ring - lighter variants
inner_colors=[
    (252,165,165),  # Light red
    (253,186,116),  # Light orange
    (253,224,71),   # Light yellow
    (134,239,172),  # Light green
    (147,197,253),  # Light blue
    (196,181,253),  # Light purple
]
inner_r=Inches(0.78)
for i,hc in enumerate(inner_colors):
    angle=math.radians(i*60 - 90)
    px=cx+inner_r*math.cos(angle)
    py=cy+inner_r*math.sin(angle)
    dot(s3,px,py,Inches(0.22),hc)

# Center white dot
dot(s3,cx,cy,Inches(0.2),WHITE)

# Label under wheel
lbl(s3,"Het kleurwiel",Inches(0.4),Inches(6.7),Inches(5.0),Inches(0.5),size=13,c=(180,190,220),italic=True,align=PP_ALIGN.CENTER)

# Right panel: HSL explanation
lbl(s3,"Het HSL model",Inches(6.3),Inches(1.5),Inches(6.5),Inches(0.7),size=26,c=DARK,bold=True)

hsl_items=[
    ("H = Hue (tint)", "0 tot 360 graden op het kleurwiel", ACCENT1),
    ("S = Saturation (verzadiging)", "0% grijs tot 100% volledig vol", GREEN),
    ("L = Lightness (helderheid)", "0% zwart tot 100% wit", BLUE),
]
for idx,(title,desc,col) in enumerate(hsl_items):
    ty=Inches(2.35)+idx*Inches(1.3)
    dot(s3,Inches(6.55),ty+Inches(0.2),Inches(0.14),col)
    lbl(s3,title,Inches(6.85),ty,Inches(5.8),Inches(0.45),size=17,c=DARK,bold=True)
    lbl(s3,desc,Inches(6.85),ty+Inches(0.42),Inches(5.8),Inches(0.5),size=14,c=BODY)

# Tip box
tip_l=Inches(6.3); tip_t=Inches(6.0); tip_w=Inches(6.65); tip_h=Inches(1.1)
box(s3,tip_l,tip_t,tip_w,tip_h,(234,88,12))
lbl(s3,"Tip:",tip_l+Inches(0.18),tip_t+Inches(0.1),Inches(0.6),Inches(0.4),size=14,c=WHITE,bold=True)
lbl(s3,"Gebruik HSL in CSS en Figma voor consistentere paletten dan met RGB",
    tip_l+Inches(0.18),tip_t+Inches(0.38),tip_w-Inches(0.36),Inches(0.58),size=13,c=WHITE,italic=True)

pnr(s3,3,light=True)
notitie(s3,"HSL is intuItiever dan RGB voor designers. Uitleggen: je kunt gewoon de L aanpassen voor lichtere/donkerdere varianten, of de S voor minder/meer levendige kleuren. Laat zien in Figma hoe de HSL sliders werken.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4: KLEURHARMONIE
# ─────────────────────────────────────────────────────────────────────────────
s4=ns(light_bg())
page_hdr(s4,"Kleurharmonie: combinaties die werken")

col_data=[
    ("Aanvullend", RED,
     ["Tegenovergestelde kleuren op het wiel","Hoog contrast, krachtig en energiek","Goed voor accenten en call-to-action"],
     "Bv. rood + groen, blauw + oranje"),
    ("Analoog", ACCENT1,
     ["Kleuren naast elkaar op het wiel","Harmonieus en kalm gevoel","Populair in natuur-gerelateerde merken"],
     "Bv. geel + oranje + rood"),
    ("Triadisch", BLUE,
     ["Drie kleuren op gelijke afstanden","Levendig maar evenwichtig","Werkt goed voor speelse merken"],
     "Bv. rood + blauw + geel"),
]

col_w=Inches(3.85); col_gap=Inches(0.22); col_start=Inches(0.82)
for ci,(name,col,bullets,example) in enumerate(col_data):
    cx=col_start+ci*(col_w+col_gap)
    col_kop(s4,cx,Inches(1.75),col_w,Inches(0.52),name,col)
    # Body card
    box(s4,cx,Inches(2.28),col_w,Inches(3.6),WHITE,border=(215,225,240))
    blist(s4,bullets,cx+Inches(0.12),Inches(2.38),col_w-Inches(0.18),Inches(2.5),size=14,c=BODY,gap=8)
    lbl(s4,example,cx+Inches(0.12),Inches(4.72),col_w-Inches(0.18),Inches(0.55),size=12,c=MUTED,italic=True)

# Bottom tip
tip_l=Inches(0.82); tip_t=Inches(6.05); tip_w=Inches(11.7); tip_h=Inches(0.92)
box(s4,tip_l,tip_t,tip_w,tip_h,(255,247,237),border=(253,186,116))
lbl(s4,"60/30/10 regel:",tip_l+Inches(0.18),tip_t+Inches(0.1),Inches(2.2),Inches(0.4),size=14,c=ACCENT1,bold=True)
lbl(s4,"60% basiskleur, 30% aanvullende kleur, 10% accent. Dit zorgt voor balans zonder saai te zijn.",
    tip_l+Inches(2.3),tip_t+Inches(0.1),tip_w-Inches(2.5),Inches(0.6),size=14,c=BODY)

pnr(s4,4)
notitie(s4,"Toon voorbeelden van bekende merken bij elke harmonievorm. Aanvullend: Netflix (rood+zwart). Analoog: veel natuur-apps. Triadisch: kindgerichte merken. De 60/30/10 regel is een praktische vuistregel voor webdesign.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5: WCAG CONTRAST
# ─────────────────────────────────────────────────────────────────────────────
s5=ns(dark_bg(False))
hero_hdr(s5,"WCAG contrast ratios","Toegankelijkheid begint bij leesbare tekst")

# Two columns: AA and AAA
col_w=Inches(4.8); left_l=Inches(0.9); right_l=Inches(6.45)

# AA column
box(s5,left_l,Inches(2.1),col_w,Inches(3.0),(20,40,90))
lbl(s5,"WCAG AA (minimaal)",left_l+Inches(0.2),Inches(2.2),col_w-Inches(0.3),Inches(0.5),size=16,c=(180,200,255),bold=True)
aa_items=["Normale tekst: 4.5:1 ratio","Grote tekst (18px+): 3:1 ratio","UI componenten: 3:1 ratio"]
blist(s5,aa_items,left_l+Inches(0.15),Inches(2.75),col_w-Inches(0.3),Inches(2.0),size=15,c=WHITE,gap=10)

# AAA column
box(s5,right_l,Inches(2.1),col_w,Inches(3.0),(20,40,90))
lbl(s5,"WCAG AAA (optimaal)",right_l+Inches(0.2),Inches(2.2),col_w-Inches(0.3),Inches(0.5),size=16,c=(180,255,200),bold=True)
aaa_items=["Normale tekst: 7:1 ratio","Grote tekst: 4.5:1 ratio","Streef hiernaar voor body text"]
blist(s5,aaa_items,right_l+Inches(0.15),Inches(2.75),col_w-Inches(0.3),Inches(2.0),size=15,c=WHITE,gap=10)

# Three contrast demo boxes stacked (right side of slide, compact)
demo_l=Inches(11.0); demo_configs=[
    ((15,23,42),(248,250,252),"Pass AAA: 15.8:1"),
    ((51,65,85),(220,230,245),"Pass AA: 4.8:1"),
    ((160,170,185),(210,218,228),"Fail: 2.1:1"),
]
demo_start_t=Inches(2.1)
for di,(tc,bc,label) in enumerate(demo_configs):
    dt=demo_start_t+di*Inches(1.0)
    box(s5,demo_l,dt,Inches(2.15),Inches(0.88),bc)
    lbl(s5,label,demo_l+Inches(0.1),dt+Inches(0.18),Inches(1.95),Inches(0.5),size=13,c=tc,bold=True)

pnr(s5,5,light=True)
notitie(s5,"WCAG staat voor Web Content Accessibility Guidelines. Leg uit dat AA verplicht is voor overheidswebsites in Nederland en dat het goed doel project minimaal AA moet halen. Gebruik webaim.org/resources/contrastchecker live in de les.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6: SEMANTISCHE KLEUREN
# ─────────────────────────────────────────────────────────────────────────────
s6=ns(light_bg())
page_hdr(s6,"Semantische kleuren","Kleur als communicatiemiddel")

sem_colors=[
    ("Success", GREEN, "Bevestiging, opgeslagen, klaar", "Formulier succesvol verstuurd"),
    ("Error", RED, "Fouten, validatie, blokkades", "Vul een geldig e-mailadres in"),
    ("Warning", (217,119,6), "Waarschuwing, let op, tijdelijk", "Je sessie verloopt over 5 minuten"),
    ("Info", BLUE, "Informatie, tips, neutraal", "Er is een nieuwe versie beschikbaar"),
]

grid_l=Inches(1.0); grid_t=Inches(1.75)
cw=Inches(5.7); ch=Inches(2.3); gap=Inches(0.22)

for idx,(label,col,desc,example) in enumerate(sem_colors):
    ci=idx%2; ri=idx//2
    cl=grid_l+ci*(cw+gap)
    ct=grid_t+ri*(ch+gap)
    inner_l,inner_t,inner_w=card(s6,cl,ct,cw,ch,accent=col)
    dot(s6,cl+Inches(0.38),ct+Inches(0.44),Inches(0.22),col)
    lbl(s6,label,cl+Inches(0.72),ct+Inches(0.2),cw-Inches(0.9),Inches(0.45),size=18,c=DARK,bold=True)
    lbl(s6,desc,inner_l,ct+Inches(0.7),inner_w,Inches(0.45),size=13,c=BODY)
    # Example text box
    box(s6,inner_l,ct+Inches(1.22),inner_w,Inches(0.62),(col[0]//6,col[1]//6,col[2]//6) if sum(col)<200 else (col[0]//5+200,col[1]//5+200,col[2]//5+200),border=col)
    lbl(s6,f'"{example}"',inner_l+Inches(0.1),ct+Inches(1.3),inner_w-Inches(0.15),Inches(0.45),size=12,c=col,italic=True)

pnr(s6,6)
notitie(s6,"Semantische kleuren zijn kleuren met een vaste betekenis. Studenten kennen dit al van verkeerslichten. Op websites gebruiken we deze conventie consistent zodat gebruikers weten wat ze kunnen verwachten. Laat zien in voorbeelden van bekende websites.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7: KLEURPALET OPZETTEN
# ─────────────────────────────────────────────────────────────────────────────
s7=ns(light_bg())
page_hdr(s7,"Kleurpalet: van 50 tot 900","Shades voor elk gebruik")

# Shade strip - 9 boxes from light to dark (orange/amber palette)
shade_values=[
    (50,  (255,240,220)),
    (100, (254,220,170)),
    (200, (253,186,116)),
    (300, (251,146,60)),
    (400, (249,115,22)),
    (500, (234,88,12)),
    (600, (194,65,12)),
    (700, (154,52,18)),
    (800, (124,42,12)),
    (900, (100,35,5)),
]

box_w=Inches(1.2); box_h=Inches(0.9); strip_t=Inches(2.35); strip_start=Inches(0.65)
for idx,(shade,color) in enumerate(shade_values):
    bx=strip_start+idx*(box_w+Inches(0.02))
    box(s7,bx,strip_t,box_w,box_h,color)
    tc=WHITE if shade>=500 else DARK
    lbl(s7,str(shade),bx,strip_t+Inches(0.58),box_w,Inches(0.3),size=11,c=tc,bold=True,align=PP_ALIGN.CENTER)

# Pointer to 500
lbl(s7,"Basis",strip_start+5*(box_w+Inches(0.02)),strip_t+box_h+Inches(0.05),box_w,Inches(0.35),size=11,c=ACCENT1,bold=True,align=PP_ALIGN.CENTER)

# Bullets
blist(s7,[
    "500 is de basiskleur, gebruik dit als referentie",
    "50 en 100 zijn goed als achtergronden en hover states",
    "700 en 800 zijn goed voor hover states van knoppen",
    "Gebruik Tailwind CSS of Radix Colors als inspiratie voor shades",
],Inches(1.0),Inches(3.65),Inches(11.3),Inches(2.2),size=16,c=BODY,gap=10)

# Tip card
tl=Inches(1.0); tt=Inches(5.95); tw=Inches(11.3); th=Inches(1.08)
cx2,ct2,cw2=card(s7,tl,tt,tw,th,accent=TEAL)
lbl(s7,"Probeer het:",cx2,tt+Inches(0.08),Inches(2.0),Inches(0.42),size=14,c=TEAL,bold=True)
lbl(s7,"Ga naar tailwindcss.com/docs/customizing-colors voor een volledig shade-systeem als referentie.",
    cx2,tt+Inches(0.5),cw2,Inches(0.5),size=13,c=BODY,italic=True)

pnr(s7,7)
notitie(s7,"Het shade systeem (50-900) komt uit Tailwind CSS en is de industriestandaard geworden. Studenten hoeven niet alle shades zelf te bedenken, ze kunnen dit als startpunt nemen en aanpassen. De 500 waarde is altijd de 'pure' merkkleur.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8: DONKERE KLEUREN EN NEUTRALS
# ─────────────────────────────────────────────────────────────────────────────
s8=ns(dark_bg(False))
hero_hdr(s8,"Neutrals: de onzichtbare helden","De ruggengraat van elk goed kleurpalet")

blist(s8,[
    "Neutrals vormen 60 procent of meer van een goed kleurpalet",
    "Kies warme grijzen (licht bruinig) of koele grijzen (licht blauwachtig)",
    "Meng je neutral altijd licht met je merkkleur voor cohesie",
    "Gebruik nooit zuiver zwart (#000000) op webpaginas: te hard contrast",
    "Donkere modus: verwissel niet gewoon licht en donker, bouw apart",
],Inches(1.0),Inches(2.1),Inches(11.3),Inches(2.8),size=17,c=WHITE,gap=11)

# Neutral scale demo
neutral_scale=[
    ("surface", (248,250,252)),
    ("border",  (203,213,225)),
    ("muted",   (148,163,184)),
    ("body",    (71,85,105)),
    ("heading", (15,23,42)),
]

demo_start=Inches(1.0); demo_t=Inches(5.3); ndemo_w=Inches(2.2); ndemo_h=Inches(0.78)
for ni,(label,nc) in enumerate(neutral_scale):
    nx=demo_start+ni*(ndemo_w+Inches(0.08))
    box(s8,nx,demo_t,ndemo_w,ndemo_h,nc,border=(100,115,130))
    tc=DARK if nc[0]>180 else WHITE
    lbl(s8,label,nx,demo_t+Inches(0.52),ndemo_w,Inches(0.3),size=11,c=tc,bold=False,align=PP_ALIGN.CENTER)

pnr(s8,8,light=True)
notitie(s8,"Neutrals worden vaak onderschat maar ze zijn net zo belangrijk als de merkkleur. Een warme neutral (licht bruin/beige tinten) past goed bij oranje of rode merkkleuren. Een koele neutral (licht blauwachtig) past goed bij blauw. Dit maakt het ontwerp samenhangend.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9: FIGMA KLEURSTIJLEN
# ─────────────────────────────────────────────────────────────────────────────
s9=ns(dark_bg(False))
hero_hdr(s9,"Figma: kleurstijlen instellen","Eenmalig instellen, overal hergebruiken")

steps=[
    ("1","Klik op een element met de gewenste kleur in Figma"),
    ("2","Klik op het vierkantige icoontje naast de kleurwaarde"),
    ("3","Klik op het plus-teken om een nieuwe stijl aan te maken"),
    ("4","Gebruik een naamstructuur: 'brand/primary', 'semantic/success', 'neutral/100'"),
    ("5","Hergebruik stijlen door op het stijlenicoontje te klikken bij elk element"),
]
for idx,(num,step) in enumerate(steps):
    ty=Inches(2.1)+idx*Inches(0.86)
    dot(s9,Inches(1.22),ty+Inches(0.2),Inches(0.22),ACCENT1)
    lbl(s9,num,Inches(1.12),ty+Inches(0.05),Inches(0.25),Inches(0.35),size=13,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s9,step,Inches(1.65),ty,Inches(10.5),Inches(0.72),size=16,c=WHITE)

# Tip card
tl=Inches(1.0); tt=Inches(6.28); tw=Inches(11.3); th=Inches(0.85)
cx2,ct2,cw2=card(s9,tl,tt,tw,th,accent=ACCENT2,bg=(20,35,80))
lbl(s9,"Minimaal voor het project:",cx2,tt+Inches(0.08),Inches(3.0),Inches(0.42),size=14,c=ACCENT2,bold=True)
lbl(s9,"primary (en shades), neutrals (5-7 stappen), semantic (success/error/warning)",
    cx2,tt+Inches(0.46),cw2,Inches(0.35),size=13,c=(200,215,240),italic=True)

pnr(s9,9,light=True)
notitie(s9,"Laat dit live zien in Figma. Open een lege file, maak een rechthoek, ga naar Fill, klik het vierkantje en maak een stijl brand/primary. Wissel dan naar een andere rechthoek en laat zien hoe je dezelfde stijl kunt hergebruiken. Als je de stijl aanpast, update alles tegelijk.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10: GOED DOEL KLEURPALET
# ─────────────────────────────────────────────────────────────────────────────
s10=ns(dark_bg(False))
hero_hdr(s10,"Goed doel: kleurpalet kiezen","Kleur die past bij het doel en de doelgroep")

blist(s10,[
    "Denk aan de emotie die het goede doel uitstraalt: warmte, hoop, urgentie?",
    "Kies een primaire kleur die past bij het thema van het goede doel",
    "Voeg een neutrale kleur toe (grijs) voor tekst en achtergronden",
    "Kies een accentkleur voor CTA knoppen en highlights",
    "Controleer direct of je kleur contrast WCAG AA haalt",
],Inches(1.0),Inches(2.1),Inches(8.5),Inches(2.8),size=16,c=WHITE,gap=10)

# Three example dots with labels
ex_colors=[(ACCENT1,"Primair"),(MUTED,"Neutraal"),(ACCENT2,"Accent")]
dot_cx=Inches(10.0)
for di,(col,lname) in enumerate(ex_colors):
    dcy=Inches(2.5)+di*Inches(1.38)
    dot(s10,dot_cx,dcy,Inches(0.38),col)
    lbl(s10,lname,dot_cx+Inches(0.52),dcy-Inches(0.18),Inches(2.0),Inches(0.5),size=15,c=WHITE)

# Tool links card
tl=Inches(1.0); tt=Inches(5.55); tw=Inches(11.3); th=Inches(1.55)
cx2,ct2,cw2=card(s10,tl,tt,tw,th,accent=BLUE,bg=(12,25,65))
lbl(s10,"Handige tools:",cx2,tt+Inches(0.08),Inches(2.5),Inches(0.42),size=15,c=BLUE2,bold=True)
lbl(s10,"coolors.co voor palette generatie en inspiratie",cx2,tt+Inches(0.5),cw2,Inches(0.38),size=14,c=WHITE)
lbl(s10,"webaim.org/resources/contrastchecker voor contrast check",cx2,tt+Inches(0.88),cw2,Inches(0.38),size=14,c=WHITE)
lbl(s10,"Probeer je primaire kleur nu al te kiezen en te controleren.",cx2,tt+Inches(1.18),cw2,Inches(0.32),size=13,c=(200,215,240),italic=True)

pnr(s10,10,light=True)
notitie(s10,"Dit is het moment waarop studenten al moeten nadenken over hun specifieke goed doel. Vraag een paar studenten wat hun goed doel is en bespreek samen welke kleur daarbij past. Laat studenten coolors.co openen op hun eigen laptop.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11: LIVE CONTRAST CHECK
# ─────────────────────────────────────────────────────────────────────────────
s11=ns(light_bg())
page_hdr(s11,"Live analyse: contrast checken","Opdrachten voor nu in de les")

opdrachten=[
    (BLUE,"Ga naar webaim.org/resources/contrastchecker. Voer je primaire kleur en wit in. Haal je WCAG AA?"),
    (GREEN,"Open een bestaande website en check de bodytekst kleur met de browser inspector. Wat is de hex waarde?"),
    (ACCENT1,"Kijk naar de CTA knop op een website. Is de knoptekst goed leesbaar, en wat is het contrast ratio?"),
    (RED,"Vergelijk twee sites: een met veel kleur en een minimalistische. Wat voelt professioneler en waarom?"),
]

for idx,(col,tekst) in enumerate(opdrachten):
    ty=Inches(1.85)+idx*Inches(1.28)
    cl,ct,cw2=card(s11,Inches(0.9),ty,Inches(11.55),Inches(1.14),accent=col)
    dot(s11,Inches(1.12),ty+Inches(0.44),Inches(0.2),col)
    lbl(s11,f"Opdracht {idx+1}",Inches(1.45),ty+Inches(0.04),Inches(2.5),Inches(0.38),size=13,c=col,bold=True)
    lbl(s11,tekst,Inches(1.45),ty+Inches(0.42),cw2-Inches(0.4),Inches(0.6),size=14,c=BODY)

pnr(s11,11)
notitie(s11,"Geef studenten 8-10 minuten voor deze opdrachten. Loop rond en help waar nodig. Bespreek daarna de bevindingen klassikaal. Vraag wie er een contrast ratio gevonden heeft die niet AA haalt - dat komt vaker voor dan je denkt.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12: KLEUR EN EMOTIE
# ─────────────────────────────────────────────────────────────────────────────
s12=ns(light_bg())
page_hdr(s12,"Kleur en emotie","Wat kleuren communiceren")

kleur_data=[
    ("Rood",   (220,38,38),  "Urgentie, passie, gevaar, energie"),
    ("Oranje", (234,88,12),  "Vriendelijkheid, enthousiasme, creativiteit"),
    ("Geel",   (234,179,8),  "Optimisme, aandacht, waarschuwing"),
    ("Groen",  (22,163,74),  "Groei, gezondheid, natuur, succes"),
    ("Blauw",  (37,99,235),  "Vertrouwen, rust, professionaliteit"),
    ("Paars",  (124,58,237), "Luxe, creativiteit, wijsheid"),
]

# Table headers
th_w=Inches(2.5); ta_w=Inches(8.1)
tl=Inches(0.9); tt=Inches(1.78)
col_kop(s12,tl,tt,th_w,Inches(0.5),"Kleur",DARK)
col_kop(s12,tl+th_w+Inches(0.08),tt,ta_w,Inches(0.5),"Associaties",DARK)

for ri,(name,col,assoc) in enumerate(kleur_data):
    ry=tt+Inches(0.5)+ri*Inches(0.78)
    row_bg=WHITE if ri%2==0 else OFFWH
    box(s12,tl,ry,th_w,Inches(0.78),row_bg,border=(215,225,240))
    dot(s12,tl+Inches(0.3),ry+Inches(0.35),Inches(0.2),col)
    lbl(s12,name,tl+Inches(0.6),ry+Inches(0.12),th_w-Inches(0.7),Inches(0.5),size=15,c=DARK,bold=True)
    box(s12,tl+th_w+Inches(0.08),ry,ta_w,Inches(0.78),row_bg,border=(215,225,240))
    lbl(s12,assoc,tl+th_w+Inches(0.22),ry+Inches(0.18),ta_w-Inches(0.3),Inches(0.46),size=14,c=BODY)

pnr(s12,12)
notitie(s12,"Kleurpsychologie is cultureel bepaald, maar in westerse context zijn deze associaties vrij consistent. Vraag studenten: welke kleur past bij jullie goed doel? Een voedselbank heeft andere behoeften dan een kinderhospitaal of een milieuorganisatie.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13: BRONNEN
# ─────────────────────────────────────────────────────────────────────────────
s13=ns(light_bg())
page_hdr(s13,"Bronnen en verdieping","Tools en referenties voor kleur en contrast")

bronnen=[
    ("Coolors kleurenpalet generator","coolors.co","Genereer en sla kleurenpaletten op",ACCENT1),
    ("WCAG contrast checker","webaim.org/resources/contrastchecker","Check contrast ratios direct online",GREEN),
    ("Tailwind CSS kleurenpalet","tailwindcss.com/docs/customizing-colors","Professioneel shade-systeem als referentie",BLUE),
    ("Radix Colors semantische kleuren","radix-ui.com/colors","Toegankelijk kleurensysteem met semantic support",(124,58,237)),
    ("Color Hunt voor inspiratie","colorhunt.co","Duizenden gecureerde kleurenpaletten",RED),
]

for bi,(title,url,desc,col) in enumerate(bronnen):
    bt=Inches(1.78)+bi*Inches(1.02)
    cl,ct,cw2=card(s13,Inches(0.9),bt,Inches(11.55),Inches(0.88),accent=col)
    lbl(s13,title,cl,ct+Inches(0.02),Inches(4.0),Inches(0.4),size=15,c=DARK,bold=True)
    lbl(s13,url,Inches(4.95),ct+Inches(0.02),Inches(3.2),Inches(0.4),size=13,c=col,italic=True)
    lbl(s13,desc,Inches(8.2),ct+Inches(0.02),cw2-Inches(4.0),Inches(0.55),size=13,c=MUTED)

pnr(s13,13)
notitie(s13,"Deel deze links ook in de Brightspace pagina voor les 4 zodat studenten thuis kunnen oefenen. Coolors en ColorHunt zijn goede starters voor studenten die moeite hebben met kleurkeuzes.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14: OPDRACHT
# ─────────────────────────────────────────────────────────────────────────────
s14=ns(dark_bg(False))
hero_hdr(s14,"Opdracht . Les 4","Goed Doel Website . Kleurpalet")

steps_opdracht=[
    ("1","Kies een primaire kleur die past bij het thema van jouw goed doel"),
    ("2","Bouw een kleurpalet met minimaal 5 tinten (50 t/m 900 schaal)"),
    ("3","Voeg neutrals toe: minimaal 5 grijstinten voor tekst en achtergronden"),
    ("4","Kies een accentkleur voor CTA knoppen en highlights"),
    ("5","Check alle tekst-achtergrond combinaties op WCAG AA contrast"),
    ("6","Sla alle kleuren op als kleurstijlen in Figma (brand/primary, neutral/x, semantic/x)"),
]

for idx,(num,step) in enumerate(steps_opdracht):
    ty=Inches(2.1)+idx*Inches(0.82)
    dot(s14,Inches(1.22),ty+Inches(0.2),Inches(0.22),ACCENT1)
    lbl(s14,num,Inches(1.12),ty+Inches(0.04),Inches(0.25),Inches(0.35),size=13,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s14,step,Inches(1.65),ty,Inches(10.5),Inches(0.72),size=16,c=WHITE)

# Deadline/deliverable hint
box(s14,Inches(1.0),Inches(7.05),Inches(11.3),Inches(0.32),(ACCENT1[0],ACCENT1[1],ACCENT1[2]))
lbl(s14,"Inleveren: Figma file met kleurstijlen klaar voor les 5, gedeeld via de Brightspace opdracht",
    Inches(1.15),Inches(7.08),Inches(11.0),Inches(0.28),size=12,c=WHITE,bold=False)

pnr(s14,14,light=True)
notitie(s14,"Dit is de huiswerkopdracht. Leg nadruk op stap 5: het contrast checken. Veel studenten vergeten dit en kiezen kleuren die niet voldoen. Stap 6 (Figma stijlen) is ook nieuw - loop hier kort doorheen aan het einde van de les.")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15: VOLGENDE LES
# ─────────────────────────────────────────────────────────────────────────────
s15=ns(dark_bg(True))
# Left accent strip
gbox(s15,Inches(0.0),Inches(0.0),Inches(0.18),Inches(7.5),ACCENT1,ACCENT2,angle=270)

lbl(s15,"Volgende les",Inches(1.0),Inches(1.0),Inches(11.0),Inches(0.65),size=20,c=(200,210,240))
lbl(s15,"Componenten\nen Design Systems",Inches(1.0),Inches(1.6),Inches(11.0),Inches(2.0),size=46,c=WHITE,bold=True)
gbox(s15,Inches(1.0),Inches(3.52),Inches(3.0),Inches(0.1),ACCENT1,ACCENT2,angle=0)

preview_items=[
    "Atomic design: atoms, molecules en organisms",
    "Componenten aanmaken in Figma",
    "Varianten en Auto Layout",
    "Hero en nav als herbruikbaar component bouwen",
]
blist(s15,preview_items,Inches(1.0),Inches(3.75),Inches(11.0),Inches(2.2),size=17,c=(220,225,250),gap=11)

# Bottom tip card
tl=Inches(1.0); tt=Inches(6.3); tw=Inches(11.3); th=Inches(0.85)
box(s15,tl,tt,tw,th,(30,50,120),border=(80,100,180))
lbl(s15,"Voorbereiding:",tl+Inches(0.18),tt+Inches(0.1),Inches(2.0),Inches(0.42),size=13,c=ACCENT2,bold=True)
lbl(s15,"Bekijk alvast bradfordzz.com/atomic-design voor de theorie achter componenten",
    tl+Inches(2.1),tt+Inches(0.12),tw-Inches(2.3),Inches(0.6),size=13,c=(200,215,250),italic=True)

pnr(s15,15,light=True)
notitie(s15,"Sluit af met enthousiasme voor les 5. Componenten en design systems zijn het fundament van moderne UI design - het wordt duidelijk waarom we al die basis-theorie nodig hadden. Zorg dat studenten hun kleurpalet af hebben want dat gaan we in les 5 direct gebruiken.")

# ─────────────────────────────────────────────────────────────────────────────
# OPSLAAN
# ─────────────────────────────────────────────────────────────────────────────
OUT = r"C:\Users\Peter Marcelis\Documents\Lesmateriaal\design\periode-1\les-4\04_kleur-contrast.pptx"
prs.save(OUT)
print(f"Presentatie opgeslagen: {OUT}")
print(f"Aantal slides: {len(prs.slides)}")
