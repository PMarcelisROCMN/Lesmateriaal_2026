"""
Les 2: Whitespace en Alignment
Design regels (zelfde systeem als les 1):
- Geen rounded containers met flat overlays
- Rounded alleen voor kleine badges
- Alles past binnen zijn container
- Geen em-dashes in zichtbare tekst
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image, ImageDraw
import numpy as np, io

NAVY  = (6,  14,  44); NAVY2 = (30, 58, 138)
BLUE  = (37, 99, 235); BLUE2 = (59,130, 246)
CYAN  = (6, 182, 212); TEAL  = (13,148, 136); TEAL2 = (5, 80, 75)
WHITE = (255,255,255); OFFWH = (248,250,252)
DARK  = (15, 23,  42); BODY  = (51, 65,  85); MUTED = (100,116,135)
GREEN = (22,163,  74); RED   = (220, 38,  38); AMBER = (180, 90,  0)

def rgb(c): return RGBColor(c[0],c[1],c[2])
SW,SH=Inches(13.33),Inches(7.5); BW,BH=1920,1080

def garr(w,h,c1,c2,d='diag'):
    x,y=np.linspace(0,1,w,dtype=np.float32),np.linspace(0,1,h,dtype=np.float32)
    if d=='h': t=np.tile(x,(h,1))
    elif d=='v': t=np.tile(y.reshape(-1,1),(1,w))
    else: xx,yy=np.meshgrid(x,y); t=xx*.55+yy*.45
    a=np.zeros((h,w,3),dtype=np.float32)
    for i in range(3): a[:,:,i]=c1[i]+(c2[i]-c1[i])*t
    return np.clip(a,0,255).astype(np.uint8)

def dark_bg(circ=True):
    img=Image.fromarray(garr(BW,BH,NAVY,NAVY2)).convert('RGBA')
    if circ:
        ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
        d.ellipse([BW-500,-180,BW+200,520],fill=(13,148,136,20))
        d.ellipse([-60,BH-360,380,BH+100],fill=(6,182,212,22))
        img=Image.alpha_composite(img,ov)
    return img.convert('RGB')

def teal_bg():
    img=Image.fromarray(garr(BW,BH,(5,55,50),(13,100,90))).convert('RGBA')
    ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
    d.ellipse([BW-420,-130,BW+130,420],fill=(255,255,255,12))
    img=Image.alpha_composite(img,ov); return img.convert('RGB')

def light_bg():
    img=Image.fromarray(garr(BW,BH,OFFWH,WHITE)).convert('RGBA')
    ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
    d.ellipse([BW-580,-220,BW+180,540],fill=(13,148,136,8))
    img=Image.alpha_composite(img,ov); return img.convert('RGB')

def buf(img):
    b=io.BytesIO(); img.save(b,format='PNG'); b.seek(0); return b

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
BLANK=prs.slide_layouts[6]

def ns(bg=None):
    s=prs.slides.add_slide(BLANK)
    if bg:
        p=s.shapes.add_picture(buf(bg),0,0,SW,SH)
        t=s.shapes._spTree; t.remove(p._element); t.insert(2,p._element)
    else:
        s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(WHITE)
    return s

def box(s,l,t,w,h,c,border=None):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c)
    if border: sh.line.color.rgb=rgb(border); sh.line.width=Pt(0.75)
    else: sh.line.fill.background()
    return sh

def gbox(s,l,t,w,h,c1,c2,angle=90):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.gradient(); st=sh.fill.gradient_stops
    st[0].position=0.0; st[0].color.rgb=rgb(c1)
    st[1].position=1.0; st[1].color.rgb=rgb(c2)
    sh.fill.gradient_angle=angle; sh.line.fill.background(); return sh

def dot(s,cx,cy,r,c):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,cx-r,cy-r,r*2,r*2)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c); sh.line.fill.background()

def lbl(s,text,l,t,w,h,size=18,c=BODY,bold=False,italic=False,align=PP_ALIGN.LEFT):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=text; r.font.name="Calibri"; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=rgb(c)

def blist(s,items,l,t,w,h,size=18,c=BODY,gap=7):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(gap); r=p.add_run()
        r.text=f"  {chr(8226)}  {item}"; r.font.name="Calibri"; r.font.size=Pt(size); r.font.color.rgb=rgb(c)

def notitie(s,t): s.notes_slide.notes_text_frame.text=t
def pnr(s,n,tot=16,light=False):
    lbl(s,f"{n} / {tot}",Inches(12.5),Inches(7.1),Inches(0.8),Inches(0.35),size=11,c=(180,180,200) if light else MUTED,align=PP_ALIGN.RIGHT)

def hero_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(1.0),size=36,c=WHITE,bold=True)
    gbox(s,Inches(1.0),Inches(1.42),Inches(2.0),Inches(0.08),TEAL,CYAN,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.55),Inches(11.3),Inches(0.42),size=15,c=(160,210,210),italic=True)

def page_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(0.95),size=34,c=DARK,bold=True)
    gbox(s,Inches(1.0),Inches(1.38),Inches(1.8),Inches(0.08),TEAL,CYAN,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.5),Inches(11.3),Inches(0.38),size=14,c=MUTED,italic=True)

def card(s,l,t,w,h,accent=TEAL,bg=OFFWH):
    STRIP=Inches(0.07)
    box(s,l,t,w,h,bg,border=(215,225,240))
    gbox(s,l,t,STRIP,h,accent,CYAN,angle=270)
    return l+STRIP+Inches(0.18), t+Inches(0.14), w-STRIP-Inches(0.35)

def col_kop(s,l,t,w,h,text,c):
    box(s,l,t,w,h,c)
    lbl(s,text,l,t+Inches(0.07),w,h-Inches(0.07),size=20,c=WHITE,bold=True,align=PP_ALIGN.CENTER)


# ── SLIDE 1: Titel ────────────────────────────────────────
s=ns(dark_bg(True)); pnr(s,1,light=True)
gbox(s,0,0,Inches(0.5),SH,TEAL,CYAN,angle=270)
lbl(s,"Whitespace\nen Alignment",Inches(1.0),Inches(1.5),Inches(9.5),Inches(2.5),size=56,c=WHITE,bold=True)
gbox(s,Inches(1.0),Inches(3.95),Inches(2.8),Inches(0.1),TEAL,CYAN,angle=0)
lbl(s,"Periode 1  ·  Les 2",Inches(1.0),Inches(4.15),Inches(6.0),Inches(0.5),size=18,c=(160,210,210))
lbl(s,"De kracht van ruimte en lijn",Inches(1.0),Inches(4.72),Inches(6.0),Inches(0.4),size=14,c=(100,180,180),italic=True)
notitie(s,"Vandaag behandelen we twee van de meest fundamentele principes. Whitespace en alignment zijn onzichtbaar als het goed is, maar direct zichtbaar als het fout gaat.")


# ── SLIDE 2: Agenda ───────────────────────────────────────
s=ns(light_bg()); pnr(s,2)
page_hdr(s,"Vandaag")
items=["Wat is whitespace en waarom werkt het","Micro en macro whitespace","Het 8pt grid systeem","Alignment en visuele lijn","Figma: layout grid, guides en align tools","Wireframe van het goed doel project"]
bkl=[TEAL,(13,148,136),BLUE,BLUE2,(99,150,246),(5,150,105)]
for i,(item,bc) in enumerate(zip(items,bkl)):
    T=Inches(1.88)+i*Inches(0.88)
    dot(s,Inches(1.35),T+Inches(0.2),Inches(0.22),bc)
    lbl(s,str(i+1),Inches(1.13),T+Inches(0.04),Inches(0.44),Inches(0.38),size=14,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s,item,Inches(1.85),T+Inches(0.02),Inches(9.0),Inches(0.5),size=20,c=DARK)
notitie(s,"Loop de agenda door. Nadruk: we maken vandaag ook echt ons eerste wireframe in Figma voor het goed doel project.")


# ── SLIDE 3: Wat is whitespace ────────────────────────────
s=ns(); pnr(s,3)
gbox(s,0,0,Inches(5.0),SH,NAVY,(20,40,100),angle=270)
gbox(s,0,0,Inches(0.5),SH,TEAL,CYAN,angle=270)
lbl(s,"Whitespace",Inches(0.75),Inches(1.5),Inches(4.0),Inches(0.8),size=38,c=WHITE,bold=True)
lbl(s,"is niet leeg,\nhet is ruimte\nmet een doel",Inches(0.75),Inches(2.32),Inches(4.0),Inches(1.8),size=24,c=(180,210,255))
blist(s,["Alles tussen en rondom ontwerpelementen","Geeft het oog rust en richting","Bepaalt wat belangrijk is via contrast door ruimte","Maakt inhoud scanbaar en leesbaar","Gebruikt door Apple, Google en Linear als designtaal"],
      Inches(5.4),Inches(1.6),Inches(7.5),Inches(3.3),size=18,c=DARK)
cl,ct,cw=card(s,Inches(5.4),Inches(5.05),Inches(7.5),Inches(2.05),accent=TEAL)
lbl(s,'"White space is to be regarded as an active element, not a passive background."',cl,ct,cw,Inches(1.3),size=15,c=DARK,italic=True)
lbl(s,"Jan Tschichold",cl,ct+Inches(1.35),cw,Inches(0.35),size=12,c=MUTED)
notitie(s,"Whitespace is geen 'niets', het is een actieve designkeuze. Vergelijk met pauzes in muziek: de stilte geeft betekenis aan de noten.")


# ── SLIDE 4: Micro vs Macro ───────────────────────────────
s=ns(light_bg()); pnr(s,4)
page_hdr(s,"Twee soorten whitespace")
KW=Inches(5.55); KH=Inches(5.25); KT=Inches(1.75); HDR=Inches(0.55)
for ci,(title_,c1,rows) in enumerate([
    ("Micro whitespace",TEAL,[("Letter-spacing","Ruimte tussen letters"),("Line-height","Ruimte tussen regels tekst"),("Padding","Ruimte binnen een element zoals een knop"),("Gap","Ruimte tussen kleine elementen")]),
    ("Macro whitespace",NAVY,[("Sectie-spacing","Ruimte tussen secties op de pagina"),("Marges","Ruimte aan de zijkanten van de layout"),("Hero-ruimte","Lucht rondom het centrale element"),("Kolom-gap","Ruimte tussen kolommen in een grid")]),
]):
    L=Inches(1.0)+ci*(KW+Inches(0.45))
    col_kop(s,L,KT,KW,HDR,title_,c1)
    box(s,L,KT+HDR,KW,KH-HDR,OFFWH,border=(215,225,240))
    for j,(rowlbl,rowdesc) in enumerate(rows):
        RY=KT+HDR+j*Inches(1.15)
        bg=OFFWH if j%2==0 else WHITE
        box(s,L,RY,KW,Inches(1.1),bg,border=(220,228,245))
        lbl(s,rowlbl,L+Inches(0.15),RY+Inches(0.12),Inches(2.1),Inches(0.42),size=15,c=DARK,bold=True)
        lbl(s,rowdesc,L+Inches(0.15),RY+Inches(0.6),KW-Inches(0.3),Inches(0.42),size=14,c=BODY)
notitie(s,"Micro zijn de kleine details, macro de grote structuur. Een ontwerp kan goede micro maar slechte macro spacing hebben.")


# ── SLIDE 5: Whitespace in de praktijk ───────────────────
s=ns(light_bg()); pnr(s,5)
page_hdr(s,"Whitespace in de praktijk")
# Twee platte panelen naast elkaar. Geen rounded corners op de containers.
PW=Inches(5.35); PH=Inches(5.1); PT=Inches(1.75); GAP=Inches(0.25)
# Slecht
box(s,Inches(1.0),PT,PW,PH,OFFWH,border=(215,225,240))
box(s,Inches(1.0),PT,PW,Inches(0.45),(200,50,50))
lbl(s,"Zonder whitespace",Inches(1.1),PT+Inches(0.07),PW-Inches(0.2),Inches(0.35),size=13,c=WHITE,bold=True)
# Inhoud dicht op elkaar
for i,(w,y) in enumerate([(PW-Inches(0.2),Inches(0.62)),(PW-Inches(0.2),Inches(0.9)),(PW-Inches(0.2),Inches(1.18)),(Inches(3.2),Inches(1.46)),(PW-Inches(0.2),Inches(1.74)),(PW-Inches(0.2),Inches(2.02)),(Inches(2.0),Inches(2.3))]):
    box(s,Inches(1.1),PT+y,w,Inches(0.22),(170,175,190))
box(s,Inches(1.1),PT+Inches(2.68),Inches(2.2),Inches(0.48),RED)
lbl(s,"KLIK HIER",Inches(1.1),PT+Inches(2.72),Inches(2.2),Inches(0.38),size=13,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
lbl(s,"Tekst plakt aan de randen\nGeen ademruimte\nAlles heeft gelijke nadruk",Inches(1.1),PT+Inches(3.38),PW-Inches(0.2),Inches(1.0),size=12,c=RED,italic=True)
# Met whitespace
L2=Inches(1.0)+PW+GAP
box(s,L2,PT,PW,PH,OFFWH,border=(215,225,240))
box(s,L2,PT,PW,Inches(0.45),(20,100,80))
lbl(s,"Met whitespace",L2+Inches(0.15),PT+Inches(0.07),PW-Inches(0.3),Inches(0.35),size=13,c=WHITE,bold=True)
# Inhoud met ruimte: ruime inset
INSET=Inches(0.38)
box(s,L2+INSET,PT+Inches(0.75),Inches(3.5),Inches(0.3),DARK)
for i in range(3):
    box(s,L2+INSET,PT+Inches(1.28)+i*Inches(0.44),Inches(3.8),Inches(0.22),(170,175,190))
box(s,L2+INSET,PT+Inches(2.45),Inches(2.0),Inches(0.52),BLUE)
lbl(s,"Lees meer",L2+INSET,PT+Inches(2.5),Inches(2.0),Inches(0.38),size=13,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
lbl(s,"Marges creeren focus\nRuimte geeft hierarchie\nKnop heeft ademruimte",L2+INSET,PT+Inches(3.38),PW-INSET-Inches(0.2),Inches(1.0),size=12,c=GREEN,italic=True)
lbl(s,"Dezelfde inhoud, andere ervaring",Inches(1.0),Inches(6.95),(SW-Inches(1.5)),Inches(0.35),size=14,c=MUTED,italic=True,align=PP_ALIGN.CENTER)
notitie(s,"Laat beide panelen zien en vraag wat het verschil is. Studenten voelen het intuïtief, nu leren ze het bewust te controleren.")


# ── SLIDE 6: 8pt grid ────────────────────────────────────
s=ns(dark_bg(False)); pnr(s,6,light=True)
gbox(s,0,0,Inches(0.5),SH,TEAL,CYAN,angle=270)
hero_hdr(s,"Het 8pt grid systeem","De meest gebruikte spacingstandaard in professioneel design")
blist(s,["Alle spacing is een veelvoud van 8: zo zijn dat 8, 16, 24, 32, 48 en 64px","Werkt moeiteloos op alle schermresoluties en pixeldichtheden","Zorgt voor visuele consistentie zonder nadenken","Ontwikkelaars kunnen je ontwerp direct vertalen naar code","Gebruikt door Google Material Design, Apple HIG en IBM Carbon"],
      Inches(1.0),Inches(2.1),Inches(11.3),Inches(3.5),size=19,c=(200,220,255))
cl,ct,cw=card(s,Inches(1.0),Inches(5.75),Inches(11.3),Inches(1.35),accent=TEAL,bg=(20,35,80))
lbl(s,"Bron: spec.fm/specifics/8-pt-grid",cl,ct,cw,Inches(0.42),size=14,c=(160,210,210),bold=True)
lbl(s,"Uitgebreide uitleg met visuele voorbeelden van het systeem",cl,ct+Inches(0.45),cw,Inches(0.42),size=14,c=(130,165,210))
notitie(s,"Waarom 8? De meeste schermen werken met een basisraster van 8. Elk veelvoud van 8 rendert altijd scherp.")


# ── SLIDE 7: 8pt grid visueel ────────────────────────────
s=ns(light_bg()); pnr(s,7)
page_hdr(s,"8pt grid in de getallen")
spacings=[(8,"8px","Binnenste padding en icon gap",TEAL),(16,"16px","Standaard component padding",BLUE),(24,"24px","Sectie-interne spacing",BLUE2),(32,"32px","Spacing tussen componenten",(99,150,246)),(48,"48px","Sectie-spacing klein",(124,58,237)),(64,"64px","Sectie-spacing groot en hero marge",NAVY2)]
MAX_BAR=Inches(6.5)
for i,(px,plbl,desc,c) in enumerate(spacings):
    T=Inches(1.72)+i*Inches(0.88)
    bw=MAX_BAR*(px/64.0)
    gbox(s,Inches(1.0),T+Inches(0.12),bw,Inches(0.5),c,CYAN,angle=0)
    lbl(s,plbl,Inches(1.0)+bw+Inches(0.12),T+Inches(0.15),Inches(1.1),Inches(0.42),size=16,c=c,bold=True)
    lbl(s,desc,Inches(1.0)+bw+Inches(1.3),T+Inches(0.17),Inches(4.8),Inches(0.4),size=15,c=BODY)
# Demo knop rechts
box(s,Inches(9.6),Inches(1.72),Inches(3.35),Inches(5.15),DARK)
gbox(s,Inches(9.6),Inches(1.72),Inches(3.35),Inches(0.08),TEAL,CYAN,angle=0)
lbl(s,"Voorbeeld:",Inches(9.78),Inches(1.92),Inches(2.95),Inches(0.38),size=12,c=(160,185,230),italic=True)
box(s,Inches(9.78),Inches(2.42),Inches(2.6),Inches(0.58),BLUE)
lbl(s,"Meer informatie",Inches(9.78),Inches(2.47),Inches(2.6),Inches(0.48),size=14,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
lbl(s,"Hoogte: 48px (6x8)\nPadding: 16px (2x8)",Inches(9.78),Inches(3.18),Inches(2.8),Inches(0.72),size=12,c=(100,130,180),italic=True)
notitie(s,"Laat studenten een knop maken. Vraag: gebruik 16px padding links/rechts en 12px boven/onder. Wacht, 12 is geen veelvoud van 8! Maak het 8 of 16.")


# ── SLIDE 8: Alignment ────────────────────────────────────
s=ns(light_bg()); pnr(s,8)
page_hdr(s,"Alignment en visuele lijn")
blist(s,["Elementen op dezelfde lijn voelen als een geheel","Links uitlijnen is de standaard voor westerse leesvolgorde","Centreren werkt goed voor koppen en CTA's, niet voor lopende tekst","Rechts uitlijnen is zelden nodig op pagina's, wel in tabellen voor getallen"],
      Inches(1.0),Inches(1.75),Inches(7.4),Inches(3.0),size=18,c=DARK)
# Demo: slechte vs betere alignment, als twee platte panelen
PW=Inches(5.4); PH=Inches(2.42); PT=Inches(4.88)
# Slecht
box(s,Inches(1.0),PT,PW,PH,OFFWH,border=(215,225,240))
box(s,Inches(1.0),PT,PW,Inches(0.4),(200,50,50))
lbl(s,"Slechte alignment",Inches(1.12),PT+Inches(0.07),PW-Inches(0.24),Inches(0.3),size=13,c=WHITE,bold=True)
box(s,Inches(1.3), PT+Inches(0.58),Inches(3.4),Inches(0.28),(180,180,200))
box(s,Inches(1.55),PT+Inches(1.0), Inches(2.8),Inches(0.28),(180,180,200))
box(s,Inches(1.15),PT+Inches(1.42),Inches(1.8),Inches(0.42),RED)
lbl(s,"Knop",Inches(1.15),PT+Inches(1.48),Inches(1.8),Inches(0.32),size=12,c=WHITE,align=PP_ALIGN.CENTER)
# Betere alignment
L2=Inches(1.0)+PW+Inches(0.22)
box(s,L2,PT,PW,PH,OFFWH,border=(215,225,240))
box(s,L2,PT,PW,Inches(0.4),(20,100,80))
lbl(s,"Betere alignment",L2+Inches(0.15),PT+Inches(0.07),PW-Inches(0.3),Inches(0.3),size=13,c=WHITE,bold=True)
AL=L2+Inches(0.35)  # vaste links-uitlijn voor alle drie
box(s,AL,PT+Inches(0.58),Inches(3.4),Inches(0.28),(180,180,200))
box(s,AL,PT+Inches(1.0), Inches(2.8),Inches(0.28),(180,180,200))
box(s,AL,PT+Inches(1.42),Inches(1.8),Inches(0.42),BLUE)
lbl(s,"Knop",AL,PT+Inches(1.48),Inches(1.8),Inches(0.32),size=12,c=WHITE,align=PP_ALIGN.CENTER)
# Uitlijningslijn
gbox(s,AL,PT+Inches(0.5),Inches(0.05),Inches(1.55),TEAL,CYAN,angle=270)
notitie(s,"Laat studenten twee layouts bouwen in Figma: een zonder bewuste alignment en een met. Ze voelen het verschil direct.")


# ── SLIDE 9: Figma layout grid ───────────────────────────
s=ns(dark_bg(False)); pnr(s,9,light=True)
gbox(s,0,0,Inches(0.5),SH,TEAL,CYAN,angle=270)
hero_hdr(s,"Figma: Layout Grid instellen")
stappen=[("1","Selecteer een frame, bijvoorbeeld Desktop 1440px"),("2","Rechtsboven in het properties panel, klik het plus-teken naast Layout grid"),("3","Kies Columns en stel in: 12 kolommen, margin 80px, gutter 24px"),("4","Klik op het oogje om de grid te tonen of te verbergen"),("5","Alle elementen snappen nu automatisch aan de grid")]
for i,(num,stap) in enumerate(stappen):
    T=Inches(2.1)+i*Inches(0.97)
    dot(s,Inches(1.2),T+Inches(0.22),Inches(0.22),TEAL)
    lbl(s,num,Inches(0.98),T+Inches(0.06),Inches(0.44),Inches(0.38),size=14,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s,stap,Inches(1.62),T+Inches(0.04),Inches(11.0),Inches(0.5),size=18,c=(210,225,255))
box(s,Inches(1.0),Inches(7.0),(SW-Inches(1.5)),Inches(0.42),(20,45,45))
lbl(s,"Sneltoets om grid te tonen of te verbergen: Ctrl+Shift+4 op Windows, Cmd+Shift+4 op Mac",Inches(1.25),Inches(7.07),(SW-Inches(2.1)),Inches(0.3),size=13,c=(160,210,210),align=PP_ALIGN.CENTER)
notitie(s,"Demo: open een desktop frame, voeg layout grid toe, laat zien hoe elementen snappen. Tip: 12 kolommen zijn deelbaar door 2, 3, 4 en 6.")


# ── SLIDE 10: Rulers, guides, snapping ───────────────────
s=ns(light_bg()); pnr(s,10)
page_hdr(s,"Figma: Rulers, Guides en Snapping")
KW=Inches(3.6); GAP=Inches(0.28); KH=Inches(5.2); KT=Inches(1.75); HDR=Inches(0.55)
data=[
    ("Rulers",TEAL,["Toon of verberg: Ctrl+R (Cmd+R)","Klik en sleep vanuit de ruler","Zet gidsen op vaste posities","Handig voor paginamarges"]),
    ("Guides",BLUE,["Sleep vanuit de ruler naar de canvas","Dubbelklik om exacte positie in te stellen","Gidsen snappen automatisch","Verwijder door terug te slepen naar de ruler"]),
    ("Snapping",(124,58,237),["Automatisch uitlijnen aan andere elementen","Rode lijnen tonen de uitlijning","Shift ingedrukt houden om te overbruggen","Aan of uit via View en dan Snap to objects"]),
]
for i,(title_,c,items) in enumerate(data):
    L=Inches(1.0)+i*(KW+GAP)
    col_kop(s,L,KT,KW,HDR,title_,c)
    box(s,L,KT+HDR,KW,KH-HDR,OFFWH,border=(215,225,240))
    for j,item in enumerate(items):
        lbl(s,f"  {chr(8226)}  {item}",L+Inches(0.15),KT+HDR+j*Inches(1.1)+Inches(0.12),KW-Inches(0.3),Inches(0.88),size=15,c=BODY)
notitie(s,"Demo: sleep een guide vanuit de linker ruler. Laat snapping uitschakelen. Studenten zien meteen dat handmatig uitlijnen veel moeilijker is.")


# ── SLIDE 11: Align tools ─────────────────────────────────
s=ns(light_bg()); pnr(s,11)
page_hdr(s,"Figma: Align gereedschappen")
lbl(s,"Selecteer meerdere elementen en kijk rechtsboven in het properties panel voor de uitlijningsopties",Inches(1.0),Inches(1.55),Inches(11.3),Inches(0.5),size=16,c=MUTED,italic=True)
tools=[("Links uitlijnen","Align left edge",TEAL),("Centreren horizontaal","Align horizontal center",BLUE),("Rechts uitlijnen","Align right edge",BLUE2),("Boven uitlijnen","Align top edge",(124,58,237)),("Centreren verticaal","Align vertical center",(140,60,200)),("Onder uitlijnen","Align bottom edge",NAVY2),("Verdelen horizontaal","Distribute horizontally",(5,150,105)),("Verdelen verticaal","Distribute vertically",GREEN)]
TW=Inches(2.98); TH=Inches(1.68); GAP=Inches(0.22)
for i,(nl,en,c) in enumerate(tools):
    row=i//4; col=i%4
    L=Inches(1.0)+col*(TW+GAP); T=Inches(2.2)+row*(TH+GAP)
    box(s,L,T,TW,TH,OFFWH,border=(215,225,240))
    box(s,L,T,TW,Inches(0.08),c)
    dot(s,L+Inches(0.38),T+Inches(0.62),Inches(0.28),c)
    lbl(s,nl,L+Inches(0.15),T+Inches(1.05),TW-Inches(0.3),Inches(0.4),size=15,c=DARK,bold=True)
    lbl(s,en,L+Inches(0.15),T+Inches(1.42),TW-Inches(0.3),Inches(0.22),size=11,c=MUTED,italic=True)
notitie(s,"Demo: maak 3 rechthoeken op willekeurige posities, selecteer alle drie, laat align left zien. Dan distribute horizontally. Studenten doen mee.")


# ── SLIDE 12: Wireframe stap voor stap ───────────────────
s=ns(dark_bg(False)); pnr(s,12,light=True)
gbox(s,0,0,Inches(0.5),SH,TEAL,CYAN,angle=270)
hero_hdr(s,"Wireframe maken, stap voor stap","Goed doel project: eerste schets in Figma")
stappen=[("1","Frame aanmaken: F, kies Desktop 1440x1024, naam 'home-wireframe'"),("2","Layout grid instellen: 12 kolommen, margin 80px, gutter 24px"),("3","Navigatiebalk: rechthoek 1440x80px bovenaan, naam 'nav'"),("4","Hero sectie: groot vlak van ongeveer 600px hoog, placeholdertekst voor kop en subtitel"),("5","CTA knop: rechthoek 180x56px binnen de hero, naam 'hero/cta-knop'"),("6","Testimonials: drie gelijke kaarten naast elkaar met 8pt spacing"),("7","Alles in grijstinten, geen kleur. Focus op structuur en verhoudingen")]
for i,(num,stap) in enumerate(stappen):
    T=Inches(1.95)+i*Inches(0.74)
    dot(s,Inches(1.2),T+Inches(0.18),Inches(0.2),TEAL)
    lbl(s,num,Inches(0.98),T+Inches(0.02),Inches(0.44),Inches(0.36),size=13,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s,stap,Inches(1.6),T,Inches(11.0),Inches(0.48),size=16,c=(210,225,255))
notitie(s,"Dit is de hands-on kern van de les. Geef 20 tot 25 minuten. Loop rond en help bij stap 3 en stap 6.")


# ── SLIDE 13: Live analyse ────────────────────────────────
s=ns(light_bg()); pnr(s,13)
page_hdr(s,"Live analyse: whitespace herkennen")
blist(s,["Open apple.com: hoeveel ruimte staat er tussen de elementen?","Open een drukke nieuwssite: wat doet het gebrek aan whitespace met je?","Kijk naar de hero sectie: hoe groot is de marge rondom de kop?","Vergelijk de knoppen: hoeveel padding zit er in de CTA?"],
      Inches(1.0),Inches(1.75),Inches(11.3),Inches(3.0),size=19,c=DARK)
cl,ct,cw=card(s,Inches(1.0),Inches(4.88),Inches(11.3),Inches(2.25),accent=TEAL)
lbl(s,"Observatieopdracht",cl,ct,cw,Inches(0.42),size=15,c=DARK,bold=True)
lbl(s,"Kijk naar de wireframe van een medestudent. Benoem 2 plekken waar meer whitespace zou helpen en 2 plekken waar het al goed werkt. Geef mondeling feedback.",cl,ct+Inches(0.48),cw,Inches(1.45),size=16,c=BODY)
notitie(s,"Klassikaal 2 tot 3 wireframes reviewen. Gebruik als kader: heb je whitespace als actief element gebruikt of is het toevallig?")


# ── SLIDE 14: Bronnen ─────────────────────────────────────
s=ns(light_bg()); pnr(s,14)
page_hdr(s,"Bronnen en verdieping")
bronnen=[("8pt grid systeem","spec.fm/specifics/8-pt-grid",TEAL),("Whitespace in UX","interaction-design.org, zoek op 'white space'",BLUE),("F-patroon scangedrag","nngroup.com/articles/f-shaped-pattern-reading-web-content",BLUE2),("Every Layout","every-layout.dev, gratis boek over CSS layout principes",(124,58,237)),("Visual hierarchy","nngroup.com/articles/visual-hierarchy-ux-definition",NAVY2)]
for i,(title_,url,c) in enumerate(bronnen):
    T=Inches(1.75)+i*Inches(0.98)
    cl,ct,cw=card(s,Inches(1.0),T,Inches(11.3),Inches(0.82),accent=c)
    lbl(s,title_,cl,ct+Inches(0.04),Inches(2.8),Inches(0.38),size=16,c=DARK,bold=True)
    lbl(s,url,cl+Inches(2.95),ct+Inches(0.06),cw-Inches(2.95),Inches(0.38),size=14,c=c)
notitie(s,"Verwijs studenten naar deze bronnen. Spec.fm/8-pt-grid is goed om voor de volgende les te lezen.")


# ── SLIDE 15: Opdracht ────────────────────────────────────
s=ns(dark_bg(False)); pnr(s,15,light=True)
gbox(s,0,0,Inches(0.5),SH,TEAL,CYAN,angle=270)
lbl(s,"Opdracht  ·  Les 2",Inches(1.0),Inches(0.45),Inches(11.3),Inches(0.7),size=32,c=WHITE,bold=True)
lbl(s,"Goed Doel Website  ·  Wireframe",Inches(1.0),Inches(1.12),Inches(8.0),Inches(0.45),size=17,c=(160,210,210),italic=True)
gbox(s,Inches(1.0),Inches(1.55),Inches(2.2),Inches(0.08),TEAL,CYAN,angle=0)
stappen=[("1","Open je Figma bestand van les 1"),("2","Maak een nieuw frame: Desktop 1440x1024, naam 'home-wireframe'"),("3","Stel een 12-koloms layout grid in met margin 80px en gutter 24px"),("4","Bouw in grijstinten: navigatiebalk, hero met CTA knop en drie testimonial-kaarten"),("5","Gebruik 8pt spacing overal, controleer alle afstanden"),("6","Benoem alle lagen goed: 'nav', 'hero/kop', 'testimonial-1' enzovoort")]
bkl2=[TEAL,(13,148,136),BLUE,BLUE2,(99,150,246),(5,150,105)]
for i,((num,stap),bc) in enumerate(zip(stappen,bkl2)):
    T=Inches(1.82)+i*Inches(0.77)
    dot(s,Inches(1.2),T+Inches(0.19),Inches(0.2),bc)
    lbl(s,num,Inches(0.98),T+Inches(0.03),Inches(0.4),Inches(0.36),size=13,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s,stap,Inches(1.6),T+Inches(0.02),Inches(11.1),Inches(0.48),size=16,c=(210,225,255))
box(s,Inches(1.0),Inches(6.7),(SW-Inches(1.5)),Inches(0.65),(20,45,45))
lbl(s,"Inleveren: Figma link via de leeromgeving voor het begin van les 3",Inches(1.25),Inches(6.85),(SW-Inches(2.1)),Inches(0.38),size=13,c=(160,210,210),align=PP_ALIGN.CENTER)
notitie(s,"Geef de laatste 20 minuten om te starten. Controleer of iedereen de layout grid goed heeft ingesteld.")


# ── SLIDE 16: Volgende les ────────────────────────────────
s=ns(dark_bg(True)); pnr(s,16,light=True)
gbox(s,0,0,Inches(0.5),SH,TEAL,CYAN,angle=270)
lbl(s,"Volgende les",Inches(1.0),Inches(1.5),Inches(11.3),Inches(0.7),size=42,c=WHITE,bold=True)
lbl(s,"Typografie en Legibility",Inches(1.0),Inches(2.22),Inches(8.0),Inches(0.6),size=26,c=CYAN)
gbox(s,Inches(1.0),Inches(2.84),Inches(2.5),Inches(0.09),TEAL,CYAN,angle=0)
for i,item in enumerate(["Type scale van H1 tot body text: hoe maak je een systeem","Line-height, letter-spacing en maximale regellengte","Lettertypes kiezen: wanneer serif en wanneer sans-serif","Tekststijlen aanmaken en beheren in Figma"]):
    box(s,Inches(1.0),Inches(3.12)+i*Inches(0.92),(SW-Inches(1.5)),Inches(0.75),(15,40,40))
    gbox(s,Inches(1.0),Inches(3.12)+i*Inches(0.92),Inches(0.07),Inches(0.75),TEAL,CYAN,angle=270)
    lbl(s,item,Inches(1.25),Inches(3.2)+i*Inches(0.92),(SW-Inches(2.0)),Inches(0.55),size=17,c=(210,225,255))
box(s,Inches(1.0),Inches(6.68),(SW-Inches(1.5)),Inches(0.65),(10,32,32))
lbl(s,"Tip: kijk op fonts.google.com/knowledge voor uitstekende achtergrondlectuur over typografie",Inches(1.25),Inches(6.83),(SW-Inches(2.1)),Inches(0.38),size=14,c=(160,210,210),align=PP_ALIGN.CENTER)
notitie(s,"Sluit de les af. Tip: vraag studenten typografie bewust te observeren. Welke lettertypes zie je op merksites en waarom voelen die anders?")

prs.save(r"C:\Users\Peter Marcelis\Documents\Lesmateriaal\design\periode-1\les-2\02_whitespace-alignment.pptx")
print(f"Klaar: {len(prs.slides)} slides")
