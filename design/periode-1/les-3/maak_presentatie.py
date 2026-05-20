"""
Les 3: Typografie en Legibility
Periode 1 - Designcursus Jaar 1
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image, ImageDraw
import numpy as np, io

NAVY=(6,14,44); NAVY2=(30,58,138); BLUE=(37,99,235); BLUE2=(59,130,246)
CYAN=(6,182,212); TEAL=(13,148,136); WHITE=(255,255,255); OFFWH=(248,250,252)
DARK=(15,23,42); BODY=(51,65,85); MUTED=(100,116,135)
GREEN=(22,163,74); RED=(220,38,38)
ACCENT1=(124,58,237); ACCENT2=(167,139,250)  # Purple/Violet for les 3

def rgb(c): return RGBColor(c[0],c[1],c[2])
SW,SH=Inches(13.33),Inches(7.5); BW,BH=1920,1080

def garr(w,h,c1,c2,d='diag'):
    x,y=np.linspace(0,1,w,dtype=np.float32),np.linspace(0,1,h,dtype=np.float32)
    if d=='h': t=np.tile(x,(h,1))
    elif d=='v': t=np.tile(y.reshape(-1,1),(1,w))
    else: xx,yy=np.meshgrid(x,y); t=xx*.55+yy*.45
    a=np.zeros((h,w,3),dtype=np.float32)
    for i in range(3): a[:,:,i]=c1[i]+(c2[i]-c1[i])*t
    return np.clip(a,0,255).astype(np.uint8)

def dark_bg(circ=True):
    img=Image.fromarray(garr(BW,BH,NAVY,NAVY2)).convert('RGBA')
    if circ:
        ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
        d.ellipse([BW-500,-180,BW+200,520],fill=(ACCENT1[0],ACCENT1[1],ACCENT1[2],20))
        d.ellipse([-60,BH-360,380,BH+100],fill=(ACCENT2[0],ACCENT2[1],ACCENT2[2],22))
        img=Image.alpha_composite(img,ov)
    return img.convert('RGB')

def light_bg():
    img=Image.fromarray(garr(BW,BH,OFFWH,WHITE)).convert('RGBA')
    ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
    d.ellipse([BW-580,-220,BW+180,540],fill=(ACCENT1[0],ACCENT1[1],ACCENT1[2],8))
    img=Image.alpha_composite(img,ov); return img.convert('RGB')

def buf(img):
    b=io.BytesIO(); img.save(b,format='PNG'); b.seek(0); return b

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
BLANK=prs.slide_layouts[6]

def ns(bg=None):
    s=prs.slides.add_slide(BLANK)
    if bg:
        p=s.shapes.add_picture(buf(bg),0,0,SW,SH)
        t=s.shapes._spTree; t.remove(p._element); t.insert(2,p._element)
    else:
        s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(WHITE)
    return s

def box(s,l,t,w,h,c,border=None):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c)
    if border: sh.line.color.rgb=rgb(border); sh.line.width=Pt(0.75)
    else: sh.line.fill.background()
    return sh

def gbox(s,l,t,w,h,c1,c2,angle=90):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.gradient(); st=sh.fill.gradient_stops
    st[0].position=0.0; st[0].color.rgb=rgb(c1)
    st[1].position=1.0; st[1].color.rgb=rgb(c2)
    sh.fill.gradient_angle=angle; sh.line.fill.background(); return sh

def dot(s,cx,cy,r,c):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,cx-r,cy-r,r*2,r*2)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c); sh.line.fill.background()

def lbl(s,text,l,t,w,h,size=18,c=BODY,bold=False,italic=False,align=PP_ALIGN.LEFT):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=text; r.font.name="Calibri"; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=rgb(c)

def blist(s,items,l,t,w,h,size=18,c=BODY,gap=7):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(gap); r=p.add_run()
        r.text=f"  {chr(8226)}  {item}"; r.font.name="Calibri"; r.font.size=Pt(size); r.font.color.rgb=rgb(c)

def notitie(s,t): s.notes_slide.notes_text_frame.text=t

def pnr(s,n,tot=15,light=False):
    lbl(s,f"{n} / {tot}",Inches(12.5),Inches(7.1),Inches(0.8),Inches(0.35),size=11,c=(180,180,200) if light else MUTED,align=PP_ALIGN.RIGHT)

def hero_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(1.0),size=36,c=WHITE,bold=True)
    gbox(s,Inches(1.0),Inches(1.42),Inches(2.0),Inches(0.08),ACCENT1,ACCENT2,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.55),Inches(11.3),Inches(0.42),size=15,c=(200,200,240),italic=True)

def page_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(0.95),size=34,c=DARK,bold=True)
    gbox(s,Inches(1.0),Inches(1.38),Inches(1.8),Inches(0.08),ACCENT1,ACCENT2,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.5),Inches(11.3),Inches(0.38),size=14,c=MUTED,italic=True)

def card(s,l,t,w,h,accent=None,bg=OFFWH):
    a=accent if accent else ACCENT1
    STRIP=Inches(0.07)
    box(s,l,t,w,h,bg,border=(215,225,240))
    gbox(s,l,t,STRIP,h,a,ACCENT2,angle=270)
    return l+STRIP+Inches(0.18), t+Inches(0.14), w-STRIP-Inches(0.35)

def col_kop(s,l,t,w,h,text,c):
    box(s,l,t,w,h,c)
    lbl(s,text,l,t+Inches(0.07),w,h-Inches(0.07),size=20,c=WHITE,bold=True,align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 1: Titel
# ─────────────────────────────────────────
s1=ns(dark_bg(True))

# Vertical accent strip left
gbox(s1,Inches(0.0),Inches(0.0),Inches(0.22),Inches(7.5),ACCENT1,ACCENT2,angle=270)

lbl(s1,"Typografie\nen Legibility",Inches(1.0),Inches(1.2),Inches(11.0),Inches(2.6),size=64,c=WHITE,bold=True)
gbox(s1,Inches(1.0),Inches(3.9),Inches(3.2),Inches(0.09),ACCENT1,ACCENT2,angle=0)
lbl(s1,"Periode 1  .  Les 3",Inches(1.0),Inches(4.05),Inches(8.0),Inches(0.5),size=20,c=(200,200,240))
lbl(s1,"De kunst van leesbare letters",Inches(1.0),Inches(4.6),Inches(8.0),Inches(0.5),size=18,c=(167,139,250),italic=True)

# Decorative dots
dot(s1,Inches(11.8),Inches(1.5),Inches(0.55),(ACCENT1[0],ACCENT1[1],ACCENT1[2]))
dot(s1,Inches(12.4),Inches(2.2),Inches(0.35),(ACCENT2[0],ACCENT2[1],ACCENT2[2]))
dot(s1,Inches(11.2),Inches(2.8),Inches(0.22),(100,60,200))

pnr(s1,1,light=True)
notitie(s1,"Slide 1: Welkom bij les 3. Typografie is een van de fundamentele vaardigheden van een designer. Vandaag leren we de regels die bepalen of tekst prettig leesbaar is of niet. Leg uit dat dit direct toepasbaar is op het goed doel project.")

# ─────────────────────────────────────────
# SLIDE 2: Agenda
# ─────────────────────────────────────────
s2=ns(light_bg())
page_hdr(s2,"Vandaag","Overzicht van de les")

items=[
    ("1","Serif versus sans-serif, wanneer gebruik je wat",(124,58,237)),
    ("2","Type scale van H1 tot caption",(109,40,217)),
    ("3","Line-height en letter-spacing",(91,33,182)),
    ("4","Maximale regellengte en leescomfort",(124,58,237)),
    ("5","Font pairs kiezen en combineren",(109,40,217)),
    ("6","Figma: tekststijlen aanmaken voor het project",(91,33,182)),
]

for i,(num,text,col) in enumerate(items):
    y=Inches(1.75)+i*Inches(0.84)
    dot(s2,Inches(1.35),y+Inches(0.22),Inches(0.22),col)
    lbl(s2,num,Inches(1.18),y+Inches(0.06),Inches(0.35),Inches(0.36),size=13,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s2,text,Inches(1.65),y,Inches(10.5),Inches(0.44),size=19,c=DARK)

pnr(s2,2)
notitie(s2,"Slide 2: Doorloop de agenda. Benadruk dat al deze onderwerpen terugkomen in het goed doel project. Studenten die al een idee hebben voor hun charity website kunnen alvast nadenken welke fonts bij hun thema passen.")

# ─────────────────────────────────────────
# SLIDE 3: Serif vs Sans-serif
# ─────────────────────────────────────────
s3=ns()
s3.background.fill.solid(); s3.background.fill.fore_color.rgb=rgb(OFFWH)

# Left panel: dark (serif)
box(s3,Inches(0.0),Inches(0.0),Inches(5.5),Inches(7.5),NAVY)
lbl(s3,"Serif",Inches(0.3),Inches(0.35),Inches(4.8),Inches(0.75),size=30,c=WHITE,bold=True)
gbox(s3,Inches(0.3),Inches(1.1),Inches(1.4),Inches(0.07),ACCENT1,ACCENT2,angle=0)

# Big Ag demo left
lbl(s3,"Ag",Inches(0.3),Inches(1.25),Inches(3.5),Inches(1.5),size=72,c=(167,139,250),bold=False,italic=False)

blist(s3,[
    "Gebruikt voor editorial en print",
    "Straalt traditie en autoriteit uit",
    "Schreven geleiden het oog langs tekst",
    "Voorbeelden: Georgia, Playfair Display",
    "Werkt goed als display heading online",
],Inches(0.3),Inches(2.75),Inches(4.8),Inches(2.6),size=15,c=(220,220,255),gap=5)

# Right panel: light (sans-serif)
lbl(s3,"Sans-serif",Inches(5.8),Inches(0.35),Inches(7.0),Inches(0.75),size=30,c=DARK,bold=True)
gbox(s3,Inches(5.8),Inches(1.1),Inches(1.4),Inches(0.07),ACCENT1,ACCENT2,angle=0)

# Big Ag demo right
lbl(s3,"Ag",Inches(5.8),Inches(1.25),Inches(3.5),Inches(1.5),size=72,c=ACCENT1,bold=False)

blist(s3,[
    "Digital-first ontworpen",
    "Modern en clean esthetiek",
    "Uitstekende leesbaarheid op schermen",
    "Voorbeelden: Inter, Roboto, Outfit",
    "Standaard keuze voor body text op web",
],Inches(5.8),Inches(2.75),Inches(7.0),Inches(2.6),size=15,c=BODY,gap=5)

# Divider line
box(s3,Inches(5.45),Inches(0.2),Inches(0.08),Inches(7.1),(210,218,235))

# Bottom tip card
cx,cy,cw=card(s3,Inches(0.7),Inches(6.25),Inches(12.0),Inches(0.95))
lbl(s3,"Tip voor websites: gebruik sans-serif voor body text, serif kan werken voor display headings.",cx,cy,cw,Inches(0.7),size=14,c=BODY)

pnr(s3,3)
notitie(s3,"Slide 3: Laat studenten zelf noemen welke fonts ze kennen. Serif versus sans-serif is een van de meest fundamentele onderscheidingen in typografie. Wijs op de schreven (kleine haakjes) die typerend zijn voor serif fonts.")

# ─────────────────────────────────────────
# SLIDE 4: Type scale
# ─────────────────────────────────────────
s4=ns(light_bg())
page_hdr(s4,"Type scale: hierarchie in grootte","Van H1 tot caption, elke stap heeft een doel")

rows=[
    ("H1","48px",9.0,"Paginatitel, hero heading"),
    ("H2","36px",6.75,"Sectietitels en paginanamen"),
    ("H3","28px",5.25,"Subtitels en kaartkoppen"),
    ("H4","22px",4.1,"Component headers"),
    ("H5","18px",3.4,"Kleine headers en labels"),
    ("Body","16px",3.0,"Lopende tekst, beschrijvingen"),
    ("Caption","12px",2.25,"Labels, metadata, bijschriften"),
]

start_y=Inches(1.65)
for i,(name,px,bw,usage) in enumerate(rows):
    y=start_y+i*Inches(0.75)
    # Name label
    lbl(s4,name,Inches(0.7),y,Inches(1.0),Inches(0.5),size=14,c=DARK,bold=True)
    # px label
    lbl(s4,px,Inches(1.75),y,Inches(0.7),Inches(0.5),size=13,c=MUTED)
    # Bar
    bar_w=min(Inches(bw),Inches(9.0))
    gbox(s4,Inches(2.5),y+Inches(0.1),bar_w,Inches(0.28),ACCENT1,ACCENT2,angle=0)
    # Usage
    lbl(s4,usage,Inches(2.6),y+Inches(0.35),Inches(9.5),Inches(0.35),size=11,c=MUTED,italic=True)

pnr(s4,4)
notitie(s4,"Slide 4: De type scale is een van de meest tijdbesparende tools voor een designer. Door een vaste schaal te definiëren hoef je nooit meer te raden welke maat je moet gebruiken. Verwijs naar het project: studenten gaan dit exact zo instellen in Figma.")

# ─────────────────────────────────────────
# SLIDE 5: Line-height en Letter-spacing
# ─────────────────────────────────────────
s5=ns(light_bg())
page_hdr(s5,"Line-height en Letter-spacing","Ruimte tussen regels en letters bepaalt leescomfort")

PANEL_W=Inches(5.3)
PANEL_H=Inches(4.2)
LEFT_X=Inches(0.8)
RIGHT_X=Inches(7.0)
PANEL_Y=Inches(1.8)

# Left panel: cramped
box(s5,LEFT_X,PANEL_Y,PANEL_W,PANEL_H,WHITE,border=(215,225,240))
# Red header bar
box(s5,LEFT_X,PANEL_Y,PANEL_W,Inches(0.45),RED)
lbl(s5,"Te weinig ruimte",LEFT_X+Inches(0.15),PANEL_Y+Inches(0.05),PANEL_W-Inches(0.2),Inches(0.38),size=15,c=WHITE,bold=True)

# Cramped text simulation
for i in range(8):
    yy=PANEL_Y+Inches(0.55)+i*Inches(0.36)
    ww=PANEL_W-Inches(0.3) if i%3!=2 else PANEL_W*0.65
    box(s5,LEFT_X+Inches(0.15),yy,ww,Inches(0.18),(210,215,225))

blist(s5,[
    "Line-height: 1.1",
    "Letter-spacing: 0",
    "Moeilijk te volgen",
    "Vermoeiend na 1 minuut",
],LEFT_X+Inches(0.1),PANEL_Y+PANEL_H+Inches(0.08),PANEL_W,Inches(0.8),size=13,c=RED,gap=3)

# Right panel: well-spaced
box(s5,RIGHT_X,PANEL_Y,PANEL_W,PANEL_H,WHITE,border=(215,225,240))
# Green header bar
box(s5,RIGHT_X,PANEL_Y,PANEL_W,Inches(0.45),GREEN)
lbl(s5,"Goed ingesteld",RIGHT_X+Inches(0.15),PANEL_Y+Inches(0.05),PANEL_W-Inches(0.2),Inches(0.38),size=15,c=WHITE,bold=True)

# Spaced text simulation
for i in range(6):
    yy=PANEL_Y+Inches(0.6)+i*Inches(0.55)
    ww=PANEL_W-Inches(0.3) if i%3!=2 else PANEL_W*0.7
    box(s5,RIGHT_X+Inches(0.15),yy,ww,Inches(0.22),(210,215,225))

blist(s5,[
    "Line-height: 1.6",
    "Letter-spacing: normaal",
    "Oog volgt moeiteloos",
    "Comfortabel leesbaar",
],RIGHT_X+Inches(0.1),PANEL_Y+PANEL_H+Inches(0.08),PANEL_W,Inches(0.8),size=13,c=GREEN,gap=3)

# Bottom info card
cx,cy,cw=card(s5,Inches(0.7),Inches(6.72),Inches(12.0),Inches(0.62))
lbl(s5,"Line-height: 1.5x fontgrootte voor body, 1.2x voor headings. Letter-spacing: 0 tot -0.02em voor headings, 0 voor body.",cx,cy,cw,Inches(0.5),size=13,c=BODY)

pnr(s5,5)
notitie(s5,"Slide 5: Demonstreer in de browser hoe line-height de leeservaring verandert. Open inspector tools en pas de CSS aan live. Laat studenten de twee kolommen vergelijken en verwoorden wat ze anders voelen.")

# ─────────────────────────────────────────
# SLIDE 6: Maximale regellengte
# ─────────────────────────────────────────
s6=ns(dark_bg(False))
gbox(s6,Inches(0.0),Inches(0.0),Inches(0.22),Inches(7.5),ACCENT1,ACCENT2,angle=270)
hero_hdr(s6,"De 45 tot 75 karakters regel","Regellengte is een van de meest onderschatte designbeslissingen")

blist(s6,[
    "Te korte regels: oog springt te vaak terug, ritme wordt verstoord",
    "Te lange regels: oog verliest de volgende regel, vermoeiend om te lezen",
    "Optimaal: 45 tot 75 karakters per regel voor body text",
    "In CSS: max-width van 65ch is een bewezen startpunt",
    "Narrowere kolommen werken goed voor zijbalken en kaartjes",
],Inches(1.0),Inches(2.1),Inches(7.0),Inches(2.8),size=17,c=(220,225,255),gap=8)

# Visual demo: three text columns of different widths
DEMO_X=Inches(8.5)
DEMO_Y=Inches(1.8)
DEMO_H=Inches(4.5)

# Narrow column (too short)
box(s6,DEMO_X,DEMO_Y,Inches(1.1),DEMO_H,(30,45,90))
lbl(s6,"Te smal",DEMO_X,DEMO_Y+Inches(0.05),Inches(1.1),Inches(0.3),size=10,c=(220,38,38),bold=True,align=PP_ALIGN.CENTER)
for i in range(9):
    yy=DEMO_Y+Inches(0.4)+i*Inches(0.42)
    box(s6,DEMO_X+Inches(0.1),yy,Inches(0.9),Inches(0.15),(70,90,160))

# Ideal column
box(s6,DEMO_X+Inches(1.3),DEMO_Y,Inches(2.0),DEMO_H,(20,35,80))
lbl(s6,"Ideaal",DEMO_X+Inches(1.3),DEMO_Y+Inches(0.05),Inches(2.0),Inches(0.3),size=10,c=ACCENT2,bold=True,align=PP_ALIGN.CENTER)
gbox(s6,DEMO_X+Inches(1.3),DEMO_Y,Inches(2.0),Inches(0.06),ACCENT1,ACCENT2,angle=0)
for i in range(9):
    yy=DEMO_Y+Inches(0.4)+i*Inches(0.42)
    box(s6,DEMO_X+Inches(1.4),yy,Inches(1.8),Inches(0.15),ACCENT1)

# Too wide column
box(s6,DEMO_X+Inches(3.5),DEMO_Y,Inches(1.1),DEMO_H,(30,45,90))
lbl(s6,"Te breed",DEMO_X+Inches(3.5),DEMO_Y+Inches(0.05),Inches(1.1),Inches(0.3),size=10,c=(220,38,38),bold=True,align=PP_ALIGN.CENTER)
for i in range(9):
    yy=DEMO_Y+Inches(0.4)+i*Inches(0.42)
    box(s6,DEMO_X+Inches(3.6),yy,Inches(0.9),Inches(0.15),(70,90,160))

lbl(s6,"max-width: 65ch",Inches(1.0),Inches(6.85),Inches(6.5),Inches(0.45),size=15,c=ACCENT2,italic=True)

pnr(s6,6,light=True)
notitie(s6,"Slide 6: Meet samen met studenten de regellengte van een populaire website. Wikipedia heeft doorgaans een te brede kolom, Medium.com heeft een bijna perfecte regellengte. Laat ze dit vergelijken.")

# ─────────────────────────────────────────
# SLIDE 7: Font pairs kiezen
# ─────────────────────────────────────────
s7=ns(light_bg())
page_hdr(s7,"Font pairs: heading plus body","Goede font combinaties versterken de merkidentiteit")

pairs=[
    ("Playfair Display + Inter","Editorial feel, goed voor lifestyle en cultuur websites. De serif heading geeft karakter, Inter als body zorgt voor uitstekende leesbaarheid.",ACCENT1),
    ("Outfit + Outfit","Dezelfde familie, varieer met weight. Modern en clean, uitstekend voor tech en startup projecten. Simpel en effectief.",BLUE),
    ("Space Grotesk + Source Sans Pro","Tech en startup gevoel. Beide geometrisch, maar Space Grotesk heeft meer karakter als heading.",GREEN),
]

for i,(pair,desc,acc) in enumerate(pairs):
    y=Inches(1.75)+i*Inches(1.7)
    cx,cy,cw=card(s7,Inches(0.8),y,Inches(11.5),Inches(1.52),accent=acc)
    lbl(s7,pair,cx,cy,cw,Inches(0.38),size=20,c=DARK,bold=True)
    lbl(s7,desc,cx,cy+Inches(0.38),cw,Inches(0.72),size=14,c=BODY)
    lbl(s7,"Aa Bb 123",cx+cw-Inches(1.8),cy,Inches(1.7),Inches(0.38),size=18,c=acc,bold=True,align=PP_ALIGN.RIGHT)

pnr(s7,7)
notitie(s7,"Slide 7: Google Fonts heeft een uitstekende 'pairings' sectie. fonts.google.com/knowledge bevat artikelen over hoe je fonts combineert. Laat studenten hun laptop pakken en twee fonts zoeken die bij hun goed doel passen.")

# ─────────────────────────────────────────
# SLIDE 8: Legibility vs Readability
# ─────────────────────────────────────────
s8=ns(light_bg())
page_hdr(s8,"Legibility versus Readability","Twee verwante maar verschillende kwaliteiten van tekst")

PANEL_W2=Inches(5.6)
PANEL_H2=Inches(3.9)
LX=Inches(0.8)
RX=Inches(6.8)
PY=Inches(1.75)

# Left: Legibility
box(s8,LX,PY,PANEL_W2,PANEL_H2,(248,245,255),border=(200,180,240))
lbl(s8,"Legibility",LX+Inches(0.2),PY+Inches(0.15),PANEL_W2-Inches(0.3),Inches(0.48),size=22,c=ACCENT1,bold=True)
lbl(s8,"Herkenbaarheid van individuele letters",LX+Inches(0.2),PY+Inches(0.6),PANEL_W2-Inches(0.3),Inches(0.4),size=13,c=MUTED,italic=True)
blist(s8,[
    "Contrast tussen lettervormen",
    "Duidelijke lettervormen en tekens",
    "Geschikt voor kleine groottes",
    "Groot impact op toegankelijkheid",
],LX+Inches(0.1),PY+Inches(1.0),PANEL_W2-Inches(0.2),Inches(2.5),size=15,c=BODY,gap=6)

# Right: Readability
box(s8,RX,PY,PANEL_W2,PANEL_H2,(245,252,248),border=(180,230,200))
lbl(s8,"Readability",RX+Inches(0.2),PY+Inches(0.15),PANEL_W2-Inches(0.3),Inches(0.48),size=22,c=GREEN,bold=True)
lbl(s8,"Leesbaarheid van langere tekst",RX+Inches(0.2),PY+Inches(0.6),PANEL_W2-Inches(0.3),Inches(0.4),size=13,c=MUTED,italic=True)
blist(s8,[
    "Line-height en letter-spacing",
    "Regellengte en kolombreedte",
    "Contrast met de achtergrond",
    "Lettergrootte in context",
],RX+Inches(0.1),PY+Inches(1.0),PANEL_W2-Inches(0.2),Inches(2.5),size=15,c=BODY,gap=6)

# Quote card at bottom
cx,cy,cw=card(s8,Inches(0.8),Inches(5.85),Inches(11.7),Inches(1.35))
lbl(s8,'"Legibility is a quality of the type. Readability is a quality of the typographic arrangement."',cx,cy,cw,Inches(0.6),size=14,c=DARK,italic=True)
lbl(s8,"Walter Tracy",cx,cy+Inches(0.6),cw,Inches(0.4),size=13,c=MUTED)

pnr(s8,8)
notitie(s8,"Slide 8: Dit onderscheid is subtiel maar belangrijk. Legibility gaat over hoe goed je een letter kunt herkennen (is de g onderscheidbaar van de q?). Readability gaat over hoe comfortabel je langere tekst kunt lezen.")

# ─────────────────────────────────────────
# SLIDE 9: Figma tekststijlen aanmaken
# ─────────────────────────────────────────
s9=ns(dark_bg(False))
gbox(s9,Inches(0.0),Inches(0.0),Inches(0.22),Inches(7.5),ACCENT1,ACCENT2,angle=270)
hero_hdr(s9,"Figma: tekststijlen instellen","Eenmalig aanmaken, overal hergebruiken")

steps=[
    ("1","Selecteer een tekstelement met de juiste instellingen in je Figma frame",(124,58,237)),
    ("2","Klik rechtsboven op de vier puntjes naast het tekstveld in het properties panel",(109,40,217)),
    ("3","Kies 'Create style' en geef het een naam zoals 'heading/h1'",(91,33,182)),
    ("4","Gebruik schuine strepen voor groepering: 'body/regular' en 'body/bold'",(124,58,237)),
    ("5","Hergebruik stijlen via het stijlenpanel door op het icoon te klikken",(109,40,217)),
]

for i,(num,text,col) in enumerate(steps):
    y=Inches(2.05)+i*Inches(0.84)
    dot(s9,Inches(1.35),y+Inches(0.22),Inches(0.26),col)
    lbl(s9,num,Inches(1.18),y+Inches(0.06),Inches(0.35),Inches(0.36),size=13,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s9,text,Inches(1.75),y,Inches(10.5),Inches(0.44),size=16,c=(220,220,255))

# Tip card (dark version)
box(s9,Inches(0.9),Inches(6.55),Inches(11.5),Inches(0.72),(20,35,85))
gbox(s9,Inches(0.9),Inches(6.55),Inches(0.07),Inches(0.72),ACCENT1,ACCENT2,angle=270)
lbl(s9,"Maak minimaal: H1, H2, H3, body, body-small, label, caption. Dit dekt 90% van de webdesign behoeften.",Inches(1.1),Inches(6.65),Inches(11.0),Inches(0.55),size=14,c=(200,210,255))

pnr(s9,9,light=True)
notitie(s9,"Slide 9: Open Figma live tijdens de les en demonstreer dit stap voor stap. Laat studenten meteen meedoen op hun eigen laptop. Het aanmaken van tekststijlen kost 10 minuten maar bespaart uren later.")

# ─────────────────────────────────────────
# SLIDE 10: Figma type system tabel
# ─────────────────────────────────────────
s10=ns(light_bg())
page_hdr(s10,"Type system in het project","Overzicht van de aanbevolen tekststijlen")

COL_W=Inches(3.7)
COL_X=[Inches(1.0),Inches(4.7),Inches(8.4)]
HDR_Y=Inches(1.85)
HDR_H=Inches(0.5)

col_kop(s10,COL_X[0],HDR_Y,COL_W,HDR_H,"Stijlnaam",ACCENT1)
col_kop(s10,COL_X[1],HDR_Y,COL_W,HDR_H,"Instellingen",BLUE)
col_kop(s10,COL_X[2],HDR_Y,COL_W,HDR_H,"Gebruik",GREEN)

rows_t=[
    ("heading/h1","48px Bold, line-height 1.2","Hero en paginatitels"),
    ("heading/h2","36px Bold, line-height 1.25","Sectietitels"),
    ("heading/h3","24px SemiBold, line-height 1.3","Kaart- en componenttitels"),
    ("body/regular","16px Regular, line-height 1.6","Beschrijvingen en paragrafen"),
    ("body/small","14px Regular, line-height 1.5","Secundaire informatie"),
    ("label","12px Medium, line-height 1.4","Knoppen en navigatie"),
]

for i,(name,settings,usage) in enumerate(rows_t):
    ry=HDR_Y+HDR_H+i*Inches(0.82)
    rbg=OFFWH if i%2==0 else WHITE
    for cx in COL_X:
        box(s10,cx,ry,COL_W,Inches(0.82),rbg,border=(220,225,235))
    PAD=Inches(0.15)
    RPAD=Inches(0.1)
    lbl(s10,name,COL_X[0]+PAD,ry+RPAD,COL_W-Inches(0.2),Inches(0.62),size=13,c=DARK,bold=True)
    lbl(s10,settings,COL_X[1]+PAD,ry+RPAD,COL_W-Inches(0.2),Inches(0.62),size=13,c=BODY)
    lbl(s10,usage,COL_X[2]+PAD,ry+RPAD,COL_W-Inches(0.2),Inches(0.62),size=13,c=BODY)

pnr(s10,10)
notitie(s10,"Slide 10: Dit is de referentietabel voor het project. Studenten kunnen dit als screenshot opslaan of de waarden direct overnemen in Figma. Benadruk dat de namen met schuine streep (heading/h1) een maphierarchie creëren in Figma.")

# ─────────────────────────────────────────
# SLIDE 11: Goed doel project typografie
# ─────────────────────────────────────────
s11=ns(dark_bg(False))
gbox(s11,Inches(0.0),Inches(0.0),Inches(0.22),Inches(7.5),ACCENT1,ACCENT2,angle=270)
hero_hdr(s11,"Goed doel: typografie kiezen","Praktische gids voor je charity website")

advice=[
    "Kies een gratis Google Font voor de heading: zoek op fonts.google.com/knowledge",
    "Kies Inter of Outfit als body font: beide zijn uitstekend leesbaar op schermen",
    "Stel je type scale in Figma in voor je begint met ontwerpen",
    "Controleer contrast: body text minimaal WCAG AA (4.5:1 ratio)",
]

for i,text in enumerate(advice):
    y=Inches(2.1)+i*Inches(0.72)
    dot(s11,Inches(1.2),y+Inches(0.18),Inches(0.14),ACCENT2)
    lbl(s11,text,Inches(1.45),y,Inches(7.2),Inches(0.44),size=16,c=(220,220,255))

# Starter pack card
box(s11,Inches(8.8),Inches(2.0),Inches(4.2),Inches(4.0),(20,30,80))
gbox(s11,Inches(8.8),Inches(2.0),Inches(4.2),Inches(0.08),ACCENT1,ACCENT2,angle=0)
lbl(s11,"Starter pack",Inches(8.95),Inches(2.15),Inches(4.0),Inches(0.45),size=16,c=ACCENT2,bold=True)

starter=[
    ("Heading:", "Outfit Bold 48/36/28px"),
    ("Body:", "Inter Regular 16px lh 1.6"),
    ("Caption:", "Inter Regular 12px"),
]
for i,(label,val) in enumerate(starter):
    yy=Inches(2.72)+i*Inches(0.9)
    lbl(s11,label,Inches(9.0),yy,Inches(1.1),Inches(0.38),size=13,c=ACCENT2,bold=True)
    lbl(s11,val,Inches(10.15),yy,Inches(2.7),Inches(0.38),size=13,c=(200,210,255))

pnr(s11,11,light=True)
notitie(s11,"Slide 11: Laat studenten nu hun goed doel thema bedenken. Een dierenasiel vraagt om andere fonts dan een mensenrechtenorganisatie. Bespreek hoe typografie gevoel en vertrouwen communiceert.")

# ─────────────────────────────────────────
# SLIDE 12: Live analyse
# ─────────────────────────────────────────
s12=ns(light_bg())
page_hdr(s12,"Live analyse: typografie herkennen","Onderzoek echte websites met de inspector")

questions=[
    ("1","Open een nieuwssite en een techblog. Welk font gebruikt elke site voor de headings?",(124,58,237)),
    ("2","Hoeveel tekststijlen zie je op een typische landingspagina?",(37,99,235)),
    ("3","Meet de regellengte van een alinea. Hoeveel woorden passen er op een regel?",(22,163,74)),
    ("4","Bekijk de letter-spacing van een heading op Apple.com. Voelt het ruimer of strakker dan normaal?",(220,38,38)),
]

for i,(num,text,acc) in enumerate(questions):
    y=Inches(1.85)+i*Inches(1.12)
    cx,cy,cw=card(s12,Inches(0.8),y,Inches(11.5),Inches(0.95),accent=acc)
    lbl(s12,f"{num}.",cx-Inches(0.12),cy,Inches(0.35),Inches(0.55),size=18,c=acc,bold=True)
    lbl(s12,text,cx+Inches(0.28),cy,cw-Inches(0.28),Inches(0.68),size=15,c=BODY)

lbl(s12,"Gebruik de browser inspector (F12) of Figma Inspect panel om lettertypewaarden te zien.",
    Inches(0.8),Inches(6.6),Inches(11.5),Inches(0.5),size=14,c=MUTED,italic=True)

pnr(s12,12)
notitie(s12,"Slide 12: Dit is een actieve opdracht voor in de les. Laat studenten de inspector openen (F12 in browser, dan tabblad Elements of Computed). Ze kunnen CSS properties lezen en font-family, font-size, line-height zien.")

# ─────────────────────────────────────────
# SLIDE 13: Bronnen
# ─────────────────────────────────────────
s13=ns(light_bg())
page_hdr(s13,"Bronnen en verdieping","Verder leren over typografie")

resources=[
    ("Google Fonts Knowledge","fonts.google.com/knowledge","Alles over font keuze en combinaties",ACCENT1),
    ("Typewolf: font inspiratie","typewolf.com","Real-world font gebruik en trending fonts",BLUE),
    ("The Elements of Typographic Style","webtypography.net","Het klassieker online, gratis te lezen",(140,100,220)),
    ("Practical Typography","practicaltypography.com","Praktische typografieregels voor het web",GREEN),
    ("Type Scale generator","typescale.com","Genereer een type scale automatisch",(124,58,237)),
]

for i,(name,url,desc,acc) in enumerate(resources):
    y=Inches(1.85)+i*Inches(1.0)
    cx,cy,cw=card(s13,Inches(0.8),y,Inches(11.5),Inches(0.85),accent=acc)
    lbl(s13,name,cx,cy,Inches(4.5),Inches(0.42),size=16,c=DARK,bold=True)
    lbl(s13,url,cx+Inches(4.6),cy,Inches(4.5),Inches(0.38),size=14,c=acc,italic=True)
    lbl(s13,desc,cx,cy+Inches(0.4),cw,Inches(0.35),size=13,c=MUTED)

pnr(s13,13)
notitie(s13,"Slide 13: Deel deze bronnen via de leeromgeving. fonts.google.com/knowledge is bijzonder toegankelijk voor beginners. Typewolf is inspirerend maar meer voor gevorderde studenten.")

# ─────────────────────────────────────────
# SLIDE 14: Opdracht
# ─────────────────────────────────────────
s14=ns(dark_bg(False))
gbox(s14,Inches(0.0),Inches(0.0),Inches(0.22),Inches(7.5),ACCENT1,ACCENT2,angle=270)

lbl(s14,"Opdracht . Les 3",Inches(1.0),Inches(0.35),Inches(11.0),Inches(0.6),size=22,c=ACCENT2,bold=True)
lbl(s14,"Goed Doel Website . Typografie",Inches(1.0),Inches(0.92),Inches(11.0),Inches(0.5),size=16,c=(180,180,220),italic=True)
gbox(s14,Inches(1.0),Inches(1.42),Inches(2.2),Inches(0.08),ACCENT1,ACCENT2,angle=0)

steps14=[
    "Kies twee Google Fonts: een heading font en een body font",
    "Maak een type scale aan in Figma: H1 t/m caption (minimaal 6 stijlen)",
    "Verwerk de typografie in je wireframe van les 2",
    "Controleer de regellengte van je bodytekst: max 75 karakters",
    "Lever je Figma link in via de leeromgeving voor les 4",
]

for i,text in enumerate(steps14):
    y=Inches(1.65)+i*Inches(0.84)
    dot(s14,Inches(1.35),y+Inches(0.22),Inches(0.24),(124,58,237))
    lbl(s14,str(i+1),Inches(1.2),y+Inches(0.08),Inches(0.28),Inches(0.32),size=12,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s14,text,Inches(1.72),y,Inches(10.5),Inches(0.44),size=17,c=(220,225,255))

# Note at bottom
box(s14,Inches(0.9),Inches(6.6),Inches(11.5),Inches(0.65),(20,30,80))
gbox(s14,Inches(0.9),Inches(6.6),Inches(0.07),Inches(0.65),ACCENT1,ACCENT2,angle=270)
lbl(s14,"Tip: bekijk fonts.google.com/knowledge voor uitleg over welke fonts goed samen werken",Inches(1.1),Inches(6.7),Inches(11.0),Inches(0.48),size=13,c=(180,200,255))

pnr(s14,14,light=True)
notitie(s14,"Slide 14: Loop de opdracht stap voor stap door. Geef studenten de rest van de les om te starten. Les 4 beginnen we met het reviewen van de font keuzes en type scales in groepjes.")

# ─────────────────────────────────────────
# SLIDE 15: Volgende les
# ─────────────────────────────────────────
s15=ns(dark_bg(True))
gbox(s15,Inches(0.0),Inches(0.0),Inches(0.22),Inches(7.5),ACCENT1,ACCENT2,angle=270)

lbl(s15,"Volgende les",Inches(1.0),Inches(0.5),Inches(11.0),Inches(0.9),size=40,c=WHITE,bold=True)
gbox(s15,Inches(1.0),Inches(1.38),Inches(2.5),Inches(0.08),ACCENT1,ACCENT2,angle=0)
lbl(s15,"Kleur en Contrast",Inches(1.0),Inches(1.52),Inches(9.0),Inches(0.65),size=28,c=ACCENT2,bold=True)

next_items=[
    ("WCAG contrast ratio's en toegankelijkheid",(124,58,237)),
    ("Kleurharmonie: aanvullend, triadisch en analoog",(109,40,217)),
    ("Semantische kleuren voor states (success, error, warning)",(91,33,182)),
    ("Kleurpaletten opzetten in Figma",(124,58,237)),
]

for i,(text,col) in enumerate(next_items):
    y=Inches(2.35)+i*Inches(0.82)
    cx2,cy2,cw2=card(s15,Inches(1.0),y,Inches(11.2),Inches(0.68),accent=col,bg=(18,28,72))
    lbl(s15,text,cx2,cy2,cw2,Inches(0.48),size=15,c=(210,215,255))

# Bottom tip
box(s15,Inches(0.9),Inches(6.55),Inches(11.5),Inches(0.7),(20,30,80))
gbox(s15,Inches(0.9),Inches(6.55),Inches(0.07),Inches(0.7),ACCENT1,ACCENT2,angle=270)
lbl(s15,"Kijk alvast naar coolors.co en kies drie kleuren die passen bij het thema van jouw goed doel",Inches(1.1),Inches(6.66),Inches(11.0),Inches(0.52),size=13,c=(180,200,255))

pnr(s15,15,light=True)
notitie(s15,"Slide 15: Geef de tip voor de volgende les expliciet mee. Studenten die thuis coolors.co verkennen komen beter voorbereid naar les 4. Herinner ze aan de deadline voor de Figma link van de typografie opdracht.")

# ─────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────
OUT=r"C:\Users\Peter Marcelis\Documents\Lesmateriaal\design\periode-1\les-3\03_typografie-legibility.pptx"
prs.save(OUT)
print(f"Presentatie opgeslagen: {OUT}")
print(f"Aantal slides: {len(prs.slides)}")
