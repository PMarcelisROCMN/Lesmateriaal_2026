"""
Les 6: Project Afronding Periode 1
Goed Doel Website - Review en Oplevering
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image, ImageDraw
import numpy as np, io

# ─── Kleurpalet ───────────────────────────────────────────────────────────────
NAVY   = (6,  14,  44)
NAVY2  = (30,  58, 138)
BLUE   = (37,  99, 235)
BLUE2  = (59, 130, 246)
CYAN   = (6,  182, 212)
TEAL   = (13, 148, 136)
WHITE  = (255, 255, 255)
OFFWH  = (248, 250, 252)
DARK   = (15,  23,  42)
BODY   = (51,  65,  85)
MUTED  = (100, 116, 135)
GREEN  = (22,  163,  74)
RED    = (220,  38,  38)
ACCENT1 = (225, 29,  72)   # Rose
ACCENT2 = (251, 113, 133)  # Pink
PURPLE  = (124,  58, 237)
ORANGE  = (234,  88,  12)
AMBER   = (217, 119,   6)

def rgb(c): return RGBColor(c[0], c[1], c[2])

SW, SH = Inches(13.33), Inches(7.5)
BW, BH = 1920, 1080

# ─── Achtergrond helpers ───────────────────────────────────────────────────────
def garr(w, h, c1, c2, d='diag'):
    x = np.linspace(0, 1, w, dtype=np.float32)
    y = np.linspace(0, 1, h, dtype=np.float32)
    if d == 'h':
        t = np.tile(x, (h, 1))
    elif d == 'v':
        t = np.tile(y.reshape(-1, 1), (1, w))
    else:
        xx, yy = np.meshgrid(x, y)
        t = xx * .55 + yy * .45
    a = np.zeros((h, w, 3), dtype=np.float32)
    for i in range(3):
        a[:, :, i] = c1[i] + (c2[i] - c1[i]) * t
    return np.clip(a, 0, 255).astype(np.uint8)

def dark_bg(circ=True):
    img = Image.fromarray(garr(BW, BH, NAVY, NAVY2)).convert('RGBA')
    if circ:
        ov = Image.new('RGBA', (BW, BH), (0, 0, 0, 0))
        d  = ImageDraw.Draw(ov)
        d.ellipse([BW-500, -180, BW+200, 520],
                  fill=(ACCENT1[0], ACCENT1[1], ACCENT1[2], 20))
        d.ellipse([-60, BH-360, 380, BH+100],
                  fill=(ACCENT2[0], ACCENT2[1], ACCENT2[2], 22))
        img = Image.alpha_composite(img, ov)
    return img.convert('RGB')

def light_bg():
    img = Image.fromarray(garr(BW, BH, OFFWH, WHITE)).convert('RGBA')
    ov  = Image.new('RGBA', (BW, BH), (0, 0, 0, 0))
    d   = ImageDraw.Draw(ov)
    d.ellipse([BW-580, -220, BW+180, 540],
              fill=(ACCENT1[0], ACCENT1[1], ACCENT1[2], 8))
    img = Image.alpha_composite(img, ov)
    return img.convert('RGB')

def buf(img):
    b = io.BytesIO()
    img.save(b, format='PNG')
    b.seek(0)
    return b

# ─── Presentatie-object ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ─── Basis teken-helpers ───────────────────────────────────────────────────────
def ns(bg=None):
    s = prs.slides.add_slide(BLANK)
    if bg:
        p = s.shapes.add_picture(buf(bg), 0, 0, SW, SH)
        t = s.shapes._spTree
        t.remove(p._element)
        t.insert(2, p._element)
    else:
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = rgb(WHITE)
    return s

def box(s, l, t, w, h, c, border=None):
    sh = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(c)
    if border:
        sh.line.color.rgb = rgb(border)
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh

def gbox(s, l, t, w, h, c1, c2, angle=90):
    sh = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, t, w, h)
    sh.fill.gradient()
    st = sh.fill.gradient_stops
    st[0].position = 0.0; st[0].color.rgb = rgb(c1)
    st[1].position = 1.0; st[1].color.rgb = rgb(c2)
    sh.fill.gradient_angle = angle
    sh.line.fill.background()
    return sh

def dot(s, cx, cy, r, c):
    sh = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, cx-r, cy-r, r*2, r*2)
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(c)
    sh.line.fill.background()

def lbl(s, text, l, t, w, h, size=18, c=BODY, bold=False, italic=False,
        align=PP_ALIGN.LEFT):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text           = text
    r.font.name      = "Calibri"
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = rgb(c)

def blist(s, items, l, t, w, h, size=18, c=BODY, gap=7):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(gap)
        r = p.add_run()
        r.text           = f"  {chr(8226)}  {item}"
        r.font.name      = "Calibri"
        r.font.size      = Pt(size)
        r.font.color.rgb = rgb(c)

def notitie(s, t):
    s.notes_slide.notes_text_frame.text = t

def pnr(s, n, tot=14, light=False):
    lbl(s, f"{n} / {tot}",
        Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.35),
        size=11, c=(180, 180, 200) if light else MUTED,
        align=PP_ALIGN.RIGHT)

def hero_hdr(s, title, sub=None):
    lbl(s, title,
        Inches(1.0), Inches(0.5), Inches(11.3), Inches(1.0),
        size=36, c=WHITE, bold=True)
    gbox(s, Inches(1.0), Inches(1.42), Inches(2.0), Inches(0.08),
         ACCENT1, ACCENT2, angle=0)
    if sub:
        lbl(s, sub,
            Inches(1.0), Inches(1.55), Inches(11.3), Inches(0.42),
            size=15, c=(255, 200, 210), italic=True)

def page_hdr(s, title, sub=None):
    lbl(s, title,
        Inches(1.0), Inches(0.5), Inches(11.3), Inches(0.95),
        size=34, c=DARK, bold=True)
    gbox(s, Inches(1.0), Inches(1.38), Inches(1.8), Inches(0.08),
         ACCENT1, ACCENT2, angle=0)
    if sub:
        lbl(s, sub,
            Inches(1.0), Inches(1.5), Inches(11.3), Inches(0.38),
            size=14, c=MUTED, italic=True)

def card(s, l, t, w, h, accent=None, bg=OFFWH):
    a    = accent if accent else ACCENT1
    STRIP = Inches(0.07)
    box(s, l, t, w, h, bg, border=(215, 225, 240))
    gbox(s, l, t, STRIP, h, a, ACCENT2, angle=270)
    return l + STRIP + Inches(0.18), t + Inches(0.14), w - STRIP - Inches(0.35)

def col_kop(s, l, t, w, h, text, c):
    box(s, l, t, w, h, c)
    lbl(s, text, l, t + Inches(0.07), w, h - Inches(0.07),
        size=20, c=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Titel
# ══════════════════════════════════════════════════════════════════════════════
s1 = ns(dark_bg(True))

# Rose accent strip left
gbox(s1, Inches(0), Inches(0), Inches(0.18), SH, ACCENT1, ACCENT2, angle=270)

# Les badge
box(s1, Inches(1.0), Inches(2.3), Inches(1.0), Inches(0.38), ACCENT1)
lbl(s1, "Les 6", Inches(1.0), Inches(2.3), Inches(1.0), Inches(0.38),
    size=14, c=WHITE, bold=True, align=PP_ALIGN.CENTER)

lbl(s1, "Project Afronding",
    Inches(1.0), Inches(2.82), Inches(11.0), Inches(0.9),
    size=52, c=WHITE, bold=True)
lbl(s1, "Periode 1",
    Inches(1.0), Inches(3.62), Inches(11.0), Inches(0.7),
    size=42, c=ACCENT2, bold=True)

gbox(s1, Inches(1.0), Inches(4.38), Inches(3.0), Inches(0.07),
     ACCENT1, ACCENT2, angle=0)

lbl(s1, "Goed Doel Website  |  Review en Oplevering",
    Inches(1.0), Inches(4.52), Inches(11.0), Inches(0.5),
    size=16, c=(255, 200, 210), italic=True)

dot(s1, Inches(11.5), Inches(1.8), Inches(0.28), (ACCENT1[0], ACCENT1[1], ACCENT1[2]))
dot(s1, Inches(12.1), Inches(2.3), Inches(0.16), (ACCENT2[0], ACCENT2[1], ACCENT2[2]))
dot(s1, Inches(11.0), Inches(5.8), Inches(0.2),  (ACCENT2[0], ACCENT2[1], ACCENT2[2]))

pnr(s1, 1, light=True)
notitie(s1, "Welkom bij de laatste les van periode 1. Vandaag ronden we het goed doel website project af: we doen een review, geven peer feedback en presenteren het werk.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Terugblik
# ══════════════════════════════════════════════════════════════════════════════
s2 = ns(light_bg())
page_hdr(s2, "Wat hebben we geleerd",
         sub="Vijf lessen design, van principes tot componenten")

lessons = [
    ("Les 1", "Wat is design, designprincipes en Figma kennismaking",        TEAL),
    ("Les 2", "Whitespace, alignment en het 8pt grid systeem",               BLUE),
    ("Les 3", "Typografie, type scale en tekststijlen in Figma",             PURPLE),
    ("Les 4", "Kleurtheorie, WCAG contrast en kleurstijlen in Figma",        ORANGE),
    ("Les 5", "Componenten, atomic design en Auto Layout",                   GREEN),
]

TOP_START = Inches(1.85)
CARD_H    = Inches(0.72)
GAP       = Inches(0.12)
CL        = Inches(1.0)
CW        = Inches(11.3)

for i, (les, txt, acc) in enumerate(lessons):
    t = TOP_START + i * (CARD_H + GAP)
    il, it, iw = card(s2, CL, t, CW, CARD_H, accent=acc)
    # Label badge
    box(s2, il, it + Inches(0.06), Inches(0.72), Inches(0.36), acc)
    lbl(s2, les, il, it + Inches(0.06), Inches(0.72), Inches(0.36),
        size=13, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
    lbl(s2, txt, il + Inches(0.82), it + Inches(0.09), iw - Inches(0.82), Inches(0.45),
        size=16, c=BODY)

pnr(s2, 2)
notitie(s2, "Loop kort langs elke les. Vraag: wat herinner je je nog van les 1? Wat was het moeilijkst?")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Project checklist
# ══════════════════════════════════════════════════════════════════════════════
s3 = ns(light_bg())
page_hdr(s3, "Goed doel website: requirements",
         sub="Controleer of je ontwerp aan alle eisen voldoet")

COL_W   = Inches(3.7)
COL_H   = Inches(5.25)
COL_TOP = Inches(1.72)
KOP_H   = Inches(0.48)
BODY_H  = COL_H - KOP_H

cols = [
    ("Inhoud",  NAVY2, Inches(0.75),
     ["Logo placeholder", "Hero heading en subtitel", "CTA knop",
      "Navigatiemenu met 3 links", "3 testimonial-kaarten", "Footer"]),
    ("Design",  ACCENT1, Inches(0.75) + COL_W + Inches(0.45),
     ["Kleurpalet met kleurstijlen", "Type scale met tekststijlen",
      "8pt grid toegepast", "WCAG AA contrast op alle tekst"]),
    ("Figma",   BLUE,   Inches(0.75) + 2*(COL_W + Inches(0.45)),
     ["Pagina gestructureerd in frames", "Lagen correct benoemd",
      "Componenten aangemaakt", "Varianten toegepast"]),
]

for label, c, cl, items in cols:
    col_kop(s3, cl, COL_TOP, COL_W, KOP_H, label, c)
    box(s3, cl, COL_TOP + KOP_H, COL_W, BODY_H, OFFWH, border=(215, 225, 240))
    blist(s3, items,
          cl + Inches(0.18), COL_TOP + KOP_H + Inches(0.16),
          COL_W - Inches(0.36), BODY_H - Inches(0.32),
          size=15, c=BODY, gap=9)

pnr(s3, 3)
notitie(s3, "Studenten kunnen hier met een pen of mentaal afvinken. Loop de drie kolommen langs en vraag of er punten zijn waarbij ze twijfelen.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Hero sectie review
# ══════════════════════════════════════════════════════════════════════════════
s4 = ns(light_bg())
page_hdr(s4, "Hero sectie: review criteria",
         sub="Wat maakt een sterke hero?")

# ── Wireframe (links) ──
WF_L = Inches(0.85)
WF_T = Inches(1.72)
WF_W = Inches(5.5)
WF_H = Inches(5.1)

# Outer frame
box(s4, WF_L, WF_T, WF_W, WF_H, (235, 240, 250), border=DARK)

# Label "1440px frame"
lbl(s4, "1440px frame",
    WF_L + Inches(0.1), WF_T + Inches(0.05), WF_W - Inches(0.2), Inches(0.28),
    size=10, c=MUTED, italic=True)

# Nav bar
NAV_H = Inches(0.48)
box(s4, WF_L, WF_T, WF_W, NAV_H, DARK)
lbl(s4, "Navigatie", WF_L + Inches(0.1), WF_T + Inches(0.1),
    WF_W - Inches(0.2), NAV_H - Inches(0.1),
    size=11, c=(200, 210, 230))

# Hero bg area
HB_T = WF_T + NAV_H + Inches(0.1)
HB_H = Inches(2.8)
box(s4, WF_L + Inches(0.1), HB_T, WF_W - Inches(0.2), HB_H, (200, 210, 230))
lbl(s4, "Hero achtergrond (foto / kleur)",
    WF_L + Inches(0.3), HB_T + Inches(0.12), WF_W - Inches(0.5), Inches(0.35),
    size=11, c=MUTED, italic=True)

# Heading placeholder
HD_T = HB_T + Inches(0.5)
box(s4, WF_L + Inches(0.3), HD_T, WF_W - Inches(0.6), Inches(0.45), (150, 170, 210))
lbl(s4, "H1 Heading",
    WF_L + Inches(0.35), HD_T + Inches(0.07), WF_W - Inches(0.7), Inches(0.3),
    size=12, c=WHITE, bold=True)

# Subtitle placeholder
ST_T = HD_T + Inches(0.54)
box(s4, WF_L + Inches(0.3), ST_T, WF_W - Inches(0.9), Inches(0.28), (180, 195, 225))
lbl(s4, "Subtitel / missieomschrijving",
    WF_L + Inches(0.35), ST_T + Inches(0.04), WF_W - Inches(1.0), Inches(0.24),
    size=10, c=(80, 90, 110), italic=True)

# CTA button
CT_T = ST_T + Inches(0.4)
box(s4, WF_L + Inches(0.3), CT_T, Inches(1.4), Inches(0.36), ACCENT1)
lbl(s4, "Doneer nu",
    WF_L + Inches(0.3), CT_T + Inches(0.06), Inches(1.4), Inches(0.28),
    size=12, c=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ── Criteria (rechts) ──
CR_L = WF_L + WF_W + Inches(0.55)
CR_T = Inches(1.85)
CR_W = Inches(6.3)

lbl(s4, "Review criteria",
    CR_L, CR_T, CR_W, Inches(0.4),
    size=16, c=DARK, bold=True)

criteria = [
    "Duidelijke H1 heading met maximaal 8 woorden",
    "Subtitel die de missie uitlegt in 1-2 zinnen",
    "CTA knop met actiewerkwoord: Doneer nu, Lees meer, Meld je aan",
    "Voldoende contrast tussen tekst en achtergrond",
    "Hero-hoogte minimaal 50vh (op Desktop: 500 tot 700px)",
]

for i, cr in enumerate(criteria):
    cy = CR_T + Inches(0.52) + i * Inches(0.8)
    dot(s4, CR_L + Inches(0.22), cy + Inches(0.2), Inches(0.12), ACCENT1)
    lbl(s4, cr, CR_L + Inches(0.52), cy, CR_W - Inches(0.6), Inches(0.68),
        size=15, c=BODY)

pnr(s4, 4)
notitie(s4, "Laat studenten hun eigen hero openen in Figma. Loop door elk criterium en laat ze aangeven of ze eraan voldoen.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Navigatie review
# ══════════════════════════════════════════════════════════════════════════════
s5 = ns(light_bg())
page_hdr(s5, "Navigatie: review criteria",
         sub="Consistente navigatie is de basis van UX")

# ── Nav wireframe (links) ──
WF_L = Inches(0.85)
WF_T = Inches(1.85)
WF_W = Inches(5.5)
NAV_H = Inches(0.72)  # 64px proportional

# Outer nav container
box(s5, WF_L, WF_T, WF_W, NAV_H, DARK)

# Logo placeholder
box(s5, WF_L + Inches(0.16), WF_T + Inches(0.14),
    Inches(0.88), Inches(0.44), (80, 100, 150))
lbl(s5, "Logo", WF_L + Inches(0.16), WF_T + Inches(0.18),
    Inches(0.88), Inches(0.36),
    size=12, c=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Nav links
link_labels = ["Home", "Over ons", "Doneer"]
for j, lk in enumerate(link_labels):
    lx = WF_L + Inches(1.3) + j * Inches(1.1)
    lbl(s5, lk, lx, WF_T + Inches(0.22), Inches(1.0), Inches(0.3),
        size=11, c=(200, 215, 240), align=PP_ALIGN.CENTER)
    # Active state underline on first link
    if j == 0:
        box(s5, lx + Inches(0.2), WF_T + NAV_H - Inches(0.08),
            Inches(0.6), Inches(0.06), ACCENT1)

# CTA button rechts in nav
box(s5, WF_L + WF_W - Inches(1.4), WF_T + Inches(0.16),
    Inches(1.22), Inches(0.4), ACCENT1)
lbl(s5, "Doneer nu",
    WF_L + WF_W - Inches(1.4), WF_T + Inches(0.22),
    Inches(1.22), Inches(0.3),
    size=10, c=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Annotation labels
lbl(s5, "Logo links",
    WF_L, WF_T + NAV_H + Inches(0.1), Inches(1.1), Inches(0.3),
    size=10, c=MUTED, italic=True)
lbl(s5, "3-5 links",
    WF_L + Inches(1.3), WF_T + NAV_H + Inches(0.1), Inches(2.0), Inches(0.3),
    size=10, c=MUTED, italic=True)
lbl(s5, "CTA rechts",
    WF_L + WF_W - Inches(1.5), WF_T + NAV_H + Inches(0.1), Inches(1.4), Inches(0.3),
    size=10, c=MUTED, italic=True)

# Hoogte annotatie
lbl(s5, "Hoogte: 64-80px",
    WF_L, WF_T + NAV_H + Inches(0.48), WF_W, Inches(0.3),
    size=11, c=BLUE, italic=True)

# ── Criteria (rechts) ──
CR_L = WF_L + WF_W + Inches(0.55)
CR_T = Inches(1.85)
CR_W = Inches(6.3)

lbl(s5, "Review criteria",
    CR_L, CR_T, CR_W, Inches(0.4),
    size=16, c=DARK, bold=True)

criteria = [
    "Logo of merkidentiteit links in de navigatie",
    "3 tot 5 navigatielinks, geen meer",
    "Actieve pagina visueel gemarkeerd",
    "CTA knop of link rechts in de navigatie",
    "Consistent op alle paginas: gebruik een component",
    "Hoogte 64 tot 80px (veelvoud van 8)",
]

for i, cr in enumerate(criteria):
    cy = CR_T + Inches(0.52) + i * Inches(0.72)
    dot(s5, CR_L + Inches(0.22), cy + Inches(0.2), Inches(0.12), ACCENT1)
    lbl(s5, cr, CR_L + Inches(0.52), cy, CR_W - Inches(0.6), Inches(0.6),
        size=15, c=BODY)

pnr(s5, 5)
notitie(s5, "Veel voorkomende fouten: te veel links, geen actieve staat, CTA ontbreekt in de nav. Wijs dit aan in de wireframe.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – CTA review
# ══════════════════════════════════════════════════════════════════════════════
s6 = ns(light_bg())
page_hdr(s6, "Call-to-Action: review criteria",
         sub="Een goede CTA is duidelijk, contrastrijk en prominent")

PANEL_T = Inches(1.85)
PANEL_H = Inches(5.1)
PANEL_W = Inches(5.6)
KOP_H   = Inches(0.52)

# ── Goede CTA (links) ──
GL = Inches(0.75)
box(s6, GL, PANEL_T, PANEL_W, PANEL_H, OFFWH, border=(215, 225, 240))
col_kop(s6, GL, PANEL_T, PANEL_W, KOP_H, "Goede CTA", GREEN)

good_items = [
    "Actief werkwoord: Doneer nu, Meld je aan",
    "Contraststerke kleur t.o.v. achtergrond",
    "Voldoende padding (min. 16x48px)",
    "Slechts een primaire actie per pagina",
    "Witruimte eromheen: geeft lucht en focus",
]
blist(s6, good_items,
      GL + Inches(0.22), PANEL_T + KOP_H + Inches(0.2),
      PANEL_W - Inches(0.44), PANEL_H - KOP_H - Inches(0.4),
      size=15, c=BODY, gap=10)

# CTA voorbeeld knop in goede panel
box(s6, GL + Inches(1.5), PANEL_T + PANEL_H - Inches(0.88),
    Inches(2.2), Inches(0.52), GREEN)
lbl(s6, "Doneer nu",
    GL + Inches(1.5), PANEL_T + PANEL_H - Inches(0.84),
    Inches(2.2), Inches(0.44),
    size=15, c=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ── Zwakke CTA (rechts) ──
RL = GL + PANEL_W + Inches(0.6)
box(s6, RL, PANEL_T, PANEL_W, PANEL_H, OFFWH, border=(215, 225, 240))
col_kop(s6, RL, PANEL_T, PANEL_W, KOP_H, "Zwakke CTA", RED)

weak_items = [
    "Vaag label zoals 'klik hier' of 'meer info'",
    "Te weinig padding, knop voelt klein aan",
    "Verloren tussen andere elementen",
    "Zelfde kleur als de achtergrond",
]
blist(s6, weak_items,
      RL + Inches(0.22), PANEL_T + KOP_H + Inches(0.2),
      PANEL_W - Inches(0.44), PANEL_H - KOP_H - Inches(0.4),
      size=15, c=BODY, gap=10)

# Zwakke CTA voorbeeld knop
box(s6, RL + Inches(1.5), PANEL_T + PANEL_H - Inches(0.88),
    Inches(2.2), Inches(0.52), (215, 225, 240))
lbl(s6, "Klik hier",
    RL + Inches(1.5), PANEL_T + PANEL_H - Inches(0.84),
    Inches(2.2), Inches(0.44),
    size=15, c=MUTED, align=PP_ALIGN.CENTER)

pnr(s6, 6)
notitie(s6, "Laat studenten in tweetallen elkaars CTA beoordelen. Is het label actief? Is er genoeg contrast en witruimte?")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Testimonials review
# ══════════════════════════════════════════════════════════════════════════════
s7 = ns(light_bg())
page_hdr(s7, "Testimonials: review criteria",
         sub="Drie consistente kaarten bouwen vertrouwen op")

# ── Wireframe 3 testimonial cards (links) ──
WF_L = Inches(0.75)
WF_T = Inches(1.85)
WF_W = Inches(5.6)
WF_H = Inches(5.0)

CARD_W = Inches(1.6)
CARD_H = Inches(3.8)
CARD_GAP = Inches(0.2)
CARD_T   = WF_T + Inches(0.6)
CARD_L   = WF_L + Inches(0.1)

for ci in range(3):
    cl = CARD_L + ci * (CARD_W + CARD_GAP)
    # Card body
    box(s7, cl, CARD_T, CARD_W, CARD_H, WHITE, border=(210, 220, 235))
    # Avatar circle
    dot(s7, cl + CARD_W/2, CARD_T + Inches(0.48), Inches(0.32), (180, 195, 220))
    lbl(s7, "AB", cl + CARD_W/2 - Inches(0.32),
        CARD_T + Inches(0.22), Inches(0.64), Inches(0.42),
        size=13, c=DARK, bold=True, align=PP_ALIGN.CENTER)
    # Name
    box(s7, cl + Inches(0.18), CARD_T + Inches(0.96),
        CARD_W - Inches(0.36), Inches(0.22), (200, 210, 230))
    # Function
    box(s7, cl + Inches(0.25), CARD_T + Inches(1.24),
        CARD_W - Inches(0.5), Inches(0.18), (215, 225, 240))
    # Quote lines
    for ql in range(3):
        box(s7, cl + Inches(0.18),
            CARD_T + Inches(1.55) + ql * Inches(0.32),
            CARD_W - Inches(0.36) - ql * Inches(0.2), Inches(0.18),
            (225, 232, 242))

# Section label
lbl(s7, "3 gelijke testimonial-kaarten",
    WF_L, WF_T + Inches(0.1), WF_W, Inches(0.4),
    size=11, c=MUTED, italic=True)

# Gap annotation
box(s7, CARD_L + CARD_W, CARD_T + Inches(1.5),
    CARD_GAP, Inches(0.06), ACCENT1)
lbl(s7, "gap\n24px",
    CARD_L + CARD_W, CARD_T + Inches(1.62), CARD_GAP + Inches(0.3), Inches(0.4),
    size=9, c=ACCENT1, italic=True)

# ── Criteria (rechts) ──
CR_L = WF_L + WF_W + Inches(0.55)
CR_T = Inches(1.85)
CR_W = Inches(6.3)

lbl(s7, "Review criteria",
    CR_L, CR_T, CR_W, Inches(0.4),
    size=16, c=DARK, bold=True)

criteria = [
    "Drie gelijke kaarten: consistente breedte en hoogte",
    "Avatar of initialen placeholder (cirkel)",
    "Naam en functie van de getuige",
    "Quote tekst in italic, max 2-3 zinnen",
    "Consistent gebruik van typografiestijlen",
    "Witruimte tussen de kaarten: 24px of 32px gap",
]

for i, cr in enumerate(criteria):
    cy = CR_T + Inches(0.52) + i * Inches(0.72)
    dot(s7, CR_L + Inches(0.22), cy + Inches(0.2), Inches(0.12), ACCENT1)
    lbl(s7, cr, CR_L + Inches(0.52), cy, CR_W - Inches(0.6), Inches(0.6),
        size=15, c=BODY)

pnr(s7, 7)
notitie(s7, "Let op: zijn de drie kaarten echt gelijk van hoogte? Zijn ze een component in Figma? Dat is de technische eis.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Peer review (FIST framework)
# ══════════════════════════════════════════════════════════════════════════════
s8 = ns(dark_bg(False))
hero_hdr(s8, "Peer review: hoe geef je goede feedback",
         sub="Gebruik het FIST-framework voor constructieve kritiek")

steps = [
    ("F", "Feit",       "Beschrijf wat je ziet zonder oordeel",                    ACCENT1),
    ("I", "Impact",     "Wat is het effect van deze designkeuze op de gebruiker",   BLUE2),
    ("S", "Suggestie",  "Geef een concrete verbetering",                            TEAL),
    ("T", "Toelichting","Leg uit waarom jouw suggestie werkt",                      GREEN),
]

STEP_T = Inches(2.0)
STEP_H = Inches(0.88)
STEP_W = Inches(11.3)
GAP    = Inches(0.12)

for i, (letter, label, desc, c) in enumerate(steps):
    t = STEP_T + i * (STEP_H + GAP)
    box(s8, Inches(1.0), t, STEP_W, STEP_H, (20, 35, 75))
    # Letter badge
    box(s8, Inches(1.0), t, Inches(0.72), STEP_H, c)
    lbl(s8, letter, Inches(1.0), t + Inches(0.2), Inches(0.72), Inches(0.5),
        size=26, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Label + desc
    lbl(s8, label, Inches(1.88), t + Inches(0.08), Inches(2.0), Inches(0.38),
        size=16, c=WHITE, bold=True)
    lbl(s8, desc,  Inches(1.88), t + Inches(0.44), Inches(9.2), Inches(0.38),
        size=14, c=(180, 195, 225))

# Bottom tip
TIP_T = STEP_T + 4 * (STEP_H + GAP) + Inches(0.18)
box(s8, Inches(1.0), TIP_T, Inches(11.3), Inches(0.64), (10, 25, 60))
lbl(s8, "Koppel je feedback altijd aan designprincipes: whitespace, alignment, contrast, typografie",
    Inches(1.2), TIP_T + Inches(0.12), Inches(10.9), Inches(0.44),
    size=14, c=ACCENT2, italic=True)

pnr(s8, 8, light=True)
notitie(s8, "Oefen het FIST-framework klassikaal: laat een student een scherm tonen en een andere student doorloopt F-I-S-T hardop.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Figma presentatiemodus
# ══════════════════════════════════════════════════════════════════════════════
s9 = ns(dark_bg(False))
hero_hdr(s9, "Figma: presenteren aan de klas",
         sub="Zet je ontwerp klaar voor de review-ronde")

steps = [
    "Zorg dat alle frames een naam hebben en in de juiste volgorde staan",
    "Klik op het Present-icoontje (pijltje rechtsboven) of druk op Ctrl+Alt+Enter",
    "Navigeer tussen frames met de pijltoetsen",
    "Gebruik de prototype-verbindingen voor een interactieve demo",
    "Deel de presentatielink zodat medestudenten kunnen meekijken",
]

STEP_T = Inches(2.0)
STEP_H = Inches(0.75)
STEP_W = Inches(10.5)
GAP    = Inches(0.1)

for i, step in enumerate(steps):
    t = STEP_T + i * (STEP_H + GAP)
    box(s9, Inches(1.0), t, STEP_W, STEP_H, (20, 35, 75))
    # Number badge
    box(s9, Inches(1.0), t, Inches(0.64), STEP_H, ACCENT1)
    lbl(s9, str(i+1), Inches(1.0), t + Inches(0.16), Inches(0.64), Inches(0.44),
        size=20, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
    lbl(s9, step, Inches(1.8), t + Inches(0.18), STEP_W - Inches(0.88), Inches(0.44),
        size=15, c=(200, 215, 240))

# Tip box
TIP_T = STEP_T + 5 * (STEP_H + GAP) + Inches(0.2)
box(s9, Inches(1.0), TIP_T, Inches(11.3), Inches(0.72), (10, 25, 60))
lbl(s9, "Tip: ",
    Inches(1.18), TIP_T + Inches(0.14), Inches(0.7), Inches(0.44),
    size=14, c=ACCENT2, bold=True)
lbl(s9, "Maak een aparte page 'presentatie' met alleen je definitieve schermen op volgorde",
    Inches(1.8), TIP_T + Inches(0.14), Inches(10.3), Inches(0.44),
    size=14, c=(180, 195, 225), italic=True)

pnr(s9, 9, light=True)
notitie(s9, "Demonstreer dit live in Figma als je een beamer hebt. Laat studenten 5 minuten nemen om hun Figma presentatiepagina klaar te zetten.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Beoordelingscriteria
# ══════════════════════════════════════════════════════════════════════════════
s10 = ns(light_bg())
page_hdr(s10, "Beoordelingscriteria",
         sub="Hoe wordt je goed doel website beoordeeld?")

criteria_data = [
    ("Design kwaliteit", "40%", ACCENT1,
     ["Whitespace en alignment consequent toegepast",
      "Typografie: type scale correct gebruikt",
      "Kleur: palet en WCAG contrast kloppen",
      "Principes uit alle 5 lessen zichtbaar"]),
    ("Technische uitvoering", "35%", BLUE,
     ["Correcte Figma frame-structuur",
      "Componenten en stijlen aangemaakt",
      "Lagen correct benoemd en gegroepeerd",
      "Auto Layout toegepast op componenten"]),
    ("Compleetheid", "25%", TEAL,
     ["Alle vereiste secties aanwezig",
      "Opdrachten van alle lessen verwerkt",
      "Figma link ingeleverd voor deadline",
      "Presentatie afgerond tijdens les 6"]),
]

COL_W = Inches(3.7)
COL_T = Inches(1.82)
COL_H = Inches(4.8)
KOP_H = Inches(0.56)

for ci, (title, pct, c, items) in enumerate(criteria_data):
    cl = Inches(0.75) + ci * (COL_W + Inches(0.42))

    # Column header with percentage badge
    box(s10, cl, COL_T, COL_W, KOP_H, c)
    lbl(s10, title, cl + Inches(0.12), COL_T + Inches(0.07),
        COL_W - Inches(1.1), KOP_H - Inches(0.1),
        size=16, c=WHITE, bold=True)
    # Percentage badge (right)
    box(s10, cl + COL_W - Inches(0.9), COL_T + Inches(0.06),
        Inches(0.8), KOP_H - Inches(0.12), WHITE)
    lbl(s10, pct, cl + COL_W - Inches(0.9), COL_T + Inches(0.1),
        Inches(0.8), KOP_H - Inches(0.12),
        size=18, c=c, bold=True, align=PP_ALIGN.CENTER)

    # Body
    box(s10, cl, COL_T + KOP_H, COL_W, COL_H - KOP_H, OFFWH, border=(215, 225, 240))
    blist(s10, items,
          cl + Inches(0.16), COL_T + KOP_H + Inches(0.18),
          COL_W - Inches(0.32), COL_H - KOP_H - Inches(0.36),
          size=14, c=BODY, gap=9)

# Bottom note
NOTE_T = COL_T + COL_H + Inches(0.22)
box(s10, Inches(0.75), NOTE_T, Inches(11.8), Inches(0.48), (240, 244, 252),
    border=(200, 210, 230))
lbl(s10,
    "Lever je Figma link in via de leeromgeving  |  Deadline: einde les 6",
    Inches(0.95), NOTE_T + Inches(0.1), Inches(11.4), Inches(0.32),
    size=13, c=NAVY2, italic=True)

pnr(s10, 10)
notitie(s10, "Bespreek de weging: design kwaliteit telt het zwaarst. Studenten hoeven niet perfect te zijn, maar moeten aantonen dat ze de principes begrijpen.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – Tips voor presenteren
# ══════════════════════════════════════════════════════════════════════════════
s11 = ns(light_bg())
page_hdr(s11, "Tips voor je presentatie",
         sub="Hoe vertel je het verhaal achter je ontwerp?")

tips = [
    "Begin met: wat is het goed doel en wie is de doelgroep",
    "Leg uit waarom je deze kleuren en lettertypes hebt gekozen",
    "Wijs de designkeuzes aan op het scherm: whitespace, alignment, contrast",
    "Vertel wat je lastig vond en hoe je het hebt opgelost",
    "Sluit af met wat je zou verbeteren als je meer tijd had",
]

TIP_T = Inches(1.85)
TIP_H = Inches(0.82)
TIP_W = Inches(11.3)
GAP   = Inches(0.14)

accents = [ACCENT1, BLUE, TEAL, GREEN, PURPLE]

for i, (tip, acc) in enumerate(zip(tips, accents)):
    t = TIP_T + i * (TIP_H + GAP)
    il, it, iw = card(s11, Inches(1.0), t, TIP_W, TIP_H, accent=acc)
    # Step number
    dot(s11, il + Inches(0.22), it + Inches(0.24), Inches(0.2), acc)
    lbl(s11, str(i+1), il + Inches(0.07), it + Inches(0.06), Inches(0.36), Inches(0.36),
        size=14, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
    lbl(s11, tip, il + Inches(0.5), it + Inches(0.08), iw - Inches(0.5), Inches(0.52),
        size=16, c=BODY)

pnr(s11, 11)
notitie(s11, "Geef studenten 2 minuten om deze tips te lezen en dan te beginnen. Presentaties: max 3 minuten per student.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – Bronnen voor zelfstudie
# ══════════════════════════════════════════════════════════════════════════════
s12 = ns(light_bg())
page_hdr(s12, "Bronnen voor zelfstudie",
         sub="Verdiep je kennis buiten de les")

resources = [
    ("Refactoring UI",         "Praktisch boek over UI design door de makers van Tailwind CSS",
     "refactoringui.com",    ACCENT1),
    ("Figma Best Practices",   "Officiele gidsen voor components, stijlen en Auto Layout",
     "figma.com/best-practices", BLUE),
    ("Nielsen Norman Group",   "Onderzoeksgebaseerde UX richtlijnen en artikelen",
     "nngroup.com",          GREEN),
    ("Designprincipes",        "Overzicht van alle bekende designprincipes met uitleg",
     "principles.design",    PURPLE),
    ("Awwwards",               "Inspiratie van de beste designsites ter wereld",
     "awwwards.com",         AMBER),
]

RC_T = Inches(1.82)
RC_H = Inches(0.9)
RC_W = Inches(11.3)
GAP  = Inches(0.1)

for i, (name, desc, url, acc) in enumerate(resources):
    t = RC_T + i * (RC_H + GAP)
    il, it, iw = card(s12, Inches(1.0), t, RC_W, RC_H, accent=acc)
    lbl(s12, name, il, it + Inches(0.02), Inches(2.6), Inches(0.38),
        size=16, c=DARK, bold=True)
    lbl(s12, desc, il + Inches(2.75), it + Inches(0.04), iw - Inches(5.0), Inches(0.38),
        size=14, c=BODY)
    lbl(s12, url,  iw - Inches(2.0),  it + Inches(0.22), Inches(2.3), Inches(0.32),
        size=12, c=acc, italic=True)

pnr(s12, 12)
notitie(s12, "Vertel kort wat je zelf het meest gebruikt. Refactoring UI is een aanrader voor visueel sterke studenten. NN/g voor studenten die geinteresseerd zijn in UX research.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – Volgende periode
# ══════════════════════════════════════════════════════════════════════════════
s13 = ns(dark_bg(True))

# Rose accent strip left
gbox(s13, Inches(0), Inches(0), Inches(0.18), SH, ACCENT1, ACCENT2, angle=270)

lbl(s13, "Volgende periode:",
    Inches(1.0), Inches(0.55), Inches(11.0), Inches(0.6),
    size=22, c=(180, 195, 230))
lbl(s13, "Portfolio Website",
    Inches(1.0), Inches(1.08), Inches(11.0), Inches(0.88),
    size=44, c=WHITE, bold=True)
gbox(s13, Inches(1.0), Inches(1.9), Inches(3.0), Inches(0.08),
     ACCENT1, ACCENT2, angle=0)

# Preview cards 2x2
preview = [
    ("Responsive design en mobile first",                    ACCENT1),
    ("JavaScript slider en interactieve elementen",          BLUE),
    ("About pagina en projectenpagina",                      TEAL),
    ("Consistent stijlsysteem over meerdere paginas",        GREEN),
]

PC_W = Inches(5.2)
PC_H = Inches(0.78)
PC_T = Inches(2.1)
GAP  = Inches(0.16)

for i, (txt, acc) in enumerate(preview):
    row = i // 2
    col = i % 2
    pl = Inches(1.0) + col * (PC_W + Inches(0.7))
    pt = PC_T + row * (PC_H + GAP)
    box(s13, pl, pt, PC_W, PC_H, (20, 35, 80))
    gbox(s13, pl, pt, Inches(0.08), PC_H, acc, ACCENT2, angle=270)
    dot(s13, pl + Inches(0.35), pt + PC_H/2, Inches(0.1), acc)
    lbl(s13, txt, pl + Inches(0.52), pt + Inches(0.2),
        PC_W - Inches(0.68), Inches(0.44),
        size=15, c=(200, 215, 240))

# Bottom note
NOTE_T = PC_T + 2*(PC_H + GAP) + Inches(0.28)
box(s13, Inches(1.0), NOTE_T, Inches(11.3), Inches(0.58), (10, 22, 60))
lbl(s13,
    "Periode 2 start over twee weken. Bekijk alvast portfoliosites van designers voor inspiratie.",
    Inches(1.2), NOTE_T + Inches(0.12), Inches(10.9), Inches(0.38),
    size=14, c=(180, 195, 225), italic=True)

pnr(s13, 13, light=True)
notitie(s13, "Geef studenten een voorproefje van periode 2. Noem een paar inspiratiebronnen: awwwards.com, dribbble.com. Vraag welk onderdeel ze het meest interessant lijkt.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – Afsluiting
# ══════════════════════════════════════════════════════════════════════════════
s14 = ns(dark_bg(True))

# Large accent strip left
gbox(s14, Inches(0), Inches(0), Inches(0.25), SH, ACCENT1, ACCENT2, angle=270)

# Decorative dots
dot(s14, Inches(11.8), Inches(1.2), Inches(0.35), (ACCENT1[0], ACCENT1[1], ACCENT1[2]))
dot(s14, Inches(12.5), Inches(2.0), Inches(0.2),  (ACCENT2[0], ACCENT2[1], ACCENT2[2]))
dot(s14, Inches(11.2), Inches(6.2), Inches(0.22), (ACCENT2[0], ACCENT2[1], ACCENT2[2]))

lbl(s14, "Gefeliciteerd!",
    Inches(1.0), Inches(0.8), Inches(11.0), Inches(1.2),
    size=60, c=WHITE, bold=True)

gbox(s14, Inches(1.0), Inches(1.88), Inches(3.5), Inches(0.1),
     ACCENT1, ACCENT2, angle=0)

lbl(s14, "Je eerste professionele designproject is klaar.",
    Inches(1.0), Inches(2.05), Inches(11.0), Inches(0.55),
    size=22, c=ACCENT2)

achievements = [
    "Je hebt ontworpen met intentie, niet toevallig",
    "Je begrijpt whitespace, typografie, kleur en componenten",
    "Je Figma skills zijn een solide basis voor alles wat volgt",
]

ACH_T = Inches(2.78)
for i, ach in enumerate(achievements):
    t = ACH_T + i * Inches(0.62)
    dot(s14, Inches(1.22), t + Inches(0.22), Inches(0.12), ACCENT1)
    lbl(s14, ach, Inches(1.5), t, Inches(10.5), Inches(0.52),
        size=18, c=(210, 220, 245))

# Bottom discussion card
DISC_T = Inches(5.45)
box(s14, Inches(1.0), DISC_T, Inches(11.3), Inches(0.88), (15, 30, 70))
gbox(s14, Inches(1.0), DISC_T, Inches(0.1), Inches(0.88),
     ACCENT1, ACCENT2, angle=270)
lbl(s14, "Vraag: ",
    Inches(1.22), DISC_T + Inches(0.16), Inches(0.9), Inches(0.5),
    size=15, c=ACCENT2, bold=True)
lbl(s14, "Welk principe vond jij het moeilijkst te toepassen en waarom?",
    Inches(2.05), DISC_T + Inches(0.16), Inches(10.1), Inches(0.5),
    size=15, c=(190, 205, 235), italic=True)

pnr(s14, 14, light=True)
notitie(s14, "Sluit de les positief af. Laat iedereen kort antwoord geven op de vraag. Dit is ook een goede evaluatie voor jou als docent: welke onderwerpen verdienen meer aandacht in periode 2?")


# ─── Opslaan ──────────────────────────────────────────────────────────────────
OUTPUT = r"C:\Users\Peter Marcelis\Documents\Lesmateriaal\design\periode-1\les-6\06_project-afronding-p1.pptx"
prs.save(OUTPUT)
print(f"Presentatie opgeslagen: {OUTPUT}")
print(f"Aantal slides: {len(prs.slides)}")
