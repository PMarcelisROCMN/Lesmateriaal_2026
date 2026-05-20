"""
Les 7: Responsive Design en Mobile First
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image, ImageDraw
import numpy as np, io

NAVY=(6,14,44); NAVY2=(30,58,138); BLUE=(37,99,235); BLUE2=(59,130,246)
CYAN=(6,182,212); TEAL=(13,148,136); WHITE=(255,255,255); OFFWH=(248,250,252)
DARK=(15,23,42); BODY=(51,65,85); MUTED=(100,116,135)
GREEN=(22,163,74); RED=(220,38,38)
ACCENT1=(79,70,229); ACCENT2=(129,140,248)  # Indigo for les 7

def rgb(c): return RGBColor(c[0],c[1],c[2])
SW,SH=Inches(13.33),Inches(7.5); BW,BH=1920,1080

def garr(w,h,c1,c2,d='diag'):
    x,y=np.linspace(0,1,w,dtype=np.float32),np.linspace(0,1,h,dtype=np.float32)
    if d=='h': t=np.tile(x,(h,1))
    elif d=='v': t=np.tile(y.reshape(-1,1),(1,w))
    else: xx,yy=np.meshgrid(x,y); t=xx*.55+yy*.45
    a=np.zeros((h,w,3),dtype=np.float32)
    for i in range(3): a[:,:,i]=c1[i]+(c2[i]-c1[i])*t
    return np.clip(a,0,255).astype(np.uint8)

def dark_bg(circ=True):
    img=Image.fromarray(garr(BW,BH,NAVY,NAVY2)).convert('RGBA')
    if circ:
        ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
        d.ellipse([BW-500,-180,BW+200,520],fill=(ACCENT1[0],ACCENT1[1],ACCENT1[2],20))
        d.ellipse([-60,BH-360,380,BH+100],fill=(ACCENT2[0],ACCENT2[1],ACCENT2[2],22))
        img=Image.alpha_composite(img,ov)
    return img.convert('RGB')

def light_bg():
    img=Image.fromarray(garr(BW,BH,OFFWH,WHITE)).convert('RGBA')
    ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
    d.ellipse([BW-580,-220,BW+180,540],fill=(ACCENT1[0],ACCENT1[1],ACCENT1[2],8))
    img=Image.alpha_composite(img,ov); return img.convert('RGB')

def buf(img):
    b=io.BytesIO(); img.save(b,format='PNG'); b.seek(0); return b

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
BLANK=prs.slide_layouts[6]

def ns(bg=None):
    s=prs.slides.add_slide(BLANK)
    if bg:
        p=s.shapes.add_picture(buf(bg),0,0,SW,SH)
        t=s.shapes._spTree; t.remove(p._element); t.insert(2,p._element)
    else:
        s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(WHITE)
    return s

def box(s,l,t,w,h,c,border=None):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c)
    if border: sh.line.color.rgb=rgb(border); sh.line.width=Pt(0.75)
    else: sh.line.fill.background()
    return sh

def gbox(s,l,t,w,h,c1,c2,angle=90):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.gradient(); st=sh.fill.gradient_stops
    st[0].position=0.0; st[0].color.rgb=rgb(c1)
    st[1].position=1.0; st[1].color.rgb=rgb(c2)
    sh.fill.gradient_angle=angle; sh.line.fill.background(); return sh

def dot(s,cx,cy,r,c):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,cx-r,cy-r,r*2,r*2)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c); sh.line.fill.background()

def lbl(s,text,l,t,w,h,size=18,c=BODY,bold=False,italic=False,align=PP_ALIGN.LEFT):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=text; r.font.name="Calibri"; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=rgb(c)

def blist(s,items,l,t,w,h,size=18,c=BODY,gap=7):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(gap); r=p.add_run()
        r.text=f"  {chr(8226)}  {item}"; r.font.name="Calibri"; r.font.size=Pt(size); r.font.color.rgb=rgb(c)

def notitie(s,t): s.notes_slide.notes_text_frame.text=t

def pnr(s,n,tot=16,light=False):
    lbl(s,f"{n} / {tot}",Inches(12.5),Inches(7.1),Inches(0.8),Inches(0.35),size=11,c=(180,180,200) if light else MUTED,align=PP_ALIGN.RIGHT)

def hero_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(1.0),size=36,c=WHITE,bold=True)
    gbox(s,Inches(1.0),Inches(1.42),Inches(2.0),Inches(0.08),ACCENT1,ACCENT2,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.55),Inches(11.3),Inches(0.42),size=15,c=(200,205,255),italic=True)

def page_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(0.95),size=34,c=DARK,bold=True)
    gbox(s,Inches(1.0),Inches(1.38),Inches(1.8),Inches(0.08),ACCENT1,ACCENT2,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.5),Inches(11.3),Inches(0.38),size=14,c=MUTED,italic=True)

def card(s,l,t,w,h,accent=None,bg=OFFWH):
    a=accent if accent else ACCENT1
    STRIP=Inches(0.07)
    box(s,l,t,w,h,bg,border=(215,225,240))
    gbox(s,l,t,STRIP,h,a,ACCENT2,angle=270)
    return l+STRIP+Inches(0.18), t+Inches(0.14), w-STRIP-Inches(0.35)

def col_kop(s,l,t,w,h,text,c):
    box(s,l,t,w,h,c)
    lbl(s,text,l,t+Inches(0.07),w,h-Inches(0.07),size=20,c=WHITE,bold=True,align=PP_ALIGN.CENTER)

# ─── SLIDE 1: Titel ───────────────────────────────────────────────────────────
s1=ns(dark_bg(True))
# Indigo accent strip left
gbox(s1,Inches(0.45),Inches(1.2),Inches(0.08),Inches(5.1),ACCENT1,ACCENT2,angle=270)
lbl(s1,"Responsive Design\nen Mobile First",Inches(0.85),Inches(1.2),Inches(10.0),Inches(2.6),size=52,c=WHITE,bold=True)
gbox(s1,Inches(0.85),Inches(3.85),Inches(3.5),Inches(0.08),ACCENT1,ACCENT2,angle=0)
lbl(s1,"Periode 2  ·  Les 7",Inches(0.85),Inches(4.0),Inches(7.0),Inches(0.5),size=18,c=(200,205,255),bold=False)
lbl(s1,"Ontwerpen voor elk scherm, elke context",Inches(0.85),Inches(4.55),Inches(9.0),Inches(0.5),size=15,c=(160,170,220),italic=True)
# Corner dots
dot(s1,Inches(12.5),Inches(1.2),Inches(0.12),ACCENT1)
dot(s1,Inches(12.8),Inches(1.5),Inches(0.07),ACCENT2)
dot(s1,Inches(12.2),Inches(1.7),Inches(0.05),(150,150,200))
pnr(s1,1,light=True)
notitie(s1,"Welkom bij les 7 van periode 2. Vandaag starten we met responsive design en mobile first. Dit is een fundamenteel principe voor moderne webdesign.")

# ─── SLIDE 2: Welkom in periode 2 ────────────────────────────────────────────
s2=ns(dark_bg(False))
hero_hdr(s2,"Periode 2: Portfolio Website")
lbl(s2,"Wat staat er op het programma:",Inches(1.0),Inches(2.0),Inches(11.3),Inches(0.4),size=14,c=(200,205,255),bold=True)
periode_items=[
    "Les 7: Responsive design en mobile first (vandaag)",
    "Les 8: Navigatiepatronen en JavaScript sliders",
    "Les 9: Portfolio design en personal branding",
    "Les 10: Design systems verdieping",
    "Les 11: Animatie en micro-interacties",
    "Les 12: Project afronding en presentatie",
]
blist(s2,periode_items,Inches(1.0),Inches(2.4),Inches(11.0),Inches(2.8),size=15,c=(220,225,255),gap=5)
# Project description card
cl,ct,cw=card(s2,Inches(1.0),Inches(5.35),Inches(11.3),Inches(1.7),accent=ACCENT1,bg=(20,30,70))
lbl(s2,"Project: ontwerp een portfolio website met about page, projectenpagina en een interactieve slider",cl,ct,cw,Inches(1.4),size=15,c=(200,210,255))
pnr(s2,2,light=True)
notitie(s2,"Geef een overzicht van periode 2. Leg uit dat het portfolio project de rode draad is door alle lessen.")

# ─── SLIDE 3: Agenda ──────────────────────────────────────────────────────────
s3=ns(light_bg())
page_hdr(s3,"Vandaag","Agenda")
agenda=[
    "Wat is responsive design en waarom is het essentieel",
    "Mobile first versus desktop first strategie",
    "Breakpoints en hoe je ze kiest",
    "Fluid typography en spacing",
    "Figma: responsive frames instellen",
    "Portfolio project introductie en eerste wireframe",
]
col_w=Inches(5.5)
col_gap=Inches(0.3)
col_l1=Inches(0.9)
col_l2=col_l1+col_w+col_gap
top=Inches(1.85)
item_h=Inches(0.72)
for i,item in enumerate(agenda):
    col=i%2; row=i//2
    l=col_l1 if col==0 else col_l2
    t=top+row*item_h
    num_c=ACCENT1 if col==0 else ACCENT2
    box(s3,l,t,col_w,Inches(0.58),(240,242,255),border=(215,225,240))
    gbox(s3,l,t,Inches(0.07),Inches(0.58),ACCENT1,ACCENT2,angle=270)
    lbl(s3,str(i+1),l+Inches(0.12),t+Inches(0.1),Inches(0.35),Inches(0.38),size=17,c=ACCENT1,bold=True)
    lbl(s3,item,l+Inches(0.45),t+Inches(0.1),col_w-Inches(0.6),Inches(0.38),size=14,c=DARK)
pnr(s3,3)
notitie(s3,"Loop kort door de agenda. Geef aan dat de praktijk (Figma) het grootste deel van de les beslaat.")

# ─── SLIDE 4: Wat is responsive design ───────────────────────────────────────
s4=ns()
# Left dark panel
box(s4,0,0,Inches(5.5),SH,NAVY2)
gbox(s4,0,0,Inches(5.5),SH,NAVY,(30,58,138),angle=90)
lbl(s4,"60%",Inches(0.3),Inches(0.8),Inches(4.8),Inches(1.8),size=96,c=WHITE,bold=True)
lbl(s4,"van alle websitebezoeken\nkomt van mobiele\napparaten",Inches(0.3),Inches(2.6),Inches(4.8),Inches(1.8),size=22,c=(200,210,255))
dot(s4,Inches(4.8),Inches(0.5),Inches(0.18),ACCENT1)
dot(s4,Inches(4.5),Inches(0.85),Inches(0.1),ACCENT2)

# Right content
page_hdr(s4,"Wat is responsive design?")
# Override header position for right panel
lbl(s4,"Responsive design",Inches(6.0),Inches(0.5),Inches(6.8),Inches(0.75),size=28,c=DARK,bold=True)
gbox(s4,Inches(6.0),Inches(1.22),Inches(1.8),Inches(0.07),ACCENT1,ACCENT2,angle=0)
bullets=[
    "Past zich aan aan de schermgrootte van de gebruiker",
    "Een codebase voor alle apparaten",
    "Layout, afbeeldingen en typografie passen mee",
    "Geen aparte mobiele site nodig",
]
blist(s4,bullets,Inches(5.7),Inches(1.4),Inches(7.2),Inches(2.8),size=16,c=BODY,gap=8)

# Bottom summary card
cl,ct,cw=card(s4,Inches(5.7),Inches(4.6),Inches(7.4),Inches(1.55),accent=ACCENT1)
lbl(s4,"Responsive design = flexibele grids + flexibele afbeeldingen + media queries",cl,ct,cw,Inches(1.3),size=15,c=DARK,bold=True)
pnr(s4,4)
notitie(s4,"Het 60% cijfer is afgerond gemiddelde. Bron: StatCounter 2024. Leg uit wat dit betekent voor ontwerpers.")

# ─── SLIDE 5: Mobile First ────────────────────────────────────────────────────
s5=ns(light_bg())
page_hdr(s5,"Mobile First: de strategie")

panel_w=Inches(5.6)
panel_h=Inches(4.0)
panel_t=Inches(1.75)

# Left panel: Mobile First
box(s5,Inches(0.7),panel_t,panel_w,panel_h,(235,240,255),border=(180,190,255))
gbox(s5,Inches(0.7),panel_t,panel_w,Inches(0.55),ACCENT1,(80,72,230),angle=90)
lbl(s5,"Mobile First",Inches(0.85),panel_t+Inches(0.08),panel_w-Inches(0.3),Inches(0.45),size=18,c=WHITE,bold=True)
mf_items=[
    "Start met het kleinste scherm",
    "Voeg complexiteit toe voor grotere schermen",
    "Dwingt tot prioriteren van content",
    "Eenvoudiger CSS structuur (min-width queries)",
    "Google geeft de voorkeur aan mobile-first sites",
]
blist(s5,mf_items,Inches(0.9),panel_t+Inches(0.65),panel_w-Inches(0.4),Inches(3.2),size=14,c=BODY,gap=6)

# Right panel: Desktop First
box(s5,Inches(6.7),panel_t,panel_w,panel_h,(245,245,250),border=(210,215,230))
gbox(s5,Inches(6.7),panel_t,panel_w,Inches(0.55),MUTED,(80,90,110),angle=90)
lbl(s5,"Desktop First",Inches(6.85),panel_t+Inches(0.08),panel_w-Inches(0.3),Inches(0.45),size=18,c=WHITE,bold=True)
df_items=[
    "Start groot, strip functionaliteit voor klein",
    "Leidt vaak tot clunky mobiele versies",
    "Max-width queries zijn lastiger te onderhouden",
    "Was de standaard voor 2010",
    "Minder geschikt voor mobile-first indexering",
]
blist(s5,df_items,Inches(6.9),panel_t+Inches(0.65),panel_w-Inches(0.4),Inches(3.2),size=14,c=BODY,gap=6)

# Bottom tip
cl,ct,cw=card(s5,Inches(0.7),Inches(6.0),Inches(12.0),Inches(1.1),accent=ACCENT1)
lbl(s5,"Les vuistregel: ontwerp altijd mobile first in Figma",cl,ct,cw,Inches(0.85),size=16,c=DARK,bold=True)
pnr(s5,5)
notitie(s5,"Bespreek waarom mobile first strategisch slimmer is. Leg uit dat min-width queries compositioneel beter zijn dan max-width.")

# ─── SLIDE 6: Breakpoints ────────────────────────────────────────────────────
s6=ns(light_bg())
page_hdr(s6,"Breakpoints: wanneer springt de layout?")

breakpoints=[
    ("375px","Mobile","Standaard\niPhone breedte"),
    ("640px","Small","Kleine tablets\nen landscape"),
    ("768px","Medium","Tablets"),
    ("1024px","Large","Kleine laptops"),
    ("1280px","Extra Large","Desktop"),
    ("1536px","2XL","Brede schermen"),
]
bar_l=Inches(0.8)
bar_t=Inches(2.2)
bar_w=Inches(11.7)
bar_h=Inches(0.45)

# Main gradient bar
gbox(s6,bar_l,bar_t,bar_w,bar_h,ACCENT1,ACCENT2,angle=0)

# Marker positions (proportional across the bar)
positions=[0.0,0.22,0.32,0.46,0.55,0.75]
marker_colors=[WHITE,WHITE,WHITE,WHITE,WHITE,WHITE]

for i,(px,label,desc) in enumerate(breakpoints):
    pos=positions[i]
    mx=bar_l+bar_w*pos
    # Vertical marker line
    box(s6,mx,bar_t-Inches(0.08),Inches(0.025),bar_h+Inches(0.16),(255,255,255))
    # Label above
    lbl(s6,px,mx-Inches(0.35),bar_t-Inches(0.55),Inches(0.75),Inches(0.35),size=11,c=ACCENT1,bold=True,align=PP_ALIGN.CENTER)
    lbl(s6,label,mx-Inches(0.45),bar_t+bar_h+Inches(0.08),Inches(0.95),Inches(0.32),size=12,c=DARK,bold=True,align=PP_ALIGN.CENTER)
    lbl(s6,desc,mx-Inches(0.55),bar_t+bar_h+Inches(0.4),Inches(1.15),Inches(0.55),size=10,c=MUTED,italic=True,align=PP_ALIGN.CENTER)
    # dot on bar
    dot(s6,mx,bar_t+bar_h/2,Inches(0.065),WHITE)

# Cards at bottom for each
card_w=Inches(1.85)
card_h=Inches(1.1)
card_t=Inches(5.6)
card_colors=[ACCENT1,BLUE,(22,163,74),NAVY2,(124,58,237),(220,38,38)]
for i,(px,label,desc) in enumerate(breakpoints):
    cl=Inches(0.6)+i*(card_w+Inches(0.12))
    box(s6,cl,card_t,card_w,card_h,card_colors[i])
    lbl(s6,label,cl+Inches(0.1),card_t+Inches(0.05),card_w-Inches(0.2),Inches(0.35),size=13,c=WHITE,bold=True)
    lbl(s6,px,cl+Inches(0.1),card_t+Inches(0.38),card_w-Inches(0.2),Inches(0.35),size=20,c=WHITE,bold=True)
    lbl(s6,desc,cl+Inches(0.1),card_t+Inches(0.72),card_w-Inches(0.2),Inches(0.35),size=10,c=(220,230,255),italic=True)

pnr(s6,6)
notitie(s6,"Leg uit dat breakpoints geen vaste regels zijn maar afgeleid worden van je content. Tailwind CSS breakpoints zijn een handige referentie.")

# ─── SLIDE 7: Figma responsive frames ────────────────────────────────────────
s7=ns(dark_bg(False))
hero_hdr(s7,"Figma: drie frames naast elkaar")

steps=[
    ("1","Maak drie frames aan","Mobile 375x812px, Tablet 768x1024px, Desktop 1440x900px"),
    ("2","Geef ze namen","'home/mobile', 'home/tablet', 'home/desktop'"),
    ("3","Begin met de mobile versie","De content is identiek maar de layout verschilt"),
    ("4","Gebruik constraints","Vastmaken op elementen om ze flexibel te maken"),
    ("5","Schakel over tussen frames","Bekijk hoe de layout verandert per schermgrootte"),
]
step_t=Inches(2.0)
step_h=Inches(0.82)
for i,(num,title,desc) in enumerate(steps):
    t=step_t+i*step_h
    # Number circle background
    box(s7,Inches(0.8),t+Inches(0.06),Inches(0.42),Inches(0.42),ACCENT1)
    lbl(s7,num,Inches(0.8),t+Inches(0.06),Inches(0.42),Inches(0.42),size=16,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s7,title,Inches(1.38),t+Inches(0.05),Inches(4.5),Inches(0.38),size=16,c=WHITE,bold=True)
    lbl(s7,desc,Inches(1.38),t+Inches(0.42),Inches(10.8),Inches(0.35),size=13,c=(180,190,230),italic=True)
    if i<len(steps)-1:
        box(s7,Inches(0.97),t+step_h-Inches(0.06),Inches(0.04),Inches(0.18),(100,110,180))

# Tip card
cl,ct,cw=card(s7,Inches(0.8),Inches(6.15),Inches(11.7),Inches(1.0),accent=ACCENT1,bg=(20,30,70))
lbl(s7,"Tip: Groepeer de drie frames in een group per pagina: zo hou je je Figma overzichtelijk",cl,ct,cw,Inches(0.75),size=14,c=(200,210,255),italic=True)
pnr(s7,7,light=True)
notitie(s7,"Demonstreer live in Figma hoe je een responsive frame aanmaakt. Laat constraints zien op tekst en afbeeldingen.")

# ─── SLIDE 8: Grid per breakpoint ────────────────────────────────────────────
s8=ns(light_bg())
page_hdr(s8,"Grid aanpassen per breakpoint")

cols_data=[
    ("Mobile","4 kolommen",ACCENT1,[
        "Margin: 16px aan beide kanten",
        "Gutter: 16px tussen kolommen",
        "Content stack: alles onder elkaar",
        "Full-width knoppen",
        "Grote touch targets (min 44px)",
    ]),
    ("Tablet","8 kolommen",BLUE,[
        "Margin: 32px aan beide kanten",
        "Gutter: 24px tussen kolommen",
        "Twee kolommen layout mogelijk",
        "Navigatie zichtbaar",
        "Cards naast elkaar (2 per rij)",
    ]),
    ("Desktop","12 kolommen",NAVY2,[
        "Margin: 80px aan beide kanten",
        "Gutter: 24px tussen kolommen",
        "Complexe multi-column layouts",
        "Sidebar navigatie mogelijk",
        "Cards naast elkaar (3-4 per rij)",
    ]),
]
col_w=Inches(3.7)
col_h=Inches(4.5)
col_t=Inches(1.72)
col_gap=Inches(0.45)
start_l=Inches(0.8)

for i,(title,subtitle,color,items) in enumerate(cols_data):
    l=start_l+i*(col_w+col_gap)
    col_kop(s8,l,col_t,col_w,Inches(0.55),title,color)
    box(s8,l,col_t+Inches(0.55),col_w,col_h-Inches(0.55),(248,249,255),border=(215,225,240))
    lbl(s8,subtitle,l+Inches(0.15),col_t+Inches(0.65),col_w-Inches(0.3),Inches(0.4),size=14,c=color,bold=True)
    blist(s8,items,l+Inches(0.1),col_t+Inches(1.05),col_w-Inches(0.2),Inches(3.2),size=13,c=BODY,gap=5)

pnr(s8,8)
notitie(s8,"Leg uit dat het 12-koloms grid standaard is in frameworks zoals Bootstrap en Tailwind. Mobiel gebruikt 4 kolommen als vereenvoudiging.")

# ─── SLIDE 9: Fluid typography ────────────────────────────────────────────────
s9=ns(dark_bg(False))
hero_hdr(s9,"Fluid typography en spacing")

typo_bullets=[
    "Gebruik rem eenheden, niet pixels: 1rem = 16px in de browser",
    "Fluid type scale: lettergrootte schaalt mee met de viewport breedte",
    "CSS clamp(): clamp(1rem, 2.5vw, 2rem) begrenst min en max",
    "Vermijd te kleine tekst op mobiel: body minimaal 16px (1rem)",
    "Headings op mobiel kleiner dan op desktop: H1 48px desktop, 32px mobiel",
]
blist(s9,typo_bullets,Inches(0.9),Inches(2.0),Inches(11.4),Inches(3.5),size=17,c=(220,225,255),gap=10)

# Comparison boxes
box_w=Inches(3.6)
box_h=Inches(0.85)
box_t=Inches(5.05)
# Mobile
box(s9,Inches(0.8),box_t,box_w,box_h,(30,40,90))
lbl(s9,"Mobiel: H1 = 32px",Inches(0.95),box_t+Inches(0.08),box_w-Inches(0.3),Inches(0.35),size=14,c=(180,190,255),bold=True)
lbl(s9,"Body = 16px (1rem)",Inches(0.95),box_t+Inches(0.46),box_w-Inches(0.3),Inches(0.3),size=12,c=(150,160,220))
# Desktop
box(s9,Inches(4.8),box_t,box_w,box_h,(20,30,80))
lbl(s9,"Desktop: H1 = 48px",Inches(4.95),box_t+Inches(0.08),box_w-Inches(0.3),Inches(0.35),size=14,c=(200,210,255),bold=True)
lbl(s9,"Body = 18px (1.125rem)",Inches(4.95),box_t+Inches(0.46),box_w-Inches(0.3),Inches(0.3),size=12,c=(170,180,240))
# Arrow
lbl(s9,">>",Inches(4.35),box_t+Inches(0.28),Inches(0.45),Inches(0.4),size=18,c=ACCENT2,bold=True)

cl,ct,cw=card(s9,Inches(8.6),box_t,Inches(4.4),box_h,accent=ACCENT2,bg=(20,30,70))
lbl(s9,"In Figma: gebruik aparte tekststijlen voor mobiel en desktop, of gebruik de tokens plugin",cl,ct,cw,Inches(0.8),size=12,c=(190,200,255))

pnr(s9,9,light=True)
notitie(s9,"Leg uit waarom relative units (rem, em, vw) beter zijn dan pixels voor schaalbaarheid. Demonstreer clamp() als je tijd hebt.")

# ─── SLIDE 10: Portfolio project intro ───────────────────────────────────────
s10=ns(light_bg())
page_hdr(s10,"Portfolio project: wat ga je bouwen")

# Left requirements
lbl(s10,"Vereisten:",Inches(0.8),Inches(1.75),Inches(5.5),Inches(0.4),size=16,c=DARK,bold=True)
requirements=[
    "Home pagina met hero en een interactieve slider",
    "About pagina met jouw verhaal en vaardigheden",
    "Projectenpagina met minimaal 3 project-kaarten",
    "Consistent stijlsysteem op alle paginas",
    "Responsive: werkt op mobiel (375px) en desktop (1440px)",
]
blist(s10,requirements,Inches(0.8),Inches(2.15),Inches(5.5),Inches(3.0),size=15,c=BODY,gap=8)

# Right: site map as flat boxes
map_l=Inches(7.0)
map_t=Inches(1.75)
map_w=Inches(5.6)

# Portfolio website (top box)
box(s10,map_l,map_t,map_w,Inches(0.55),ACCENT1)
lbl(s10,"Portfolio Website",map_l,map_t+Inches(0.08),map_w,Inches(0.42),size=15,c=WHITE,bold=True,align=PP_ALIGN.CENTER)

# Connecting line down
box(s10,map_l+map_w/2-Inches(0.025),map_t+Inches(0.55),Inches(0.05),Inches(0.3),(150,160,220))

# Three sub-pages
sub_w=Inches(1.65)
sub_t=map_t+Inches(0.85)
sub_gap=Inches(0.12)
sub_colors=[ACCENT1,BLUE,NAVY2]
sub_labels=["Home","About","Projecten"]
sub_descs=["hero, slider, intro","foto, bio, skills","grid met kaarten"]

# Horizontal line
box(s10,map_l+sub_w/2,sub_t-Inches(0.08),map_w-sub_w,Inches(0.04),(150,160,220))

for i,(label,desc,color) in enumerate(zip(sub_labels,sub_descs,sub_colors)):
    sl=map_l+i*(sub_w+sub_gap)
    # Vertical connector
    box(s10,sl+sub_w/2-Inches(0.025),sub_t-Inches(0.08),Inches(0.05),Inches(0.12),(150,160,220))
    box(s10,sl,sub_t,sub_w,Inches(0.5),color)
    lbl(s10,label,sl,sub_t+Inches(0.04),sub_w,Inches(0.3),size=14,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    # Description box below
    box(s10,sl,sub_t+Inches(0.5),sub_w,Inches(0.65),(240,242,255),border=(200,210,240))
    lbl(s10,desc,sl+Inches(0.05),sub_t+Inches(0.55),sub_w-Inches(0.1),Inches(0.55),size=12,c=BODY,italic=True,align=PP_ALIGN.CENTER)

# Assessment criteria bottom
cl,ct,cw=card(s10,Inches(0.8),Inches(5.5),Inches(11.7),Inches(1.6),accent=ACCENT1)
lbl(s10,"Inleverdatum en beoordeling",cl,ct,cw,Inches(0.35),size=14,c=DARK,bold=True)
lbl(s10,"Lever je Figma link in via de ELO voor les 8. Beoordeeld op: mobile wireframe, desktop wireframe, grid gebruik en responsive aanpak.",cl,ct+Inches(0.38),cw,Inches(0.9),size=13,c=BODY)
pnr(s10,10)
notitie(s10,"Introducer het portfolio project. Laat een voorbeeld portfolio zien voor inspiratie. Bespreek de beoordelingscriteria.")

# ─── SLIDE 11: Mobile wireframe stap voor stap ───────────────────────────────
s11=ns(dark_bg(False))
hero_hdr(s11,"Eerste wireframe: mobiele home")

steps11=[
    ("1","Maak een Mobile frame aan","375x812px, naam 'home/mobile'"),
    ("2","Navigatie","Hamburger menu icoon rechts, logo links, hoogte 56px"),
    ("3","Hero sectie","Volle breedte, grote heading, 1 CTA knop eronder"),
    ("4","Slider placeholder","Rechthoek over de breedte, hoogte 240px"),
    ("5","Intro sectie","Naam en korte bio tekst"),
    ("6","Footer","Drie icoonlinks voor social media"),
]
step_t=Inches(2.0)
step_h=Inches(0.8)
for i,(num,title,desc) in enumerate(steps11):
    t=step_t+i*step_h
    box(s11,Inches(0.8),t+Inches(0.08),Inches(0.38),Inches(0.38),ACCENT1)
    lbl(s11,num,Inches(0.8),t+Inches(0.08),Inches(0.38),Inches(0.38),size=14,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s11,title,Inches(1.32),t+Inches(0.06),Inches(3.5),Inches(0.35),size=15,c=WHITE,bold=True)
    lbl(s11,desc,Inches(1.32),t+Inches(0.42),Inches(10.8),Inches(0.32),size=13,c=(170,180,230),italic=True)

# Mobile frame mockup on right
frame_l=Inches(8.5)
frame_t=Inches(1.7)
frame_w=Inches(2.4)
frame_h=Inches(4.8)
# Phone outline
box(s11,frame_l,frame_t,frame_w,frame_h,(25,35,80),border=(100,110,200))
# Nav bar
box(s11,frame_l,frame_t,frame_w,Inches(0.42),(40,50,100))
lbl(s11,"Logo  [=]",frame_l+Inches(0.08),frame_t+Inches(0.08),frame_w-Inches(0.16),Inches(0.28),size=9,c=(180,190,255))
# Hero
box(s11,frame_l,frame_t+Inches(0.42),frame_w,Inches(1.1),(50,60,120))
lbl(s11,"HERO",frame_l,frame_t+Inches(0.72),frame_w,Inches(0.35),size=11,c=(150,160,220),align=PP_ALIGN.CENTER)
# Slider
box(s11,frame_l,frame_t+Inches(1.52),frame_w,Inches(1.55),(35,45,100))
lbl(s11,"SLIDER",frame_l,frame_t+Inches(2.2),frame_w,Inches(0.35),size=11,c=(120,130,200),align=PP_ALIGN.CENTER)
# Intro
box(s11,frame_l,frame_t+Inches(3.07),frame_w,Inches(0.95),(30,40,90))
lbl(s11,"Intro tekst",frame_l,frame_t+Inches(3.4),frame_w,Inches(0.28),size=9,c=(140,150,210),align=PP_ALIGN.CENTER)
# Footer
box(s11,frame_l,frame_t+Inches(4.02),frame_w,Inches(0.38),(20,30,80))
lbl(s11,"[ ] [ ] [ ]",frame_l,frame_t+Inches(4.08),frame_w,Inches(0.28),size=9,c=(120,130,200),align=PP_ALIGN.CENTER)

lbl(s11,"375px",frame_l,frame_t+Inches(4.85),frame_w,Inches(0.28),size=10,c=(130,140,200),align=PP_ALIGN.CENTER)

cl,ct,cw=card(s11,Inches(0.8),Inches(6.65),Inches(7.5),Inches(0.65),accent=ACCENT2,bg=(20,30,70))
lbl(s11,"Focus op structuur en volgorde van content, geen kleur of stijl",cl,ct,cw,Inches(0.5),size=14,c=(190,200,255),italic=True)
pnr(s11,11,light=True)
notitie(s11,"Laat studenten stap voor stap hun eerste wireframe opbouwen. Controleer dat iedereen het juiste frame formaat gebruikt.")

# ─── SLIDE 12: Desktop wireframe ─────────────────────────────────────────────
s12=ns(dark_bg(False))
hero_hdr(s12,"Desktop versie: hoe ziet het er anders uit")

comparisons=[
    ("Nav","Hamburger menu icoon\n(mobiel)","Volledige horizontale\nnavigatie (desktop)"),
    ("Hero","Gestapelde content\n(mobiel)","Content naast\nafbeelding (desktop)"),
    ("Slider","Volle breedte smal\n(mobiel)","Centered breed\n(desktop)"),
    ("Grid","1 kolom\n(mobiel)","2 tot 3 kolommen\n(desktop)"),
]
comp_t=Inches(2.1)
comp_h=Inches(0.92)
comp_gap=Inches(0.16)
label_w=Inches(1.5)
panel_w=Inches(4.5)

# Headers
lbl(s12,"Element",Inches(0.8),comp_t-Inches(0.45),label_w,Inches(0.38),size=13,c=(180,190,230),bold=True)
lbl(s12,"Mobiel",Inches(0.8)+label_w+Inches(0.1),comp_t-Inches(0.45),panel_w,Inches(0.38),size=13,c=ACCENT2,bold=True,align=PP_ALIGN.CENTER)
lbl(s12,"Desktop",Inches(0.8)+label_w+panel_w+Inches(0.3),comp_t-Inches(0.45),panel_w,Inches(0.38),size=13,c=ACCENT1,bold=True,align=PP_ALIGN.CENTER)
# divider
gbox(s12,Inches(0.8),comp_t-Inches(0.08),Inches(11.7),Inches(0.04),ACCENT1,ACCENT2,angle=0)

for i,(elem,mob,desk) in enumerate(comparisons):
    t=comp_t+i*(comp_h+comp_gap)
    # Element label
    box(s12,Inches(0.8),t,label_w,comp_h,(30,40,90))
    lbl(s12,elem,Inches(0.85),t+Inches(0.28),label_w-Inches(0.1),Inches(0.38),size=14,c=WHITE,bold=True)
    # Mobile panel
    ml=Inches(0.8)+label_w+Inches(0.1)
    box(s12,ml,t,panel_w,comp_h,(35,45,100))
    lbl(s12,mob,ml+Inches(0.15),t+Inches(0.12),panel_w-Inches(0.3),comp_h-Inches(0.24),size=13,c=(190,200,240))
    # Arrow
    lbl(s12,">>",ml+panel_w+Inches(0.05),t+Inches(0.28),Inches(0.28),Inches(0.38),size=14,c=ACCENT2,bold=True)
    # Desktop panel
    dl=ml+panel_w+Inches(0.3)
    box(s12,dl,t,panel_w,comp_h,(25,35,90))
    lbl(s12,desk,dl+Inches(0.15),t+Inches(0.12),panel_w-Inches(0.3),comp_h-Inches(0.24),size=13,c=(210,220,255))

cl,ct,cw=card(s12,Inches(0.8),Inches(6.4),Inches(11.7),Inches(0.75),accent=ACCENT2,bg=(18,28,68))
lbl(s12,"Hetzelfde content, andere rangschikking",cl,ct,cw,Inches(0.55),size=16,c=(200,210,255),bold=True,italic=True)
pnr(s12,12,light=True)
notitie(s12,"Bespreek dat responsive design niet over het toevoegen of verwijderen van content gaat, maar over het herschikken ervan.")

# ─── SLIDE 13: Live analyse ───────────────────────────────────────────────────
s13=ns(light_bg())
page_hdr(s13,"Live analyse: responsiviteit testen")

tasks=[
    ("1","Browser verkleinen",
     "Open je favoriete website op mobiel of verklein je browser. Klapt de navigatie in op een hamburgermenu?",
     ACCENT1),
    ("2","Typografie checken",
     "Bekijk de typografie op mobiel. Zijn de headings kleiner dan op desktop? Is de body tekst leesbaar?",
     BLUE),
    ("3","Slechte mobile site",
     "Zoek een site die er op mobiel slecht uitziet. Wat gaat er mis? Te kleine tekst, overlappende elementen?",
     RED),
    ("4","Browser DevTools",
     "Gebruik de browser devtools (F12, device toolbar) om een desktop site op 375px te bekijken.",
     GREEN),
]
card_w=Inches(5.6)
card_h=Inches(2.1)
card_gap=Inches(0.25)
positions_13=[(Inches(0.7),Inches(1.85)),(Inches(6.55),Inches(1.85)),
               (Inches(0.7),Inches(4.1)),(Inches(6.55),Inches(4.1))]

for i,(num,title,desc,color) in enumerate(tasks):
    l,t=positions_13[i]
    cl,ct,cw=card(s13,l,t,card_w,card_h,accent=color)
    lbl(s13,f"Taak {num}: {title}",cl,ct,cw,Inches(0.38),size=15,c=color,bold=True)
    lbl(s13,desc,cl,ct+Inches(0.42),cw,Inches(1.5),size=13,c=BODY)
pnr(s13,13)
notitie(s13,"Geef studenten 10 minuten voor deze analyse. Bespreek daarna klassikaal wat ze gevonden hebben.")

# ─── SLIDE 14: Bronnen ────────────────────────────────────────────────────────
s14=ns(light_bg())
page_hdr(s14,"Bronnen")

bronnen=[
    ("Responsive design basics","web.dev/responsive-web-design-basics","Google Web.dev gids: de fundamenten van responsive design",ACCENT1),
    ("CSS media queries","developer.mozilla.org/docs/Web/CSS/CSS_media_queries","MDN documentatie: volledige referentie voor media queries",BLUE),
    ("Every Layout","every-layout.dev","Responsive layout patronen zonder media queries",GREEN),
    ("Figma responsive design","figma.com/resource-library/responsive-design","Officiele Figma gids voor responsive frames en constraints",NAVY2),
    ("Mobile First Design","abookapart.com/products/mobile-first","Boek van Luke Wroblewski: de mobile first aanpak",(124,58,237)),
]
bron_h=Inches(0.92)
bron_t=Inches(1.85)
bron_w=Inches(11.7)
bron_gap=Inches(0.1)

for i,(title,url,desc,color) in enumerate(bronnen):
    t=bron_t+i*(bron_h+bron_gap)
    cl,ct,cw=card(s14,Inches(0.8),t,bron_w,bron_h,accent=color)
    lbl(s14,title,cl,ct,Inches(3.8),Inches(0.38),size=15,c=color,bold=True)
    lbl(s14,url,cl+Inches(3.9),ct+Inches(0.04),cw-Inches(3.9),Inches(0.32),size=12,c=MUTED,italic=True)
    lbl(s14,desc,cl,ct+Inches(0.42),cw,Inches(0.42),size=13,c=BODY)
pnr(s14,14)
notitie(s14,"Verwijs studenten naar deze bronnen voor verdieping. Every Layout is bijzonder interessant omdat het responsive layouts zonder media queries laat zien.")

# ─── SLIDE 15: Opdracht ───────────────────────────────────────────────────────
s15=ns(dark_bg(False))
# Accent left strip
gbox(s15,Inches(0.45),Inches(0.4),Inches(0.08),Inches(6.7),ACCENT1,ACCENT2,angle=270)
lbl(s15,"Opdracht  ·  Les 7",Inches(0.72),Inches(0.5),Inches(11.0),Inches(0.55),size=14,c=ACCENT2,bold=True)
lbl(s15,"Portfolio Website  ·  Responsive Wireframe",Inches(0.72),Inches(1.0),Inches(11.0),Inches(0.65),size=26,c=WHITE,bold=True)
gbox(s15,Inches(0.72),Inches(1.65),Inches(3.5),Inches(0.07),ACCENT1,ACCENT2,angle=0)

steps15=[
    ("1","Maak twee Figma frames","'home/mobile' (375x812) en 'home/desktop' (1440x900)"),
    ("2","Wireframe beide versies","In grijstinten, geen kleur of stijl"),
    ("3","Verplichte elementen","Navigatie, hero, slider placeholder, intro sectie, footer"),
    ("4","Grid gebruiken","4-koloms grid op mobiel, 12-koloms op desktop"),
    ("5","Inleveren","Lever je Figma link in via de ELO voor les 8"),
]
step_t=Inches(1.85)
step_h=Inches(0.88)
for i,(num,title,desc) in enumerate(steps15):
    t=step_t+i*step_h
    gbox(s15,Inches(0.72),t+Inches(0.1),Inches(0.38),Inches(0.38),ACCENT1,ACCENT2,angle=135)
    lbl(s15,num,Inches(0.72),t+Inches(0.1),Inches(0.38),Inches(0.38),size=15,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s15,title,Inches(1.25),t+Inches(0.08),Inches(4.2),Inches(0.38),size=16,c=WHITE,bold=True)
    lbl(s15,desc,Inches(1.25),t+Inches(0.45),Inches(11.0),Inches(0.35),size=13,c=(180,190,230))
pnr(s15,15,light=True)
notitie(s15,"Geef duidelijk aan wat er ingeleverd moet worden. Herinner studenten aan de deadline voor les 8.")

# ─── SLIDE 16: Volgende les ───────────────────────────────────────────────────
s16=ns(dark_bg(True))
gbox(s16,Inches(0.45),Inches(0.5),Inches(0.08),Inches(6.5),ACCENT1,ACCENT2,angle=270)
lbl(s16,"Volgende les",Inches(0.72),Inches(0.6),Inches(11.0),Inches(0.55),size=16,c=(180,190,255))
lbl(s16,"Navigatiepatronen en Sliders",Inches(0.72),Inches(1.15),Inches(11.0),Inches(1.1),size=40,c=WHITE,bold=True)
gbox(s16,Inches(0.72),Inches(2.25),Inches(4.0),Inches(0.08),ACCENT1,ACCENT2,angle=0)

preview_items=[
    "Hamburger menu en mobiele navigatiepatronen",
    "Slider en carousel design principes",
    "Toegankelijkheid van interactieve elementen",
    "Figma: prototype van een werkende slider",
]
blist(s16,preview_items,Inches(0.72),Inches(2.45),Inches(11.0),Inches(2.5),size=18,c=(200,210,255),gap=10)

# Dots decoration
dot(s16,Inches(11.5),Inches(2.0),Inches(0.15),ACCENT1)
dot(s16,Inches(11.9),Inches(2.5),Inches(0.09),ACCENT2)
dot(s16,Inches(11.3),Inches(2.8),Inches(0.06),(150,160,220))
dot(s16,Inches(12.1),Inches(1.8),Inches(0.05),(100,110,200))

cl,ct,cw=card(s16,Inches(0.72),Inches(5.5),Inches(11.7),Inches(1.0),accent=ACCENT2,bg=(18,28,70))
lbl(s16,"Kijk alvast naar sliders op Dribbble of Awwwards voor designinspiratie",cl,ct,cw,Inches(0.75),size=15,c=(190,200,255),italic=True)
pnr(s16,16,light=True)
notitie(s16,"Sluit af met enthousiasme voor les 8. Herinner studenten aan de inleveropdracht en geef ze de link naar Dribbble/Awwwards mee.")

# ─── SAVE ────────────────────────────────────────────────────────────────────
OUT=r"C:\Users\Peter Marcelis\Documents\Lesmateriaal\design\periode-2\les-7\07_responsive-mobile-first.pptx"
prs.save(OUT)
print(f"Opgeslagen: {OUT}")
print(f"Aantal slides: {len(prs.slides)}")
