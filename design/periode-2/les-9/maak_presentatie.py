"""
Les 9: Portfolio Design en Personal Branding
Periode 2, Jaar 1
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image, ImageDraw
import numpy as np, io

NAVY=(6,14,44); NAVY2=(30,58,138); BLUE=(37,99,235); BLUE2=(59,130,246)
CYAN=(6,182,212); TEAL=(13,148,136); WHITE=(255,255,255); OFFWH=(248,250,252)
DARK=(15,23,42); BODY=(51,65,85); MUTED=(100,116,135)
GREEN=(22,163,74); RED=(220,38,38)
ACCENT1=(192,38,211); ACCENT2=(232,121,249)

def rgb(c): return RGBColor(c[0],c[1],c[2])
SW,SH=Inches(13.33),Inches(7.5); BW,BH=1920,1080

def garr(w,h,c1,c2,d='diag'):
    x,y=np.linspace(0,1,w,dtype=np.float32),np.linspace(0,1,h,dtype=np.float32)
    if d=='h': t=np.tile(x,(h,1))
    elif d=='v': t=np.tile(y.reshape(-1,1),(1,w))
    else: xx,yy=np.meshgrid(x,y); t=xx*.55+yy*.45
    a=np.zeros((h,w,3),dtype=np.float32)
    for i in range(3): a[:,:,i]=c1[i]+(c2[i]-c1[i])*t
    return np.clip(a,0,255).astype(np.uint8)

def dark_bg(circ=True):
    img=Image.fromarray(garr(BW,BH,NAVY,NAVY2)).convert('RGBA')
    if circ:
        ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
        d.ellipse([BW-500,-180,BW+200,520],fill=(ACCENT1[0],ACCENT1[1],ACCENT1[2],20))
        d.ellipse([-60,BH-360,380,BH+100],fill=(ACCENT2[0],ACCENT2[1],ACCENT2[2],22))
        img=Image.alpha_composite(img,ov)
    return img.convert('RGB')

def light_bg():
    img=Image.fromarray(garr(BW,BH,OFFWH,WHITE)).convert('RGBA')
    ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
    d.ellipse([BW-580,-220,BW+180,540],fill=(ACCENT1[0],ACCENT1[1],ACCENT1[2],8))
    img=Image.alpha_composite(img,ov); return img.convert('RGB')

def buf(img):
    b=io.BytesIO(); img.save(b,format='PNG'); b.seek(0); return b

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
BLANK=prs.slide_layouts[6]

def ns(bg=None):
    s=prs.slides.add_slide(BLANK)
    if bg:
        p=s.shapes.add_picture(buf(bg),0,0,SW,SH)
        t=s.shapes._spTree; t.remove(p._element); t.insert(2,p._element)
    else:
        s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(WHITE)
    return s

def box(s,l,t,w,h,c,border=None):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c)
    if border: sh.line.color.rgb=rgb(border); sh.line.width=Pt(0.75)
    else: sh.line.fill.background()
    return sh

def gbox(s,l,t,w,h,c1,c2,angle=90):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.gradient(); st=sh.fill.gradient_stops
    st[0].position=0.0; st[0].color.rgb=rgb(c1)
    st[1].position=1.0; st[1].color.rgb=rgb(c2)
    sh.fill.gradient_angle=angle; sh.line.fill.background(); return sh

def dot(s,cx,cy,r,c):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,cx-r,cy-r,r*2,r*2)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c); sh.line.fill.background()

def lbl(s,text,l,t,w,h,size=18,c=BODY,bold=False,italic=False,align=PP_ALIGN.LEFT):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=text; r.font.name="Calibri"; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=rgb(c)

def blist(s,items,l,t,w,h,size=18,c=BODY,gap=7):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(gap); r=p.add_run()
        r.text=f"  {chr(8226)}  {item}"; r.font.name="Calibri"; r.font.size=Pt(size); r.font.color.rgb=rgb(c)

def notitie(s,t): s.notes_slide.notes_text_frame.text=t

def pnr(s,n,tot=15,light=False):
    lbl(s,f"{n} / {tot}",Inches(12.5),Inches(7.1),Inches(0.8),Inches(0.35),size=11,
        c=(180,180,200) if light else MUTED,align=PP_ALIGN.RIGHT)

def hero_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(1.0),size=36,c=WHITE,bold=True)
    gbox(s,Inches(1.0),Inches(1.42),Inches(2.0),Inches(0.08),ACCENT1,ACCENT2,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.55),Inches(11.3),Inches(0.42),size=15,c=(240,180,255),italic=True)

def page_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(0.95),size=34,c=DARK,bold=True)
    gbox(s,Inches(1.0),Inches(1.38),Inches(1.8),Inches(0.08),ACCENT1,ACCENT2,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.5),Inches(11.3),Inches(0.38),size=14,c=MUTED,italic=True)

def card(s,l,t,w,h,accent=None,bg=OFFWH):
    a=accent if accent else ACCENT1
    STRIP=Inches(0.07)
    box(s,l,t,w,h,bg,border=(215,225,240))
    gbox(s,l,t,STRIP,h,a,ACCENT2,angle=270)
    return l+STRIP+Inches(0.18), t+Inches(0.14), w-STRIP-Inches(0.35)

def col_kop(s,l,t,w,h,text,c):
    box(s,l,t,w,h,c)
    lbl(s,text,l,t+Inches(0.07),w,h-Inches(0.07),size=20,c=WHITE,bold=True,align=PP_ALIGN.CENTER)

# ─── SLIDE 1: TITEL ───────────────────────────────────────────────────────────
s1=ns(dark_bg(True))
# Fuchsia accent strip left
gbox(s1,0,0,Inches(0.18),SH,ACCENT1,ACCENT2,angle=270)
# Title
lbl(s1,"Portfolio Design\nen Personal Branding",
    Inches(1.0),Inches(1.6),Inches(9.5),Inches(2.0),size=48,c=WHITE,bold=True)
# Subtitle
lbl(s1,"Periode 2  ·  Les 9",
    Inches(1.0),Inches(3.7),Inches(8.0),Inches(0.55),size=18,c=(200,160,240))
# Gradient accent line
gbox(s1,Inches(1.0),Inches(4.3),Inches(3.0),Inches(0.08),ACCENT1,ACCENT2,angle=0)
# Tagline
lbl(s1,"Hoe presenteer je jezelf als designer?",
    Inches(1.0),Inches(4.45),Inches(9.0),Inches(0.5),size=15,c=(220,190,255),italic=True)
# Les number badge
gbox(s1,Inches(10.8),Inches(0.35),Inches(1.7),Inches(0.55),ACCENT1,ACCENT2,angle=0)
lbl(s1,"Les 9",Inches(10.8),Inches(0.35),Inches(1.7),Inches(0.55),size=16,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
pnr(s1,1,light=True)
notitie(s1,"Les 9: Portfolio Design en Personal Branding. Studenten leren hoe ze een overtuigend portfolio opzetten en een persoonlijke designstijl ontwikkelen.")

# ─── SLIDE 2: AGENDA ──────────────────────────────────────────────────────────
s2=ns(light_bg())
page_hdr(s2,"Vandaag","Wat gaan we behandelen")
items=[
    "Wat maakt een sterk portfolio",
    "Personal branding: wie ben jij als designer",
    "About pagina: structuur en inhoud",
    "Projectenpagina: projecten presenteren",
    "Consistente stijl over meerdere paginas",
    "Figma: multi-pagina ontwerp opzetten",
]
dot_colors=[ACCENT1,ACCENT2,BLUE,BLUE2,GREEN,TEAL]
for i,(item,dc) in enumerate(zip(items,dot_colors)):
    row_t=Inches(1.9)+i*Inches(0.78)
    dot(s2,Inches(1.35),row_t+Inches(0.18),Inches(0.13),dc)
    lbl(s2,item,Inches(1.6),row_t,Inches(9.5),Inches(0.6),size=18,c=BODY)
pnr(s2,2)
notitie(s2,"Loop de agenda door. Benoem dat vandaag de focus ligt op inhoud en structuur van de portfolio website die studenten aan het bouwen zijn.")

# ─── SLIDE 3: WAT MAAKT EEN STERK PORTFOLIO ──────────────────────────────────
s3=ns()
# Left dark panel
box(s3,0,0,Inches(5.2),SH,DARK)
gbox(s3,0,0,Inches(0.18),SH,ACCENT1,ACCENT2,angle=270)
lbl(s3,"Wat maakt een\nsterk portfolio?",
    Inches(0.35),Inches(0.6),Inches(4.6),Inches(1.5),size=28,c=WHITE,bold=True)
gbox(s3,Inches(0.35),Inches(2.1),Inches(1.6),Inches(0.06),ACCENT1,ACCENT2,angle=0)
lbl(s3,"Je portfolio is je eerste indruk bij een opdrachtgever of werkgever. Het vertelt in 30 seconden wie je bent.",
    Inches(0.35),Inches(2.25),Inches(4.55),Inches(1.8),size=14,c=(200,210,230),italic=False)
# Right bullets
bullets=[
    "Maximaal 5 tot 8 projecten: kies kwaliteit boven kwantiteit",
    "Elk project heeft een duidelijk verhaal: wat, hoe en waarom",
    "Consistent ontwerp: hetzelfde kleurpalet en typografie op alle paginas",
    "Snel laden: geen overbodige afbeeldingen",
    "Werkende links en prototypes laten het leven zien",
    "Een duidelijke contactmogelijkheid of CTA",
]
blist(s3,bullets,Inches(5.5),Inches(1.9),Inches(7.5),Inches(4.8),size=16,c=BODY,gap=10)
pnr(s3,3)
notitie(s3,"Bespreek waarom selectiviteit belangrijk is: een recruiteur kijkt gemiddeld 30 seconden naar een portfolio. Kwaliteit wint van kwantiteit.")

# ─── SLIDE 4: PERSONAL BRANDING ───────────────────────────────────────────────
s4=ns(light_bg())
page_hdr(s4,"Personal branding: jouw unieke stem")
# Three panels
panels=[
    ("Stijl",ACCENT1,"Welke esthetiek past bij jou? Minimalistisch, kleurrijk, editorial? Kies en wees consistent."),
    ("Toon",BLUE,"Hoe schrijf je? Formeel en professioneel, of casual en persoonlijk? Kies een toon en houd die aan."),
    ("Focus",GREEN,"Wat zijn je specialiteiten? UI design, branding, illustratie? Benoem het duidelijk."),
]
pw=Inches(3.6); gap=Inches(0.22); pt=Inches(1.9); ph=Inches(2.8)
for i,(title,col,desc) in enumerate(panels):
    pl=Inches(1.0)+i*(pw+gap)
    box(s4,pl,pt,pw,ph,col)
    lbl(s4,title,pl,pt+Inches(0.18),pw,Inches(0.55),size=22,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s4,desc,pl+Inches(0.2),pt+Inches(0.75),pw-Inches(0.4),ph-Inches(0.95),size=14,c=WHITE)
# Card below
cx,cy,cw=card(s4,Inches(1.0),Inches(4.9),Inches(11.3),Inches(0.88))
lbl(s4,"Tip: kijk naar vijf portfolios van designers die je bewondert. Wat maakt hun stijl herkenbaar?",
    cx,cy,cw,Inches(0.6),size=14,c=BODY,italic=True)
pnr(s4,4)
notitie(s4,"Laat studenten nadenken over hun eigen stijl. Vraag: welke drie woorden beschrijven jouw designsmaak? Dit helpt bij het bepalen van kleur, typografie en toon.")

# ─── SLIDE 5: ABOUT PAGINA STRUCTUUR ─────────────────────────────────────────
s5=ns(light_bg())
page_hdr(s5,"About pagina: wat zet je erin")
# Left wireframe simulation
wf_l=Inches(1.0); wf_t=Inches(1.85); wf_w=Inches(4.6); wf_h=Inches(5.2)
box(s5,wf_l,wf_t,wf_w,wf_h,DARK)
lbl(s5,"Wireframe preview",wf_l,wf_t+Inches(0.08),wf_w,Inches(0.35),size=10,c=MUTED,align=PP_ALIGN.CENTER)
# Header section
box(s5,wf_l+Inches(0.15),wf_t+Inches(0.45),wf_w-Inches(0.3),Inches(0.5),(50,70,110))
lbl(s5,"Header: naam + functietitel",wf_l+Inches(0.2),wf_t+Inches(0.5),wf_w-Inches(0.4),Inches(0.38),size=10,c=WHITE)
# Photo placeholder
box(s5,wf_l+Inches(0.15),wf_t+Inches(1.05),Inches(1.3),Inches(1.3),MUTED)
lbl(s5,"Foto",wf_l+Inches(0.15),wf_t+Inches(1.55),Inches(1.3),Inches(0.3),size=10,c=WHITE,align=PP_ALIGN.CENTER)
# Bio text (3 thin rectangles)
for bi in range(3):
    box(s5,wf_l+Inches(1.6),wf_t+Inches(1.05)+bi*Inches(0.32),wf_w-Inches(1.85),Inches(0.22),(60,80,120))
# Skills row
skill_colors=[(192,38,211),(37,99,235),(22,163,74)]
for si in range(3):
    box(s5,wf_l+Inches(0.15)+si*Inches(1.4),wf_t+Inches(2.55),Inches(1.22),Inches(0.38),skill_colors[si])
    lbl(s5,["UI Design","Branding","Web"][si],
        wf_l+Inches(0.15)+si*Inches(1.4),wf_t+Inches(2.58),Inches(1.22),Inches(0.3),size=9,c=WHITE,align=PP_ALIGN.CENTER)
# CTA button
box(s5,wf_l+Inches(1.3),wf_t+Inches(3.15),Inches(2.0),Inches(0.42),ACCENT1)
lbl(s5,"Neem contact op",wf_l+Inches(1.3),wf_t+Inches(3.18),Inches(2.0),Inches(0.35),size=10,c=WHITE,align=PP_ALIGN.CENTER)
# Right checklist
checklist=[
    "Naam en functietitel bovenaan",
    "Professionele foto of illustratie",
    "Bio van 3 tot 5 zinnen: wie je bent, wat je doet, waar je naartoe wilt",
    "Vaardigheden als badges of lijst",
    "Contactknop of e-mailadres",
    "Optioneel: tijdlijn met ervaring of opleiding",
]
blist(s5,checklist,Inches(6.0),Inches(2.0),Inches(6.8),Inches(4.8),size=15,c=BODY,gap=10)
pnr(s5,5)
notitie(s5,"Laat de wireframe zien als voorbeeld. Bespreek de hiërarchie: naam bovenaan, dan visueel, dan tekst. De CTA moet altijd zichtbaar zijn.")

# ─── SLIDE 6: PROJECTENPAGINA STRUCTUUR ──────────────────────────────────────
s6=ns(light_bg())
page_hdr(s6,"Projectenpagina: je werk tonen")
# Page title area at top of wireframe
box(s6,Inches(1.0),Inches(1.85),Inches(6.6),Inches(0.4),(220,225,235))
lbl(s6,"Projecten",Inches(1.05),Inches(1.88),Inches(6.5),Inches(0.32),size=12,c=MUTED)
# Three project cards
card_w=Inches(1.95); card_h=Inches(2.4); card_t=Inches(2.35); card_gap=Inches(0.15)
card_cols=[ACCENT1,BLUE,GREEN]
for ci in range(3):
    cl=Inches(1.0)+ci*(card_w+card_gap)
    box(s6,cl,card_t,card_w,card_h,OFFWH,border=(215,225,240))
    # Dark header strip
    box(s6,cl,card_t,card_w,Inches(0.6),DARK)
    lbl(s6,f"Project {ci+1}",cl+Inches(0.08),card_t+Inches(0.08),card_w-Inches(0.16),Inches(0.42),size=11,c=WHITE,bold=True)
    # Thumbnail placeholder
    box(s6,cl+Inches(0.1),card_t+Inches(0.68),card_w-Inches(0.2),Inches(0.7),MUTED)
    # Category badge
    box(s6,cl+Inches(0.1),card_t+Inches(1.46),Inches(0.9),Inches(0.28),card_cols[ci])
    lbl(s6,["UI","Branding","Web"][ci],
        cl+Inches(0.1),card_t+Inches(1.47),Inches(0.9),Inches(0.25),size=9,c=WHITE,align=PP_ALIGN.CENTER)
    # Short description lines
    box(s6,cl+Inches(0.1),card_t+Inches(1.82),card_w-Inches(0.2),Inches(0.14),(200,210,225))
    box(s6,cl+Inches(0.1),card_t+Inches(2.0),card_w-Inches(0.35),Inches(0.14),(200,210,225))
# Right bullets
r_bullets=[
    "Thumbnail afbeelding of mockup",
    "Projectnaam als heading",
    "Categorie als badge (UI, Branding, Web)",
    "Korte omschrijving: max 2 zinnen",
    "Link of knop naar het project",
]
blist(s6,r_bullets,Inches(7.6),Inches(2.1),Inches(5.3),Inches(4.0),size=16,c=BODY,gap=10)
pnr(s6,6)
notitie(s6,"Bespreek dat een thumbnail de aandacht trekt. Goede mockups (telefoon/laptop) maken het project direct herkenbaar als echt werk.")

# ─── SLIDE 7: PROJECTPRESENTATIE ─────────────────────────────────────────────
s7=ns(dark_bg(False))
hero_hdr(s7,"Een project goed presenteren","Van probleem naar resultaat")
steps=[
    "Begin met het probleem: wat moest er opgelost worden",
    "Toon je proces: schetsen, wireframes, iteraties",
    "Laat het eindresultaat zien: meerdere schermen of mockups",
    "Sluit af met het resultaat: wat heb je geleerd of bereikt",
    "Minder tekst, meer visuals: ratio 80 procent beeld, 20 procent tekst",
]
for i,step in enumerate(steps):
    row_t=Inches(2.0)+i*Inches(0.68)
    dot(s7,Inches(1.22),row_t+Inches(0.2),Inches(0.18),ACCENT1)
    lbl(s7,str(i+1),Inches(1.22)-Inches(0.1),row_t+Inches(0.06),Inches(0.2),Inches(0.28),size=11,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s7,step,Inches(1.55),row_t,Inches(10.5),Inches(0.6),size=16,c=WHITE)
# Card at bottom
cx,cy,cw=card(s7,Inches(1.0),Inches(5.7),Inches(11.3),Inches(1.1),bg=(30,40,70))
lbl(s7,"Een projectpagina met alleen een eindresultaat zonder context is minder overtuigend dan een pagina die het proces laat zien.",
    cx,cy,cw,Inches(0.8),size=13,c=(220,200,255),italic=True)
pnr(s7,7,light=True)
notitie(s7,"Laat voorbeeldportfolio's zien van Behance. Vergelijk een 'slechte' (alleen eindresultaat) met een 'goede' (process + eindresultaat) projectpagina.")

# ─── SLIDE 8: CONSISTENTIE ────────────────────────────────────────────────────
s8=ns(light_bg())
page_hdr(s8,"Consistentie: hetzelfde systeem overal")
# Three columns
col_data=[
    ("Kleur",ACCENT1,"Dezelfde primaire kleur, neutrals en accentkleur op elke pagina."),
    ("Typografie",BLUE,"Dezelfde heading- en bodyfonts, zelfde groottes en regelafstand."),
    ("Componenten",GREEN,"Nav, footer, knoppen en kaarten zijn overal hetzelfde component."),
]
cw2=Inches(3.6); cg=Inches(0.22)
for i,(title,col,desc) in enumerate(col_data):
    cl=Inches(1.0)+i*(cw2+cg)
    col_kop(s8,cl,Inches(1.85),cw2,Inches(0.55),title,col)
    box(s8,cl,Inches(2.4),cw2,Inches(1.6),OFFWH,border=(215,225,240))
    lbl(s8,desc,cl+Inches(0.18),Inches(2.54),cw2-Inches(0.36),Inches(1.35),size=14,c=BODY)
# Mini page wireframes row
mini_w=Inches(3.2); mini_h=Inches(1.9); mini_t=Inches(4.2); mini_gap=Inches(0.42)
for mi in range(3):
    ml=Inches(1.15)+mi*(mini_w+mini_gap)
    box(s8,ml,mini_t,mini_w,mini_h,(235,238,245),border=(200,210,220))
    # Consistent colored header bar
    box(s8,ml,mini_t,mini_w,Inches(0.3),ACCENT1)
    # Content placeholder lines
    for li in range(3):
        box(s8,ml+Inches(0.2),mini_t+Inches(0.45)+li*Inches(0.38),mini_w-Inches(0.4),Inches(0.18),(210,215,225))
    lbl(s8,["Home","About","Projecten"][mi],ml,mini_t+mini_h+Inches(0.05),mini_w,Inches(0.28),
        size=11,c=MUTED,align=PP_ALIGN.CENTER)
pnr(s8,8)
notitie(s8,"Laat zien hoe inconsistentie eruit ziet: verschillende kleuren per pagina, wisselende lettertypes. Dan het verschil met een consistent systeem.")

# ─── SLIDE 9: FIGMA MULTI-PAGINA ──────────────────────────────────────────────
s9=ns(dark_bg(False))
hero_hdr(s9,"Figma: meerdere paginas opzetten","Shared styles en componenten")
steps9=[
    "Maak in Figma Pages aan: 'Home', 'About', 'Projecten' (via de tabs linksonder)",
    "Elke page heeft de bijhorende frames: mobile en desktop versie",
    "Gebruik shared styles: kleur- en tekststijlen gelden voor het hele bestand",
    "Hergebruik componenten: navigatie en footer als instance op elke page",
    "Maak een overzichtspage 'Sitemap' met alle frames in miniatuur",
]
for i,step in enumerate(steps9):
    row_t=Inches(2.0)+i*Inches(0.68)
    dot(s9,Inches(1.22),row_t+Inches(0.2),Inches(0.18),ACCENT2)
    lbl(s9,str(i+1),Inches(1.22)-Inches(0.1),row_t+Inches(0.06),Inches(0.2),Inches(0.28),size=11,c=DARK,bold=True,align=PP_ALIGN.CENTER)
    lbl(s9,step,Inches(1.55),row_t,Inches(10.5),Inches(0.6),size=16,c=WHITE)
# Tip card
cx,cy,cw=card(s9,Inches(1.0),Inches(5.7),Inches(11.3),Inches(1.1),bg=(30,40,70))
lbl(s9,"Tip: lock de achtergrondlaag en grid op elke frame via het layer panel om per ongeluk aanpassen te voorkomen.",
    cx,cy,cw,Inches(0.8),size=13,c=(220,200,255),italic=True)
pnr(s9,9,light=True)
notitie(s9,"Demonstreer live in Figma: toon hoe Pages werken, hoe je een component maakt van de navigatie en die instantie op meerdere paginas plaatst.")

# ─── SLIDE 10: INSPIRATIEBRONNEN ─────────────────────────────────────────────
s10=ns(light_bg())
page_hdr(s10,"Inspiratie: goede portfolio's")
insp=[
    ("Dribbble","Visuele shots van UI werk. Kijk naar trending portfolio ontwerpen en laat je inspireren door kleurgebruik.",ACCENT1),
    ("Behance","Volledige case studies met proces, goed voor storytelling ideeen en projectpresentatie.",BLUE),
    ("Webflow Showcase","Live websites met indrukwekkende animaties en interactie. Analyseer de structuur.",GREEN),
    ("Readymag","Portfolio templates die je kunt analyseren op structuur en opmaak.",(217,119,6)),
]
crd_w=Inches(5.4); crd_h=Inches(1.3); crd_gap=Inches(0.22)
positions=[(Inches(1.0),Inches(1.9)),(Inches(6.9),Inches(1.9)),
           (Inches(1.0),Inches(3.38)),(Inches(6.9),Inches(3.38))]
for (cl,ct),(title,desc,col) in zip(positions,insp):
    ix,iy,iw=card(s10,cl,ct,crd_w,crd_h,accent=col)
    lbl(s10,title,ix,iy,iw,Inches(0.32),size=15,c=DARK,bold=True)
    lbl(s10,desc,ix,iy+Inches(0.34),iw,crd_h-Inches(0.55),size=12,c=BODY)
# Bottom note
cx,cy,cw=card(s10,Inches(1.0),Inches(4.88),Inches(11.3),Inches(0.88))
lbl(s10,"Analyseer 2 portfolio's: wat werkt goed in het ontwerp? Waarom trek je je oog naar dit werk?",
    cx,cy,cw,Inches(0.65),size=14,c=BODY,italic=True)
pnr(s10,10)
notitie(s10,"Laat studenten in tweetallen 2 portfolio's analyseren (5 minuten). Vraag daarna: wat was het eerste dat je opviel? Was de hiërarchie duidelijk?")

# ─── SLIDE 11: STIJLGIDS ──────────────────────────────────────────────────────
s11=ns(light_bg())
page_hdr(s11,"Portfolio stijlgids: jouw regels")
style_items=[
    ("Primaire kleur","1 krachtige kleur die past bij jouw stijl",ACCENT1),
    ("Neutrale kleuren","Wit of lichtgrijs voor achtergronden, donkergrijs voor tekst",MUTED),
    ("Heading font","1 expressief lettertype, mag een serif zijn",BLUE),
    ("Body font","1 goed leesbaar sans-serif lettertype",DARK),
    ("Kaart-stijl","Vlak of met subtiele schaduw, consistent op alle paginas",GREEN),
    ("Afbeeldingsverhouding","Kies 16:9 of 4:3 en houd dat aan voor alle thumbnails",(217,119,6)),
]
item_w=Inches(5.4); item_h=Inches(0.7); item_gap=Inches(0.18)
positions11=[(Inches(1.0),Inches(1.9)),(Inches(6.9),Inches(1.9)),
             (Inches(1.0),Inches(2.78)),(Inches(6.9),Inches(2.78)),
             (Inches(1.0),Inches(3.66)),(Inches(6.9),Inches(3.66))]
for (cl,ct),(label,desc,col) in zip(positions11,style_items):
    ix,iy,iw=card(s11,cl,ct,item_w,item_h,accent=col)
    lbl(s11,label,ix,iy,iw,Inches(0.28),size=13,c=DARK,bold=True)
    lbl(s11,desc,ix,iy+Inches(0.3),iw,item_h-Inches(0.38),size=11,c=BODY)
# Tip card
cx,cy,cw=card(s11,Inches(1.0),Inches(4.6),Inches(11.3),Inches(0.78))
lbl(s11,"Schrijf deze regels op in je Figma als een aparte 'stijlgids'-pagina.",
    cx,cy,cw,Inches(0.55),size=14,c=BODY,italic=True)
pnr(s11,11)
notitie(s11,"Laat studenten hun eigen stijlgids opstellen. Dit vormt de basis voor hun portfolio. Ze gebruiken deze regels in de opdracht van vandaag.")

# ─── SLIDE 12: LIVE OEFENING ──────────────────────────────────────────────────
s12=ns(dark_bg(False))
hero_hdr(s12,"Oefening: about pagina wireframen","Doe het nu in Figma")
steps12=[
    "Maak een nieuwe page aan in Figma: 'About'",
    "Voeg frames toe: 'about/mobile' (375x812) en 'about/desktop' (1440x900)",
    "Wireframe de mobile versie: naam, foto, bio, skills, contact",
    "Pas de desktop versie aan: zet de foto naast de bio tekst (2-koloms layout)",
    "Hergebruik je nav- en footercomponent van de homepage",
]
for i,step in enumerate(steps12):
    row_t=Inches(2.0)+i*Inches(0.68)
    dot(s12,Inches(1.22),row_t+Inches(0.2),Inches(0.18),ACCENT1)
    lbl(s12,str(i+1),Inches(1.22)-Inches(0.1),row_t+Inches(0.06),Inches(0.2),Inches(0.28),size=11,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s12,step,Inches(1.55),row_t,Inches(10.5),Inches(0.6),size=16,c=WHITE)
# Goal box
cx,cy,cw=card(s12,Inches(1.0),Inches(5.7),Inches(11.3),Inches(1.05),bg=(30,40,70))
lbl(s12,"Doel: na deze les heb je een werkende about pagina wireframe in Figma.",
    cx,cy,cw,Inches(0.75),size=14,c=(220,200,255),bold=True)
pnr(s12,12,light=True)
notitie(s12,"Geef studenten 20-25 minuten voor deze oefening. Loop langs en geef feedback op de hiërarchie. Vraag: is het direct duidelijk wie deze persoon is?")

# ─── SLIDE 13: BRONNEN ────────────────────────────────────────────────────────
s13=ns(light_bg())
page_hdr(s13,"Bronnen","Verder lezen en verkennen")
bronnen=[
    ("Dribbble","dribbble.com/tags/portfolio","Portfolio inspiratie: trending UI werk en stijlen",ACCENT1),
    ("Behance","behance.net","Volledige case studies met procesweergave",BLUE),
    ("Smashing Magazine","smashingmagazine.com, zoek 'designer portfolio'","Artikelen over personal branding voor designers",GREEN),
    ("Webflow Showcase","webflow.com/made-in-webflow","Live portfolio websites met interactie en animatie",(217,119,6)),
    ("Lapa.ninja","lapa.ninja","Curated portfolio inspiratie, gesorteerd op categorie",(124,58,237)),
]
bw2=Inches(10.8); bh2=Inches(0.75); bg2=Inches(0.16)
bt=Inches(1.9)
for i,(title,url,desc,col) in enumerate(bronnen):
    ix,iy,iw=card(s13,Inches(1.0),bt+i*(bh2+bg2),bw2,bh2,accent=col)
    lbl(s13,title,ix,iy,Inches(2.2),Inches(0.3),size=14,c=DARK,bold=True)
    lbl(s13,url,ix+Inches(2.3),iy,iw-Inches(2.3),Inches(0.3),size=12,c=MUTED)
    lbl(s13,desc,ix,iy+Inches(0.32),iw,bh2-Inches(0.38),size=12,c=BODY)
pnr(s13,13)
notitie(s13,"Verwijs studenten naar deze bronnen voor de opdracht thuis. Dribbble en Behance zijn het meest direct bruikbaar als inspiratie.")

# ─── SLIDE 14: OPDRACHT ───────────────────────────────────────────────────────
s14=ns(dark_bg(False))
gbox(s14,0,0,Inches(0.18),SH,ACCENT1,ACCENT2,angle=270)
lbl(s14,"Opdracht  ·  Les 9",Inches(0.35),Inches(0.5),Inches(10.0),Inches(0.6),size=14,c=(180,180,210))
lbl(s14,"Portfolio Website\nAbout en Projectenpagina",
    Inches(0.35),Inches(0.95),Inches(11.0),Inches(1.6),size=30,c=WHITE,bold=True)
gbox(s14,Inches(0.35),Inches(2.55),Inches(2.5),Inches(0.07),ACCENT1,ACCENT2,angle=0)
steps14=[
    "Ontwerp de About pagina in Figma: mobile en desktop versie",
    "Ontwerp de Projectenpagina: minimaal 3 projectkaarten in een grid",
    "Schrijf een stijlgids-pagina in Figma: kleur, typografie en kaartregels",
    "Zorg voor consistente navigatie en footer op alle paginas",
    "Lever je Figma link in voor les 10",
]
for i,step in enumerate(steps14):
    row_t=Inches(2.8)+i*Inches(0.72)
    dot(s14,Inches(0.52),row_t+Inches(0.2),Inches(0.16),ACCENT1)
    lbl(s14,str(i+1),Inches(0.52)-Inches(0.09),row_t+Inches(0.06),Inches(0.18),Inches(0.28),size=11,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    lbl(s14,step,Inches(0.82),row_t,Inches(11.5),Inches(0.6),size=17,c=WHITE)
pnr(s14,14,light=True)
notitie(s14,"Inlevertermijn: begin van les 10. Herinner studenten aan de requirements: mobile + desktop wireframe voor elke pagina, plus de stijlgids-pagina in Figma.")

# ─── SLIDE 15: VOLGENDE LES ───────────────────────────────────────────────────
s15=ns(dark_bg(True))
gbox(s15,0,0,Inches(0.18),SH,ACCENT1,ACCENT2,angle=270)
lbl(s15,"Volgende les",Inches(1.0),Inches(1.4),Inches(9.0),Inches(0.65),size=18,c=(200,160,240))
lbl(s15,"Design Systems\nVerdieping",
    Inches(1.0),Inches(2.0),Inches(10.5),Inches(1.8),size=44,c=WHITE,bold=True)
gbox(s15,Inches(1.0),Inches(3.85),Inches(2.5),Inches(0.08),ACCENT1,ACCENT2,angle=0)
preview=[
    "Design tokens: kleur, spacing en typografie als variabelen",
    "Component library organiseren en documenteren",
    "Figma variables en modes",
    "Handoff naar developers via het Inspect panel",
]
for i,item in enumerate(preview):
    row_t=Inches(4.05)+i*Inches(0.55)
    dot(s15,Inches(1.22),row_t+Inches(0.16),Inches(0.1),ACCENT2)
    lbl(s15,item,Inches(1.45),row_t,Inches(10.0),Inches(0.48),size=15,c=(220,200,255))
# Tip box at bottom
cx,cy,cw=card(s15,Inches(1.0),Inches(6.35),Inches(11.3),Inches(0.75),bg=(25,35,65))
lbl(s15,"Bekijk alvast tokens.studio voor een blik op hoe design tokens werken.",
    cx,cy,cw,Inches(0.52),size=13,c=(220,200,255),italic=True)
pnr(s15,15,light=True)
notitie(s15,"Sluit de les af. Vraag een student om samen te vatten wat ze vandaag hebben geleerd. Herinner aan de opdracht: Figma link inleveren voor les 10.")

# ─── OPSLAAN ──────────────────────────────────────────────────────────────────
OUT=r"C:\Users\Peter Marcelis\Documents\Lesmateriaal\design\periode-2\les-9\09_portfolio-personal-branding.pptx"
prs.save(OUT)
print(f"Presentatie opgeslagen: {OUT}")
print(f"Aantal slides: {len(prs.slides)}")
