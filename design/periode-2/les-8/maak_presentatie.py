"""
Les 8: Navigatiepatronen en Sliders
Periode 2, Jaar 1 - Portfolio Website project
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image, ImageDraw
import numpy as np, io

NAVY=(6,14,44); NAVY2=(30,58,138); BLUE=(37,99,235); BLUE2=(59,130,246)
CYAN=(6,182,212); TEAL=(13,148,136); WHITE=(255,255,255); OFFWH=(248,250,252)
DARK=(15,23,42); BODY=(51,65,85); MUTED=(100,116,135)
GREEN=(22,163,74); RED=(220,38,38)
ACCENT1=(14,165,233); ACCENT2=(56,189,248)  # Sky Blue for les 8

def rgb(c): return RGBColor(c[0],c[1],c[2])
SW,SH=Inches(13.33),Inches(7.5); BW,BH=1920,1080

def garr(w,h,c1,c2,d='diag'):
    x,y=np.linspace(0,1,w,dtype=np.float32),np.linspace(0,1,h,dtype=np.float32)
    if d=='h': t=np.tile(x,(h,1))
    elif d=='v': t=np.tile(y.reshape(-1,1),(1,w))
    else: xx,yy=np.meshgrid(x,y); t=xx*.55+yy*.45
    a=np.zeros((h,w,3),dtype=np.float32)
    for i in range(3): a[:,:,i]=c1[i]+(c2[i]-c1[i])*t
    return np.clip(a,0,255).astype(np.uint8)

def dark_bg(circ=True):
    img=Image.fromarray(garr(BW,BH,NAVY,NAVY2)).convert('RGBA')
    if circ:
        ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
        d.ellipse([BW-500,-180,BW+200,520],fill=(ACCENT1[0],ACCENT1[1],ACCENT1[2],20))
        d.ellipse([-60,BH-360,380,BH+100],fill=(ACCENT2[0],ACCENT2[1],ACCENT2[2],22))
        img=Image.alpha_composite(img,ov)
    return img.convert('RGB')

def light_bg():
    img=Image.fromarray(garr(BW,BH,OFFWH,WHITE)).convert('RGBA')
    ov=Image.new('RGBA',(BW,BH),(0,0,0,0)); d=ImageDraw.Draw(ov)
    d.ellipse([BW-580,-220,BW+180,540],fill=(ACCENT1[0],ACCENT1[1],ACCENT1[2],8))
    img=Image.alpha_composite(img,ov); return img.convert('RGB')

def buf(img):
    b=io.BytesIO(); img.save(b,format='PNG'); b.seek(0); return b

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
BLANK=prs.slide_layouts[6]

def ns(bg=None):
    s=prs.slides.add_slide(BLANK)
    if bg:
        p=s.shapes.add_picture(buf(bg),0,0,SW,SH)
        t=s.shapes._spTree; t.remove(p._element); t.insert(2,p._element)
    else:
        s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(WHITE)
    return s

def box(s,l,t,w,h,c,border=None):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c)
    if border: sh.line.color.rgb=rgb(border); sh.line.width=Pt(0.75)
    else: sh.line.fill.background()
    return sh

def gbox(s,l,t,w,h,c1,c2,angle=90):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,l,t,w,h)
    sh.fill.gradient(); st=sh.fill.gradient_stops
    st[0].position=0.0; st[0].color.rgb=rgb(c1)
    st[1].position=1.0; st[1].color.rgb=rgb(c2)
    sh.fill.gradient_angle=angle; sh.line.fill.background(); return sh

def dot(s,cx,cy,r,c):
    sh=s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,cx-r,cy-r,r*2,r*2)
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(c); sh.line.fill.background()

def lbl(s,text,l,t,w,h,size=18,c=BODY,bold=False,italic=False,align=PP_ALIGN.LEFT):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=text; r.font.name="Calibri"; r.font.size=Pt(size)
    r.font.bold=bold; r.font.italic=italic; r.font.color.rgb=rgb(c)

def blist(s,items,l,t,w,h,size=18,c=BODY,gap=7):
    bx=s.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(gap); r=p.add_run()
        r.text=f"  {chr(8226)}  {item}"; r.font.name="Calibri"; r.font.size=Pt(size); r.font.color.rgb=rgb(c)

def notitie(s,t): s.notes_slide.notes_text_frame.text=t

def pnr(s,n,tot=15,light=False):
    lbl(s,f"{n} / {tot}",Inches(12.5),Inches(7.1),Inches(0.8),Inches(0.35),size=11,c=(180,180,200) if light else MUTED,align=PP_ALIGN.RIGHT)

def hero_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(1.0),size=36,c=WHITE,bold=True)
    gbox(s,Inches(1.0),Inches(1.42),Inches(2.0),Inches(0.08),ACCENT1,ACCENT2,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.55),Inches(11.3),Inches(0.42),size=15,c=(180,230,255),italic=True)

def page_hdr(s,title,sub=None):
    lbl(s,title,Inches(1.0),Inches(0.5),Inches(11.3),Inches(0.95),size=34,c=DARK,bold=True)
    gbox(s,Inches(1.0),Inches(1.38),Inches(1.8),Inches(0.08),ACCENT1,ACCENT2,angle=0)
    if sub: lbl(s,sub,Inches(1.0),Inches(1.5),Inches(11.3),Inches(0.38),size=14,c=MUTED,italic=True)

def card(s,l,t,w,h,accent=None,bg=OFFWH):
    a=accent if accent else ACCENT1
    STRIP=Inches(0.07)
    box(s,l,t,w,h,bg,border=(215,225,240))
    gbox(s,l,t,STRIP,h,a,ACCENT2,angle=270)
    return l+STRIP+Inches(0.18), t+Inches(0.14), w-STRIP-Inches(0.35)

def col_kop(s,l,t,w,h,text,c):
    box(s,l,t,w,h,c)
    lbl(s,text,l,t+Inches(0.07),w,h-Inches(0.07),size=20,c=WHITE,bold=True,align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 1: Titel
# ─────────────────────────────────────────────
s1=ns(dark_bg(True))
# Accent strip left
gbox(s1,Inches(0.0),Inches(0.0),Inches(0.18),SH,ACCENT1,ACCENT2,angle=270)
lbl(s1,"Navigatiepatronen\nen Sliders",Inches(1.0),Inches(1.5),Inches(9.5),Inches(2.5),size=52,c=WHITE,bold=True)
gbox(s1,Inches(1.0),Inches(3.88),Inches(2.5),Inches(0.1),ACCENT1,ACCENT2,angle=0)
lbl(s1,"Periode 2  |  Les 8",Inches(1.0),Inches(4.05),Inches(8.0),Inches(0.5),size=18,c=(180,230,255),bold=False)
lbl(s1,"Van hamburgermenu tot interactieve carousel",Inches(1.0),Inches(4.6),Inches(9.5),Inches(0.5),size=15,c=(140,200,255),italic=True)
pnr(s1,1,light=True)
notitie(s1,"Welkom bij les 8. Vandaag behandelen we navigatiepatronen voor desktop en mobiel, sliderdesign en Figma prototyping. Dit sluit direct aan bij het portfolio website project.")

# ─────────────────────────────────────────────
# SLIDE 2: Agenda
# ─────────────────────────────────────────────
s2=ns(light_bg())
page_hdr(s2,"Vandaag","Overzicht van de les")
agenda=[
    "1.  Navigatiepatronen: desktop versus mobiel",
    "2.  Hamburger menu, tabs en bottom navigation",
    "3.  Sliders en carousels: design principes",
    "4.  Toegankelijkheid van navigatie en sliders",
    "5.  Figma: prototype van een slider bouwen",
    "6.  Portfolio: navigatie en slider uitwerken",
]
cx,cy=Inches(1.0),Inches(1.8)
cw,ch=Inches(11.3),Inches(5.2)
box(s2,cx,cy,cw,ch,OFFWH,border=(215,225,240))
gbox(s2,cx,cy,Inches(0.07),ch,ACCENT1,ACCENT2,angle=270)
blist(s2,agenda,cx+Inches(0.32),cy+Inches(0.18),cw-Inches(0.5),ch-Inches(0.25),size=19,c=BODY,gap=9)
pnr(s2,2)
notitie(s2,"Loop de agenda door. Leg uit dat navigatie en sliders kernelementen zijn van elke portfoliowebsite. Studenten gaan vandaag zowel theorie leren als in Figma prototypes bouwen.")

# ─────────────────────────────────────────────
# SLIDE 3: Desktop navigatiepatronen
# ─────────────────────────────────────────────
s3=ns(light_bg())
page_hdr(s3,"Desktop: navigatiepatronen","Drie veelgebruikte patronen")
panels=[
    ("Horizontale nav", ACCENT1, [
        "Bovenaan de pagina",
        "Logo links, links midden of rechts",
        "Sticky gedrag mogelijk",
        "Meest voorkomend patroon",
        "Max 5-7 links",
    ]),
    ("Mega menu", NAVY2, [
        "Voor sites met veel content",
        "Hover toont een groot dropdown-panel",
        "Gebruikt door grote retailers",
        "en nieuwssites",
    ]),
    ("Sidebar nav", BLUE, [
        "Aan de linkerkant van de pagina",
        "Goed voor dashboards",
        "Meerdere niveaus mogelijk",
        "Neemt horizontale ruimte",
    ]),
]
PW=Inches(3.7); PH=Inches(4.5); PT=Inches(1.72); GAP=Inches(0.22)
starts=[Inches(0.62), Inches(0.62)+PW+GAP, Inches(0.62)+2*(PW+GAP)]
for i,(titel,kleur,items) in enumerate(panels):
    l=starts[i]
    col_kop(s3,l,PT,PW,Inches(0.52),titel,kleur)
    box(s3,l,PT+Inches(0.52),PW,PH-Inches(0.52),OFFWH,border=(215,225,240))
    blist(s3,items,l+Inches(0.18),PT+Inches(0.62),PW-Inches(0.22),PH-Inches(0.7),size=15,c=BODY,gap=6)
pnr(s3,3)
notitie(s3,"Bespreek de drie patronen. Horizontale nav is de standaard voor portfoliosites. Mega menu is relevant voor grote sites zoals wehkamp.nl. Sidebar past goed bij admin-interfaces en dashboards.")

# ─────────────────────────────────────────────
# SLIDE 4: Mobiele navigatiepatronen
# ─────────────────────────────────────────────
s4=ns(light_bg())
page_hdr(s4,"Mobiel: navigatiepatronen","Drie veelgebruikte patronen")
mob_panels=[
    ("Hamburger menu", ACCENT1, [
        "Drie streepjes als icoon",
        "Verbergt nav achter een overlay",
        "Universeel herkend",
        "Goed voor veel links",
        "Nadeel: lagere discoverability",
    ]),
    ("Bottom navigation", GREEN, [
        "3-5 icoonlinks onderaan het scherm",
        "Duimvriendelijk bereikbaar",
        "Gebruikt door apps",
        "Minder geschikt voor websites",
    ]),
    ("Tab bar", BLUE, [
        "Bovenaan of als tabs in content",
        "Goed voor 3-5 secties",
        "Duidelijke actieve staat",
        "is belangrijk",
    ]),
]
for i,(titel,kleur,items) in enumerate(mob_panels):
    l=starts[i]
    col_kop(s4,l,PT,PW,Inches(0.52),titel,kleur)
    box(s4,l,PT+Inches(0.52),PW,PH-Inches(0.52),OFFWH,border=(215,225,240))
    blist(s4,items,l+Inches(0.18),PT+Inches(0.62),PW-Inches(0.22),PH-Inches(0.7),size=15,c=BODY,gap=6)
pnr(s4,4)
notitie(s4,"Voor een portfolio website is een hamburger menu op mobiel de meest voor de hand liggende keuze. Bottom navigation zie je vooral in native apps zoals Instagram en YouTube. Tab bars werken goed voor beperkte navigatiesecties.")

# ─────────────────────────────────────────────
# SLIDE 5: Hamburger menu design checklist
# ─────────────────────────────────────────────
s5=ns(light_bg())
page_hdr(s5,"Hamburger menu: design checklist","Wat maakt een goed hamburger menu?")

# Left: wireframe simulation
WL=Inches(1.0); WT=Inches(1.8); WW=Inches(4.2); WH=Inches(5.0)
box(s5,WL,WT,WW,WH,(230,240,255),border=(180,200,230))
lbl(s5,"Mobiele viewport",WL,WT+Inches(0.08),WW,Inches(0.35),size=11,c=MUTED,align=PP_ALIGN.CENTER)

# Phone header bar
box(s5,WL+Inches(0.12),WT+Inches(0.42),WW-Inches(0.24),Inches(0.48),NAVY)
lbl(s5,"Mijn Portfolio",WL+Inches(0.22),WT+Inches(0.47),Inches(2.0),Inches(0.38),size=13,c=WHITE,bold=True)
# Hamburger icon (three horizontal bars)
bar_l=WL+Inches(3.2); bar_top=WT+Inches(0.52); bar_w=Inches(0.5); bar_h=Inches(0.06); bar_gap=Inches(0.12)
for bi in range(3):
    box(s5,bar_l,bar_top+bi*bar_gap*1.5,bar_w,bar_h,WHITE)

# Open overlay panel
box(s5,WL+Inches(0.12),WT+Inches(0.9),WW-Inches(0.24),Inches(3.5),(30,40,80))
lbl(s5,"X  Sluiten",WL+Inches(0.25),WT+Inches(0.96),Inches(1.5),Inches(0.35),size=13,c=WHITE)
for li,link in enumerate(["Home","About","Projecten","Contact"]):
    box(s5,WL+Inches(0.12),WT+Inches(1.42)+li*Inches(0.6),WW-Inches(0.24),Inches(0.5),(40,55,100))
    lbl(s5,link,WL+Inches(0.3),WT+Inches(1.47)+li*Inches(0.6),Inches(2.5),Inches(0.4),size=14,c=WHITE,bold=(li==0))
    if li==0:
        box(s5,WL+Inches(0.12),WT+Inches(1.42),Inches(0.06),Inches(0.5),ACCENT1)

lbl(s5,"Wireframe: overlay open",WL,WT+WH-Inches(0.28),WW,Inches(0.25),size=10,c=MUTED,italic=True,align=PP_ALIGN.CENTER)

# Right: checklist
RL=Inches(5.5); RT=Inches(1.8); RW=Inches(7.5); RH=Inches(5.0)
checklist=[
    "Icoon duidelijk zichtbaar: minimaal 44x44px taptarget",
    "Label 'Menu' toevoegen verbetert herkenbaarheid",
    "Overlay dekt het hele scherm of schuift van links",
    "Sluitknop zichtbaar bovenin het overlay",
    "Actieve link is visueel gemarkeerd",
    "Animatie: smooth slide of fade, niet te snel",
]
cx2,_,cw2=card(s5,RL,RT,RW,RH,accent=ACCENT1)
lbl(s5,"Checklist",RL+Inches(0.32),RT+Inches(0.1),RW,Inches(0.35),size=14,c=ACCENT1,bold=True)
blist(s5,checklist,cx2,RT+Inches(0.5),cw2,RH-Inches(0.6),size=15,c=BODY,gap=8)
pnr(s5,5)
notitie(s5,"Loop de checklist punt voor punt door. Demonstreer op een echte site (bv. nu.nl op mobiel) hoe het hamburger menu werkt. Let op de animatie en de sluitknop.")

# ─────────────────────────────────────────────
# SLIDE 6: Slider design principes
# ─────────────────────────────────────────────
s6=ns(dark_bg(False))
hero_hdr(s6,"Slider en carousel: design principes","Wat maakt een slider effectief en gebruiksvriendelijk?")
bullets=[
    "Toon altijd indicatoren hoeveel slides er zijn (puntjes of nummers)",
    "Pijlknoppen links en rechts, ook op mobiel",
    "Autoplay is controversieel: schakel het uit of geef een pauzeknop",
    "Maximaal 5 slides, anders verliest de gebruiker het overzicht",
    "Elke slide heeft een duidelijke focus: 1 afbeelding of 1 boodschap",
    "Touch swipe op mobiel is een must",
]
blist(s6,bullets,Inches(1.0),Inches(2.0),Inches(11.3),Inches(3.5),size=17,c=(220,235,255),gap=8)

# Portfolio tip card
cx3,cy3,cw3=card(s6,Inches(1.0),Inches(5.6),Inches(11.3),Inches(1.5),accent=ACCENT1,bg=(20,30,70))
lbl(s6,"Portfolio tip",Inches(1.32),Inches(5.7),Inches(2.5),Inches(0.35),size=12,c=ACCENT1,bold=True)
lbl(s6,"Voor een portfolio slider: toon je projecten als slides, met projectnaam, categorie en een korte omschrijving.",cx3,cy3,cw3,Inches(1.1),size=14,c=(200,220,255))
pnr(s6,6,light=True)
notitie(s6,"Bespreek onderzoek van Nielsen Norman Group: autoplay carousels worden grotendeels genegeerd door gebruikers. Statische of gebruiker-gestuurde sliders presteren beter. Verwijs naar de bronnen op slide 14.")

# ─────────────────────────────────────────────
# SLIDE 7: Slider anatomie
# ─────────────────────────────────────────────
s7=ns(light_bg())
page_hdr(s7,"Slider: anatomie van een goed ontwerp","Elk onderdeel heeft een functie")

SL=Inches(1.5); ST=Inches(1.75); SW2=Inches(10.3); SH2=Inches(3.7)
# Main slide area
box(s7,SL,ST,SW2,SH2,DARK)
# Image placeholder area inside slide
box(s7,SL+Inches(0.08),ST+Inches(0.08),Inches(4.5),SH2-Inches(0.16),(35,50,80))
lbl(s7,"[Projectafbeelding]",SL+Inches(0.2),ST+SH2/2-Inches(0.2),Inches(4.3),Inches(0.4),size=13,c=MUTED,align=PP_ALIGN.CENTER)

# Text area inside slide
box(s7,SL+Inches(4.72),ST+Inches(0.08),Inches(5.5),SH2-Inches(0.16),(25,35,65))
lbl(s7,"Projectnaam",SL+Inches(4.9),ST+Inches(0.25),Inches(5.2),Inches(0.45),size=18,c=WHITE,bold=True)
box(s7,SL+Inches(4.72),ST+Inches(0.75),Inches(2.2),Inches(0.28),ACCENT1)
lbl(s7,"Categorie",SL+Inches(4.82),ST+Inches(0.77),Inches(2.0),Inches(0.24),size=11,c=WHITE,bold=True)
box(s7,SL+Inches(4.72),ST+Inches(1.15),Inches(5.2),Inches(0.12),(50,70,110))
box(s7,SL+Inches(4.72),ST+Inches(1.4),Inches(5.2),Inches(0.12),(50,70,110))
box(s7,SL+Inches(4.72),ST+Inches(1.65),Inches(3.5),Inches(0.12),(50,70,110))
lbl(s7,"Omschrijving tekst",SL+Inches(4.9),ST+Inches(1.1),Inches(5.0),Inches(0.28),size=11,c=MUTED)

# Prev arrow button
PAL=SL-Inches(0.58); PAT=ST+SH2/2-Inches(0.3); PAW=Inches(0.45); PAH=Inches(0.6)
box(s7,PAL,PAT,PAW,PAH,ACCENT1)
lbl(s7,"<",PAL,PAT+Inches(0.06),PAW,PAH-Inches(0.06),size=18,c=WHITE,bold=True,align=PP_ALIGN.CENTER)

# Next arrow button
NAL=SL+SW2+Inches(0.13); NAT=PAT
box(s7,NAL,NAT,PAW,PAH,ACCENT1)
lbl(s7,">",NAL,NAT+Inches(0.06),PAW,PAH-Inches(0.06),size=18,c=WHITE,bold=True,align=PP_ALIGN.CENTER)

# Dot indicators
dot_y=ST+SH2+Inches(0.35); dot_cx=SL+SW2/2-Inches(0.5)
for di in range(5):
    dot_color=ACCENT1 if di==0 else MUTED
    dot(s7,dot_cx+di*Inches(0.28),dot_y,Inches(0.07),dot_color)

# Annotations
lbl(s7,"Vorige slide",PAL-Inches(0.9),PAT-Inches(0.08),Inches(0.85),Inches(0.35),size=10,c=MUTED,italic=True,align=PP_ALIGN.RIGHT)
lbl(s7,"Volgende slide",NAL+PAW+Inches(0.08),NAT-Inches(0.08),Inches(0.9),Inches(0.35),size=10,c=MUTED,italic=True)
lbl(s7,"Slide indicator (actief puntje = huidige slide)",dot_cx-Inches(0.2),dot_y+Inches(0.2),Inches(3.0),Inches(0.35),size=10,c=MUTED,italic=True,align=PP_ALIGN.CENTER)
lbl(s7,"Slide inhoud",SL+Inches(0.2),ST-Inches(0.32),Inches(2.5),Inches(0.28),size=10,c=MUTED,italic=True)
lbl(s7,"Swipe area (mobiel)",SL+Inches(4.0),ST+SH2+Inches(0.05),Inches(3.0),Inches(0.28),size=10,c=MUTED,italic=True,align=PP_ALIGN.CENTER)
pnr(s7,7)
notitie(s7,"Wijs elk onderdeel aan op de wireframe. Bespreek waarom pijlknoppen ook op mobiel zichtbaar moeten zijn, naast swipe-support. De dot-indicatoren communiceren hoeveel slides er zijn.")

# ─────────────────────────────────────────────
# SLIDE 8: Toegankelijkheid
# ─────────────────────────────────────────────
s8=ns(dark_bg(False))
hero_hdr(s8,"Toegankelijkheid: navigatie en sliders","WCAG-richtlijnen en ARIA patterns")

# Left column: Navigatie
NL=Inches(1.0); NT=Inches(2.05); NW=Inches(5.6); NH=Inches(4.7)
cx4,cy4,cw4=card(s8,NL,NT,NW,NH,accent=ACCENT1,bg=(20,30,70))
lbl(s8,"Navigatie",cx4,cy4,cw4,Inches(0.35),size=16,c=ACCENT1,bold=True)
nav_items=[
    "Alle links bereikbaar via de Tab-toets",
    "aria-label=\"Hoofdnavigatie\" op het nav-element",
    "Actieve pagina: aria-current=\"page\"",
]
blist(s8,nav_items,cx4,cy4+Inches(0.4),cw4,NH-Inches(0.55),size=14,c=(200,220,255),gap=7)

# Right column: Slider
SRL=Inches(7.0); SRT=Inches(2.05); SRW=Inches(5.6); SRH=Inches(4.7)
cx5,cy5,cw5=card(s8,SRL,SRT,SRW,SRH,accent=CYAN,bg=(20,30,70))
lbl(s8,"Slider",cx5,cy5,cw5,Inches(0.35),size=16,c=CYAN,bold=True)
slider_items=[
    "Pijlknoppen hebben aria-label (\"Vorige\"/\"Volgende\")",
    "Pauzeknop voor autoplay-sliders",
    "prefers-reduced-motion: geen animatie",
    "Keyboard: pijltoetsen navigeren door slides",
]
blist(s8,slider_items,cx5,cy5+Inches(0.4),cw5,SRH-Inches(0.55),size=14,c=(200,220,255),gap=7)

# Bottom quote
box(s8,Inches(1.0),Inches(6.9),Inches(11.3),Inches(0.42),(20,30,70))
lbl(s8,"Toegankelijkheid is geen extra, het is een eis voor professioneel design.",Inches(1.2),Inches(6.93),Inches(11.0),Inches(0.38),size=13,c=ACCENT2,italic=True,bold=True)
pnr(s8,8,light=True)
notitie(s8,"Bespreek het belang van WCAG 2.1 AA. In Nederland zijn overheidswebsites verplicht toegankelijk. Voor portfolio-websites is het een teken van professionaliteit. Laat zien hoe Tab-navigatie werkt in de browser.")

# ─────────────────────────────────────────────
# SLIDE 9: Figma prototype: slider
# ─────────────────────────────────────────────
s9=ns(dark_bg(False))
hero_hdr(s9,"Figma: slider prototype bouwen","Stap voor stap een interactieve slider maken")
steps=[
    "1.  Maak 3 slide-frames: 'slide-1', 'slide-2', 'slide-3' naast elkaar",
    "2.  Maak een hoofdframe 'slider' dat de slider omhult",
    "3.  Klik op het pijl-icoon (Next knop), verbind het met slide-2 via On Click",
    "4.  Stel de animatie in: Smart Animate, Ease Out, 300ms",
    "5.  Test in de preview (Ctrl+Alt+P) of de slider werkt",
    "6.  Voeg dot-indicators toe als aparte componenten met active variant",
]
blist(s9,steps,Inches(1.0),Inches(2.05),Inches(11.3),Inches(4.5),size=16,c=(220,235,255),gap=10)

# Tip box
cx6,cy6,cw6=card(s9,Inches(1.0),Inches(6.2),Inches(11.3),Inches(0.95),accent=ACCENT1,bg=(20,30,70))
lbl(s9,"Tip: gebruik 'After Delay' interactie voor een autoplay-effect (stel in op 3000ms).",cx6,cy6,cw6,Inches(0.75),size=13,c=(200,220,255),italic=True)
pnr(s9,9,light=True)
notitie(s9,"Demonstreer live in Figma hoe je frames koppelt voor een prototype. Laat zien hoe Smart Animate werkt voor een vloeiende overgang. Studenten volgen mee op hun eigen laptop.")

# ─────────────────────────────────────────────
# SLIDE 10: Figma prototype: hamburger menu
# ─────────────────────────────────────────────
s10=ns(dark_bg(False))
hero_hdr(s10,"Figma: hamburger menu prototype","Open en sluit animatie bouwen in Figma")
hsteps=[
    "1.  Maak twee frames: 'nav/closed' en 'nav/open'",
    "2.  In 'nav/open': voeg een donker overlay-frame toe met de links erin",
    "3.  Verbind het hamburger-icoon in 'nav/closed' met 'nav/open' (On Click, overlay)",
    "4.  Voeg een sluitknop toe in 'nav/open' die teruglinkt naar 'nav/closed'",
    "5.  Animatie: Slide in from right, 300ms, Ease Out",
]
blist(s10,hsteps,Inches(1.0),Inches(2.05),Inches(11.3),Inches(4.2),size=16,c=(220,235,255),gap=10)

cx7,cy7,cw7=card(s10,Inches(1.0),Inches(6.2),Inches(11.3),Inches(0.95),accent=ACCENT2,bg=(20,30,70))
lbl(s10,"Tip: gebruik 'Open overlay' in Figma prototype voor het menu-overlay patroon.",cx7,cy7,cw7,Inches(0.75),size=13,c=(200,220,255),italic=True)
pnr(s10,10,light=True)
notitie(s10,"Laat zien hoe de overlay verbinding werkt in Figma. De student moet begrijpen dat 'nav/closed' en 'nav/open' twee aparte frames zijn die via een interactie verbonden worden.")

# ─────────────────────────────────────────────
# SLIDE 11: Portfolio navigatie uitwerken
# ─────────────────────────────────────────────
s11=ns(light_bg())
page_hdr(s11,"Portfolio: navigatie ontwerpen","Stap voor stap de navigatiecomponent uitwerken")
stappen11=[
    "Desktop: horizontale nav met logo links, 3 links (Home, About, Projecten) en een contact-CTA rechts",
    "Mobiel: hamburger icoon rechts, bij klik fullscreen overlay met dezelfde 3 links",
    "Gebruik je navigatiecomponent uit periode 1 als basis",
    "Voeg een active state toe: welke pagina is actief",
    "Maak het component responsive met twee varianten: desktop en mobiel",
]
cx8,cy8,cw8=card(s11,Inches(1.0),Inches(1.72),Inches(11.3),Inches(5.2),accent=ACCENT1)
blist(s11,stappen11,cx8,cy8,cw8,Inches(4.8),size=17,c=BODY,gap=9)

# Small callout
lbl(s11,"Lever beide varianten op als component in je Figma-bestand.",Inches(1.3),Inches(7.05),Inches(10.5),Inches(0.32),size=12,c=MUTED,italic=True)
pnr(s11,11)
notitie(s11,"Herinner studenten aan de navigatiecomponent die ze in periode 1 hebben gemaakt. Zo bouwen ze voort op bestaand werk. De twee varianten (desktop/mobiel) zijn een goede oefening in Figma-componenten met varianten.")

# ─────────────────────────────────────────────
# SLIDE 12: Portfolio slider uitwerken
# ─────────────────────────────────────────────
s12=ns(light_bg())
page_hdr(s12,"Portfolio: slider ontwerpen","De slider voor de homepage uitwerken")
stappen12=[
    "Kies 3 projecten als slidecontent",
    "Elke slide: projectafbeelding placeholder, projectnaam, categorie en korte omschrijving",
    "Navigatie: pijlen links en rechts, puntjes onderaan",
    "Desktop: slider breed (80% van de paginabreedte)",
    "Mobiel: slider volledig breed, swipe gesture",
    "Maak een Figma prototype met 3 werkende slides",
]
cx9,cy9,cw9=card(s12,Inches(1.0),Inches(1.72),Inches(11.3),Inches(5.2),accent=CYAN)
blist(s12,stappen12,cx9,cy9,cw9,Inches(4.8),size=17,c=BODY,gap=9)
lbl(s12,"Gebruik placeholder-rechthoeken voor afbeeldingen totdat je echte projecten hebt.",Inches(1.3),Inches(7.05),Inches(10.5),Inches(0.32),size=12,c=MUTED,italic=True)
pnr(s12,12)
notitie(s12,"Studenten kiezen 3 van hun eigen projecten uit periode 1 als inhoud voor de slider. Als ze nog geen projecten hebben, gebruiken ze fictieve namen. De focus ligt op het ontwerp, niet op de echte inhoud.")

# ─────────────────────────────────────────────
# SLIDE 13: Live analyse
# ─────────────────────────────────────────────
s13=ns(light_bg())
page_hdr(s13,"Live analyse: navigatie en sliders testen","Klassikale verkenning - 10 minuten")
taken=[
    ("01", "Open je portfolio wireframe van les 7. Voeg een hamburger menu toe aan de mobiele versie.", ACCENT1),
    ("02", "Zoek een website met een goede slider. Wat zijn de indicatoren? Zijn er pijlknoppen?", BLUE),
    ("03", "Test je favoriete site met de Tab-toets in de browser. Kun je alle links bereiken?", GREEN),
    ("04", "Bekijk Apple.com op mobiel versus desktop. Wat verandert er aan de navigatie?", (124,58,237)),
]
TW=Inches(5.5); TH=Inches(1.45); TT=Inches(1.72); TGAP=Inches(0.12)
for i,(nr,tekst,kleur) in enumerate(taken):
    col=i%2; row=i//2
    tl=Inches(0.62)+col*(TW+Inches(0.22)); tt=TT+row*(TH+TGAP)
    box(s13,tl,tt,TW,TH,OFFWH,border=(215,225,240))
    gbox(s13,tl,tt,Inches(0.07),TH,kleur,ACCENT2,angle=270)
    lbl(s13,nr,tl+Inches(0.22),tt+Inches(0.1),Inches(0.5),Inches(0.4),size=20,c=kleur,bold=True)
    lbl(s13,tekst,tl+Inches(0.22),tt+Inches(0.48),TW-Inches(0.35),Inches(0.85),size=13,c=BODY)
pnr(s13,13)
notitie(s13,"Geef studenten 10 minuten voor de taken. Daarna klassikale bespreking: wat viel op? Welke slider was het beste? Was de navigatie via Tab-toets bruikbaar? Dit verbindt theorie met de eigen ervaring als gebruiker.")

# ─────────────────────────────────────────────
# SLIDE 14: Bronnen
# ─────────────────────────────────────────────
s14=ns(light_bg())
page_hdr(s14,"Bronnen","Verdieping en documentatie")
bronnen=[
    ("Nielsen Norman: Hamburger menus","nngroup.com/articles/hamburger-menus", ACCENT1),
    ("Smashing Magazine: Carousel design","smashingmagazine.com  |  zoek 'carousel design'", BLUE),
    ("WAI-ARIA navigation patterns","w3.org/WAI/ARIA/apg/patterns", GREEN),
    ("Figma prototype documentation","figma.com/prototyping", (124,58,237)),
    ("Swiper.js (slider bibliotheek)","swiperjs.com", NAVY2),
]
BRW=Inches(11.3); BRT=Inches(1.72); BRH=Inches(0.92); BRGAP=Inches(0.1)
for i,(titel,url,kleur) in enumerate(bronnen):
    bt=BRT+i*(BRH+BRGAP)
    cx10,cy10,cw10=card(s14,Inches(1.0),bt,BRW,BRH,accent=kleur)
    lbl(s14,titel,cx10,cy10,cw10*0.5,BRH-Inches(0.18),size=14,c=DARK,bold=True)
    lbl(s14,url,cx10+cw10*0.5,cy10,cw10*0.5,BRH-Inches(0.18),size=13,c=MUTED,italic=True,align=PP_ALIGN.RIGHT)
pnr(s14,14)
notitie(s14,"Deel deze bronnen via de digitale leeromgeving. Nielsen Norman is de belangrijkste referentie voor navigatie en sliders. Swiper.js is de meest gebruikte slider-bibliotheek voor webprojecten.")

# ─────────────────────────────────────────────
# SLIDE 15: Opdracht
# ─────────────────────────────────────────────
s15=ns(dark_bg(False))
# Accent strip left
gbox(s15,Inches(0.0),Inches(0.0),Inches(0.18),SH,ACCENT1,ACCENT2,angle=270)
lbl(s15,"Opdracht  |  Les 8",Inches(0.5),Inches(0.35),Inches(12.0),Inches(0.5),size=14,c=ACCENT1,bold=True)
lbl(s15,"Portfolio Website  |  Navigatie en Slider",Inches(0.5),Inches(0.8),Inches(12.0),Inches(0.42),size=13,c=(160,200,240),italic=True)

opdracht=[
    "1.  Ontwerp de desktop navigatie: horizontaal, logo + 3 links + CTA",
    "2.  Ontwerp de mobiele navigatie: hamburger + overlay",
    "3.  Ontwerp de slider voor de homepage met 3 slides",
    "4.  Bouw een Figma prototype van de slider (3 werkende overgangsanimaties)",
    "5.  Maak een prototype van het hamburger menu (open en sluit animatie)",
    "6.  Lever je Figma link in voor les 9",
]
blist(s15,opdracht,Inches(0.5),Inches(1.42),Inches(12.3),Inches(5.5),size=17,c=(220,235,255),gap=10)

box(s15,Inches(0.5),Inches(6.8),Inches(12.3),Inches(0.45),(20,30,70))
lbl(s15,"Inleveren via de leeromgeving voor de volgende les. Figma-link delen is voldoende.",Inches(0.7),Inches(6.84),Inches(12.0),Inches(0.38),size=13,c=ACCENT2,italic=True)
pnr(s15,15,light=True)
notitie(s15,"Sluit de les af door de opdracht toe te lichten. Studenten werken thuis verder aan het Figma prototype. Volgende les (les 9) gaan we in op portfolio layout en typografie. Zorg dat studenten hun Figma-link klaar hebben.")

# ─────────────────────────────────────────────
# Opslaan
# ─────────────────────────────────────────────
OUT=r"C:\Users\Peter Marcelis\Documents\Lesmateriaal\design\periode-2\les-8\08_navigatie-sliders.pptx"
prs.save(OUT)
print(f"Presentatie opgeslagen: {OUT}")
print(f"Aantal slides: {len(prs.slides)}")
